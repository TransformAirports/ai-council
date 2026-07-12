#!/usr/bin/env python
"""Build the companion executive deck for
'Future-Proofing the Dulles Terminal Rebuild' (The Stopwatch and the Shell).
13 slides, 16:9, Georgia headlines / Calibri body, navy #143C6E on white.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

NAVY = RGBColor(0x14, 0x3C, 0x6E)
GRAY = RGBColor(0x5A, 0x5A, 0x5A)
LIGHT = RGBColor(0x8A, 0x9A, 0xB5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SERIF = "Georgia"
SANS = "Calibri"

SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def box(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    return tb, tf


def set_run(r, text, size, font=SANS, color=NAVY, bold=False, italic=False):
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic


def headline(slide, text, size=30, top=0.55, left=0.8, width=11.7):
    tb, tf = box(slide, left, top, width, 1.4)
    p = tf.paragraphs[0]
    set_run(p.add_run(), text, size, font=SERIF, color=NAVY, bold=True)
    return tb


def kicker(slide, text, top=0.28, left=0.82):
    tb, tf = box(slide, left, top, 8, 0.3)
    p = tf.paragraphs[0]
    r = p.add_run()
    set_run(r, text.upper(), 11, font=SANS, color=LIGHT, bold=True)
    return tb


def rule(slide, top=1.85, left=0.85, width=3.2):
    ln = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Pt(2.5))
    ln.fill.solid()
    ln.fill.fore_color.rgb = NAVY
    ln.line.fill.background()
    ln.shadow.inherit = False
    return ln


def big_number(slide, number, label, left=0.85, top=2.25, width=4.4, num_size=66):
    tb, tf = box(slide, left, top, width, 2.6)
    p = tf.paragraphs[0]
    set_run(p.add_run(), number, num_size, font=SERIF, color=NAVY, bold=True)
    p2 = tf.add_paragraph()
    p2.space_before = Pt(6)
    set_run(p2.add_run(), label, 14, font=SANS, color=GRAY)
    return tb


def body(slide, items, left=5.8, top=2.25, width=6.7, size=15, gap=10):
    """items: list of (lead, rest) tuples or plain strings."""
    tb, tf = box(slide, left, top, width, 4.6)
    first = True
    for it in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap)
        if isinstance(it, tuple):
            lead, rest = it
            set_run(p.add_run(), lead, size, font=SANS, color=NAVY, bold=True)
            set_run(p.add_run(), " " + rest, size, font=SANS, color=NAVY)
        else:
            set_run(p.add_run(), it, size, font=SANS, color=NAVY)
    return tb


def footnote(slide, text, top=6.75, left=0.85, width=10.8):
    tb, tf = box(slide, left, top, width, 0.55)
    p = tf.paragraphs[0]
    set_run(p.add_run(), text, 9.5, font=SANS, color=GRAY, italic=True)
    return tb


def page_number(slide, n):
    tb, tf = box(slide, 12.55, 7.02, 0.6, 0.35)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    set_run(p.add_run(), str(n), 11, font=SANS, color=LIGHT)
    return tb


HEADLINES = []  # collected for validation


# ---------------------------------------------------------------- slide 1: title
s = add_slide()
bg = s.shapes.add_shape(1, 0, 0, SW, SH)
bg.fill.solid()
bg.fill.fore_color.rgb = NAVY
bg.line.fill.background()
bg.shadow.inherit = False

tb, tf = box(s, 1.0, 2.15, 11.3, 2.6)
p = tf.paragraphs[0]
set_run(p.add_run(), "The Stopwatch and the Shell", 48, font=SERIF, color=WHITE, bold=True)
p2 = tf.add_paragraph()
p2.space_before = Pt(10)
set_run(p2.add_run(), "What actually survives to 2050 at Dulles", 26, font=SERIF, color=RGBColor(0xC9, 0xD6, 0xE8))

tb, tf = box(s, 1.0, 5.6, 11.3, 1.2)
p = tf.paragraphs[0]
set_run(p.add_run(), "A Transform Airports AI Council report on the $22 billion Dulles rebuild",
        15, font=SANS, color=RGBColor(0xC9, 0xD6, 0xE8))
p2 = tf.add_paragraph()
p2.space_before = Pt(6)
set_run(p2.add_run(), "Metropolitan Washington Airports Authority  |  July 2026", 13,
        font=SANS, color=RGBColor(0x8A, 0x9A, 0xB5))
HEADLINES.append("The Stopwatch and the Shell")

# ------------------------------------------------------- slide 2: thesis
s = add_slide()
kicker(s, "The thesis")
h = "The concept died in twelve years. The shell is protected at sixty-four."
headline(s, h, size=30)
rule(s)
big_number(s, "12 yrs / 64 yrs", "Life of Saarinen's mobile-lounge processing concept vs. age of his\nNational Register headhouse — the split is written into Dulles's own concrete.",
           left=0.85, top=2.3, width=5.0, num_size=54)
body(s, [
    ("The danger is not a missed forecast.", "It is pouring fifty-year concrete around any passenger-processing sequence — today's or tomorrow's — converting a 5-to-15-year decision into a fifty-year liability."),
    ("Dulles has already paid twice", "to rip out a processing model it hard-wired into the building: the ~$1.4–1.5B AeroTrain (2010) and the $3.75B extension that retires the mobile lounges for good."),
    ("The move that survives to 2050 is a refusal to predict:", "lock the durable geometry hard and oversized, keep the processing layer soft and richly provisioned, treat every sequence commitment as reversible."),
], left=6.3, top=2.25, width=6.2)
page_number(s, 2)
HEADLINES.append(h)

# ------------------------------------------- slide 3: cheapest asset
s = add_slide()
kicker(s, "The argument  ·  1 of 5")
h = "Dulles is spending its cheapest asset to buy the most expensive kind of certainty."
headline(s, h, size=28)
rule(s)
big_number(s, "$9.56", "FY2025 cost per enplanement — among the three or four cheapest US large\nhubs, in a field from $3.94 (Atlanta) to $30.16 (LAX).*",
           left=0.85, top=2.3, width=4.6)
body(s, [
    ("$21.8B in new bonds", "— four to five times MWAA's entire existing $4.7–4.9B aviation-enterprise debt — issued against that advantage."),
    ("No contractual cap:", "under MWAA's residual rate structure, incremental debt service passes mechanically to the airlines."),
    ("The $22B headline is a political anchor, not an engineering estimate.", "The same DOT RFI drew responses from $14.4B (Ferrovial) to $35–50B (Phoenix/Ironbridge); read the figure at ±50 percent,** and the program's own materials exclude \"any potential delays or cost overruns.\""),
    ("Any future-proofing move that permanently raises CPE is not future-proofing.", "It is spending the airport's core commercial asset."),
], left=5.8, top=2.2, width=6.7, size=14, gap=9)
footnote(s, "* CPE figures compiled from MWAA filings by a secondary aggregator (DWU), not read from the audited ACFR — verify before board use.   ** ±50% is analyst inference from the RFI spread, not a sourced figure.")
page_number(s, 3)
HEADLINES.append(h)

# ------------------------------------------- slide 4: eroding advantage
s = add_slide()
kicker(s, "The argument  ·  2 of 5")
h = "The cheap-hub advantage is eroding before the first program bond is issued."
headline(s, h, size=28)
rule(s)
big_number(s, "$9.56 → $12.77", "Signatory CPE, FY2025 to first-half 2026 — on the ~$5.5B of debt MWAA\nalready plans through 2028, before the program adds a dollar.*",
           left=0.85, top=2.3, width=5.4, num_size=54)
body(s, [
    ("Debt per O&D enplanement heads toward $400,", "from $223 in 2024; coverage moves toward 1.3x.*"),
    ("The program does not start the clock on Dulles's cost problem.", "It accelerates one already running."),
    ("The self-help money is frozen:", "the PFC cap has sat at $4.50 since 2000 — roughly $2.45 in today's dollars — and BIL terminal grants sunset September 30, 2026, before the first concourse is poured."),
], left=6.6, top=2.3, width=5.9, size=14.5)
footnote(s, "* Same secondary-aggregator source (DWU) as the CPE series — verify $12.77, 1.3x coverage, and $400/$223 debt-per-enplanement against MWAA financials before board use.")
page_number(s, 4)
HEADLINES.append(h)

# ------------------------------------------- slide 5: failure taxonomy
s = add_slide()
kicker(s, "The argument  ·  3 of 5")
h = "Mega-terminals break at the novel system on the critical path — not the poured structure."
headline(s, h, size=27)
rule(s)
body(s, [
    ("Berlin Brandenburg:", "nine years late at 2.3–2.6x budget — detonated by a bespoke smoke-extraction system, not the structure."),
    ("Denver automated baggage:", "$193M contract past $400M, ~$560M over in total, never ran at scale, decommissioned 2005."),
    ("Heathrow T5:", "structure on time and on budget; the baggage software filter lost some 42,000 bags in the first ten days."),
    ("LAX people mover:", "+76 percent and three years late."),
    ("Kansas City (2023):", "forty gates, $1.5B, on time and on budget — no bespoke system on its critical path."),
], left=0.85, top=2.15, width=6.6, size=14.5, gap=8)
tb, tf = box(s, 7.9, 2.15, 4.6, 4.2)
p = tf.paragraphs[0]
set_run(p.add_run(), "The one structural exception proves the rule:", 14.5, font=SANS, color=NAVY, bold=True)
p2 = tf.add_paragraph()
p2.space_before = Pt(4)
set_run(p2.add_run(), "Denver's Great Hall blew up on weak 1990s concrete under the existing terminal — an unknown-existing-conditions risk that lives in Dulles's renovation packages for its 1960s–1980s concourses, not the new-build ones.", 13.5, font=SANS, color=NAVY)
p3 = tf.add_paragraph()
p3.space_before = Pt(14)
set_run(p3.add_run(), "Spending instruction: overbuild the durable, unretrofittable layer. Keep security regime, sensing, and fit-out deliberately soft — those are the systems that dated fastest everywhere else.", 14, font=SERIF, color=NAVY, bold=True, italic=True)
page_number(s, 5)
HEADLINES.append(h)

# ------------------------------------------- slide 6: coordination
s = add_slide()
kicker(s, "The argument  ·  4 of 5")
h = "For a decade, coordination has beaten concrete — and Dulles has not exhausted it."
headline(s, h, size=28)
rule(s)
big_number(s, "+20%", "Memphis capacity gain from FAA wake-turbulence recategorization —\nup to 22 more arrivals an hour. No runway, no taxiway. A rule.",
           left=0.85, top=2.3, width=4.6)
body(s, [
    ("A-CDM cut European taxi-out times 10–20 percent", "with a data protocol, not a taxiway."),
    ("Legacy hubs carry an estimated 30–50 percent of latent gate capacity,*", "recoverable within a banked operation through common-use gating and turn discipline."),
    ("Dulles: 139 gates serving 29M passengers", "against a master-plan ambition toward 90M. The instinct is to multiply gates. The utilization data says look harder first."),
    ("A $22B program sized to today's uncoordinated peak", "is funding coordination debt in structural steel. Pull the coordination levers first; build the right, smaller amount of the expensive thing."),
], left=5.8, top=2.2, width=6.7, size=14, gap=9)
footnote(s, "* Circa-2012 utilization study, compiled before Delta left Memphis in 2013 — treat the 30–50% range as directional pending a current source.")
page_number(s, 6)
HEADLINES.append(h)

# ------------------------------------------- slide 7: durable vs volatile
s = add_slide()
kicker(s, "The argument  ·  5 of 5")
h = "The durable commitments are geometric; the volatile ones are procedural."
headline(s, h, size=28)
rule(s)
# two-column table effect
col1 = s.shapes.add_shape(1, Inches(0.85), Inches(2.15), Inches(5.7), Inches(3.9))
col1.fill.solid(); col1.fill.fore_color.rgb = NAVY; col1.line.fill.background(); col1.shadow.inherit = False
tf = col1.text_frame; tf.word_wrap = True
tf.margin_left = Inches(0.3); tf.margin_right = Inches(0.3); tf.margin_top = Inches(0.25)
p = tf.paragraphs[0]
set_run(p.add_run(), "50-YEAR PHYSICS — LOCK HARD, OVERSIZED", 15, font=SERIF, color=WHITE, bold=True)
for t in ["Walking distance and circulation geometry",
          "Apron and taxilane clearance",
          "Baggage-system topology",
          "Curb and roadway geometry",
          "Column grid, floor-to-floor, utility and vertical spines"]:
    pp = tf.add_paragraph(); pp.space_before = Pt(9)
    set_run(pp.add_run(), "•  " + t, 14, font=SANS, color=WHITE)
col2 = s.shapes.add_shape(1, Inches(6.85), Inches(2.15), Inches(5.7), Inches(3.9))
col2.fill.solid(); col2.fill.fore_color.rgb = RGBColor(0xE9, 0xEE, 0xF5); col2.line.fill.background(); col2.shadow.inherit = False
tf = col2.text_frame; tf.word_wrap = True
tf.margin_left = Inches(0.3); tf.margin_right = Inches(0.3); tf.margin_top = Inches(0.25)
p = tf.paragraphs[0]
set_run(p.add_run(), "5-TO-15-YEAR LAYERS — LEAVE SOFT", 15, font=SERIF, color=NAVY, bold=True)
for t in ["Checkpoint lane count and machine mix",
          "Gate assignment and turn discipline",
          "Check-in fit-out — replaced 3–5x inside the building's life*",
          "Every operational-intelligence product (A-CDM, digital twin, biometrics)",
          "Security regime and sensing"]:
    pp = tf.add_paragraph(); pp.space_before = Pt(9)
    set_run(pp.add_run(), "•  " + t, 14, font=SANS, color=NAVY)
tb, tf = box(s, 0.85, 6.2, 11.7, 0.5)
p = tf.paragraphs[0]
set_run(p.add_run(), "When these constraints are wrong, no software rescues them: Istanbul locked 20–30-minute walks on day one; Kansas City's curb backed up to the highway.",
        13, font=SANS, color=NAVY, italic=True)
footnote(s, "* 3–5x fit-out replacement cadence is a chief-engineer estimate, no external source.", top=6.82)
page_number(s, 7)
HEADLINES.append(h)

# ------------------------------------------- slide 8: counter-case
s = add_slide()
kicker(s, "The counter-case, honestly presented")
h = "The serious opposing view: build the conventional thing superbly and pocket the difference."
headline(s, h, size=26)
rule(s)
body(s, [
    ("Screening is consolidating, not distributing.", "TSA's 2025 solicitation seeks fewer officers per passenger; self-service screening is one PreCheck-only prototype at Las Vegas with no deployment timeline; CT machines push the central checkpoint to grow first."),
    ("Off-airport check-in has had thirty years to scale and never did.", "Hong Kong and Vienna have offered it since the 1990s; both remain optional, partial, some facilities suspended."),
    ("Flexibility is an option you pay for and rarely exercise.", "The real-options literature's own headline case (Zurich) found the flexible plan worth roughly 5 percent more than the best conventional one — under favorable assumptions."),
    ("Execution, not architecture, is what failed at BER, T5, and Denver", "— and LaGuardia Terminal B, fully centralized, delivered on time and on budget and was named best new terminal in the world."),
    ("The real risk is a departed carrier, not a dated design.", "United carries ~68 percent of Dulles traffic. Pittsburgh built to US Airways' spec and was dehubbed at a CPE of $6."),
], left=0.85, top=2.15, width=11.7, size=14, gap=8)
page_number(s, 8)
HEADLINES.append(h)

# ------------------------------------------- slide 9: why it falls short
s = add_slide()
kicker(s, "Why the counter-case is insufficient")
h = "Its best point — dehubbing — is the strongest argument for convertible geometry."
headline(s, h, size=27)
rule(s)
big_number(s, "40M vs 29M", "Handling capacity Dulles built in 1997 vs. actual traffic twenty-eight years\nlater. \"Demand fills whatever we build\" is a dangerous prior at this airport.",
           left=0.85, top=2.3, width=5.2, num_size=50)
body(s, [
    ("Concede screening and off-airport processing", "— the thesis never asked to build for either. Provision; do not pour concrete. And do not oversize the central checkpoint: CT and self-service raise passengers per square foot."),
    ("Daxing is the thesis executed, not its refutation:", "it committed hard to geometry — farthest gate under 8 minutes — and hard-wired no check-in typology or screening regime."),
    ("The flexibility-cost objection holds against fit-out softness and fails against deep-service provisioning:", "conduit, power density, risers, floor loading are cheap in new slab, ruinous in an occupied terminal."),
    ("A single-tenant concourse strands when the tenant's math turns (Pittsburgh).", "Convertible, MARS-capable, re-lettable structure is a recoverable asset. Cincinnati fell 74 percent and demolished a concourse — there was no tenant to convert to."),
], left=6.4, top=2.2, width=6.1, size=13.5, gap=8)
page_number(s, 9)
HEADLINES.append(h)

# ------------------------------------------- slide 10: design FOR now
s = add_slide()
kicker(s, "Recommendations  ·  1 of 3")
h = "Design FOR now: lock the geometry operations can never fix."
headline(s, h, size=28)
rule(s)
body(s, [
    ("1.  CBP-driven international geometry, binding at the April 2027 construction start:", "one-way sterile corridors, arriving-passenger separation, the FIS baggage hall. Build FIS geometry for the 2050 ceiling; phase FIS fit-out and gate count against CBP staffing reality — officers are appropriated annually."),
    ("2.  The central checkpoint's location and secure boundary,", "traced to the TSA Checkpoint Requirements and Planning Guide (May 2025) — the governing document. Leave lane count and machine mix soft."),
    ("3.  Apron and taxilane geometry, vertical-circulation and utility spines,", "generous floor-to-floor heights, and wide column grids — 30–50 m clear spans* buying column-free, re-partitionable floor plates."),
    ("This is the widebody international gateway function", "— anchored by United's immunized Atlantic joint venture and DCA's perimeter and slot cage — the least-reversible part of United's commitment and the part with the longest half-life."),
], left=0.85, top=2.15, width=11.7, size=14.5, gap=10)
footnote(s, "* 30–50 m clear-span families are a chief-engineer standard, no external source.")
page_number(s, 10)
HEADLINES.append(h)

# ------------------------------------------- slide 11: accommodate later / bet against
s = add_slide()
kicker(s, "Recommendations  ·  2 of 3")
h = "Accommodate LATER; bet AGAINST the sequence."
headline(s, h, size=28)
rule(s)
tb, tf = box(s, 0.85, 2.15, 5.9, 4.4)
p = tf.paragraphs[0]
set_run(p.add_run(), "ACCOMMODATE LATER (ranked by retrofit-cost asymmetry)", 13.5, font=SERIF, color=NAVY, bold=True)
for lead, rest in [
    ("1. The substrate:", " apron-edge power, transformer vaults, conduit density, fast-charge power density, comms risers, an open operational-data architecture. Plausibly single-digit percent of the program* — the last scope to cut, not the first."),
    ("2. MARS / swing-stand clearances:", " one position serves a widebody or two narrowbodies; the clearance is fifty-year physics, the daily decision is soft."),
    ("3. Checkpoint structural bays", " that can grow with CT, then thin with self-service."),
    ("4. Check-in halls convertible", " to dwell, lounge, and revenue space — and Concourse B's 33 United regional gates built to convert to common use."),
]:
    pp = tf.add_paragraph(); pp.space_before = Pt(7)
    set_run(pp.add_run(), lead, 13, font=SANS, color=NAVY, bold=True)
    set_run(pp.add_run(), rest, 13, font=SANS, color=NAVY)
tb, tf = box(s, 7.1, 2.15, 5.4, 4.4)
p = tf.paragraphs[0]
set_run(p.add_run(), "BET AGAINST, EXPLICITLY", 13.5, font=SERIF, color=NAVY, bold=True)
for t in [
    "Sizing the building for today's uncoordinated peak before the coordination levers are pulled",
    "Distributed or off-airport screening as a built-form organizing principle — TSA has scheduled nothing",
    "eVTOL-at-scale apron geometry as a load-bearing assumption",
    "Any single-vendor, first-of-a-kind life-safety or baggage system with no fallback — the BER and Denver failure mode",
    "A permanent, hard-walled row of check-in counters — the one piece of processing geometry already known to be dying",
]:
    pp = tf.add_paragraph(); pp.space_before = Pt(8)
    set_run(pp.add_run(), "•  " + t, 13, font=SANS, color=NAVY)
footnote(s, "* Single-digit-percent provisioning cost is an analyst estimate — program team should cost precisely.")
page_number(s, 11)
HEADLINES.append(h)

# ------------------------------------------- slide 12: above the architecture
s = add_slide()
kicker(s, "Recommendations  ·  3 of 3")
h = "Two decisions sit above the architecture — and the design cannot solve either."
headline(s, h, size=27)
rule(s)
big_number(s, "$26.47", "CPE at Dulles's 2013 peak — the last time IAD sat near the cost basis this\nprogram would recreate, United weighed leaving. Newark's physical\nconstraint saved the hub, not anything MWAA built.*",
           left=0.85, top=2.3, width=5.2, num_size=60)
body(s, [
    ("The financial decision:", "condition the connecting-hub scope on a matched-duration United commitment — a use-and-lease term materially beyond 2039 with majority-in-interest co-underwriting — and stress-test bond coverage against a 30 percent United pull-down."),
    ("The political decision:", "the program is federally sponsored and federally timed; the sponsor's mandate ends with the administration, and the 2034 date rests on a NEPA posture already being probed in the D.C. Circuit. Design the durable moves to be defensible if the federal push evaporates."),
    ("Never plan around a favorable act of Congress.", "The perimeter rule took 58 years to move — then moved against MWAA over the unanimous objection of both home-state delegations."),
], left=6.4, top=2.25, width=6.1, size=14, gap=9)
footnote(s, "* Historical CPE from the same secondary aggregator (DWU) — verify against the audited ACFR before board use.")
page_number(s, 12)
HEADLINES.append(h)

# ------------------------------------------- slide 13: close
s = add_slide()
bg = s.shapes.add_shape(1, 0, 0, SW, SH)
bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background(); bg.shadow.inherit = False
tb, tf = box(s, 1.0, 2.3, 11.3, 3.0)
p = tf.paragraphs[0]
set_run(p.add_run(),
        "Draw the boundary between what moves and what pours.",
        38, font=SERIF, color=WHITE, bold=True)
p2 = tf.add_paragraph()
p2.space_before = Pt(18)
set_run(p2.add_run(),
        "The sequence will change again. The question in the program meeting is which lines on the drawing Dulles will still be able to move when it does — and which it is about to pour into the ground for fifty years.",
        18, font=SANS, color=RGBColor(0xC9, 0xD6, 0xE8))
p3 = tf.add_paragraph()
p3.space_before = Pt(18)
set_run(p3.add_run(), "Everything else is fit-out.", 22, font=SERIF, color=WHITE, italic=True, bold=True)
tb, tf = box(s, 12.55, 7.02, 0.6, 0.35)
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
set_run(p.add_run(), "13", 11, font=SANS, color=RGBColor(0x8A, 0x9A, 0xB5))
HEADLINES.append("Draw the boundary between what moves and what pours.")

OUT = "/Users/christiankessleriv/Repos/ai-council-mwaa/outputs/stage4/future-proofing-the-dulles-terminal-rebuild.pptx"
prs.save(OUT)
print(f"Saved {OUT} with {len(prs.slides.slides if hasattr(prs.slides,'slides') else prs.slides._sldIdLst)} slides")

# ---- validation pass ----
from pptx import Presentation as P2
check = P2(OUT)
n = len(check.slides)
assert n == 13, f"expected 13 slides, got {n}"
problems = []
for i, (slide, expect) in enumerate(zip(check.slides, HEADLINES), 1):
    texts = " || ".join(sh.text_frame.text for sh in slide.shapes if sh.has_text_frame)
    if expect not in texts:
        problems.append((i, expect))
if problems:
    for i, e in problems:
        print(f"MISSING headline on slide {i}: {e}")
    raise SystemExit(1)
print(f"VALIDATED: {n} slides, every slide carries its headline.")
