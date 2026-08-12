#!/usr/bin/env python
"""Build the technical read-ahead deck for
'Quiet-by-Design Standards for MWAA Terminals'
(reader-facing title: a communications-integrity standard for MWAA terminals).

18 slides, 16:9. Georgia display / Calibri body. Brand tokens from
assets/brand/design-system.json. Slide contract, headlines, exhibits, and
source notes follow outputs/stage4/visual-brief.json exactly.
"""
import json
import math
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_PATTERN, MSO_LINE_DASH_STYLE

ROOT = Path("/Users/christiankessleriv/Repos/ai-council-mwaa")
BRIEF = json.loads((ROOT / "outputs/stage4/visual-brief.json").read_text())
OUT = ROOT / "outputs/stage4/quiet-by-design-standards-for-mwaa-terminals-2.pptx"

# ------------------------------------------------------------------ brand
NAVY = RGBColor(0x0B, 0x2D, 0x4D)      # runway navy
BLUE = RGBColor(0x2E, 0x84, 0xA5)      # terminal blue
GOLD = RGBColor(0xD4, 0xA2, 0x4C)      # guidance gold
SLATE = RGBColor(0x41, 0x56, 0x69)     # operations slate
FOG = RGBColor(0xED, 0xF3, 0xF6)       # apron fog
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x17, 0x23, 0x2D)
RED = RGBColor(0xA6, 0x41, 0x3A)       # alert red
GOLD_TINT = RGBColor(0xF6, 0xEA, 0xD3)
LIGHT_ON_NAVY = RGBColor(0xC9, 0xD9, 0xE6)
MUTED_ON_NAVY = RGBColor(0x9F, 0xB6, 0xC9)
# stepped tints of terminal blue for the STI bands (light -> saturated)
BAND_TINTS = [
    RGBColor(0xE4, 0xEF, 0xF4),
    RGBColor(0xC6, 0xDE, 0xE9),
    RGBColor(0x9E, 0xC7, 0xD9),
    RGBColor(0x6C, 0xA9, 0xC3),
    RGBColor(0x2E, 0x84, 0xA5),
]

SERIF = "Georgia"
SANS = "Calibri"
SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]

HEADLINES = {int(s["slide_number"]): s["headline"] for s in BRIEF["slides"]}
NOTES = {int(s["slide_number"]): s.get("speaker_note", "") for s in BRIEF["slides"]}
SIG_SLIDE = int(BRIEF["signature_visual"]["slide_number"])


def add_slide():
    return prs.slides.add_slide(BLANK)


def set_alt(shape, desc):
    for attr in ("nvSpPr", "nvGrpSpPr", "nvGraphicFramePr", "nvPicPr"):
        el = getattr(shape._element, attr, None)
        if el is not None:
            el.cNvPr.set("descr", desc)
            return


def txbox(shapes, l, t, w, h, name=None):
    tb = shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    if name:
        tb.name = name
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tb, tf


def run(p, text, size, font=SANS, color=INK, bold=False, italic=False):
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    return r


def para(tf, first=False, align=None, before=None, after=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    if align is not None:
        p.alignment = align
    if before is not None:
        p.space_before = Pt(before)
    if after is not None:
        p.space_after = Pt(after)
    return p


def rect(shapes, l, t, w, h, fill=None, line_color=None, line_w=None,
         dash=None, shape=MSO_SHAPE.RECTANGLE, name=None):
    sp = shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    if name:
        sp.name = name
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line_color is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line_color
        sp.line.width = Pt(line_w or 1.0)
        if dash:
            sp.line.dash_style = dash
    return sp


def panel_text(sp, pad=0.18):
    tf = sp.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(pad)
    tf.margin_top = tf.margin_bottom = Inches(0.12)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    return tf


def headline(slide, n, size=35, top=0.40, height=1.32, color=NAVY,
             left=0.55, width=12.23):
    tb, tf = txbox(slide.shapes, left, top, width, height,
                   name=f"s{n:02d}-title")
    p = para(tf, first=True)
    run(p, HEADLINES[n], size, font=SERIF, color=color, bold=True)
    return tb


def gold_rule(slide, top=1.82, left=0.57, width=2.4):
    rect(slide.shapes, left, top, width, 0.05, fill=GOLD)


def source_note(slide, n, text, color=SLATE):
    tb, tf = txbox(slide.shapes, 0.55, 6.98, 11.5, 0.44,
                   name=f"s{n:02d}-SourceNote")
    p = para(tf, first=True)
    run(p, text, 10, color=color, italic=True)
    return tb


def page_num(slide, n, color=SLATE):
    tb, tf = txbox(slide.shapes, 12.55, 7.00, 0.5, 0.35)
    p = para(tf, first=True, align=PP_ALIGN.RIGHT)
    run(p, str(n), 11, color=color)


def notes(slide, n):
    text = NOTES.get(n, "")
    if text:
        slide.notes_slide.notes_text_frame.text = text


def caption(slide, text, top, size=16, left=0.55, width=12.23, height=0.7,
            italic=False, bold=False, color=NAVY):
    tb, tf = txbox(slide.shapes, left, top, width, height)
    p = para(tf, first=True)
    run(p, text, size, color=color, italic=italic, bold=bold)
    return tb


def gtext(g, l, t, w, h, entries, align=None):
    """Textbox inside a group. entries: list of (text, size, kw-dict) paragraphs
    where each entry may itself be a list of run tuples."""
    tb, tf = txbox(g.shapes, l, t, w, h)
    first = True
    for entry in entries:
        p = para(tf, first=first)
        first = False
        if align is not None:
            p.alignment = align
        runs = entry if isinstance(entry, list) else [entry]
        for text, size, kw in runs:
            run(p, text, size, **kw)
    return tb


def bullets(tf, items, size, color, first_used=False, gap=6, lead_bold=True):
    for lead, rest in items:
        p = para(tf, first=not first_used, before=gap)
        first_used = True
        if lead:
            run(p, lead, size, color=color, bold=lead_bold)
        if rest:
            run(p, rest, size, color=color)


# =================================================================== slide 1
s = add_slide()
rect(s.shapes, 0, 0, 13.333, 7.5, fill=NAVY, name="cover-field")
rect(s.shapes, 1.0, 2.12, 2.2, 0.06, fill=GOLD, name="cover-rule")
headline(s, 1, size=50, top=2.38, height=2.05, color=WHITE, left=1.0,
         width=11.4)
tb, tf = txbox(s.shapes, 1.32, 4.55, 10.9, 1.25)
p = para(tf, first=True)
run(p, "A communications-integrity standard for MWAA terminals — argued as "
       "safety engineering, delivered through the Design Manual.",
    24, font=SERIF, color=LIGHT_ON_NAVY)
rect(s.shapes, 1.0, 4.50, 0.07, 1.15, fill=GOLD, name="cover-accent")
tb, tf = txbox(s.shapes, 1.0, 6.12, 11.2, 0.4)
p = para(tf, first=True)
run(p, "Technical read-ahead  ·  Metropolitan Washington Airports Authority"
       "  ·  August 2026", 16, color=MUTED_ON_NAVY)
source_note(s, 1, "Sources: Transform Airports AI Council technical "
            "read-ahead, August 2026. Companion to the verified report "
            "‘The Loud Terminal Is the Unintelligible One.’",
            color=MUTED_ON_NAVY)
notes(s, 1)

# =================================================================== slide 2
s = add_slide()
headline(s, 2)
gold_rule(s)
g = s.shapes.add_group_shape()
g.name = "Exhibit — capital passing through the Design Manual"
X0, PER_B = 3.05, 0.405   # $1B -> inches
rows = [
    ("IAD — committed", 6.99, "$6.99B", NAVY, False, 2.30),
    ("DCA — committed", 2.39, "$2.39B", NAVY, False, 3.12),
    ("Dulles transformation", 20.0, ">$20B", GOLD, True, 3.94),
]
for label, val, vtxt, color, hatched, y in rows:
    gtext(g, 0.62, y + 0.05, 2.32, 0.5,
          [(label, 12.5, dict(color=INK, bold=True))])
    bar = rect(g.shapes, X0, y, PER_B * val, 0.52,
               fill=None if hatched else color)
    if hatched:
        bar.fill.patterned()
        bar.fill.pattern = MSO_PATTERN.LIGHT_DOWNWARD_DIAGONAL
        bar.fill.fore_color.rgb = GOLD
        bar.fill.back_color.rgb = WHITE
        bar.line.color.rgb = GOLD
        bar.line.width = Pt(1.25)
        bar.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    gtext(g, X0 + PER_B * val + 0.12, y + 0.06, 1.1, 0.45,
          [(vtxt, 14, dict(color=NAVY, bold=True))])
gtext(g, X0, 2.94, 8.6, 0.32,
      [("Committed capital construction program — 15-year airline use and "
        "lease agreement effective Jan 1, 2025", 11, dict(color=SLATE))])
gtext(g, X0, 4.58, 8.6, 0.32,
      [("Provisional pending Board action — announced July 29, 2026 "
        "(hatched)", 11, dict(color=SLATE))])
rect(g.shapes, X0 - 0.03, 2.20, 0.03, 2.75, fill=SLATE)
gtext(g, X0 - 0.1, 5.02, 4.0, 0.3,
      [("USD billions", 11, dict(color=SLATE, italic=True))])
set_alt(g, "Horizontal bars: committed capital of $6.99 billion at IAD and "
        "$2.39 billion at DCA under the 2025 use-and-lease agreement, and a "
        "hatched provisional bar above $20 billion for the announced Dulles "
        "transformation.")
caption(s, "Every square foot of that work passes through the MWAA Design "
        "Manual — “a mandatory guide with the force of law on the "
        "airport property.”", 5.72, height=0.85)
source_note(s, 2, "Sources: MWAA use-and-lease press release (2024); July 29,"
            " 2026 joint announcement as reported by Travel Agent Central and"
            " Simple Flying ($20B+ provisional pending Board action); MWAA "
            "Design Manual description, mwaa.com.")
page_num(s, 2)
notes(s, 2)

# =================================================================== slide 3
s = add_slide()
headline(s, 3)
gold_rule(s)
left = rect(s.shapes, 0.55, 2.10, 5.85, 4.35, fill=FOG)
tf = panel_text(left, pad=0.28)
p = para(tf, first=True)
run(p, "THE CONSTRAINT", 16, color=SLATE, bold=True)
p = para(tf, before=10)
run(p, "The NTSB’s final report on the January 29, 2025 midair "
       "collision over the Potomac — 67 fatalities — resets public "
       "tolerance for any policy paraphrasable as “reduce safety "
       "announcements.”", 17, color=INK)
rect(s.shapes, 6.62, 2.35, 0.07, 3.85, fill=GOLD)
right = rect(s.shapes, 6.92, 2.10, 5.86, 4.35, fill=None)
tf = panel_text(right, pad=0.28)
p = para(tf, first=True)
run(p, "THE INVERSION", 16, color=NAVY, bold=True)
p = para(tf, before=10)
run(p, "The standard is defensible only as the opposite of subtraction: a "
       "floor that raises the probability the announcement that matters — "
       "the gate change, the shelter-in-place, the IROPS diversion — "
       "reaches the passenger who needs it.", 17, color=NAVY)
source_note(s, 3, "Sources: NTSB Aircraft Accident Report AIR-26/02 (final, "
            "2026, ntsb.gov); contemporaneous coverage (CNN).")
page_num(s, 3)
notes(s, 3)

# =================================================================== slide 4
s = add_slide()
headline(s, 4)
gold_rule(s)
caption(s, "A terminal fails these floors by being loud, reverberant, and "
        "cross-broadcast — precisely the unspecified default.", 1.98,
        height=0.60, italic=True)
g = s.shapes.add_group_shape()
g.name = "SIGNATURE VISUAL — The Intelligibility Ruler"

RX0, RSPAN = 0.90, 11.50            # STI 0.0 at RX0, 1.0 at RX0+RSPAN


def sti_x(v):
    return RX0 + RSPAN * v


BAND_T, BAND_H = 3.82, 0.80
zones = [
    (0.00, 0.30, "Bad", "<0.30", INK),
    (0.30, 0.45, "Poor", "0.30–0.45", INK),
    (0.45, 0.60, "Fair", "0.45–0.60", INK),
    (0.60, 0.75, "Good", "0.60–0.75", INK),
    (0.75, 1.00, "Excellent", ">0.75", WHITE),
]
for i, (a, b, zname, zval, tcol) in enumerate(zones):
    z = rect(g.shapes, sti_x(a), BAND_T, RSPAN * (b - a), BAND_H,
             fill=BAND_TINTS[i], line_color=WHITE, line_w=0.75)
    tf = panel_text(z, pad=0.06)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = para(tf, first=True, align=PP_ALIGN.CENTER)
    run(p, zname, 13, color=tcol, bold=True)
    p = para(tf, align=PP_ALIGN.CENTER)
    run(p, zval, 11, color=tcol)
for v in (0.0, 0.30, 0.45, 0.60, 0.75, 1.0):
    gtext(g, sti_x(v) - 0.35, BAND_T + BAND_H + 0.06, 0.7, 0.26,
          [(f"{v:.2f}", 11, dict(color=SLATE))], align=PP_ALIGN.CENTER)
gtext(g, RX0 + RSPAN - 3.0, BAND_T + BAND_H + 0.34, 3.0, 0.26,
      [("Speech Transmission Index (0–1)", 11,
        dict(color=SLATE, italic=True))], align=PP_ALIGN.RIGHT)

# NFPA 72 marker pair (code) at 0.45 and 0.50
rect(g.shapes, sti_x(0.45) - 0.015, 2.98, 0.03, BAND_T - 2.98, fill=NAVY)
rect(g.shapes, sti_x(0.50) - 0.015, 2.98, 0.03, BAND_T - 2.98, fill=NAVY)
rect(g.shapes, sti_x(0.45) - 0.015, 2.98, sti_x(0.50) - sti_x(0.45) + 0.03,
     0.03, fill=NAVY)
gtext(g, sti_x(0.45) - 1.4, 2.52, 3.6, 0.34,
      [("NFPA 72 emergency-voice floor (code)", 12.5,
        dict(color=NAVY, bold=True))])
gtext(g, sti_x(0.45) - 2.95, 3.06, 2.85, 0.66,
      [("0.45 per location", 12, dict(color=NAVY, bold=True)),
       (" — at ≥90% of locations per acoustically distinguishable "
        "space", 11, dict(color=SLATE))], align=PP_ALIGN.RIGHT)
gtext(g, sti_x(0.50) + 0.12, 3.06, 2.9, 0.66,
      [("0.50 average", 12, dict(color=NAVY, bold=True)),
       (" — per acoustically distinguishable space", 11,
        dict(color=SLATE))])

# ACRP 175 marker (guidance) at 0.45 — dot on the ruler, legend note at right
rect(g.shapes, sti_x(0.45) - 0.045, 3.70, 0.09, 0.09, fill=BLUE,
     shape=MSO_SHAPE.OVAL)
rect(g.shapes, 6.72, 5.42, 0.11, 0.11, fill=BLUE, shape=MSO_SHAPE.OVAL)
gtext(g, 6.95, 5.28, 5.45, 0.60,
      [[("ACRP 175 pier-style PA minimum (guidance): ", 12,
         dict(color=BLUE, bold=True)),
        ("STI 0.45 with 10 dB(A) signal-to-noise", 12, dict(color=BLUE))]])

# default zone of the unspecified terminal (Bad–Poor)
dz = rect(g.shapes, RX0, 5.28, sti_x(0.45) - RX0, 0.62, fill=FOG,
          line_color=RED, line_w=1.25, dash=MSO_LINE_DASH_STYLE.DASH)
tf = panel_text(dz, pad=0.12)
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = para(tf, first=True)
run(p, "Unspecified reverberant, cross-broadcast terminal defaults here "
       "(Bad–Poor)", 12, color=RED, bold=True)
rect(g.shapes, sti_x(0.22) - 0.015, 5.02, 0.03, 0.26, fill=RED)

gtext(g, RX0, 6.10, RSPAN, 0.30,
      [("Code: NFPA 72 (adopted edition governs)   ·   Guidance: ACRP "
        "Research Report 175 (TRB, 2017)   ·   Measurement method: "
        "IEC 60268-16:2020", 11, dict(color=SLATE, italic=True))])
set_alt(g, "The Intelligibility Ruler: a horizontal Speech Transmission "
        "Index scale from 0 to 1 banded Bad, Poor, Fair, Good, Excellent, "
        "with NFPA 72 code markers at 0.45 per location and 0.50 average, "
        "an ACRP 175 guidance marker at 0.45 with 10 dB(A) signal-to-noise, "
        "and a shaded low-end zone where the unspecified reverberant "
        "terminal defaults.")
source_note(s, 4, "Sources: NFPA 72 (code paywalled; thresholds corroborated"
            " across 10+ records reproducing code language); ACRP Research "
            "Report 175 (TRB, 2017); IEC 60268-16:2020 (standard paywalled; "
            "bands corroborated in secondary records).")
page_num(s, 4)
notes(s, 4)

# =================================================================== slide 5
s = add_slide()
headline(s, 5)
gold_rule(s)
g = s.shapes.add_group_shape()
g.name = "Exhibit — Joint Commission alarm-governance timeline"
rect(g.shapes, 1.05, 3.32, 11.1, 0.035, fill=SLATE)
nodes = [
    (2.05, "The problem, documented",
     "One academic medical center logged 74,535 alarms across eight units "
     "in one week (illustrative magnitude). Clinicians had stopped hearing "
     "them."),
    (6.35, "Jan 1, 2014 — NPSG.06.01.01 phase one",
     "Every accredited hospital must inventory and prioritize its clinical "
     "alarms."),
    (10.45, "Jan 1, 2016 — phase two",
     "Governance policies, accountability, and proof on the record — not "
     "silence."),
]
for cx, head, body in nodes:
    rect(g.shapes, cx - 0.11, 3.22, 0.22, 0.22, fill=NAVY,
         shape=MSO_SHAPE.OVAL)
    gtext(g, cx - 1.55, 2.28, 3.25, 0.80,
          [(head, 12.5, dict(color=NAVY, bold=True))])
    gtext(g, cx - 1.55, 3.72, 3.25, 1.55, [(body, 12, dict(color=INK))])
set_alt(g, "Three-node timeline: alarm fatigue documented at illustrative "
        "magnitude, then Joint Commission NPSG.06.01.01 phase one in 2014 "
        "requiring alarm inventory and prioritization, then phase two in "
        "2016 requiring governance and accountability.")
caption(s, "The pattern maps directly to airports: inventory, rank, govern, "
        "log.", 5.82, height=0.45, bold=True)
source_note(s, 5, "Sources: The Joint Commission, National Patient Safety "
            "Goal NPSG.06.01.01 (phased 2014/2016); Respiratory Therapy "
            "Magazine alarm-fatigue review (per-bed figures vary by study; "
            "magnitude illustrative).")
page_num(s, 5)
notes(s, 5)

# =================================================================== slide 6
s = add_slide()
headline(s, 6)
gold_rule(s)
frame = s.shapes.add_table(5, 4, Inches(0.55), Inches(2.05), Inches(12.23),
                           Inches(4.15))
tbl = frame.table
tbl.first_row = True
tbl.horz_banding = False
for i, w in enumerate((1.85, 4.85, 3.05, 2.48)):
    tbl.columns[i].width = Inches(w)
headers = ["Airport", "Announcement policy (operator-described)",
           "Structural profile vs. IAD/DCA",
           "Published intelligibility or missed-boarding data"]
data = [
    ("London City (LCY)",
     "Silent policy since 2008; no flight or gate announcements since "
     "Aug 18, 2016, except essential and emergency",
     "Small, point-to-point, business-traveler-dominant"),
    ("Helsinki-Vantaa (HEL)",
     "Silent regime since 2015; gate-area and emergency overrides retained",
     "Small transfer airport, high English proficiency"),
    ("Copenhagen (CPH)",
     "Gate and boarding information moved to screens and app",
     "Point-to-point-weighted Northern European hub"),
    ("Amsterdam Schiphol (AMS)",
     "Operator-described “Silent Airport”; information on screens "
     "and airport app",
     "Large transfer hub — the closest structural analogue; rests on screen "
     "density and app-carrying passengers"),
]


def cell_text(cell, text, size, color, bold=False, fill=None,
              align=None, anchor=MSO_ANCHOR.MIDDLE):
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    cell.vertical_anchor = anchor
    cell.margin_left = cell.margin_right = Inches(0.1)
    cell.margin_top = cell.margin_bottom = Inches(0.05)
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    run(p, text, size, color=color, bold=bold)


for c, h in enumerate(headers):
    cell_text(tbl.cell(0, c), h, 12, WHITE, bold=True, fill=NAVY)
for r, (a, b, c) in enumerate(data, 1):
    fill = WHITE if r % 2 else FOG
    cell_text(tbl.cell(r, 0), a, 12, NAVY, bold=True, fill=fill)
    cell_text(tbl.cell(r, 1), b, 12, INK, fill=fill)
    cell_text(tbl.cell(r, 2), c, 12, INK, fill=fill)
merged = tbl.cell(1, 3)
merged.merge(tbl.cell(4, 3))
cell_text(merged, "None identified in this research pass", 13, INK,
          bold=True, fill=GOLD_TINT, align=PP_ALIGN.CENTER)
set_alt(frame, "Comparison table of London City, Helsinki-Vantaa, "
        "Copenhagen, and Schiphol silent-airport policies and structural "
        "profiles; a merged cell notes that none published intelligibility "
        "or missed-boarding measurement.")
caption(s, "Policy analogues, not design precedents.", 6.42, height=0.4,
        italic=True)
source_note(s, 6, "Sources: operator self-descriptions — "
            "londoncityairport.com; Finavia (Helsinki-Vantaa); Copenhagen "
            "Airport passenger-communication policy; Amsterdam Schiphol "
            "Silent Airport pages. No published intelligibility or "
            "missed-boarding measurement identified for any of the four.")
page_num(s, 6)
notes(s, 6)

# =================================================================== slide 7
s = add_slide()
headline(s, 7)
gold_rule(s)
g = s.shapes.add_group_shape()
g.name = "Exhibit — SFO reported before/after paging volumes"
stamp = rect(g.shapes, 0.62, 2.10, 4.1, 0.42, fill=GOLD_TINT,
             line_color=GOLD, line_w=1.0)
tf = panel_text(stamp, pad=0.1)
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = para(tf, first=True)
run(p, "SFO-reported; methodology not published", 12, color=INK, bold=True)

BASE_Y, MAX_H = 5.55, 2.35
charts = [
    ("Individual paging occurrences per day", 492, 261, 1.05),
    ("Cumulative announcement time, minutes per day", 145, 58, 4.45),
]
for title, before, after, x0 in charts:
    gtext(g, x0 - 0.25, 2.75, 3.3, 0.55,
          [(title, 12, dict(color=NAVY, bold=True))])
    scale = MAX_H / before
    for i, (lab, val, color) in enumerate(
            (("before", before, SLATE), ("after", after, BLUE))):
        bx = x0 + i * 1.45
        bh = val * scale
        rect(g.shapes, bx, BASE_Y - bh, 0.95, bh, fill=color)
        gtext(g, bx - 0.15, BASE_Y - bh - 0.30, 1.25, 0.28,
              [(str(val), 13, dict(color=NAVY, bold=True))],
              align=PP_ALIGN.CENTER)
        gtext(g, bx - 0.15, BASE_Y + 0.06, 1.25, 0.26,
              [(lab, 11.5, dict(color=SLATE))], align=PP_ALIGN.CENTER)
rect(g.shapes, 0.85, BASE_Y, 6.7, 0.025, fill=SLATE)
set_alt(g, "Two paired before-and-after bar charts of SFO-reported Quiet "
        "Airport outcomes: individual paging occurrences fell from 492 to "
        "261 per day and cumulative announcement time from 145 to 58 "
        "minutes per day; stamped SFO-reported, methodology not published.")
cw = rect(s.shapes, 8.25, 2.10, 4.53, 4.35, fill=FOG)
tf = panel_text(cw, pad=0.24)
p = para(tf, first=True)
run(p, "THE SAME RECORD", 16, color=SLATE, bold=True)
p = para(tf, before=10)
run(p, "Documented passenger complaints and missed-flight accounts "
       "attributed to the reduced-announcement practice, in contemporaneous "
       "reporting.", 16, color=INK)
p = para(tf, before=10)
run(p, "Announcement reduction without intelligibility engineering and "
       "channel redundancy is subtraction, not design.", 16, color=NAVY,
    bold=True)
source_note(s, 7, "Sources: SFO Quiet Airport outcomes as reported by "
            "International Airport Review (2023) and View from the Wing, "
            "citing SFO; all figures SFO-reported, methodology not "
            "published. Counterweight: Simple Flying, documented "
            "missed-flight and confusion accounts.")
page_num(s, 7)
notes(s, 7)

# =================================================================== slide 8
s = add_slide()
headline(s, 8)
gold_rule(s)
g = s.shapes.add_group_shape()
g.name = "Exhibit — accessibility regulatory stack"
bands = [
    (NAVY, WHITE, "28 CFR 35.160 — ADA Title II",
     "Effective communication with people with disabilities; auxiliary "
     "aids and services.", None),
    (BLUE, WHITE, "2010 ADA Standards §219 / §706",
     "Assistive-listening systems in each assembly area where audible "
     "communication is integral to use.", None),
    (SLATE, WHITE, "14 CFR Part 382",
     "Prompt access to the same gate information for passengers needing "
     "visual or hearing assistance.", None),
    (FOG, INK, "ACRP Research Report 239",
     "Digital-first substitution is partial, not sufficient, for travelers "
     "with disabilities and older adults.", "guidance"),
]
y = 2.10
for fill, tcol, head, body, chip in bands:
    b = rect(g.shapes, 0.62, y, 7.55, 1.02, fill=fill)
    tf = panel_text(b, pad=0.2)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = para(tf, first=True)
    run(p, head + "  ", 13, color=tcol, bold=True)
    run(p, body, 12, color=tcol)
    if chip:
        ch = rect(g.shapes, 7.05, y + 0.08, 1.0, 0.30, fill=GOLD_TINT,
                  line_color=GOLD, line_w=0.75)
        ctf = panel_text(ch, pad=0.04)
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        cp = para(ctf, first=True, align=PP_ALIGN.CENTER)
        run(cp, "guidance", 10.5, color=INK, bold=True)
    y += 1.12
set_alt(g, "Vertical stack of four regulation bands: 28 CFR 35.160 "
        "effective communication, 2010 ADA Standards sections 219 and 706 "
        "assistive listening, 14 CFR Part 382 equal gate information, and "
        "ACRP Research Report 239 guidance that digital substitution is "
        "partial.")
tb, tf = txbox(s.shapes, 8.55, 2.30, 4.2, 4.0)
p = para(tf, first=True)
run(p, "A standard that mandates redundant visual, digital, staffed, and "
       "hearing-loop channels binds MWAA harder than current practice, "
       "not softer.", 17, color=NAVY, bold=True)
p = para(tf, before=12)
run(p, "Reducing broadcast without those redundancies compounds exclusion.",
    16, color=INK)
source_note(s, 8, "Sources: 28 C.F.R. § 35.160 (ada.gov); DOJ 2010 ADA "
            "Standards §§219, 706; 14 C.F.R. Part 382 (ecfr.gov); "
            "ACRP Research Report 239 (TRB, 2023).")
page_num(s, 8)
notes(s, 8)

# =================================================================== slide 9
s = add_slide()
headline(s, 9)
gold_rule(s)
g = s.shapes.add_group_shape()
g.name = "Exhibit — two-tier standard framework"
colA = rect(g.shapes, 0.62, 2.02, 5.95, 3.60, fill=NAVY)
tf = panel_text(colA, pad=0.22)
p = para(tf, first=True)
run(p, "MANDATORY AT BOTH AIRPORTS", 13, font=SERIF, color=WHITE, bold=True)
bullets(tf, [
    ("•  ", "Acoustically-distinguishable-space (ADS) map as a design "
     "deliverable"),
    ("•  ", "Intelligibility floors per NFPA 72 and ACRP 175"),
    ("•  ", "Measurement per IEC 60268-16:2020, commissioned under "
     "representative traffic, with cure period"),
    ("•  ", "Accessibility floors: 28 CFR 35.160, 2010 ADA Standards, "
     "14 CFR Part 382"),
    ("•  ", "Hearing conservation per 29 CFR 1910.95"),
    ("•  ", "Audit-logged announcement governance"),
    ("•  ", "Concessionaire TVs and gate-lounge audio kept in scope"),
], 11.5, WHITE, first_used=True, gap=5, lead_bold=False)
colB = rect(g.shapes, 6.82, 2.02, 5.95, 3.60, fill=FOG)
tf = panel_text(colB, pad=0.22)
p = para(tf, first=True)
run(p, "CONTEXT-DEPENDENT BY TERMINAL", 13, font=SERIF, color=NAVY,
    bold=True)
bullets(tf, [
    ("•  ", "Tier levels for PA zoning and announcement density"),
    ("•  ", "Quiet and low-stimulation space scope"),
    ("•  ", "Enforcement mechanics by lease type"),
], 12.5, INK, first_used=True, gap=8, lead_bold=False)
band = rect(g.shapes, 0.62, 5.78, 12.15, 0.98, fill=GOLD_TINT,
            line_color=GOLD, line_w=1.0)
tf = panel_text(band, pad=0.22)
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = para(tf, first=True)
run(p, "IROPS OVERRIDE — ", 12.5, color=INK, bold=True)
run(p, "during declared irregular operations, the Airport Duty Manager "
       "holds named authority to exceed announcement-density templates; "
       "every deviation is logged and reconciled afterward.", 12.5,
    color=INK)
set_alt(g, "Two-column framework: mandatory floors at both airports and "
        "context-dependent tiers by terminal, over an amber band defining "
        "the Airport Duty Manager IROPS override with logged deviations.")
source_note(s, 9, "Sources: framework per the verified report; floors per "
            "NFPA 72, ACRP 175, IEC 60268-16:2020, 28 CFR 35.160, 2010 ADA "
            "Standards, 14 CFR Part 382, 29 CFR 1910.95; governance "
            "capabilities per ACI-NA (Feb 27, 2026).")
page_num(s, 9)
notes(s, 9)

# ================================================================== slide 10
s = add_slide()
headline(s, 10)
gold_rule(s)
g = s.shapes.add_group_shape()
g.name = "Exhibit — how the standard enters force"
steps = [
    ("1", "Office of Engineering drafts the amendment", False),
    ("2", "Primary-text verification: Design Manual acoustic content; AUL "
     "majority-in-interest", True),
    ("3", "Airline and tenant consultation; formal comment window", False),
    ("4", "Board information item — resolution only on material objection",
     False),
    ("5", "Design Manual publication binds consultants and contractors",
     False),
    ("6", "Commissioning under representative traffic, with cure period",
     False),
]
BOXW, GAP, Y, H = 1.86, 0.21, 2.30, 2.15
for i, (num, text, contingent) in enumerate(steps):
    x = 0.62 + i * (BOXW + GAP)
    if contingent:
        b = rect(g.shapes, x, Y, BOXW, H, fill=GOLD_TINT, line_color=GOLD,
                 line_w=1.25, dash=MSO_LINE_DASH_STYLE.DASH)
        tcol = INK
    else:
        b = rect(g.shapes, x, Y, BOXW, H, fill=FOG, line_color=SLATE,
                 line_w=0.75)
        tcol = INK
    tf = panel_text(b, pad=0.12)
    p = para(tf, first=True)
    run(p, num, 14, font=SERIF, color=NAVY, bold=True)
    p = para(tf, before=4)
    run(p, text, 11.5, color=tcol)
    if i < 5:
        rect(g.shapes, x + BOXW + 0.035, Y + H / 2 - 0.07, 0.14, 0.14,
             fill=SLATE, shape=MSO_SHAPE.RIGHT_TRIANGLE)
set_alt(g, "Six-step flow from Office of Engineering drafting through "
        "primary-text verification, airline consultation, Board "
        "information item, Design Manual publication, and commissioning "
        "under representative traffic; the verification step is flagged "
        "as the contingency.")
side = rect(s.shapes, 0.55, 4.90, 12.23, 1.55, fill=None, line_color=GOLD,
            line_w=1.5)
tf = panel_text(side, pad=0.22)
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = para(tf, first=True)
run(p, "First 90 days: ", 16, color=NAVY, bold=True)
run(p, "commission the Concourse E intelligibility and ambient baseline at "
       "the fall 2026 opening. Acceptance floor: STI ≥ 0.45 at "
       "≥ 90% of locations per acoustically distinguishable space, "
       "average ≥ 0.50, measured per IEC 60268-16:2020 under "
       "representative traffic.", 16, color=INK)
source_note(s, 10, "Sources: path per the verified report; Design Manual "
            "authority per MWAA (mwaa.com); Concourse E per MWAA project "
            "disclosures (flydulles.com), verified Aug 11, 2026. AUL "
            "majority-in-interest provisions not verified from primary "
            "text — the path's flagged contingency.")
page_num(s, 10)
notes(s, 10)

# ================================================================== slide 11
s = add_slide()
headline(s, 11)
gold_rule(s)
g = s.shapes.add_group_shape()
g.name = "Exhibit — debt-service coverage trajectory"


def cov_y(v):                     # 1.0x -> 5.85in, 1.8x -> 2.35in
    return 5.85 - (v - 1.0) * (3.50 / 0.8)


for v in (1.0, 1.3, 1.63):
    rect(g.shapes, 2.10, cov_y(v), 9.6, 0.02,
         fill=SLATE if v == 1.0 else FOG if False else RGBColor(0xD5, 0xDF,
                                                                0xE6))
    gtext(g, 0.95, cov_y(v) - 0.14, 1.05, 0.28,
          [(f"{v:.2f}×", 11.5, dict(color=SLATE))],
          align=PP_ALIGN.RIGHT)
x24, x29 = 3.30, 10.50
_dx, _dy = x29 - x24, cov_y(1.30) - cov_y(1.63)
_chord = math.hypot(_dx, _dy)
_cx, _cy = (x24 + x29) / 2.0, (cov_y(1.63) + cov_y(1.30)) / 2.0
conn = g.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(_cx - _chord / 2.0),
                          Inches(_cy), Inches(_chord), Emu(1))
conn.rotation = math.degrees(math.atan2(_dy, _dx))
conn.fill.background()
conn.line.color.rgb = SLATE
conn.line.width = Pt(1.5)
conn.line.dash_style = MSO_LINE_DASH_STYLE.DASH
conn.shadow.inherit = False
rect(g.shapes, x24 - 0.09, cov_y(1.63) - 0.09, 0.18, 0.18, fill=NAVY,
     shape=MSO_SHAPE.OVAL)
rect(g.shapes, x29 - 0.09, cov_y(1.30) - 0.09, 0.18, 0.18, fill=WHITE,
     line_color=GOLD, line_w=2.0, shape=MSO_SHAPE.OVAL)
gtext(g, x24 - 1.0, cov_y(1.63) - 0.62, 3.0, 0.35,
      [("1.63× (2024, reported)", 13, dict(color=NAVY, bold=True))])
gtext(g, x29 - 3.3, cov_y(1.30) + 0.22, 3.6, 0.35,
      [("≈1.3× (projected over five years)", 13,
        dict(color=NAVY, bold=True))], align=PP_ALIGN.RIGHT)
gtext(g, 5.55, 2.35, 4.3, 0.55,
      [("Projected range — only the two reported points are marked; no "
        "intermediate values reported", 11, dict(color=SLATE,
                                                 italic=True))])
ann = rect(g.shapes, 2.55, 4.95, 4.25, 0.55, fill=FOG)
tf = panel_text(ann, pad=0.12)
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = para(tf, first=True)
run(p, "≈$5.5B of new debt phases in, 2025–2028", 12,
    color=INK, bold=True)
gtext(g, 2.9, 6.0, 2.0, 0.28, [("2024", 11.5, dict(color=SLATE))])
gtext(g, 9.7, 6.0, 3.0, 0.28,
      [("+5 years (projection)", 11.5, dict(color=SLATE))])
gtext(g, 0.80, 2.15, 4.6, 0.30,
      [("Net-revenue debt-service coverage (×)", 11.5,
        dict(color=SLATE, italic=True))])
set_alt(g, "Two-point coverage trajectory: net-revenue debt-service "
        "coverage reported at 1.63 times in 2024, projected to about 1.3 "
        "times over five years as roughly 5.5 billion dollars of new debt "
        "phases in; only the two reported points are marked.")
caption(s, "Added specification cost inside programs already in design cuts "
        "against coverage; specification at programming stage does not.",
        6.38, height=0.5)
source_note(s, 11, "Sources: Moody's projection as reported by The Bond "
            "Buyer, ‘MWAA goes to market’ (2025; "
            "subscription-gated). Underlying rating-agency reports are the "
            "primary citation-grade source — confirm before board-facing "
            "publication.")
page_num(s, 11)
notes(s, 11)

# ================================================================== slide 12
s = add_slide()
headline(s, 12)
gold_rule(s)
g = s.shapes.add_group_shape()
g.name = "Exhibit — design-time vs retrofit acoustic treatment ranges"
DX0, PER_D = 3.30, 0.172            # $1/sq ft -> inches


def dx(v):
    return DX0 + PER_D * v


for label, a, b, color, y in (
        ("New construction — acoustic ceiling treatment", 5, 15, NAVY, 2.60),
        ("Retrofit — acoustic panel treatment", 18, 45, SLATE, 3.75)):
    gtext(g, 0.62, y - 0.42, 5.8, 0.35,
          [(label, 12.5, dict(color=NAVY, bold=True))])
    rect(g.shapes, dx(a), y, PER_D * (b - a), 0.55, fill=color)
    gtext(g, dx(a) - 0.75, y + 0.10, 0.65, 0.32,
          [(f"${a}", 12.5, dict(color=INK, bold=True))],
          align=PP_ALIGN.RIGHT)
    gtext(g, dx(b) + 0.10, y + 0.10, 0.85, 0.32,
          [(f"${b}", 12.5, dict(color=INK, bold=True))])
rect(g.shapes, DX0, 4.62, PER_D * 50, 0.02, fill=SLATE)
for v in (0, 10, 20, 30, 40, 50):
    rect(g.shapes, dx(v) - 0.01, 4.62, 0.02, 0.08, fill=SLATE)
    gtext(g, dx(v) - 0.35, 4.74, 0.7, 0.26,
          [(f"${v}", 11, dict(color=SLATE))], align=PP_ALIGN.CENTER)
gtext(g, dx(50) - 2.4, 5.02, 2.4, 0.28,
      [("USD per square foot", 11, dict(color=SLATE, italic=True))],
      align=PP_ALIGN.RIGHT)
face = rect(g.shapes, 0.62, 5.45, 12.15, 0.55, fill=GOLD_TINT,
            line_color=GOLD, line_w=1.0)
tf = panel_text(face, pad=0.14)
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = para(tf, first=True)
run(p, "Non-airport commercial benchmarks, medium confidence — airside "
       "premiums would raise both ranges.", 12, color=INK, bold=True)
set_alt(g, "Two horizontal range bars in dollars per square foot: "
        "new-construction acoustic ceiling treatment 5 to 15 dollars and "
        "retrofit panel treatment 18 to 45 dollars, labeled as non-airport "
        "medium-confidence benchmarks.")
caption(s, "The ranges point one direction: specify at design time.", 6.28,
        height=0.42, bold=True)
source_note(s, 12, "Sources: commercial pricing guides — "
            "designtransitionstudio.com (new-construction acoustic ceiling, "
            "2026 guide); acousticmod.com (retrofit vs. new construction; "
            "page since relocated). Non-airport, directional benchmarks.")
page_num(s, 12)
notes(s, 12)

# ================================================================== slide 13
s = add_slide()
headline(s, 13)
gold_rule(s)
g = s.shapes.add_group_shape()
g.name = "Exhibit — IAD vs DCA anchor-tenant profile"
iad = rect(g.shapes, 0.62, 2.02, 5.95, 3.30, fill=NAVY)
tf = panel_text(iad, pad=0.22)
p = para(tf, first=True)
run(p, "IAD — INTERNATIONAL-CONNECTION HUB", 13, font=SERIF, color=WHITE,
    bold=True)
bullets(tf, [
    ("•  ", "29.01M passengers in 2025; 10.53M international"),
    ("•  ", "Fastest 2025 growth among the 50 largest US airports"),
    ("•  ", "United: ≈67% of passenger traffic; 50 "
     "jetbridge-equipped gates in Concourses C and D"),
], 12.5, WHITE, first_used=True, gap=8, lead_bold=False)
dca = rect(g.shapes, 6.82, 2.02, 5.95, 3.30, fill=SLATE)
tf = panel_text(dca, pad=0.22)
p = para(tf, first=True)
run(p, "DCA — SLOT-CONTROLLED DOMESTIC O&D", 13, font=SERIF, color=WHITE,
    bold=True)
bullets(tf, [
    ("•  ", "American: ≈53% market share (end-2024)"),
    ("•  ", "2025 traffic declined — the January collision, softened "
     "government-related travel, a weeks-long federal shutdown"),
    ("•  ", "Podium PA effectively airline-owned at preferential-use "
     "gates"),
], 12.5, WHITE, first_used=True, gap=8, lead_bold=False)
strip = rect(g.shapes, 0.62, 5.50, 12.15, 1.10, fill=FOG)
tf = panel_text(strip, pad=0.22)
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = para(tf, first=True)
run(p, "Airline briefing reported by Cranky Flier (Aug 4, 2026): ", 12.5,
    color=INK)
run(p, "fully loaded IAD transformation cost per enplanement ≈$90.64 "
       "vs. $12.88 today", 12.5, color=INK, bold=True)
run(p, " — analyst-reported airline positioning, not an MWAA figure.",
    12.5, color=INK)
set_alt(g, "Side-by-side profile of IAD, a growing international hub where "
        "United holds about 67 percent of traffic, and slot-controlled DCA "
        "where American holds about 53 percent; a footer strip carries the "
        "analyst-reported cost-per-enplanement positioning.")
source_note(s, 13, "Sources: Simple Flying (United IAD share and gates, "
            "verified Aug 11, 2026; American DCA share, March 2025) — "
            "directional trade-press figures; IAD traffic per MWAA "
            "reporting via Simple Flying and FFXnow; CPE figures per "
            "Cranky Flier (Aug 4, 2026), attributed as analyst commentary.")
page_num(s, 13)
notes(s, 13)

# ================================================================== slide 14
s = add_slide()
headline(s, 14)
gold_rule(s)
g = s.shapes.add_group_shape()
g.name = "Exhibit — holdroom audio governance boundaries"
outer = rect(g.shapes, 0.62, 2.05, 7.35, 4.55, fill=None, line_color=SLATE,
             line_w=1.25)
gtext(g, 0.82, 2.15, 6.9, 0.30,
      [("HOLDROOM — conceptual, not a floorplan", 11.5,
        dict(color=SLATE, bold=True, italic=True))])
mwaa = rect(g.shapes, 0.90, 2.60, 6.80, 1.55, fill=BAND_TINTS[1])
tf = panel_text(mwaa, pad=0.16)
p = para(tf, first=True)
run(p, "MWAA-controlled overhead PA zones", 12.5, color=NAVY, bold=True)
p = para(tf, before=3)
run(p, "Governed by the Design Manual today; the amendment adds "
       "intelligibility floors, ADS mapping, and audit-logged governance.",
    11.5, color=INK)
pod = rect(g.shapes, 0.90, 4.35, 3.30, 1.95, fill=NAVY)
tf = panel_text(pod, pad=0.16)
p = para(tf, first=True)
run(p, "Airline podium PA (preferential-use gate)", 12.5, color=WHITE,
    bold=True)
p = para(tf, before=3)
run(p, "Amendment written to reach it — contingent on AUL "
       "majority-in-interest verification.", 11.5, color=LIGHT_ON_NAVY)
conc = rect(g.shapes, 4.40, 4.35, 3.30, 1.95, fill=GOLD_TINT,
            line_color=GOLD, line_w=1.0)
tf = panel_text(conc, pad=0.16)
p = para(tf, first=True)
run(p, "Concessionaire TVs and audio", 12.5, color=INK, bold=True)
p = para(tf, before=3)
run(p, "Lease terms; ambient thresholds set in the amendment.", 11.5,
    color=INK)
gov = rect(g.shapes, 8.25, 2.05, 4.52, 2.85, fill=FOG)
tf = panel_text(gov, pad=0.18)
p = para(tf, first=True)
run(p, "AUDITABLE GOVERNANCE THE AMENDMENT CAN SPECIFY", 12, color=NAVY,
    bold=True)
bullets(tf, [
    ("•  ", "Centralized management"),
    ("•  ", "Replay and audit logs of every broadcast"),
    ("•  ", "Elimination of manual-initiation errors"),
    ("•  ", "Templated multi-language delivery"),
], 12, INK, first_used=True, gap=5, lead_bold=False)
enf = rect(g.shapes, 8.25, 5.05, 4.52, 1.55, fill=None, line_color=SLATE,
           line_w=1.0)
tf = panel_text(enf, pad=0.18)
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = para(tf, first=True)
run(p, "Enforcement: ", 12, color=NAVY, bold=True)
run(p, "annual acoustic inspections keyed to lease renewal · cure "
       "notices · lease-management escalation for pattern "
       "noncompliance.", 12, color=INK)
set_alt(g, "Schematic holdroom boundary diagram, labeled conceptual and "
        "not a floorplan, separating MWAA-controlled overhead PA, "
        "airline podium PA at preferential-use gates, and concessionaire "
        "audio, with panels listing auditable governance capabilities and "
        "enforcement mechanics.")
source_note(s, 14, "Sources: airline gate control per trade-press reporting "
            "(Simple Flying); AUL structure per MWAA press release, "
            "majority-in-interest provisions not verified from primary "
            "text; governance platform capabilities per ACI-NA (Feb 27, "
            "2026) — industry advocacy, directional.")
page_num(s, 14)
notes(s, 14)

# ================================================================== slide 15
s = add_slide()
headline(s, 15)
gold_rule(s)
g = s.shapes.add_group_shape()
g.name = "Exhibit — first-cycle implementation timeline"
rect(g.shapes, 0.90, 3.42, 11.55, 0.035, fill=SLATE)
miles = [
    (1.65, "Days 1–30", "Primary-text verification",
     "Design Manual acoustic and PA content; AUL majority-in-interest "
     "provisions (Office of Engineering).", False),
    (4.75, "Fall 2026", "Concourse E baseline",
     "≈435,000 sq ft, 14 gates. Commission the intelligibility and "
     "ambient baseline under representative traffic at opening.", False),
    (7.85, "This capital cycle", "Amendment issued",
     "Board information item; escalation to resolution only on material "
     "airline objection.", False),
    (10.95, "Phase-dependent", "Reference application",
     "DCA Terminal 1 concourse replacement — subject to phase "
     "verification, not verified in this research pass; fallback: earliest "
     "IAD C/D element still in programming.", True),
]
for cx, period, action, detail, contingent in miles:
    if contingent:
        rect(g.shapes, cx - 0.13, 3.31, 0.26, 0.26, fill=GOLD_TINT,
             line_color=GOLD, line_w=1.5, dash=MSO_LINE_DASH_STYLE.DASH,
             shape=MSO_SHAPE.OVAL)
    else:
        rect(g.shapes, cx - 0.13, 3.31, 0.26, 0.26, fill=NAVY,
             shape=MSO_SHAPE.OVAL)
    gtext(g, cx - 1.42, 2.42, 2.85, 0.80,
          [[(period + "\n", 12.5, dict(color=NAVY, bold=True)),
            (action, 12, dict(color=SLATE, bold=True))]])
    card = rect(g.shapes, cx - 1.42, 3.85, 2.85, 2.35,
                fill=GOLD_TINT if contingent else FOG,
                line_color=GOLD if contingent else None,
                line_w=1.25 if contingent else None,
                dash=MSO_LINE_DASH_STYLE.DASH if contingent else None)
    tf = panel_text(card, pad=0.12)
    p = para(tf, first=True)
    run(p, detail, 11.5, color=INK)
set_alt(g, "Four-milestone timeline: days 1 to 30 primary-text "
        "verification, fall 2026 Concourse E baseline commissioning, "
        "amendment issued this capital cycle, and a phase-dependent "
        "reference application at DCA Terminal 1 marked contingent in "
        "amber with dashed borders.")
caption(s, "The Concourse E baseline becomes the number every future MWAA "
        "facility is measured against.", 6.40, height=0.45, bold=True)
source_note(s, 15, "Sources: Concourse E per MWAA project disclosures "
            "(flydulles.com, verified Aug 11, 2026); reference-application "
            "phasing per the verified report, flagged unverified; open "
            "items listed on the final page of this packet.")
page_num(s, 15)
notes(s, 15)

# ================================================================== slide 16
s = add_slide()
headline(s, 16)
gold_rule(s)
g = s.shapes.add_group_shape()
g.name = "Exhibit — commissioning protocol"
protocol = [
    ("1", "Empty-building STI baseline",
     "Measured per IEC 60268-16:2020 — necessary, but not a commissioning "
     "result."),
    ("2", "Representative-traffic verification",
     "Required because the standard's own scope (§7.13, §8.9.3) "
     "does not cover fluctuating background noise — and terminals are the "
     "paradigm fluctuating-noise environment."),
    ("3", "Cure period and re-measurement",
     "Verification-with-cure-period, not an occupancy condition — keeps "
     "acceptance off the occupancy critical path."),
]
for i, (num, head, body) in enumerate(protocol):
    x = 0.62 + i * 4.18
    b = rect(g.shapes, x, 2.15, 3.85, 2.60, fill=FOG if i != 1 else None,
             line_color=BLUE if i == 1 else None,
             line_w=1.5 if i == 1 else None)
    tf = panel_text(b, pad=0.16)
    p = para(tf, first=True)
    run(p, num + "  ", 15, font=SERIF, color=NAVY, bold=True)
    run(p, head, 12.5, color=NAVY, bold=True)
    p = para(tf, before=5)
    run(p, body, 12, color=INK)
    if i < 2:
        rect(g.shapes, x + 3.90, 3.35, 0.14, 0.14, fill=SLATE,
             shape=MSO_SHAPE.RIGHT_TRIANGLE)
strip = rect(g.shapes, 0.62, 5.15, 12.15, 1.30, fill=None, line_color=SLATE,
             line_w=1.0)
tf = panel_text(strip, pad=0.2)
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = para(tf, first=True)
run(p, "Related floors:  ", 12, color=NAVY, bold=True)
run(p, "29 CFR 1910.95 — hearing-conservation program at 85 dB(A) "
       "eight-hour TWA   ·   ANSI/ASA S12.60 — 35 dB(A) "
       "core-learning-space background, the reference benchmark for "
       "MWAA-defined low-stimulation zones only (reference, not terminal "
       "ambient).", 12, color=INK)
set_alt(g, "Three-step commissioning protocol: empty-building STI "
        "baseline, representative-traffic verification required by the "
        "measurement standard's scope limits, then cure period and "
        "re-measurement, with a reference strip of related occupational "
        "and classroom acoustic floors.")
source_note(s, 16, "Sources: IEC 60268-16:2020 scope limitation (standard "
            "paywalled; scope confirmed in records citing §7.13, "
            "§8.9.3); 29 C.F.R. § 1910.95 (OSHA); ANSI/ASA S12.60 "
            "(paywalled; 35 dB(A) figure corroborated).")
page_num(s, 16)
notes(s, 16)

# ================================================================== slide 17
s = add_slide()
headline(s, 17)
gold_rule(s)
g = s.shapes.add_group_shape()
g.name = "Exhibit — existing floors vs amendment marginal content"
lp = rect(g.shapes, 0.62, 2.02, 5.95, 3.35, fill=FOG)
tf = panel_text(lp, pad=0.22)
p = para(tf, first=True)
run(p, "WHAT EXISTING FLOORS REQUIRE", 13, font=SERIF, color=NAVY,
    bold=True)
bullets(tf, [
    ("NFPA 72 — ", "emergency voice intelligibility only; no ADS mapping "
     "as a general deliverable, no routine-PA governance"),
    ("2010 ADA Standards — ", "receiver counts, not intelligibility floors "
     "for daily announcements"),
    ("ACRP 175 — ", "guidance, not code"),
], 12, INK, first_used=True, gap=8)
rp = rect(g.shapes, 6.82, 2.02, 5.95, 3.35, fill=NAVY)
tf = panel_text(rp, pad=0.22)
p = para(tf, first=True)
run(p, "WHAT ONLY THE AMENDMENT ADDS", 13, font=SERIF, color=WHITE,
    bold=True)
bullets(tf, [
    ("•  ", "ADS map at each design milestone"),
    ("•  ", "Representative-traffic commissioning"),
    ("•  ", "Audit-logged governance of routine PA"),
    ("•  ", "Channel-redundancy floors"),
    ("•  ", "IROPS override"),
    ("•  ", "Tenant-audio boundaries"),
], 12, WHITE, first_used=True, gap=5, lead_bold=False)
cav = rect(g.shapes, 0.62, 5.52, 12.15, 1.28, fill=GOLD_TINT,
           line_color=GOLD, line_w=1.0)
tf = panel_text(cav, pad=0.2)
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = para(tf, first=True)
run(p, "The serious objection: ", 12, color=INK, bold=True)
run(p, "premature specification before program-brief maturity drives "
       "change orders, and hard commissioning gates on the critical path "
       "in this coverage window would be a bond-rating conversation. "
       "Answered by cure-period commissioning and milestone — not "
       "occupancy — deliverables.", 12, color=INK)
set_alt(g, "Gap comparison: left panel lists what NFPA 72, the ADA "
        "Standards, and ACRP 175 already require; right panel lists the "
        "amendment's marginal content; an amber band states the premature-"
        "specification objection and its answer.")
source_note(s, 17, "Sources: counter-case and rebuttal per the verified "
            "report; regulatory scope per NFPA 72, 2010 ADA Standards, and "
            "ACRP Research Report 175 as corroborated in the evidence "
            "base.")
page_num(s, 17)
notes(s, 17)

# ================================================================== slide 18
s = add_slide()
headline(s, 18)
gold_rule(s)
frame = s.shapes.add_table(6, 3, Inches(0.55), Inches(2.02), Inches(12.23),
                           Inches(3.95))
tbl = frame.table
tbl.first_row = True
tbl.horz_banding = False
for i, w in enumerate((5.55, 2.85, 3.83)):
    tbl.columns[i].width = Inches(w)
for c, h in enumerate(("Open verification item", "Owner", "Closure route")):
    cell_text(tbl.cell(0, c), h, 12, WHITE, bold=True, fill=NAVY)
checklist = [
    ("Current Design Manual acoustic, PA-zoning, and STI content",
     "Office of Engineering",
     "Read Volume 2 and applicable specification sections (public, "
     "mwaa.com)"),
    ("2025 AUL majority-in-interest thresholds and cost-recovery "
     "categorization", "Office of Engineering with airline relations",
     "Primary agreement text"),
    ("DCA Terminal 1 concourse replacement project phase",
     "Capital program office", "Capital program status reporting"),
    ("SFO Quiet Airport measurement methodology", "Research follow-up",
     "SFO Airport Commission minutes or program report"),
    ("Rating-agency primary reports behind the Bond Buyer coverage "
     "figures", "Finance", "Moody's, Fitch, and S&P MWAA reports"),
]
for r, (item, owner, route) in enumerate(checklist, 1):
    fill = WHITE if r % 2 else FOG
    cell_text(tbl.cell(r, 0), item, 12, INK, bold=True, fill=fill)
    cell_text(tbl.cell(r, 1), owner, 12, INK, fill=fill)
    cell_text(tbl.cell(r, 2), route, 12, INK, fill=fill)
set_alt(frame, "Checklist table of five open verification items with "
        "owners and closure routes, from Design Manual content and "
        "use-and-lease provisions to rating-agency primary reports.")
caption(s, "Every figure in this packet carries a documented source; the "
        "full source appendix is in the companion report. Nothing on this "
        "page is asserted as verified.", 6.18, height=0.7)
source_note(s, 18, "Sources: open items per the report's evidence review. "
            "Closing them is the first 90 days of assigned work.")
page_num(s, 18)
notes(s, 18)

# ------------------------------------------------------------------ save
prs.save(OUT)
print(f"Saved {OUT}")

# ------------------------------------------------------------- validation
import re


def norm(t):
    return " ".join(re.findall(r"[a-z0-9]+", t.casefold()))


check = Presentation(OUT)
assert len(check.slides) == 18, len(check.slides)
sig = []
for n, slide in enumerate(check.slides, 1):
    titles = [sh for sh in slide.shapes
              if getattr(sh, "name", "").endswith("-title")]
    assert titles, f"slide {n}: no title shape"
    got = titles[0].text_frame.text
    want = HEADLINES[n]
    assert norm(got) == norm(want), f"slide {n}: {got!r} != {want!r}"
    for sh in slide.shapes:
        if str(getattr(sh, "name", "")).casefold().startswith(
                "signature visual —".casefold()):
            sig.append(n)
assert sig == [SIG_SLIDE], sig
print("VALIDATED: 18 slides, headlines match brief, signature marker on "
      f"slide {SIG_SLIDE}.")
