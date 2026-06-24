#!/usr/bin/env python
"""Build the companion executive deck for 'From Compliance to Capability'."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- palette & type ----------------------------------------------------------
NAVY   = RGBColor(0x14, 0x3C, 0x6E)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREY   = RGBColor(0x55, 0x5B, 0x66)
LIGHT  = RGBColor(0xE7, 0xEC, 0xF2)
ACCENT = RGBColor(0xB1, 0x8A, 0x3E)  # restrained gold for the one big number
RULE   = RGBColor(0xC8, 0xD2, 0xDF)

SERIF = "Georgia"
SANS  = "Calibri"

OUT = "from-compliance-to-capability-a-strategic-review-of-the-mwaa-contracting-manual.pptx"

# 16:9
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

MARGIN = Inches(0.9)
CW = SW - 2 * MARGIN


def add_slide():
    return prs.slides.add_slide(BLANK)


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def box(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tb, tf


def style(run, *, font=SANS, size=18, color=NAVY, bold=False, italic=False):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic


def para(tf, text, *, font=SANS, size=18, color=NAVY, bold=False, italic=False,
         align=PP_ALIGN.LEFT, space_after=6, space_before=0, leading=1.05, first=False):
    p = tf.paragraphs[0] if first and tf.paragraphs[0].text == "" else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    p.line_spacing = leading
    r = p.add_run()
    r.text = text
    style(r, font=font, size=size, color=color, bold=bold, italic=italic)
    return p


def rect(slide, l, t, w, h, color, line=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(0.75)
    sp.shadow.inherit = False
    return sp


def slide_number(slide, n):
    tb, tf = box(slide, SW - Inches(0.9), SH - Inches(0.55), Inches(0.6), Inches(0.35))
    para(tf, str(n), font=SANS, size=11, color=GREY, align=PP_ALIGN.RIGHT, first=True, space_after=0)


def kicker(slide, text):
    tb, tf = box(slide, MARGIN, Inches(0.55), CW, Inches(0.4))
    p = para(tf, text.upper(), font=SANS, size=12.5, color=ACCENT, bold=True, first=True, space_after=0)
    # letter spacing
    rPr = p.runs[0]._r.get_or_add_rPr()
    rPr.set("spc", "180")


def headline(slide, text, *, top=Inches(1.0), size=30, w=CW, color=NAVY):
    tb, tf = box(slide, MARGIN, top, w, Inches(1.7))
    para(tf, text, font=SERIF, size=size, color=color, bold=False, leading=1.04, first=True, space_after=0)
    return tb


def hrule(slide, top, w=Inches(2.2), color=ACCENT, l=MARGIN):
    rect(slide, l, top, w, Pt(2.4), color)


# =============================================================================
# 1 — TITLE
# =============================================================================
s = add_slide()
bg(s, NAVY)
rect(s, 0, 0, SW, Inches(0.18), ACCENT)
tb, tf = box(s, MARGIN, Inches(1.05), CW, Inches(0.5))
p = para(tf, "METROPOLITAN WASHINGTON AIRPORTS AUTHORITY", font=SANS, size=14, color=LIGHT, bold=True, first=True, space_after=0)
p.runs[0]._r.get_or_add_rPr().set("spc", "240")

tb, tf = box(s, MARGIN, Inches(2.15), CW, Inches(2.6))
para(tf, "From Compliance to Capability", font=SERIF, size=52, color=WHITE, bold=True, first=True, space_after=8, leading=1.0)
para(tf, "A Strategic Review of the MWAA Contracting Manual", font=SERIF, size=26, color=LIGHT, italic=True, space_after=0, leading=1.05)

rect(s, MARGIN, Inches(5.05), Inches(2.2), Pt(2.4), ACCENT)
tb, tf = box(s, MARGIN, Inches(5.3), CW, Inches(1.2))
para(tf, "The binding constraint is not what the manual forbids.", font=SANS, size=18, color=WHITE, first=True, space_after=2)
para(tf, "It is what the manual never measures.", font=SANS, size=18, color=WHITE, space_after=0)

tb, tf = box(s, MARGIN, SH - Inches(0.75), CW, Inches(0.4))
para(tf, "Transform Airports AI Council  ·  Executive Companion Deck  ·  June 2026", font=SANS, size=12, color=RGBColor(0x9F,0xB0,0xC6), first=True, space_after=0)

# =============================================================================
# 2 — THESIS
# =============================================================================
s = add_slide(); bg(s, WHITE)
kicker(s, "The thesis in one sentence")
hrule(s, Inches(1.05))
tb, tf = box(s, MARGIN, Inches(1.9), CW, Inches(3.4))
para(tf, "The manual is not a cage. It is a modern toolbox whose tools are rarely the default and whose clock is never timed.",
     font=SERIF, size=38, color=NAVY, leading=1.12, first=True, space_after=0)
tb, tf = box(s, MARGIN, Inches(5.4), CW, Inches(1.3))
para(tf, "MWAA's binding constraint is not what the manual forbids — it forbids remarkably little. The constraint is what it fails to measure, to presume, and to instrument. The reform that matters is not less control. It is the same accountability, delivered faster and proven cleaner.",
     font=SANS, size=17, color=GREY, leading=1.25, first=True, space_after=0)
slide_number(s, 2)

# =============================================================================
# 3 — THE CLOCK  (big number: 9 months)
# =============================================================================
s = add_slide(); bg(s, WHITE)
kicker(s, "The clock the manual forgot to keep")
headline(s, "A $4.9 million buy took nine months to award — and the manual sets a clock for everyone except itself.",
         top=Inches(1.0), size=27)
hrule(s, Inches(2.55))

# big number block
rect(s, MARGIN, Inches(3.0), Inches(4.4), Inches(3.4), LIGHT)
tb, tf = box(s, MARGIN, Inches(3.35), Inches(4.4), Inches(2.6))
para(tf, "9", font=SERIF, size=120, color=NAVY, bold=True, align=PP_ALIGN.CENTER, first=True, space_after=0, leading=0.9)
para(tf, "MONTHS FROM SOLICITATION TO AWARD", font=SANS, size=13, color=NAVY, bold=True, align=PP_ALIGN.CENTER, space_after=0)
tb, tf = box(s, MARGIN+Inches(0.3), Inches(5.85), Inches(3.8), Inches(0.5))
para(tf, "Issued Aug 25, 2025 · closed Dec 10 · awarded May 22, 2026", font=SANS, size=11.5, color=GREY, align=PP_ALIGN.CENTER, first=True, space_after=0)

# right column points
rx = MARGIN + Inches(4.9)
rw = SW - MARGIN - rx
tb, tf = box(s, rx, Inches(3.0), rw, Inches(3.6))
para(tf, "Every clock that constrains an outside party is instrumented:", font=SANS, size=16, color=NAVY, bold=True, first=True, space_after=8)
for t in ["30-day minimum bid advertisement",
          "15-day sole-source notice, posted for 30",
          "7 days for a protester to file"]:
    para(tf, "—  " + t, font=SANS, size=15.5, color=GREY, space_after=6, leading=1.1)
para(tf, "The one interval the requesting office actually feels — time-to-award — has no target, no commitment, no measurement.",
     font=SANS, size=15.5, color=NAVY, bold=True, space_before=8, space_after=0, leading=1.15)
slide_number(s, 3)

# =============================================================================
# 4 — TOOLKIT MODERN (big: 50-75%)
# =============================================================================
s = add_slide(); bg(s, WHITE)
kicker(s, "Finding 1 · The document is not behind. It is asleep.")
headline(s, "The toolkit is already modern; the deficit is dormancy, not absence.",
         top=Inches(1.0), size=28)
hrule(s, Inches(2.2))

tb, tf = box(s, MARGIN, Inches(2.6), Inches(6.7), Inches(4.2))
para(tf, "Edition 6.0 already authorizes:", font=SANS, size=16, color=NAVY, bold=True, first=True, space_after=8)
for t in ["Construction manager-at-risk",
          "All three flavors of design-build — including progressive design-build, named outright",
          "Multiple Award Construction Contracts",
          "Qualifications-based selection and best-value tradeoff",
          "Pilot programs, unsolicited proposals, stipends",
          "Cooperative purchasing off other governments' contracts"]:
    para(tf, "•  " + t, font=SANS, size=15, color=GREY, space_after=7, leading=1.12)
para(tf, "A 2026 manual that names progressive design-build is ahead of most public owners. The reform target is how these tools get defaulted — not whether they exist.",
     font=SANS, size=15, color=NAVY, italic=True, space_before=6, space_after=0, leading=1.18)

# big stat
rect(s, MARGIN + Inches(7.2), Inches(2.6), Inches(3.3), Inches(3.5), NAVY)
tb, tf = box(s, MARGIN + Inches(7.2), Inches(2.95), Inches(3.3), Inches(2.8))
para(tf, "50–75%", font=SERIF, size=58, color=WHITE, bold=True, align=PP_ALIGN.CENTER, first=True, space_after=4, leading=0.95)
para(tf, "of major airport programs now run CMAR or progressive design-build as their baseline", font=SANS, size=14, color=LIGHT, align=PP_ALIGN.CENTER, space_after=0, leading=1.18)
tb, tf = box(s, MARGIN + Inches(7.2), Inches(6.25), Inches(3.3), Inches(0.4))
para(tf, "Yet Chapter 3 still treats sealed bid as the default.", font=SANS, size=12, color=GREY, align=PP_ALIGN.CENTER, italic=True, first=True, space_after=0)
slide_number(s, 4)

# =============================================================================
# 5 — OWN LAW (big: ~90%)
# =============================================================================
s = add_slide(); bg(s, WHITE)
kicker(s, "Finding 2 · The freedom MWAA declines to use")
headline(s, "MWAA writes its own procurement law — so roughly 90% of any reform is a Board vote, not a legislative act.",
         top=Inches(1.0), size=26)
hrule(s, Inches(2.55))

rect(s, MARGIN, Inches(3.05), Inches(4.4), Inches(3.3), LIGHT)
tb, tf = box(s, MARGIN, Inches(3.5), Inches(4.4), Inches(2.6))
para(tf, "~90%", font=SERIF, size=96, color=NAVY, bold=True, align=PP_ALIGN.CENTER, first=True, space_after=0, leading=0.9)
para(tf, "OF REFORM IS WITHIN THE BOARD'S OWN GIFT", font=SANS, size=13, color=NAVY, bold=True, align=PP_ALIGN.CENTER, space_after=0)

rx = MARGIN + Inches(4.9); rw = SW - MARGIN - rx
tb, tf = box(s, rx, Inches(3.05), rw, Inches(3.6))
para(tf, "Section 1.7 exempts the Authority from the Federal Acquisition Regulation and from Virginia and D.C. procurement codes.",
     font=SANS, size=16.5, color=NAVY, bold=True, first=True, space_after=10, leading=1.15)
para(tf, "The most common excuse in public procurement — the state code won't let us — is unavailable here. Only two constraints bind:",
     font=SANS, size=15.5, color=GREY, space_after=8, leading=1.18)
para(tf, "—  The federal strings on grant money", font=SANS, size=15.5, color=NAVY, space_after=5)
para(tf, "—  Political memory", font=SANS, size=15.5, color=NAVY, space_after=10)
para(tf, "The Authority spends that freedom imitating the conservative defaults of leagues it does not have to play in.",
     font=SANS, size=15.5, color=GREY, italic=True, space_after=0, leading=1.18)
slide_number(s, 5)

# =============================================================================
# 6 — THRESHOLD ($3M vs $9B)
# =============================================================================
s = add_slide(); bg(s, WHITE)
kicker(s, "Finding 3 · The threshold has not kept pace")
headline(s, "The $3 million Board threshold is a rounding error against a $9 billion capital program.",
         top=Inches(1.0), size=27)
hrule(s, Inches(2.3))

# two-number comparison
cw = Inches(5.3)
gap = Inches(0.55)
lft = MARGIN
rgt = MARGIN + cw + gap

rect(s, lft, Inches(2.9), cw, Inches(2.5), LIGHT)
tb, tf = box(s, lft, Inches(3.15), cw, Inches(2.1))
para(tf, "$3M", font=SERIF, size=72, color=NAVY, bold=True, align=PP_ALIGN.CENTER, first=True, space_after=2, leading=0.95)
para(tf, "Board reserves approval of any contract at or above this line", font=SANS, size=14, color=GREY, align=PP_ALIGN.CENTER, space_after=0, leading=1.15)

rect(s, rgt, Inches(2.9), cw, Inches(2.5), NAVY)
tb, tf = box(s, rgt, Inches(3.15), cw, Inches(2.1))
para(tf, "$9B", font=SERIF, size=72, color=WHITE, bold=True, align=PP_ALIGN.CENTER, first=True, space_after=2, leading=0.95)
para(tf, "capital program — $6.99B of it at Dulles; a Dulles master plan reportedly scoped up to $22B through 2034",
     font=SANS, size=14, color=LIGHT, align=PP_ALIGN.CENTER, space_after=0, leading=1.15)

tb, tf = box(s, MARGIN, Inches(5.75), CW, Inches(1.4))
para(tf, "MWAA already fully delegates ALL competitively offered construction to the CEO — proof it trusts delegation where schedule matters most. Re-indexing the threshold and moving routine awards to a consent agenda frees transaction speed while the Board still governs at the program tier.",
     font=SANS, size=16, color=NAVY, first=True, space_after=0, leading=1.25)
slide_number(s, 6)

# =============================================================================
# 7 — DESIGNED IN AT AWARD (25 years)
# =============================================================================
s = add_slide(); bg(s, WHITE)
kicker(s, "Finding 4 · The default")
headline(s, "The largest operational losses are designed in at award — where the manual goes silent.",
         top=Inches(1.0), size=28)
hrule(s, Inches(2.2))

rect(s, MARGIN, Inches(2.65), Inches(4.0), Inches(3.5), NAVY)
tb, tf = box(s, MARGIN, Inches(3.0), Inches(4.0), Inches(2.9))
para(tf, "25", font=SERIF, size=120, color=WHITE, bold=True, align=PP_ALIGN.CENTER, first=True, space_after=0, leading=0.85)
para(tf, "YEARS OPERATIONS WILL RUN THE ASSET THE BID WAS OPENED ON",
     font=SANS, size=13, color=LIGHT, bold=True, align=PP_ALIGN.CENTER, space_after=0, leading=1.2)

rx = MARGIN + Inches(4.5); rw = SW - MARGIN - rx
tb, tf = box(s, rx, Inches(2.65), rw, Inches(4.0))
para(tf, "The manual requires none of:", font=SANS, size=16, color=NAVY, bold=True, first=True, space_after=8)
for t in ["Total-cost-of-ownership scoring",
          "An operations representative on the source-selection committee",
          "Operational acceptance as a condition of final payment"]:
    para(tf, "—  " + t, font=SANS, size=15.5, color=GREY, space_after=7, leading=1.15)
para(tf, "The jet bridge bought on first cost throws faults in year seven — and the operator knew it the day the bid was opened.",
     font=SANS, size=15.5, color=NAVY, italic=True, space_before=6, space_after=8, leading=1.2)
para(tf, "This bundle adds days at the front and saves a decade of write-ups. Off-the-shelf ACRP practice.",
     font=SANS, size=15.5, color=NAVY, bold=True, space_after=0, leading=1.2)
slide_number(s, 7)

# =============================================================================
# 8 — INSTRUMENT (8 -> 4 days)
# =============================================================================
s = add_slide(); bg(s, WHITE)
kicker(s, "Finding 5 · The instrument")
headline(s, "The technology gap is a binder that should be software — the real prize is Tier 1, not AI.",
         top=Inches(1.0), size=28)
hrule(s, Inches(2.2))

tb, tf = box(s, MARGIN, Inches(2.6), Inches(6.6), Inches(4.0))
para(tf, "The manual defines its “contract management system” as one that “consists of procurement policies and processes” — a governance binder, not software.",
     font=SANS, size=16, color=NAVY, first=True, space_after=10, leading=1.22)
para(tf, "In 2026, protests “sent electronically will not be accepted” — certified mail or hand delivery only — while the GAO runs its entire 1,803-protest docket online.",
     font=SANS, size=15.5, color=GREY, space_after=10, leading=1.22)
para(tf, "The proven wins are electronic sourcing, structured supplier data, spend analytics. Buy the foundation; discount the autonomous-AI pitch.",
     font=SANS, size=15.5, color=NAVY, bold=True, space_after=0, leading=1.22)

rect(s, MARGIN + Inches(7.1), Inches(2.6), Inches(3.4), Inches(3.7), LIGHT)
tb, tf = box(s, MARGIN + Inches(7.1), Inches(2.95), Inches(3.4), Inches(3.0))
para(tf, "8 → 4", font=SERIF, size=66, color=NAVY, bold=True, align=PP_ALIGN.CENTER, first=True, space_after=2, leading=0.95)
para(tf, "DAYS", font=SANS, size=18, color=NAVY, bold=True, align=PP_ALIGN.CENTER, space_after=8)
para(tf, "Peer-reviewed cycle time at an airport company after digital procurement.", font=SANS, size=13.5, color=GREY, align=PP_ALIGN.CENTER, space_after=8, leading=1.2)
para(tf, "Halving turnaround is the realistic prize — not the vendor's 60%.", font=SANS, size=13.5, color=NAVY, align=PP_ALIGN.CENTER, italic=True, space_after=0, leading=1.2)
slide_number(s, 8)

# =============================================================================
# 9 — COUNTER-CASE (2012)
# =============================================================================
s = add_slide(); bg(s, NAVY)
tb, tf = box(s, MARGIN, Inches(0.55), CW, Inches(0.4))
p = para(tf, "THE COUNTER-CASE, HONESTLY PRESENTED", font=SANS, size=12.5, color=ACCENT, bold=True, first=True, space_after=0)
p.runs[0]._r.get_or_add_rPr().set("spc", "180")
tb = headline(s, "The controls were bought with a scandal — and the credit window makes execution risk maximally expensive right now.",
              top=Inches(1.05), size=26, color=WHITE)
rect(s, MARGIN, Inches(2.75), Inches(2.2), Pt(2.4), ACCENT)

cards = [
    ("2012", "DOT IG found a culture, not a paperwork problem: competition limited on ~two-thirds of contracts over $200K (117 of 190), ~$83.6M spent without required Board approval. The $200K trigger and quarterly report are a settlement with Congress."),
    ("$5.5B", "New debt 2025–2028. Coverage falls toward ~1.3x. Rating agencies reward predictability of delivery, not velocity — “faster procurement” can read as “weakened controls.”"),
    ("$288M", "Denver's Great Hall: signed at ~30% design, terminated 14 months later, the partner seeking ~$288M in claims. Speed without a design-maturity gate authors the most expensive failure a hub can."),
]
cx = MARGIN
cw3 = (CW - Inches(0.8)) / 3
for i, (big, body) in enumerate(cards):
    x = cx + i * (cw3 + Inches(0.4))
    rect(s, x, Inches(3.15), cw3, Inches(3.5), RGBColor(0x1C,0x4A,0x82))
    tb, tf = box(s, x + Inches(0.25), Inches(3.4), cw3 - Inches(0.5), Inches(3.1))
    para(tf, big, font=SERIF, size=44, color=WHITE, bold=True, first=True, space_after=8, leading=0.95)
    para(tf, body, font=SANS, size=13, color=LIGHT, space_after=0, leading=1.2)
slide_number(s, 9)

# =============================================================================
# 10 — WHY INSUFFICIENT (62%)
# =============================================================================
s = add_slide(); bg(s, WHITE)
kicker(s, "Why the counter-case is insufficient")
headline(s, "Control is not pre-approval — so reform must substitute, never subtract.",
         top=Inches(1.0), size=28)
hrule(s, Inches(2.05))

tb, tf = box(s, MARGIN, Inches(2.5), Inches(6.7), Inches(4.2))
para(tf, "The 2012 IG punished uncompeted, unaccountable spend. The manual's response is a great deal of sequential pre-review that touches neither failure:",
     font=SANS, size=15.5, color=NAVY, first=True, space_after=8, leading=1.2)
for t in ["Hard-copy evaluations returned in paper",
          "Serial approval chains that could run in parallel",
          "A mail-only protest channel",
          "A monthly Board cadence on a $3M goods buy"]:
    para(tf, "—  " + t, font=SANS, size=15, color=GREY, space_after=5, leading=1.12)
para(tf, "These touch order of operations, not competition or accountability. Tiering controls to risk is what Sarbanes-Oxley did after Enron — scaled testing, not repeal. Net control quality rises.",
     font=SANS, size=15, color=NAVY, bold=True, space_before=8, space_after=0, leading=1.2)

rect(s, MARGIN + Inches(7.2), Inches(2.5), Inches(3.3), Inches(3.9), LIGHT)
tb, tf = box(s, MARGIN + Inches(7.2), Inches(2.8), Inches(3.3), Inches(3.4))
para(tf, "62%", font=SERIF, size=72, color=NAVY, bold=True, align=PP_ALIGN.CENTER, first=True, space_after=2, leading=0.95)
para(tf, "of a $527.7M contingency burned on Silver Line Phase 2 — one package's change orders alone topped $41.9M.",
     font=SANS, size=14, color=GREY, align=PP_ALIGN.CENTER, space_after=10, leading=1.2)
para(tf, "Every change order that waits on a Board cycle turns cost risk into schedule risk. The credit-protective reform is faster, not slower.",
     font=SANS, size=13.5, color=NAVY, align=PP_ALIGN.CENTER, italic=True, space_after=0, leading=1.2)
slide_number(s, 10)

# =============================================================================
# 11 — THE AGENDA (matrix)
# =============================================================================
s = add_slide(); bg(s, WHITE)
kicker(s, "Implications · a ranked agenda")
headline(s, "The center of gravity is execution discipline and Board self-governance — not deregulation.",
         top=Inches(1.0), size=26)
hrule(s, Inches(2.25))

rows = [
    ("Re-index the $3M Board threshold; route routine awards to a consent agenda", "9", "4", "20"),
    ("Mandate TCO scoring, an operations seat, and an acceptance-gate payment", "8", "3", "20"),
    ("Invert the burden of proof: alternative delivery presumptive for major capital", "8", "4", "20"),
    ("Install a design-maturity gate — never sign a contract at immature design", "8", "3", "20"),
    ("Risk-tiered control architecture (MOONSHOT) — the SOX-scaling move", "9", "7", "21"),
]
# header
ty = Inches(2.7)
col_x = [MARGIN, MARGIN + Inches(8.5), MARGIN + Inches(9.6), MARGIN + Inches(10.7)]
col_w = [Inches(8.4), Inches(1.0), Inches(1.0), Inches(1.0)]
hdrs = ["REFORM (TOP-RANKED)", "IMPACT", "EFFORT", "PRIORITY"]
rect(s, MARGIN, ty, CW, Inches(0.45), NAVY)
for j, h in enumerate(hdrs):
    tb, tf = box(s, col_x[j] + Inches(0.1), ty + Inches(0.07), col_w[j], Inches(0.35))
    al = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
    para(tf, h, font=SANS, size=11.5, color=WHITE, bold=True, align=al, first=True, space_after=0)

rh = Inches(0.62)
for i, (name, imp, eff, pri) in enumerate(rows):
    ry = ty + Inches(0.45) + i * rh
    if i % 2 == 0:
        rect(s, MARGIN, ry, CW, rh, LIGHT)
    moon = "MOONSHOT" in name
    nm = name.replace(" (MOONSHOT)", "")
    tb, tf = box(s, col_x[0] + Inches(0.1), ry + Inches(0.13), col_w[0], rh)
    para(tf, nm, font=SANS, size=13.5, color=(ACCENT if moon else NAVY), bold=moon, first=True, space_after=0, leading=1.0)
    for j, val in zip((1, 2, 3), (imp, eff, pri)):
        tb, tf = box(s, col_x[j], ry + Inches(0.1), col_w[j], rh)
        para(tf, val, font=SERIF, size=20, color=NAVY, bold=(j == 3), align=PP_ALIGN.CENTER, first=True, space_after=0)

tb, tf = box(s, MARGIN, Inches(6.5), CW, Inches(0.7))
para(tf, "Priority = Impact + Innovation + (10 − Effort), 1–10, analyst judgment. Four reforms tie at the top at 20; the moonshot scores 21. The flat top is the honest result — there is no single heroic fix.",
     font=SANS, size=12, color=GREY, italic=True, first=True, space_after=0, leading=1.15)
slide_number(s, 11)

# =============================================================================
# 12 — QUICK WINS
# =============================================================================
s = add_slide(); bg(s, WHITE)
kicker(s, "Quick wins · near-zero political risk")
headline(s, "Five quick wins cost almost nothing and can start now.",
         top=Inches(1.0), size=29)
hrule(s, Inches(1.95))

wins = [
    ("Adopt a cycle-time standard", "One KPI — days from requisition to award, by method and dollar band — plus quarterly variance review. Free, and it measures the band where most volume lives."),
    ("Publish a cooperative-purchasing playbook", "Approved-vehicle list and a one-page piggyback decision tree with a price-reasonableness check. Authority already in §3.7."),
    ("Convert the §1.5 forecast into an industry-and-teaming day", "A compliance artifact doubles as a supplier-engagement engine that raises competition and seeds SLBE teaming."),
    ("Digitize the analog channels", "Accept electronic protest filing; end hard-copy evaluation. Both are pure dwell with no integrity value."),
    ("Correct the FTA circular citation (4220.1F → 4220.1G)", "A currency defect a federal reviewer will catch first if MWAA does not."),
]
y = Inches(2.45)
for i, (t, b) in enumerate(wins):
    yy = y + i * Inches(0.92)
    # number chip
    rect(s, MARGIN, yy, Inches(0.55), Inches(0.55), NAVY)
    tb, tf = box(s, MARGIN, yy + Inches(0.07), Inches(0.55), Inches(0.45))
    para(tf, str(i + 1), font=SERIF, size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER, first=True, space_after=0)
    tb, tf = box(s, MARGIN + Inches(0.8), yy - Inches(0.02), CW - Inches(0.8), Inches(0.9))
    para(tf, t, font=SANS, size=16.5, color=NAVY, bold=True, first=True, space_after=2, leading=1.0)
    para(tf, b, font=SANS, size=13, color=GREY, space_after=0, leading=1.12)
slide_number(s, 12)

# =============================================================================
# 13 — CLOSING TAKEAWAY
# =============================================================================
s = add_slide(); bg(s, NAVY)
rect(s, 0, 0, SW, Inches(0.18), ACCENT)
tb, tf = box(s, MARGIN, Inches(2.0), CW, Inches(2.6))
para(tf, "Integrity is not what slows MWAA down.", font=SERIF, size=46, color=WHITE, bold=True, first=True, space_after=6, leading=1.05)
para(tf, "Sequence is.", font=SERIF, size=46, color=ACCENT, bold=True, space_after=0, leading=1.05)

rect(s, MARGIN, Inches(4.5), Inches(2.2), Pt(2.4), ACCENT)
tb, tf = box(s, MARGIN, Inches(4.8), CW, Inches(1.6))
para(tf, "The integrity the 2012 scandal bought stays untouched. The reform MWAA needs for the multi-billion-dollar decade ahead is one built to make sure the right thing happens fast enough to matter — and provably, auditable to the last signature.",
     font=SANS, size=18, color=LIGHT, first=True, space_after=0, leading=1.3)
tb, tf = box(s, MARGIN, SH - Inches(0.8), CW, Inches(0.4))
para(tf, "Same accountability. Less cycle time.", font=SANS, size=15, color=RGBColor(0x9F,0xB0,0xC6), bold=True, first=True, space_after=0)

prs.save(OUT)
print("saved", OUT, "slides:", len(prs.slides._sldIdLst))
