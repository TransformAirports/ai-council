#!/usr/bin/env python3
"""Companion executive deck: Visitor Access at Houston Airports."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

NAVY = RGBColor(0x14, 0x3C, 0x6E)
GREY = RGBColor(0x5A, 0x64, 0x72)
LIGHT = RGBColor(0x8A, 0x96, 0xA6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_BG = RGBColor(0xF2, 0xF5, 0xF9)

SERIF = "Georgia"
SANS = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

MARGIN = Inches(0.75)
CONTENT_W = SLIDE_W - Inches(1.5)


def add_slide():
    return prs.slides.add_slide(BLANK)


def textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    return box, tf


def set_run(run, text, font=SANS, size=16, color=NAVY, bold=False, italic=False):
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic


def headline(slide, text, size=30, top=Inches(0.55), width=CONTENT_W, color=NAVY):
    box, tf = textbox(slide, MARGIN, top, width, Inches(1.5))
    p = tf.paragraphs[0]
    set_run(p.add_run(), text, font=SERIF, size=size, color=color, bold=True)
    p.line_spacing = 1.05
    return box


def kicker(slide, text, top=Inches(0.28)):
    box, tf = textbox(slide, MARGIN, top, CONTENT_W, Inches(0.35))
    p = tf.paragraphs[0]
    set_run(p.add_run(), text.upper(), font=SANS, size=11, color=LIGHT, bold=True)
    return box


def rule(slide, top=Inches(1.85)):
    ln = slide.shapes.add_shape(1, MARGIN, top, Inches(1.2), Emu(0))
    ln.height = Pt(3)
    ln.fill.solid()
    ln.fill.fore_color.rgb = NAVY
    ln.line.fill.background()
    ln.shadow.inherit = False
    return ln


def page_number(slide, n):
    box, tf = textbox(slide, SLIDE_W - Inches(1.0), SLIDE_H - Inches(0.45),
                      Inches(0.6), Inches(0.3))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    set_run(p.add_run(), str(n), font=SANS, size=10, color=LIGHT)


def body_bullets(slide, items, left=MARGIN, top=Inches(2.15), width=CONTENT_W,
                 height=Inches(4.8), size=15, space_after=12):
    """items: list of (lead, rest) tuples or plain strings."""
    box, tf = textbox(slide, left, top, width, height)
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(space_after)
        p.line_spacing = 1.12
        if isinstance(item, tuple):
            lead, rest = item
            set_run(p.add_run(), lead, font=SANS, size=size, color=NAVY, bold=True)
            set_run(p.add_run(), rest, font=SANS, size=size, color=GREY)
        else:
            set_run(p.add_run(), item, font=SANS, size=size, color=GREY)
    return box


def big_stat(slide, number, label, left, top, num_size=54, width=Inches(3.6),
             label_size=12, num_color=NAVY):
    box, tf = textbox(slide, left, top, width, Inches(2.0))
    p = tf.paragraphs[0]
    set_run(p.add_run(), number, font=SERIF, size=num_size, color=num_color, bold=True)
    p2 = tf.add_paragraph()
    p2.space_before = Pt(4)
    p2.line_spacing = 1.1
    set_run(p2.add_run(), label, font=SANS, size=label_size, color=GREY)
    return box


def footnote(slide, text, top=SLIDE_H - Inches(0.62)):
    box, tf = textbox(slide, MARGIN, top, CONTENT_W - Inches(0.7), Inches(0.45))
    p = tf.paragraphs[0]
    set_run(p.add_run(), text, font=SANS, size=9.5, color=LIGHT, italic=True)
    return box


# ---------------------------------------------------------------- Slide 1: Title
s = add_slide()
bg = s.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
bg.fill.solid()
bg.fill.fore_color.rgb = NAVY
bg.line.fill.background()
bg.shadow.inherit = False

box, tf = textbox(s, Inches(1.0), Inches(2.15), Inches(11.3), Inches(2.2))
p = tf.paragraphs[0]
set_run(p.add_run(), "The Velvet Rope at the Oversubscribed Door",
        font=SERIF, size=44, color=WHITE, bold=True)
p.line_spacing = 1.05
p2 = tf.add_paragraph()
p2.space_before = Pt(18)
set_run(p2.add_run(),
        "A visitor pass program for Houston Airports — what it costs, what it buys, and where it should live",
        font=SANS, size=20, color=RGBColor(0xC9, 0xD6, 0xE8))

box, tf = textbox(s, Inches(1.0), Inches(6.35), Inches(11.3), Inches(0.8))
p = tf.paragraphs[0]
set_run(p.add_run(), "Executive briefing  |  Houston Airport System  |  July 2026",
        font=SANS, size=13, color=RGBColor(0x9A, 0xB0, 0xCC))

ln = s.shapes.add_shape(1, Inches(1.0), Inches(2.0), Inches(1.4), Pt(4))
ln.fill.solid()
ln.fill.fore_color.rgb = WHITE
ln.line.fill.background()
ln.shadow.inherit = False

# ---------------------------------------------------------------- Slide 2: Thesis
s = add_slide()
kicker(s, "The thesis")
headline(s, "The answer is a narrow yes — and the narrowness is the answer.", size=32)
rule(s)
box, tf = textbox(s, MARGIN, Inches(2.4), Inches(11.0), Inches(3.4))
p = tf.paragraphs[0]
p.line_spacing = 1.25
set_run(p.add_run(),
        "A hard-capped, off-peak, advance-vetted, twelve-month pilot at Hobby — funded outside the "
        "airline rate base and sold as community relations.",
        font=SERIF, size=26, color=NAVY)
p2 = tf.add_paragraph()
p2.space_before = Pt(20)
p2.line_spacing = 1.2
set_run(p2.add_run(),
        "Anything larger — at Bush, at scale, justified by revenue or efficiency, layered onto active "
        "construction — will be switched off in front of angry families on the days they most want it, "
        "and it deserves to fail.",
        font=SANS, size=17, color=GREY)
page_number(s, 2)

# ---------------------------------------------------------------- Slide 3: Money
s = add_slide()
kicker(s, "The evidence — cost")
headline(s, "The program is financially invisible, so money cannot decide it.")
rule(s)
big_stat(s, "$778.4M", "Houston Airport System annual revenue-fund budget (FY2027)",
         MARGIN, Inches(2.3))
big_stat(s, "$2.92B", "Five-year capital plan, FY2026–2030",
         Inches(4.95), Inches(2.3))
big_stat(s, "$0.3–1.5M", "Estimated all-in annual program cost — analyst estimate; no peer airport "
                         "publishes one; to be confirmed with a bottom-up model",
         Inches(9.0), Inches(2.3), num_size=44)
body_bullets(s, [
    ("Less than two-tenths of one percent of the annual budget. ",
     "It moves neither the 1.8x debt-service coverage ratio, the AA- credit, nor cost per enplanement. "
     "“We can’t afford it” is not an available answer."),
    ("Neither is “it pays for itself.” ",
     "At $10.16 of non-airline revenue per passenger, even 300 visitors a day yields $0.5–1.1M gross "
     "before costs — most of it fictional, because greeters get dropped off and buy nothing."),
    ("Seattle already ran this experiment. ",
     "Its program launched to grow non-aeronautical revenue and, in the Port’s own words, “took on "
     "a life of its own, turning into more of a customer experience initiative.”"),
], top=Inches(4.55), size=14, space_after=10)
page_number(s, 3)

# ---------------------------------------------------------------- Slide 4: Scarce asset
s = add_slide()
kicker(s, "The evidence — the real constraint")
headline(s, "The scarce asset is a checkpoint lane at 6:10 a.m. — and Houston just ran out of them.")
rule(s)
big_stat(s, "42.4%", "TSA officer callout rate at Bush during the March 2026 shutdown — with Hobby at "
                     "47.4%, the two worst in the nation", MARGIN, Inches(2.3), width=Inches(3.9))
big_stat(s, "4+ hrs", "Peak security waits at Bush during the shutdown, with roughly half of all "
                      "lanes dark", Inches(5.0), Inches(2.3), width=Inches(3.7))
big_stat(s, "45 min", "Terminal E peak waits on ordinary days — even after the new International "
                      "Central Processor checkpoints came online", Inches(9.0), Inches(2.3), width=Inches(3.6))
body_bullets(s, [
    ("The cost of this program is not on the balance sheet. ",
     "It is peak-hour lane-minutes, concentrated at specific lanes and hours — at a system that "
     "periodically has none to give."),
    ("The airport does not own the lane. ",
     "TSA owns the lane, staffed to a federal budget the airport neither sets nor funds. The airport "
     "owns the permission slip: the cap, the hours, the checkpoint, the blackout dates, the off switch."),
], top=Inches(4.7), size=14, space_after=10)
page_number(s, 4)

# ---------------------------------------------------------------- Slide 5: Peer programs
s = add_slide()
kicker(s, "The evidence — the peer record")
headline(s, "Every functioning peer program is small, rationed, and built to be switched off.")
rule(s)

rows = [
    ("Seattle (SEA)", "300/day", "One checkpoint only; visitors told to arrive “two or more hours” early; suspended for World Cup and summer 2026 peaks"),
    ("Philadelphia (PHL)", "100/day", "No same-day sign-up; TSA-adjudicated 1–7 days in advance — just 10,300 visitors in ~2 years against a 36,500/yr ceiling"),
    ("Detroit (DTW)", "75/day", "~25,000 passes since Oct 2023 — roughly 40 a day, about half the cap"),
    ("New Orleans (MSY)", "50/day wk", "Opens at 11 a.m., skipping the entire morning bank; one visit per person per month"),
    ("Pittsburgh (PIT)", "Dark since 2020", "The 2017 pioneer. COVID ended it; terminal construction kept it dark. Its slack was a dehubbed terminal — capacity Houston does not have"),
]
top = Inches(2.2)
for i, (apt, cap, note) in enumerate(rows):
    y = top + Inches(0.92) * i
    if i % 2 == 0:
        band = s.shapes.add_shape(1, MARGIN, y - Inches(0.06), CONTENT_W, Inches(0.9))
        band.fill.solid()
        band.fill.fore_color.rgb = ACCENT_BG
        band.line.fill.background()
        band.shadow.inherit = False
    box, tf = textbox(s, MARGIN + Inches(0.15), y, Inches(2.5), Inches(0.8))
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    set_run(tf.paragraphs[0].add_run(), apt, font=SANS, size=15, color=NAVY, bold=True)
    box, tf = textbox(s, Inches(3.5), y, Inches(2.1), Inches(0.8))
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    set_run(tf.paragraphs[0].add_run(), cap, font=SERIF, size=19, color=NAVY, bold=True)
    box, tf = textbox(s, Inches(5.7), y, Inches(6.8), Inches(0.86))
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.line_spacing = 1.05
    set_run(p.add_run(), note, font=SANS, size=12.5, color=GREY)
footnote(s, "Realized demand runs far below cap everywhere — which shrinks the congestion risk and the benefit case together. The cap is the design, not a caveat.")
page_number(s, 5)

# ---------------------------------------------------------------- Slide 6: Trivial load
s = add_slide()
kicker(s, "The evidence — sizing the pilot")
headline(s, "The recommended load is trivial: six visitors an hour at the smaller airport.")
rule(s)
big_stat(s, "~6/hr", "50 passes a day at Hobby across a midday-to-evening window — about six "
                     "visitors an hour through one designated checkpoint", MARGIN, Inches(2.3), width=Inches(3.8))
big_stat(s, "2–4%", "Share of a single real-world screening lane (150–250 passengers/hour planning "
                    "range) that load represents", Inches(4.95), Inches(2.3), width=Inches(3.7))
big_stat(s, "14.6M", "Hobby’s annual passengers — a single 24-hour terminal, versus Bush’s "
                     "48.45 million across nine checkpoints", Inches(9.0), Inches(2.3), width=Inches(3.6))
body_bullets(s, [
    ("The load is not annual and it is not average. ",
     "It is a handful of extra bodies at one checkpoint during one hour — which is why the design "
     "confines it to the trough between departure banks, as all five peer operators do."),
    ("One number the public record does not supply: ",
     "Hobby’s checkpoint count. The final cap must be set from Hobby’s own by-hour, by-checkpoint "
     "throughput before launch — a stated launch condition, not an assumption."),
], top=Inches(4.7), size=14, space_after=10)
page_number(s, 6)

# ---------------------------------------------------------------- Slide 7: Security
s = add_slide()
kicker(s, "The evidence — security")
headline(s, "The security question is federally settled; the accountability question is not.")
rule(s)
col_w = Inches(5.75)
box, tf = textbox(s, MARGIN, Inches(2.2), col_w, Inches(0.5))
set_run(tf.paragraphs[0].add_run(), "Settled by regulation", font=SERIF, size=18, color=NAVY, bold=True)
body_bullets(s, [
    ("Identical screening, watch-list vetted. ",
     "Every visitor is screened exactly like a ticketed passenger and vetted against Secure Flight "
     "watch lists before entry, under 49 CFR Parts 1542 and 1560."),
    ("A vetted visitor is arguably a lower individual risk ",
     "than a walk-up passenger — a bad actor can already buy a $39 fare, same-day, vetted only at booking."),
    ("No breach on the record. ",
     "Roughly two dozen programs over a decade; no reported mass-casualty event or catastrophic "
     "breach in the public record."),
], top=Inches(2.75), width=col_w, size=13.5, space_after=10)
box, tf = textbox(s, Inches(6.85), Inches(2.2), col_w, Inches(0.5))
set_run(tf.paragraphs[0].add_run(), "Unsettled — and owned by the airport", font=SERIF, size=18, color=NAVY, bold=True)
body_bullets(s, [
    ("Fort Lauderdale, 2017: ",
     "a ticketed passenger killed five in baggage claim in under 80 seconds — in the public zone no "
     "pass touches — then a phantom-shooter rumor drove a self-evacuation onto the tarmac."),
    ("Visitors appear on no manifest. ",
     "Concourse density turns “is everyone out?” into an unanswerable question — hence the live "
     "registry, the single location, and the public-safety veto."),
    ("TSA takes no stance on program design, ",
     "so the airport absorbs every residual risk the federal government declines to opine on."),
], left=Inches(6.85), top=Inches(2.75), width=col_w, size=13.5, space_after=10)
page_number(s, 7)

# ---------------------------------------------------------------- Slide 8: Airlines
s = add_slide()
kicker(s, "The evidence — governance")
headline(s, "In Houston the airlines are the governance structure. Keep every dollar off the rate base.")
rule(s)
big_stat(s, "58.7%", "United’s market share at Bush — roughly 80% of flights, its third-largest hub",
         MARGIN, Inches(2.3), width=Inches(3.8))
big_stat(s, ">93%", "Southwest’s share at Hobby — a near-monopoly signatory that punishes holdroom "
                    "crowding on every 30–40-minute turn", Inches(4.95), Inches(2.3), width=Inches(3.7))
big_stat(s, "Jan 2030", "Terminal E lease expiration — the next real ratemaking window with United",
         Inches(9.0), Inches(2.3), num_size=40, width=Inches(3.6))
body_bullets(s, [
    ("Fund any part of this through the airline rate base ",
     "and it enters the Majority-in-Interest clause — handing United, at ~59% activity, a near-decisive "
     "procedural lever right before the 2030 negotiation. A bad trade for a rounding error."),
    ("Kept on the non-airline side of the ledger, ",
     "it stays the airport’s decision. Whether visitor concession revenue flows to the airport or "
     "offsets airline costs is unconfirmed in the public record — resolving it is a 90-day action item, "
     "not a rhetorical crutch."),
], top=Inches(4.7), size=14, space_after=10)
page_number(s, 8)

# ---------------------------------------------------------------- Slide 9: Counter-case
s = add_slide()
kicker(s, "The counter-case, honestly presented")
headline(s, "The strongest case against: marginal load on a saturated checkpoint, marginal benefit of a rounding error.")
rule(s)
body_bullets(s, [
    ("You do not add discretionary bodies to a queue that cannot clear paying passengers. ",
     "Houston’s 2026 shutdown waits were the worst in the nation — and visitors load the slowest lanes "
     "specifically, since PreCheck and CLEAR are unavailable to pass holders."),
    ("The benefit vanishes exactly when it is most wanted. ",
     "Seattle suspends for the World Cup and summer weekends. The program is engineered to disappear on "
     "the days families most want to walk someone to a gate — a promise built to be revoked."),
    ("The “avoided escort labor” line is close to fictional. ",
     "The categories that consume labor — unaccompanied minors, reduced-mobility passengers, badged "
     "employees — are airline-run or SIDA-controlled and do not convert to a public pass."),
    ("The staffing cannot follow the launch. ",
     "49 CFR 1542.217 requires “adequate” law-enforcement coverage; sworn hiring at HPD runs 12–18 "
     "months. The program can launch in 90 days. The staffing to support it cannot."),
    ("Houston is choosing the hard version of an easy program: ",
     "no stranded Pittsburgh capacity, a dominant carrier absorbing rising cost per enplanement "
     "($10.66 → $11.17), multi-year construction, and a City-department governance structure that "
     "turns every incident into a council hearing."),
], top=Inches(2.35), size=14, space_after=11)
page_number(s, 9)

# ---------------------------------------------------------------- Slide 10: Rebuttal
s = add_slide()
kicker(s, "Why the counter-case is insufficient")
headline(s, "Every objection indicts a specific bad program — and the good program is defined by conceding them.")
rule(s)
rows = [
    ("Saturation", "An objection to United’s peak-bank checkpoint at Bush on a shutdown morning — not to 50 vetted visitors at one Hobby checkpoint on an ordinary Tuesday, auto-paused the instant waits cross a published threshold."),
    ("Suspension", "The ability to switch it off unilaterally is the entire reason it is safe to run. The failure mode is unmanaged suspension — cured by a blackout calendar published before launch, as Seattle did for the World Cup."),
    ("Security", "The Fort Lauderdale logic leads not to “no program” but to its shape: capped hard, one location, a live registry of who is inside, and a public-safety veto written into the rules."),
    ("Staffing lag", "It does not kill the pilot; it sizes it — 50 a day at one checkpoint at the smaller airport, with a written HPD/HFD understanding as a precondition, not a follow-up."),
    ("The hard version", "Standing — and it is why the yes is narrow. Refuse the hard version; build the easy one on purpose, at Hobby, where the downside is smallest."),
]
top = Inches(2.2)
for i, (lead, rest) in enumerate(rows):
    y = top + Inches(0.95) * i
    box, tf = textbox(s, MARGIN, y, Inches(2.1), Inches(0.85))
    set_run(tf.paragraphs[0].add_run(), lead, font=SERIF, size=15.5, color=NAVY, bold=True)
    box, tf = textbox(s, Inches(3.05), y, Inches(9.5), Inches(0.9))
    p = tf.paragraphs[0]
    p.line_spacing = 1.08
    set_run(p.add_run(), rest, font=SANS, size=12.5, color=GREY)
page_number(s, 10)

# ---------------------------------------------------------------- Slide 11: Recommendation
s = add_slide()
kicker(s, "The recommendation")
headline(s, "Yes — as a pilot, not a program, and at Hobby, not Bush.")
rule(s)
left_items = [
    ("Site: ", "Hobby. One terminal, one signatory, one checkpoint to protect — but confirm West "
               "Concourse phasing first; if geometry is compromised, launch after 2027 completion."),
    ("Cap: ", "50 a day, one designated checkpoint, with an operations-center throttle that pauses "
              "issuance automatically — cap set from Hobby’s own hourly checkpoint data."),
    ("Hours: ", "Midday to evening only. Morning bank and every peak date blacked out on a calendar "
                "published before launch."),
    ("Vetting: ", "The Philadelphia firewall — advance application 1–7 days out, TSA-adjudicated "
                  "before arrival. No same-day curbside kiosks."),
]
right_items = [
    ("Eligibility: ", "Narrow and sympathetic — military send-offs, elder and first-time-traveler "
                      "assistance, sanctioned events and tours."),
    ("Funding: ", "Non-airline sources only. The moment any cost touches the rate base, United gets "
                  "a grievance for free."),
    ("Stop criteria, published at launch: ", "pause at a 20-minute standard-lane wait (provisional); "
                  "suspend on any security event; do not renew below 40% cap utilization (provisional)."),
    ("Duration: ", "Twelve months — a full seasonal cycle including one peak blackout — with renewal "
                   "contingent on published criteria, not automatic."),
]
body_bullets(s, left_items, top=Inches(2.25), width=Inches(5.8), size=13, space_after=10)
body_bullets(s, right_items, left=Inches(6.85), top=Inches(2.25), width=Inches(5.75), size=13, space_after=10)
footnote(s, "Estimated all-in cost: $0.3M–$1.5M/year — analyst estimate, not a sourced total; to be confirmed with a bottom-up Houston cost model. Phase 0: landside programming needs no ASP amendment and can start now.")
page_number(s, 11)

# ---------------------------------------------------------------- Slide 12: First 90 days
s = add_slide()
kicker(s, "Implementation")
headline(s, "The first 90 days: the security-program amendment, not vendor selection, is the critical path.")
rule(s)
big_stat(s, "6–12 mo", "Realistic TSA timeline for the Airport Security Program amendment — a "
                       "discretionary, revocable permission on the Federal Security Director’s calendar, "
                       "which no vendor selection can shorten", MARGIN, Inches(2.3), num_size=44, width=Inches(4.4))
body_bullets(s, [
    ("File the ASP amendment with the Federal Security Director ", "in week one, and manage to its calendar."),
    ("Brief City Council before they read the line item ",
     "as “spending so non-passengers can shop” — and hand communications the labor rebuttal in "
     "advance: identical screening, watch-list vetted, capped, sheddable."),
    ("Confirm how the use-and-lease agreement treats visitor concession revenue, ",
     "so the funding decision rests on fact rather than assumption."),
    ("Draft the throttle protocol, blackout calendar, and published stop criteria ",
     "with the operations center and the public-safety chair."),
    ("Negotiate the written HPD/HFD staffing understanding ",
     "before booking a benefit whose cost lands on departments the airport does not command."),
], left=Inches(5.6), top=Inches(2.3), width=Inches(7.0), size=13.5, space_after=10)
page_number(s, 12)

# ---------------------------------------------------------------- Slide 13: Closing
s = add_slide()
bg = s.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
bg.fill.solid()
bg.fill.fore_color.rgb = NAVY
bg.line.fill.background()
bg.shadow.inherit = False
box, tf = textbox(s, Inches(1.0), Inches(2.2), Inches(11.3), Inches(3.2))
p = tf.paragraphs[0]
p.line_spacing = 1.15
set_run(p.add_run(),
        "The question is not whether Houston can afford the velvet rope.",
        font=SERIF, size=32, color=WHITE, bold=True)
p2 = tf.add_paragraph()
p2.space_before = Pt(20)
p2.line_spacing = 1.2
set_run(p2.add_run(),
        "It is whether it has the discipline to keep the rope small, to hang it at the right door, and "
        "to take it down without apology on the mornings the line already runs to the arrivals hall.",
        font=SERIF, size=24, color=RGBColor(0xC9, 0xD6, 0xE8))
box, tf = textbox(s, Inches(1.0), Inches(6.2), Inches(11.3), Inches(0.6))
p = tf.paragraphs[0]
set_run(p.add_run(),
        "The test, in 2030: does this sit in a footnote — or across the negotiating table?",
        font=SANS, size=15, color=RGBColor(0x9A, 0xB0, 0xCC), italic=True)
ln = s.shapes.add_shape(1, Inches(1.0), Inches(2.0), Inches(1.4), Pt(4))
ln.fill.solid()
ln.fill.fore_color.rgb = WHITE
ln.line.fill.background()
ln.shadow.inherit = False
page_number_box, tf = textbox(s, SLIDE_W - Inches(1.0), SLIDE_H - Inches(0.45), Inches(0.6), Inches(0.3))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.RIGHT
set_run(p.add_run(), "13", font=SANS, size=10, color=RGBColor(0x9A, 0xB0, 0xCC))

OUT = ("/Users/christiankessleriv/Repos/ai-council-mwaa/runs/"
       "2026-07-23-visitor-access-at-houston-airports-2/stage4/"
       "visitor-access-at-houston-airports-2.pptx")
prs.save(OUT)
print(f"Saved {OUT} with {len(prs.slides.__iter__.__self__._sldIdLst)} slides")
