"""Build the 5-slide argument brief: Data Center Concerns on Aircraft Approach.

Canonical contract: outputs/stage4/visual-brief.json (argument_brief, 5 slides).
Brand tokens: assets/brand/design-system.json.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

NAVY = RGBColor(0x0B, 0x2D, 0x4D)
NAVY_LIGHT = RGBColor(0x14, 0x3F, 0x63)
BLUE = RGBColor(0x2E, 0x84, 0xA5)
GOLD = RGBColor(0xD4, 0xA2, 0x4C)
SLATE = RGBColor(0x41, 0x56, 0x69)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FOG = RGBColor(0xED, 0xF3, 0xF6)
GREEN = RGBColor(0x24, 0x74, 0x5C)
RED = RGBColor(0xA6, 0x41, 0x3A)

DISPLAY = "Georgia"
BODY = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def new_slide():
    slide = prs.slides.add_slide(BLANK)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY
    return slide


def set_alt(shape, text):
    """Set alt text (descr) on any shape or group."""
    el = shape._element
    cNvPr = el.find(qn("p:nvSpPr") + "/" + qn("p:cNvPr"))
    if cNvPr is None:
        cNvPr = el.find(qn("p:nvGrpSpPr") + "/" + qn("p:cNvPr"))
    if cNvPr is None:
        cNvPr = el.find(qn("p:nvCxnSpPr") + "/" + qn("p:cNvPr"))
    if cNvPr is not None:
        cNvPr.set("descr", text)


def text_box(shapes, name, x, y, w, h, paras, *, size=16, font=BODY,
             color=WHITE, bold=False, italic=False, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, space_after=None, wrap=True):
    """paras: str, or list of (text, overrides-dict)."""
    box = shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.name = name
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(0)
    if isinstance(paras, str):
        paras = [(paras, {})]
    for i, (text, ov) in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ov.get("align", align)
        if space_after is not None:
            p.space_after = Pt(space_after)
        run = p.add_run()
        run.text = text
        f = run.font
        f.name = ov.get("font", font)
        f.size = Pt(ov.get("size", size))
        f.color.rgb = ov.get("color", color)
        f.bold = ov.get("bold", bold)
        f.italic = ov.get("italic", italic)
    return box


def rect(shapes, name, x, y, w, h, fill, line=None, line_w=1.0,
         shape_type=MSO_SHAPE.RECTANGLE):
    s = shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    s.name = name
    s.shadow.inherit = False
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    tf = s.text_frame
    tf.margin_left = tf.margin_right = Pt(4)
    tf.margin_top = tf.margin_bottom = Pt(2)
    return s


def shape_text(s, paras, *, size=16, font=BODY, color=WHITE, bold=False,
               align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE):
    tf = s.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(paras, str):
        paras = [(paras, {})]
    for i, (text, ov) in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ov.get("align", align)
        run = p.add_run()
        run.text = text
        f = run.font
        f.name = ov.get("font", font)
        f.size = Pt(ov.get("size", size))
        f.color.rgb = ov.get("color", color)
        f.bold = ov.get("bold", bold)


def line(shapes, name, x1, y1, x2, y2, color, w=1.5, dash=None):
    c = shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1),
                             Inches(x2), Inches(y2))
    c.name = name
    c.line.color.rgb = color
    c.line.width = Pt(w)
    c.shadow.inherit = False
    if dash:
        ln = c.line._get_or_add_ln()
        d = ln.makeelement(qn("a:prstDash"), {"val": dash})
        ln.append(d)
    return c


def title(slide, text, *, size, y=0.30, h=1.35):
    return text_box(slide.shapes, "slide-title", 0.55, y, 12.23, h, text,
                    size=size, font=DISPLAY, color=WHITE, bold=True)


def source_footer(slide, text, y=7.06, h=0.36):
    return text_box(slide.shapes, "source-note", 0.55, y, 12.23, h, text,
                    size=10, font=BODY, color=FOG)


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ---------------------------------------------------------------- slide 1
s1 = new_slide()
title(s1, "Do not build a data center inside the approach surfaces of a "
          "Cat II/III ILS runway.", size=44, y=0.28, h=2.30)

# Approach-surface fan (freeform, narrow at runway threshold, wide to the left)
fb = s1.shapes.build_freeform(Emu(Inches(9.30)), Emu(Inches(4.36)), scale=1.0)
fb.add_line_segments(
    [(Emu(Inches(9.30)), Emu(Inches(4.84))),
     (Emu(Inches(1.30)), Emu(Inches(6.25))),
     (Emu(Inches(1.30)), Emu(Inches(2.95)))],
    close=True,
)
fan = fb.convert_to_shape()
fan.name = "approach-surface-fan"
fan.shadow.inherit = False
fan.fill.solid()
fan.fill.fore_color.rgb = NAVY_LIGHT
fan.line.color.rgb = GOLD
fan.line.width = Pt(1.5)
set_alt(fan, "Schematic fan of the RWY 1L approach and RWY 19R departure "
             "surfaces extending from the runway end.")

g1 = s1.shapes.add_group_shape()
g1.name = "site-plan-diagram"
line(g1.shapes, "centerline", 1.30, 4.60, 9.30, 4.60, GOLD, 1.25, dash="dash")
rect(g1.shapes, "runway", 9.30, 4.38, 3.10, 0.44, FOG)
for i in range(5):
    rect(g1.shapes, f"alsf-bar-{i}", 7.90 + i * 0.30, 4.44, 0.05, 0.32, FOG)
dc = rect(g1.shapes, "data-center-site", 3.60, 4.02, 1.80, 1.16, RED)
shape_text(dc, "Proposed data center", size=16, bold=True)
text_box(g1.shapes, "label-runway", 9.30, 3.98, 3.05, 0.34, "RWY 1L / 19R",
         size=16, color=GOLD, bold=True)
text_box(g1.shapes, "label-fan", 1.30, 2.42, 5.30, 0.44,
         "RWY 1L approach / RWY 19R departure surfaces", size=16, color=GOLD)
text_box(g1.shapes, "label-alsf", 7.30, 4.98, 2.90, 0.32,
         "ALSF-2 approach lights", size=16, color=WHITE)
set_alt(g1, "Plan-view schematic: the proposed data center footprint sits "
            "inside the RWY 1L approach and RWY 19R departure surface fan, "
            "under the arrival and departure corridor of a Cat II/III ILS "
            "runway with ALSF-2 lighting.")
text_box(s1.shapes, "caption", 0.55, 6.55, 12.23, 0.36,
         "9,400-ft Cat II/III ILS runway with ALSF-2 approach lighting.",
         size=16, color=WHITE, italic=True)
source_footer(s1, "Schematic, not to scale. Source: FAA instrument approach "
                  "data, KIAD RWY 01L/19R — Cat II/III ILS, ALSF-2 "
                  "(FlightAware approach plate).")
notes(s1, "Open with the claim, not background. The site sits inside the "
          "RWY 1L approach and RWY 19R departure surfaces of IAD's Cat II/III "
          "ILS runway — the most instrument-critical corridor on the "
          "airfield. This is strategic capacity land; the recommendation is "
          "do not proceed at this location, not do not proceed with the "
          "developer.")

# ---------------------------------------------------------------- slide 2
s2 = new_slide()
title(s2, "A data center puts three interference vectors under one glide "
          "path.", size=36, y=0.28, h=1.25)

g2 = s2.shapes.add_group_shape()
g2.name = "SIGNATURE VISUAL — approach corridor cross-section"
gs = g2.shapes
line(gs, "ground", 0.70, 5.95, 12.65, 5.95, FOG, 2.0)
rect(gs, "runway-profile", 0.80, 5.79, 2.60, 0.16, FOG)
for i in range(5):
    rect(gs, f"alsf-tick-{i}", 3.50 + i * 0.30, 5.78, 0.045, 0.17, GOLD)
gp = line(gs, "glide-path", 3.45, 5.90, 12.30, 1.95, GOLD, 2.25, dash="dash")
gp_ln = gp.line._get_or_add_ln()
head = gp_ln.makeelement(qn("a:headEnd"),
                         {"type": "triangle", "w": "lg", "len": "lg"})
gp_ln.append(head)
text_box(gs, "label-glide-path", 9.70, 1.58, 2.95, 0.34,
         "Cat II/III glide path", size=16, color=GOLD, bold=True,
         align=PP_ALIGN.RIGHT)
bld = rect(gs, "data-center-block", 6.60, 4.95, 2.90, 1.00, SLATE)
shape_text(bld, "Data center", size=16, bold=True,
           anchor=MSO_ANCHOR.BOTTOM)
rect(gs, "tower-1", 7.00, 4.68, 0.42, 0.27, SLATE)
rect(gs, "tower-2", 7.85, 4.68, 0.42, 0.27, SLATE)
plumes = [(6.98, 2.55, 2.13), (7.62, 2.30, 2.38), (8.26, 2.60, 2.08)]
for i, (px, py, ph) in enumerate(plumes):
    a = rect(gs, f"plume-arrow-{i}", px, py, 0.42, ph, RED,
             shape_type=MSO_SHAPE.UP_ARROW)
basin = gs.add_shape(MSO_SHAPE.OVAL, Inches(9.95), Inches(5.62),
                     Inches(1.45), Inches(0.33))
basin.name = "retention-basin"
basin.shadow.inherit = False
basin.fill.solid()
basin.fill.fore_color.rgb = BLUE
basin.line.fill.background()


def badge(shapes, name, x, y, num):
    b = shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y),
                         Inches(0.34), Inches(0.34))
    b.name = name
    b.shadow.inherit = False
    b.fill.solid()
    b.fill.fore_color.rgb = RED
    b.line.color.rgb = WHITE
    b.line.width = Pt(1.0)
    shape_text(b, num, size=16, bold=True)
    return b


badge(gs, "badge-1", 7.66, 1.86, "1")
badge(gs, "badge-2", 10.50, 5.14, "2")
badge(gs, "badge-3", 6.12, 5.14, "3")
legend = [
    ("legend-1", 2.00, "1",
     "Plume turbulence to 1,000 ft above the tower (AIM § 7-6-16)"),
    ("legend-2", 2.98, "2",
     "Open-water basin draws wildlife (AC 150/5200-33C)"),
    ("legend-3", 3.96, "3",
     "Metallic mass and emissions degrade the ILS signal "
     "(Order 6750.16E)"),
]
for name, ly, num, txt in legend:
    badge(gs, name + "-badge", 0.75, ly, num)
    text_box(gs, name + "-text", 1.24, ly - 0.02, 3.45, 0.92, txt,
             size=16, color=WHITE)
set_alt(g2, "Cross-section of the approach corridor: an aircraft on the "
            "Cat II/III glide path descends over a data center whose thermal "
            "plumes rise through the flight path, an open-water basin sits "
            "beside it, and its metallic mass lies under the ILS signal — "
            "three interference vectors under one glide path.")
text_box(s2.shapes, "ils-quote", 0.90, 6.14, 11.53, 0.70,
         "“Placing an object outside the critical area does not "
         "guarantee non-interference with the ILS signal in space.” "
         "— FAA Order 6750.16E",
         size=16, italic=True, align=PP_ALIGN.CENTER)
source_footer(s2, "Schematic, not to scale. Sources: FAA AIM § 7-6-16 "
                  "(Change 3, Sept. 5, 2024); FAA AC 150/5200-33C; FAA Order "
                  "6750.16E (Apr. 10, 2014).")
notes(s2, "This is the whole mechanism in one picture. The hazards are "
          "physical, not procedural: turbulence a pilot cannot see, wildlife "
          "the corridor cannot tolerate, and signal interference no lease "
          "term can waive. FAA Technical Operations certifies the ILS and is "
          "not a party to any lease.")

# ---------------------------------------------------------------- slide 3
s3 = new_slide()
title(s3, "The FAA has already called this hazard class incompatible — "
          "and it holds the enforcement pen.", size=35, y=0.28, h=1.45)

g3 = s3.shapes.add_group_shape()
g3.name = "authority-stack"
cols = [
    (0.55, "PRECEDENT", [
        "Thermal plumes “incompatible with airport operations” "
        "— FAA guidance, 2015",
        "Cooling-tower plumes a hazard to Long Beach departures — "
        "Puente docket",
    ]),
    (6.95, "ENFORCEMENT", [
        "Section 743 preserves FAA review: aircraft safety, ground safety, "
        "federal investment",
        "Grant Assurances 19, 20, 29 attach independently",
        "14 CFR Part 16: penalties up to 3x diverted revenue",
    ]),
]
for cx, header, entries in cols:
    text_box(g3.shapes, f"hdr-{header}", cx, 2.05, 5.80, 0.42, header,
             size=20, font=BODY, color=GOLD, bold=True)
    line(g3.shapes, f"rule-{header}", cx, 2.52, cx + 5.80, 2.52, GOLD, 2.0)
    ey = 2.78
    for j, entry in enumerate(entries):
        text_box(g3.shapes, f"entry-{header}-{j}", cx + 0.05, ey, 5.70, 0.80,
                 entry, size=17, color=WHITE)
        if j < len(entries) - 1:
            line(g3.shapes, f"sep-{header}-{j}", cx + 0.05, ey + 0.92,
                 cx + 5.75, ey + 0.92, SLATE, 1.0)
        ey += 1.16
set_alt(g3, "Two-column authority stack: FAA precedent findings on thermal "
            "plume hazards beside the FAA's independent enforcement hooks "
            "over this project.")
source_footer(s3, "Sources: FAA Technical Guidance Memorandum, Sept. 24, "
                  "2015; CEC Docket 15-AFC-01; FAA Reauthorization Act of "
                  "2024 § 743; FAA Airport Sponsor Assurances, Apr. "
                  "2025; 14 CFR Part 16; 49 U.S.C. § 46301(a)(3).",
              y=6.92, h=0.52)
notes(s3, "This is precedent, not prediction. The FAA classified thermal "
          "plumes as incompatible with airport operations in 2015 and "
          "applied that finding against a power plant near Long Beach "
          "Airport. Section 743 preserves FAA jurisdiction on exactly the "
          "safety and federal-investment grounds this project implicates, "
          "and Assurances 19, 20, and 29 each attach independently.")

# ---------------------------------------------------------------- slide 4
s4 = new_slide()
title(s4, "MWAA's own data-center precedent argues for a different parcel, "
          "not this one.", size=35, y=0.28, h=1.40)

g4 = s4.shapes.add_group_shape()
g4.name = "site-comparison"
panels = [
    (0.55, GREEN, "2018 Western Lands", [
        "424 acres to Digital Realty",
        "ALP change subject to NEPA",
        "Environmental Assessment in preparation",
    ]),
    (6.93, RED, "Proposed site", [
        "Inside the 1L approach and 19R departure surfaces",
        "Aircraft-safety test is implicated",
        "Federal-investment test is implicated",
    ]),
]
for px, accent, header, entries in panels:
    hdr = rect(g4.shapes, f"panel-hdr-{header}", px, 1.95, 5.85, 0.55, accent)
    shape_text(hdr, header, size=18, bold=True, align=PP_ALIGN.LEFT)
    hdr.text_frame.margin_left = Pt(12)
    body = rect(g4.shapes, f"panel-body-{header}", px, 2.50, 5.85, 2.95,
                NAVY_LIGHT)
    tf = body.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(14)
    for j, entry in enumerate(entries):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(12)
        run = p.add_run()
        run.text = entry
        run.font.name = BODY
        run.font.size = Pt(17)
        run.font.color.rgb = WHITE
bar = rect(g4.shapes, "section-743-bar", 0.55, 5.72, 12.23, 0.80, None,
           line=GOLD, line_w=1.75)
shape_text(bar, "Section 743 preserves FAA review on safety and federal "
                "investment — assume review; do not lean on the 45-day "
                "window.", size=17, color=WHITE)
set_alt(g4, "Side-by-side comparison: the 2018 Western Lands data-center "
            "sale required an ALP change subject to NEPA, with an "
            "Environmental Assessment in preparation; the proposed "
            "site sits inside the RWY 1L approach and 19R departure "
            "surfaces, implicating aircraft-safety and federal-investment "
            "review tests.")
source_footer(s4, "Sources: MWAA Western Lands release, 2018; FAA "
                  "Reauthorization Act of 2024 § 743; FAA ALP "
                  "Preliminary Instructions Memorandum, Oct. 3, 2024.",
              y=6.98, h=0.40)
notes(s4, "Take the objection at full strength: data centers can work at "
          "IAD — MWAA sold 424 acres of Western Lands to Digital Realty in "
          "2018. The proposed site, unlike that documented case, lies "
          "inside the RWY 1L approach and 19R departure surfaces. "
          "Section 743 makes aircraft safety, ground safety, and prior "
          "federal investment explicit review tests. The precedent "
          "validates the use, not this parcel.")

# ---------------------------------------------------------------- slide 5
s5 = new_slide()
title(s5, "Halt, file the 7460-1, and move the deal — the aeronautical "
          "study governs, not the lease.", size=35, y=0.28, h=1.45)

g5 = s5.shapes.add_group_shape()
g5.name = "action-sequence"
steps = [
    ("1", "Halt advancement pending FAA coordination", "PLANNING"),
    ("2", "File Form 7460-1; require a No Hazard determination", "PLANNING"),
    ("3", "Move the deal to master-plan non-aeronautical zones",
     "COMMERCIAL REAL ESTATE"),
    ("4", "Route Grant Assurance exposure; ALP re-enters the master plan",
     "GENERAL COUNSEL"),
]
for i, (num, action, owner) in enumerate(steps):
    bx = 0.55 + i * 3.10
    panel = rect(g5.shapes, f"step-{num}", bx, 2.20, 2.92, 3.55, NAVY_LIGHT)
    text_box(g5.shapes, f"step-{num}-num", bx + 0.18, 2.38, 2.56, 0.70, num,
             size=40, font=DISPLAY, color=GOLD, bold=True)
    text_box(g5.shapes, f"step-{num}-action", bx + 0.18, 3.20, 2.56, 1.55,
             action, size=17, color=WHITE)
    line(g5.shapes, f"step-{num}-rule", bx + 0.18, 4.95, bx + 1.10, 4.95,
         GOLD, 2.0)
    text_box(g5.shapes, f"step-{num}-owner", bx + 0.18, 5.06, 2.56, 0.62,
             owner, size=16, color=GOLD, bold=True)
set_alt(g5, "Four sequenced actions with owners: halt advancement, file "
            "Form 7460-1, move the deal to non-aeronautical zones, and "
            "route Grant Assurance exposure to General Counsel.")
source_footer(s5, "Sources: FAA OE/AAA process, 14 CFR Part 77 and Form "
                  "7460-1; FAA ALP SOP 2.00; FAA Airport Sponsor "
                  "Assurances, Apr. 2025.",
              y=7.02, h=0.36)
notes(s5, "Close on the decision anchor: the governing event is the Form "
          "7460-1 filing and the aeronautical study that follows — "
          "every commitment made before that study returns is made against "
          "an undefined regulatory outcome. This redirects the deal, not "
          "the developer. The ask is a halt-and-resite, not a kill.")

# ---------------------------------------------------------------- save
prs.core_properties.title = ("Data Center Concerns on Aircraft Approach "
                             "— RWY 19R/1L, Dulles International")
prs.core_properties.author = "Transform Airports Council"
OUT = ("/Users/christiankessleriv/Repos/ai-council-mwaa/outputs/stage4/"
       "argument-data-center-concerns-on-aircraft-approach.pptx")
prs.save(OUT)
print("saved", OUT, "slides:", len(prs.slides.__iter__.__self__._sldIdLst))
