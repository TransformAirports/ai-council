#!/usr/bin/env python
"""Build the executive companion deck for 'Data Centers at Airports: Airside Risks'."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

NAVY = RGBColor(0x14, 0x3C, 0x6E)
MIDGRAY = RGBColor(0x5A, 0x5A, 0x5A)
LIGHTGRAY = RGBColor(0x8A, 0x8A, 0x8A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x9A, 0x6A, 0x1E)  # muted bronze for sparing accents

SERIF = "Georgia"
SANS = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def txbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    return box, tf


def set_run(run, text, font=SANS, size=14, color=NAVY, bold=False, italic=False):
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic


def headline(slide, text, size=28, top=Inches(0.45), width=Inches(11.9)):
    box, tf = txbox(slide, Inches(0.7), top, width, Inches(1.35))
    p = tf.paragraphs[0]
    set_run(p.add_run(), text, font=SERIF, size=size, color=NAVY, bold=True)
    p.line_spacing = 1.05
    return box


def rule(slide, top=Inches(1.72), left=Inches(0.72), width=Inches(2.2)):
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(2.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def big_number(slide, number, label, left=Inches(0.7), top=Inches(2.1),
               width=Inches(4.6), num_size=72, label_size=15):
    box, tf = txbox(slide, left, top, width, Inches(3.6))
    p = tf.paragraphs[0]
    set_run(p.add_run(), number, font=SERIF, size=num_size, color=NAVY, bold=True)
    p2 = tf.add_paragraph()
    p2.space_before = Pt(8)
    set_run(p2.add_run(), label, font=SANS, size=label_size, color=MIDGRAY)
    p2.line_spacing = 1.15
    return box


def body_bullets(slide, items, left=Inches(5.7), top=Inches(2.1),
                 width=Inches(6.9), size=15, space=12):
    """items: list of (lead, rest) or plain strings."""
    box, tf = txbox(slide, left, top, width, Inches(4.6))
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(space)
        p.line_spacing = 1.12
        if isinstance(item, tuple):
            lead, rest = item
            set_run(p.add_run(), lead, font=SANS, size=size, color=NAVY, bold=True)
            r2 = p.add_run()
            set_run(r2, " " + rest, font=SANS, size=size, color=MIDGRAY)
        else:
            set_run(p.add_run(), item, font=SANS, size=size, color=MIDGRAY)
    return box


def source_line(slide, text):
    box, tf = txbox(slide, Inches(0.7), Inches(6.85), Inches(10.8), Inches(0.4))
    p = tf.paragraphs[0]
    set_run(p.add_run(), text, font=SANS, size=9.5, color=LIGHTGRAY, italic=True)


def page_number(slide, n):
    box, tf = txbox(slide, Inches(12.45), Inches(6.95), Inches(0.7), Inches(0.4))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    set_run(p.add_run(), str(n), font=SANS, size=11, color=LIGHTGRAY)


# ---------------------------------------------------------------- Slide 1: Title
s = add_slide()
# navy field
from pptx.enum.shapes import MSO_SHAPE
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
bg.fill.solid()
bg.fill.fore_color.rgb = NAVY
bg.line.fill.background()
bg.shadow.inherit = False

box, tf = txbox(s, Inches(1.0), Inches(2.1), Inches(11.3), Inches(2.2))
p = tf.paragraphs[0]
set_run(p.add_run(), "The Building Is Already Finished", font=SERIF, size=48, color=WHITE, bold=True)
p2 = tf.add_paragraph()
p2.space_before = Pt(16)
set_run(p2.add_run(),
        "Data centers near flight paths, and the gap between clearing a surface and proving compatibility",
        font=SANS, size=20, color=RGBColor(0xC9, 0xD6, 0xE8))
p2.line_spacing = 1.2

box, tf = txbox(s, Inches(1.0), Inches(5.9), Inches(11.3), Inches(1.0))
p = tf.paragraphs[0]
set_run(p.add_run(), "Transform Airports AI Council", font=SANS, size=14, color=WHITE, bold=True)
p2 = tf.add_paragraph()
set_run(p2.add_run(), "Data Centers at Airports: Airside Risks  |  July 2026", font=SANS, size=12,
        color=RGBColor(0xC9, 0xD6, 0xE8))

# ---------------------------------------------------------------- Slide 2: Thesis
s = add_slide()
headline(s, "The test should be proven compatibility with aircraft operations — not a roofline below a surface.", size=26)
rule(s, top=Inches(1.95))
box, tf = txbox(s, Inches(0.9), Inches(2.5), Inches(11.4), Inches(3.6))
p = tf.paragraphs[0]
set_run(p.add_run(),
        "A data center near an approach or departure path should have to prove compatibility before it is approved.",
        font=SERIF, size=24, color=NAVY, bold=False, italic=True)
p.line_spacing = 1.25
p2 = tf.add_paragraph()
p2.space_before = Pt(24)
set_run(p2.add_run(),
        "The controlling risk is not the steam wisp on a cold morning. It is the process plant running under stress — "
        "a diesel generator fleet in prolonged islanded operation — and the permanent foreclosure of airfield capacity "
        "that no operator can ever un-build. The screening tools are mature; the first-pass ones are free. "
        "What is missing is the requirement to use them, and the discipline to use them on the right day.",
        font=SANS, size=16, color=MIDGRAY)
p2.line_spacing = 1.3
page_number(s, 2)

# ---------------------------------------------------------------- Slide 3: Part 77 blind
s = add_slide()
headline(s, "The review that governs the decision cannot see the hazard.")
rule(s)
big_number(s, "1,000 ft", "FAA's own guidance (AIM 7-6-16): plume turbulence “can extend to heights of over 1,000 feet above the top of the stack or cooling tower” — worst in calm, cold, stable air, in and around approach and departure corridors.")
body_bullets(s, [
    ("Part 77 measures a shape, not a behavior.", "It asks whether the structure penetrates an airspace surface. A data center is short; it passes easily."),
    ("The hazard lives in the machinery.", "Heat rejection, backup generation, standing water, EMI — none of them is a Part 77 surface. A “No Hazard” determination can coexist with every airside risk in this report."),
    ("It has happened to an airplane.", "An aircraft rolled 50–60 degrees off level at roughly 550 ft AGL over cooling towers — a power-station field, not a data center, but exactly the mechanism a generator hall reproduces under load."),
])
source_line(s, "Sources: FAA AIM 7-6-16 (Chief-Engineer brief); Regulatory-Political brief; Technology-Scout brief, citing WACAZ.")
page_number(s, 3)

# ---------------------------------------------------------------- Slide 4: Controlling scenario
s = add_slide()
headline(s, "The controlling scenario is prolonged emergency generation, not normal cooling.")
rule(s)
big_number(s, "4,700", "diesel generators in Loudoun County alone. Statewide counts range from ~9,000 (VPM) to 10,500+ (Virginia Mercury) — two published source counts, not one reconciled total.")
body_bullets(s, [
    ("The gap is the risk.", "Routine testing runs 10–30 minutes a month per unit; the non-emergency cap is ~100 hours a year; the emergency-run allowance is unlimited."),
    ("Scale is industrial.", "One Colorado campus sought 98 generators on top of 40 installed — roughly 345 MW of prime movers on a single property."),
    ("The envelope is trending against the airport.", "Transformer lead times exceed 160 weeks; Northern Virginia grid waits exceed seven years. Virginia DEQ is being asked to expand permitted run-hours. “Emergency backup” is becoming de-facto primary generation."),
])
source_line(s, "Sources: Chief-Engineer, Emergency-Management, Operations-Analyst, Airport-COO briefs; CPR; Construction Owners; Inside Climate News.")
page_number(s, 4)

# ---------------------------------------------------------------- Slide 5: Worst-day correlation
s = add_slide()
headline(s, "The data center's worst day is the airport's worst day.")
rule(s)
big_number(s, "22 hrs", "to extinguish the September 2025 lithium-ion fire at a government data center in Daejeon, South Korea — 200+ firefighters, 60 engines, 647 government IT systems dark.")
body_bullets(s, [
    ("The airport owes a 3-minute clock.", "Part 139 indexes ARFF to reach a burning aircraft in about three minutes, with foam and dry chemical — not the sustained water or submersion a lithium-ion fire demands."),
    ("The fuel is already at the fence.", "Tier III/IV designs drive 72–96 hours of on-site diesel — 25,000 to 63,400 gallons for a single 10 MW plant. An evaporative campus draws 1–5 million gallons of water a day against the airport's fire flow."),
    ("Correlation, not coincidence.", "The grid failure or ice storm that degrades the airfield is the same event that puts the diesel fleet into continuous generation a few thousand feet from the runway. A steady-state obstruction review cannot represent that composite day."),
])
source_line(s, "Sources: Director-of-Public-Safety, Emergency-Management briefs; Korea Herald; NetworkWorld; PowerMag.")
page_number(s, 5)

# ---------------------------------------------------------------- Slide 6: Wildlife
s = add_slide()
headline(s, "Wildlife and water is the one codified, enforceable federal hook.")
rule(s)
big_number(s, "78%", "of reported bird strikes occur below 1,000 feet; roughly 90% below 3,000 feet — the approach and departure band a stormwater pond sits under.")
body_bullets(s, [
    ("The bright line already exists.", "AC 150/5200-33C bars new water impoundments of a quarter-acre or larger within a runway approach and within 5,000 feet of a runway end; hazardous attractants are presumptively incompatible within 10,000 feet at turbine airports and reach to 5 statute miles."),
    ("A hyperscale campus imports two attractants.", "Large stormwater basins holding open water, and expansive flat roofs that make ideal loafing and nesting habitat."),
    ("Hardest risk to argue away at a zoning board.", "No borrowed foreign threshold, no bespoke CFD — existing strike data and existing FAA distance standards. It anchors the screen; it does not govern it."),
])
source_line(s, "Sources: Regulatory-Political, Operations-Analyst, Deep-Research briefs; FAA AC 150/5200-33C; FAA strike data.")
page_number(s, 6)

# ---------------------------------------------------------------- Slide 7: Section 743
s = add_slide()
headline(s, "Congress narrowed federal review just as the land rush arrived.")
rule(s)
big_number(s, "45 days", "for the FAA to assert jurisdiction over a non-aeronautical project under §743 of the 2024 FAA Reauthorization. Silence forfeits it.")
body_bullets(s, [
    ("The bar is now “materially impacts.”", "The FAA cannot condition a non-aeronautical project — even through a grant assurance — unless it materially impacts the safe and efficient operation of aircraft."),
    ("One door stays open.", "Frame a plume, wildlife, or EMI risk credibly and the FAA stays in the room. Fail, and the clock runs out. The evidentiary standard is what keeps the last federal door open."),
    ("Unused authority exists.", "FAA Order 7400.2K §6-3-3 already permits a hazard finding for “physical, electromagnetic, or line-of-sight interference” with flight operations — a hook beyond geometry, rarely exercised."),
])
source_line(s, "Sources: Regulatory-Political brief, citing Kaplan Kirsch and AirTAP; Deep-Research brief.")
page_number(s, 7)

# ---------------------------------------------------------------- Slide 8: Irreversibility
s = add_slide()
headline(s, "The one risk that cannot be mitigated is capacity foreclosure — and it governs.")
rule(s)
big_number(s, "40–50 yrs", "the ground lease under a commissioned data center — a 20-to-30-year fixed asset with contractual uptime obligations. There is no phase-two removal.")
body_bullets(s, [
    ("Every other risk can be engineered.", "Dry coolers, lined basins, bunded fuel, EMI shielding, run-hour caps. Foreclosure is a property of the land and the permanence, not the machinery — it has no design fix."),
    ("What is sterilized stays sterilized.", "A future runway, parallel taxiway, displaced threshold, or lower-minimums approach not yet designed — foreclosed for a generation. The obstruction check tests today's surfaces, not 2050's airfield."),
    ("The asymmetry decides.", "Reversible mistakes are survivable; irreversible ones define careers. That asymmetry, not any single engineering finding, should govern the decision."),
])
source_line(s, "Sources: Chief-Engineer, Procurement, Airport-CEO, Airport-COO briefs.")
page_number(s, 8)

# ---------------------------------------------------------------- Slide 9: Incentives
s = add_slide()
headline(s, "Every incentive points to “yes” — and the airport holds the residual risk.")
rule(s)
big_number(s, "$236.5M", "MWAA's one-time check for 424 acres of Dulles Western Lands (2018; ~$207M net). Loudoun County's data-center tax annuity: $60M in FY2013 to $800M+ in FY2026, projected at ~45% of county revenue by FY2027.")
body_bullets(s, [
    ("The county gets the annuity; the airport keeps the airspace.", "Proceeds were legally confined to lowering Dulles's cost per enplanement. The land now carries a campus master-planned for 11.7M sq ft and ~1 GW of IT load."),
    ("Power scarcity locks it in.", "Dominion's contracted data-center load jumped from ~21 GW to ~40 GW in six months of 2024; PJM capacity prices in the zone rose roughly tenfold. Dry cooling — the airside-safe choice — costs 3–4× to install with a 25–35% power penalty."),
    ("Even the carriers say yes.", "The lease lowers CPE, so signatory airlines approve — but majority-in-interest clauses give them no vote on non-aeronautical land use. The decision that can hurt them most is made in the one room where they are absent."),
])
source_line(s, "Sources: Infrastructure-Economist (NetChoice, IEEFA, DOE), Airline-Commercial-Strategist, Airport-CEO, Procurement briefs.")
page_number(s, 9)

# ---------------------------------------------------------------- Slide 10: Counter-case
s = add_slide()
headline(s, "The counter-case is real: the frightening numbers come from smokestacks, not chiller yards.")
rule(s)
big_number(s, "0", "documented cases, across this entire panel, of a data center degrading an ILS, VOR, or GPS approach — and no documented pattern of turbulence upsets or strikes at the densest clusters in the country.")
body_bullets(s, [
    ("The physics was borrowed.", "The 1,000-ft turbulence figure and the 4.3 m/s threshold come from combustion stacks and Australian power-station guidance, not data-center cooling. The FAA's own 2006 risk analysis judged overflight risk “insignificant.”"),
    ("The natural experiment points the other way.", "Elk Grove Village abuts O'Hare; Ashburn sits under Dulles departures. Decades of coexistence, hundreds of facilities — no documented airside incident."),
    ("The technology is de-risking itself.", "The industry is moving to closed-loop and “zero-water” cooling; the marginal new campus throws less visible plume than the last. A dry data center often lowers the wildlife attractant profile versus the farmland it replaces."),
])
source_line(s, "Sources: Contrarian, Technology-Scout, Chief-Engineer, Regulatory-Political briefs.")
page_number(s, 10)

# ---------------------------------------------------------------- Slide 11: Rebuttal
s = add_slide()
headline(s, "The counter-case falls short: the skeptic conceded the controlling scenario in writing.")
rule(s)
big_number(s, "“That is a\ncombustion\nsource”", "— the Contrarian brief, conceding that in a multi-day grid failure the campus runs its entire diesel fleet continuously, in exactly the calm, cold, stable window the plume literature flags.", num_size=40)
body_bullets(s, [
    ("Absence of evidence was never evidence.", "Turbulence upsets are logged as wake, weather, or “unknown”; plume causation is rarely investigated. The Dulles compatibility study was never run — its environmental review modeled wetlands, not plumes or generators. Silence from a study never run proves nothing."),
    ("“Lighter than a hotel” is a design claim.", "True with clean-agent suppression, bunded fuel, and a rehearsed mutual-aid plan; false without them. The screen converts that assertion from marketing into proof — before approval."),
    ("Irreversibility survives every rebuttal.", "Grant every engineering fix. Capacity foreclosure remains — conceded by the Contrarian: “Land-use permanence is a real cost the height-only review does miss.”"),
])
source_line(s, "Sources: Contrarian, Deep-Research, Emergency-Management, Director-of-Public-Safety briefs.")
page_number(s, 11)

# ---------------------------------------------------------------- Slide 12: The screen
s = add_slide()
headline(s, "The answer is a two-stage, configuration-first screen — not a blanket gate.")
rule(s)
box, tf = txbox(s, Inches(0.7), Inches(2.05), Inches(5.7), Inches(4.4))
p = tf.paragraphs[0]
set_run(p.add_run(), "Stage 1 — universal disclosure (cheap)", font=SERIF, size=18, color=NAVY, bold=True)
for t in ["Cooling type", "Generator fleet size and permitted run-hours",
          "Stormwater design", "Capacity-foreclosure geometry"]:
    p2 = tf.add_paragraph()
    p2.space_before = Pt(6)
    set_run(p2.add_run(), "—  " + t, font=SANS, size=15, color=MIDGRAY)
p3 = tf.add_paragraph()
p3.space_before = Pt(14)
set_run(p3.add_run(), "Four answers set the tier. Required of every applicant inside the 5-statute-mile approach/departure zone.",
        font=SANS, size=13, color=LIGHTGRAY, italic=True)
p3.line_spacing = 1.2

box, tf = txbox(s, Inches(6.9), Inches(2.05), Inches(5.7), Inches(4.6))
p = tf.paragraphs[0]
set_run(p.add_run(), "Stage 2 — studies scale to the disclosure", font=SERIF, size=18, color=NAVY, bold=True)
items = [
    ("Universal pass/fail:", "future-capacity finding (ALP-referenced) and wildlife/stormwater compliance to AC 150/5200-33C."),
    ("Scaled:", "thermal-plume modeling across three envelopes (islanded generation, peak-ambient, cold-humid); emergency-response package with fire-flow re-run and signed mutual aid; NAVAID/RFI and lighting review."),
    ("Governance:", "developer-funded, airport-directed. Never developer-marked homework."),
]
for lead, rest in items:
    p2 = tf.add_paragraph()
    p2.space_before = Pt(8)
    p2.line_spacing = 1.15
    set_run(p2.add_run(), lead, font=SANS, size=14, color=NAVY, bold=True)
    set_run(p2.add_run(), " " + rest, font=SANS, size=14, color=MIDGRAY)
source_line(s, "Tools exist and first-pass ones are free: MITRE Exhaust Plume Analyzer, ACRP Report 108, CASA screening threshold, AC 150/5200-33C. Full package: 9–15 months, on the developer.")
page_number(s, 12)

# ---------------------------------------------------------------- Slide 13: MWAA levers
s = add_slide()
headline(s, "MWAA holds real levers on its own land: the lease, the ALP, and the §743 record.")
rule(s)
big_number(s, "18–30 mo", "the honest schedule from board direction to a defensible, FAA-consented instrument — not the 6–12 months a real-estate pro forma assumes.")
body_bullets(s, [
    ("Choose the lease, not the sale.", "Fee-simple conveyance — the 2018 instrument — is the weakest way to enforce restrictions over a multi-decade asset. A ground lease with reopeners and a capacity-reservation clause keeps enforcement, reversion, and runway optionality in the Authority's hands."),
    ("Be honest about jurisdiction.", "On MWAA land, the screen is a binding lease and ALP condition. Off the fence, the operator is a commenter — publish the same standard as MWAA's formal position in 7460 comments, §743 framings, and county proceedings."),
    ("The upside is real too.", "Grid interconnection is the scarce asset. The airport that can say “yes, here, on these terms, wired into our loads” is worth more to a developer than the airport that can only say no. Cooling type is a lease clause, not a law of physics."),
])
source_line(s, "Sources: Procurement, Regulatory-Political, Airport-COO, Chief-Engineer, Virtual-Chris briefs.")
page_number(s, 13)

# ---------------------------------------------------------------- Slide 14: Close
s = add_slide()
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
bg.fill.solid()
bg.fill.fore_color.rgb = NAVY
bg.line.fill.background()
bg.shadow.inherit = False

box, tf = txbox(s, Inches(1.2), Inches(2.3), Inches(10.9), Inches(3.0))
p = tf.paragraphs[0]
set_run(p.add_run(), "Ask the questions before the concrete cures.",
        font=SERIF, size=40, color=WHITE, bold=True)
p.line_spacing = 1.1
p2 = tf.add_paragraph()
p2.space_before = Pt(24)
set_run(p2.add_run(),
        "Some data centers belong near airports — engineered right, sited off the centerline. Some do not. "
        "What separates them is whether the airport tests compatibility before approval or negotiates occupancy "
        "from weakness after the building is finished — and whether, when it finally reserves land for the runway "
        "it may need in 2050, it still owns the ground to put it on.",
        font=SANS, size=17, color=RGBColor(0xC9, 0xD6, 0xE8))
p2.line_spacing = 1.35
page_number_box, tf2 = txbox(s, Inches(12.45), Inches(6.95), Inches(0.7), Inches(0.4))
pp = tf2.paragraphs[0]
pp.alignment = PP_ALIGN.RIGHT
set_run(pp.add_run(), "14", font=SANS, size=11, color=RGBColor(0xC9, 0xD6, 0xE8))

OUT = "/Users/christiankessleriv/Repos/ai-council-mwaa/runs/2026-07-21-data-centers-at-airports-airside-risks/stage4/data-centers-at-airports-airside-risks.pptx"
prs.save(OUT)
print(f"Saved {len(prs.slides.__iter__.__self__._sldIdLst)} slides to {OUT}")
