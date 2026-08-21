"""Build the Terminal F executive briefing deck (16 slides) per the visual brief."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE

NAVY = RGBColor(0x0B, 0x2D, 0x4D)
BLUE = RGBColor(0x2E, 0x84, 0xA5)
GOLD = RGBColor(0xD4, 0xA2, 0x4C)
SLATE = RGBColor(0x41, 0x56, 0x69)
RED = RGBColor(0xA6, 0x41, 0x3A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FOG = RGBColor(0xED, 0xF3, 0xF6)
INK = RGBColor(0x17, 0x23, 0x2D)
GRAY = RGBColor(0x9A, 0xA8, 0xB3)

DISPLAY = "Georgia"
BODY = "Calibri"

prs = Presentation()
prs.slide_width = Emu(12192000)
prs.slide_height = Emu(6858000)
BLANK = prs.slide_layouts[6]


def R(text, size=16, font=BODY, color=INK, bold=False, italic=False):
    return (text, {"size": size, "font": font, "color": color, "bold": bold,
                   "italic": italic})


def P(*runs, align=PP_ALIGN.LEFT, space_before=None, space_after=None, line=None):
    return {"runs": list(runs), "align": align, "space_before": space_before,
            "space_after": space_after, "line": line}


def add_text(sh, x, y, w, h, paras, name=None, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = sh.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if name:
        tb.name = name
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = para["align"]
        if para.get("space_before") is not None:
            p.space_before = Pt(para["space_before"])
        if para.get("space_after") is not None:
            p.space_after = Pt(para["space_after"])
        if para.get("line") is not None:
            p.line_spacing = para["line"]
        for text, fmt in para["runs"]:
            r = p.add_run()
            r.text = text
            r.font.size = Pt(fmt["size"])
            r.font.name = fmt["font"]
            r.font.bold = fmt["bold"]
            r.font.italic = fmt["italic"]
            r.font.color.rgb = fmt["color"]
    return tb


def box(sh, x, y, w, h, fill=None, line=None, lw=1.0, dash=False,
        shape=MSO_SHAPE.RECTANGLE, name=None):
    s = sh.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if name:
        s.name = name
    s.shadow.inherit = False
    for child in list(s._element):
        if child.tag.endswith("}style"):
            s._element.remove(child)
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(lw)
        if dash:
            s.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return s


def set_alt(shape, desc):
    for el in shape._element.iter():
        if el.tag.endswith("}cNvPr"):
            el.set("descr", desc)
            break


def new_slide():
    return prs.slides.add_slide(BLANK)


def add_title(slide, n, text, size=35, color=NAVY, w=12.23, y=0.30, h=1.02):
    return add_text(slide.shapes, 0.55, y, w, h,
                    [P(R(text, size=size, font=DISPLAY, color=color, bold=True),
                       line=1.02)],
                    name=f"slide-{n:02d}-title")


def add_source(slide, n, text, color=SLATE, w=11.9):
    return add_text(slide.shapes, 0.55, 7.00, w, 0.44,
                    [P(R("Source: " + text, size=9.5, color=color), line=1.05)],
                    name=f"slide-{n:02d}-source")


def add_page(slide, n, color=SLATE):
    add_text(slide.shapes, 12.66, 7.06, 0.35, 0.32,
             [P(R(f"{n:02d}", size=10, color=color), align=PP_ALIGN.RIGHT)],
             name=f"slide-{n:02d}-page")


def add_chip(slide, n):
    box(slide.shapes, 11.13, 0.34, 1.65, 0.40, fill=SLATE,
        name=f"slide-{n:02d}-chip-bg")
    add_text(slide.shapes, 11.13, 0.40, 1.65, 0.30,
             [P(R("READ-AHEAD", size=16, color=WHITE, bold=True),
                align=PP_ALIGN.CENTER)],
             name=f"slide-{n:02d}-chip")


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def lock_glyph(gs, x, y, closed=True, color=SLATE):
    """Small padlock: body square plus shackle arc; open lock offsets shackle."""
    body_w, body_h = 0.20, 0.15
    sx = x + (0.13 if not closed else 0.03)
    box(gs, sx, y - 0.11, 0.14, 0.13, fill=None, line=color, lw=1.75,
        shape=MSO_SHAPE.BLOCK_ARC if False else MSO_SHAPE.OVAL)
    box(gs, x, y, body_w, body_h, fill=color)


def style_cell(cell, paras, fill=None, anchor=MSO_ANCHOR.MIDDLE):
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    else:
        cell.fill.solid()
        cell.fill.fore_color.rgb = WHITE
    cell.vertical_anchor = anchor
    cell.margin_left = Inches(0.08)
    cell.margin_right = Inches(0.08)
    cell.margin_top = Inches(0.03)
    cell.margin_bottom = Inches(0.03)
    tf = cell.text_frame
    tf.word_wrap = True
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = para["align"]
        if para.get("line") is not None:
            p.line_spacing = para["line"]
        for text, fmt in para["runs"]:
            r = p.add_run()
            r.text = text
            r.font.size = Pt(fmt["size"])
            r.font.name = fmt["font"]
            r.font.bold = fmt["bold"]
            r.font.italic = fmt["italic"]
            r.font.color.rgb = fmt["color"]


def make_table(slide, x, y, w, h, col_widths, header, rows, name,
               cell_size=12.0, header_size=12.5, first_col_bold=True):
    gf = slide.shapes.add_table(len(rows) + 1, len(header),
                                Inches(x), Inches(y), Inches(w), Inches(h))
    gf.name = name
    tbl = gf.table
    tbl.first_row = False
    tbl.horz_banding = False
    for j, cw in enumerate(col_widths):
        tbl.columns[j].width = Inches(cw)
    for j, label in enumerate(header):
        style_cell(tbl.cell(0, j),
                   [P(R(label, size=header_size, color=WHITE, bold=True),
                      line=1.0)], fill=NAVY)
    for i, row in enumerate(rows):
        fill = FOG if i % 2 == 0 else WHITE
        for j, val in enumerate(row):
            bold = first_col_bold and j == 0
            color = NAVY if bold else INK
            style_cell(tbl.cell(i + 1, j),
                       [P(R(val, size=cell_size, color=color, bold=bold),
                          line=1.02)], fill=fill)
    return gf


# ---------------------------------------------------------------- slide 1
def slide_01():
    s = new_slide()
    box(s.shapes, 0, 0, 13.333, 7.5, fill=NAVY, name="cover-field")
    add_text(s.shapes, 0.55, 0.62, 12.2, 1.75,
             [P(R("The concrete is a decoy. The register is the deadline.",
                  size=44, font=DISPLAY, color=WHITE, bold=True), line=1.04)],
             name="slide-01-title")
    box(s.shapes, 0.57, 2.42, 2.6, 0.045, fill=GOLD, name="cover-rule")
    add_text(s.shapes, 0.55, 2.62, 12.2, 0.4,
             [P(R("An executive briefing on Terminal F  ·  Dallas Fort Worth "
                  "International Airport  ·  August 2026",
                  size=18, color=FOG))],
             name="cover-subtitle")
    g = s.shapes.add_group_shape()
    g.name = "cover-module-diagram"
    gs = g.shapes
    add_text(gs, 1.35, 3.28, 10.6, 0.32,
             [P(R("Factory-built concourse module — its design decisions closed "
                  "before it arrived", size=14, color=GOLD, italic=True),
                align=PP_ALIGN.CENTER)])
    # module deck
    box(gs, 2.3, 3.75, 8.7, 1.45, fill=SLATE, line=WHITE, lw=1.5)
    box(gs, 2.3, 3.75, 8.7, 0.07, fill=GOLD)
    for i in range(1, 6):
        box(gs, 2.3 + i * 1.45, 3.88, 0.012, 1.25, fill=WHITE)
    # transporters
    for i in range(6):
        tx = 2.75 + i * 1.38
        box(gs, tx, 5.34, 0.95, 0.30, fill=FOG)
        for k in range(3):
            box(gs, tx + 0.08 + k * 0.30, 5.60, 0.20, 0.20, fill=NAVY,
                line=FOG, lw=1.0, shape=MSO_SHAPE.OVAL)
    # ground line
    box(gs, 1.55, 5.84, 10.3, 0.018, fill=FOG)
    # dimension line
    box(gs, 2.3, 6.06, 8.7, 0.014, fill=WHITE)
    box(gs, 2.3, 5.99, 0.014, 0.15, fill=WHITE)
    box(gs, 10.986, 5.99, 0.014, 0.15, fill=WHITE)
    add_text(gs, 2.3, 6.16, 8.7, 0.3,
             [P(R("Largest module 278 ft × 136 ft", size=14, color=WHITE,
                  bold=True), align=PP_ALIGN.CENTER)])
    add_text(gs, 1.35, 6.50, 10.6, 0.3,
             [P(R("3,320 tons  ·  moved by self-propelled modular transporter  ·  "
                  "set to ~3/4-inch tolerance  ·  completed August 8, 2025",
                  size=14, color=FOG), align=PP_ALIGN.CENTER)])
    set_alt(g, "Elevation diagram of a prefabricated Terminal F concourse module "
               "on self-propelled modular transporters; the largest module "
               "measures 278 by 136 feet, weighs 3,320 tons, and was set to "
               "about three-quarter-inch tolerance on August 8, 2025.")
    add_source(s, 1, "Innovation Next+ JV / Walsh Group release, module set "
                     "completed August 8, 2025; dimensions per Dallas Innovates "
                     "(2025); tolerance corroborated by International Airport "
                     "Review (Aug 2026).", color=GRAY)
    add_page(s, 1, color=GRAY)
    add_notes(s, "The photographs traveled well. What they hide: every "
                 "consequential decision about that first module was finalized "
                 "months earlier, in a factory, at a design-lock date the public "
                 "record does not name. This briefing is about the decisions "
                 "still open — and when each one closes.")


# ---------------------------------------------------------------- slide 2
def slide_02():
    s = new_slide()
    add_title(s, 2, "A $4 billion, 31-gate building committed to one carrier "
                    "through 2043")
    g = s.shapes.add_group_shape()
    g.name = "slide-02-stats"
    gs = g.shapes
    blocks = [
        ("$1.6B → $4B", "Scope growth, May 2025 — from 15 gates to 31 gates"),
        ("2043", "Use & Lease term; American Airlines is the sole occupant"),
        ("82.6%", "American's share of DFW passengers, 2025"),
        ("~100,000", "peak-day American customers, on ~930 peak-day departures"),
    ]
    for i, (num, label) in enumerate(blocks):
        bx = 0.55 + i * 3.12
        box(gs, bx, 2.05, 2.92, 0.05, fill=GOLD)
        add_text(gs, bx, 2.35, 2.92, 1.0,
                 [P(R(num, size=30, font=DISPLAY, color=NAVY, bold=True))])
        add_text(gs, bx, 3.35, 2.92, 1.9,
                 [P(R(label, size=16, color=INK), line=1.12)])
    set_alt(g, "Four statistics: Terminal F scope grew from 1.6 billion dollars "
               "and 15 gates to 4 billion dollars and 31 gates in May 2025; the "
               "use and lease agreement runs through 2043 with American as sole "
               "occupant; American carried 82.6 percent of DFW passengers in "
               "2025; about 100,000 peak-day American customers travel on about "
               "930 peak-day departures.")
    add_text(s.shapes, 0.55, 5.75, 12.23, 0.5,
             [P(R("Phase 1 opens in 2027; the full program completes in 2030. "
                  "The plan is sound and the carrier is committed — the question "
                  "is what the building should absorb before each module's "
                  "decisions close.", size=16, color=SLATE, italic=True),
                line=1.1)],
             name="slide-02-context")
    add_source(s, 2, "DFW Airport / American Airlines joint release, May 2025; "
                     "DFW airport statistics via secondary reporting (2025); "
                     "American Airlines newsroom, July 2026 and December 2025.")
    add_page(s, 2)
    add_notes(s, "Phase 1 opens 2027; the full program completes 2030. The "
                 "premise throughout: the plan is sound and American is "
                 "committed. The question is what the building should absorb "
                 "before each module's decisions close.")


# ---------------------------------------------------------------- slide 3
def slide_03():
    s = new_slide()
    add_title(s, 3, "The deadline is release-for-fabrication — and for six "
                    "modules it has already passed")
    g = s.shapes.add_group_shape()
    g.name = "SIGNATURE VISUAL — The Reversibility Clock"
    gs = g.shapes
    ax_l, ax_r, ax_y = 2.85, 12.55, 2.28
    pitch = (ax_r - ax_l) / 5.0

    def yr_x(year_frac):
        return ax_l + (year_frac - 2025.0) * pitch

    # axis + ticks
    box(gs, ax_l, ax_y, ax_r - ax_l, 0.022, fill=SLATE)
    for i in range(6):
        tx = ax_l + i * pitch
        box(gs, tx - 0.008, ax_y - 0.05, 0.016, 0.12, fill=SLATE)
        add_text(gs, tx - 0.4, ax_y + 0.12, 0.8, 0.24,
                 [P(R(str(2025 + i), size=12.5, color=SLATE),
                    align=PP_ALIGN.CENTER)])
    # fixed milestones
    for xf, lab, lx, lw_, ly, al in [
        (2025 + 7 / 12, "Modules set — Aug 2025", 2.85, 2.2, 1.92,
         PP_ALIGN.CENTER),
        (2027.0, "Phase 1 opens — 2027", 6.00, 1.95, 1.92, PP_ALIGN.CENTER),
        (2030.0, "Program completes — 2030; window closes", 10.0, 2.55, 1.60,
         PP_ALIGN.RIGHT),
    ]:
        mx = yr_x(xf)
        box(gs, mx - 0.09, ax_y - 0.08, 0.18, 0.18, fill=NAVY,
            shape=MSO_SHAPE.DIAMOND)
        add_text(gs, lx, ly, lw_, 0.62 if ly < 1.9 else 0.3,
                 [P(R(lab, size=12.5, color=NAVY, bold=True), align=al,
                    line=1.05)])
    # today rule
    tx = yr_x(2026 + 7 / 12)
    box(gs, tx - 0.02, 1.80, 0.04, 3.35, fill=GOLD)
    add_text(gs, tx - 1.0, 1.50, 2.0, 0.28,
             [P(R("TODAY — AUG 2026", size=12.5, color=GOLD, bold=True),
                align=PP_ALIGN.CENTER)])
    # swimlane 1: set modules (closed)
    add_text(gs, 0.55, 2.78, 2.15, 0.62,
             [P(R("Six modules set Aug 2025", size=12.5, color=NAVY, bold=True),
                line=1.05)])
    box(gs, ax_l, 2.80, 0.70, 0.34, fill=SLATE)
    lock_glyph(gs, ax_l + 0.77, 2.90, closed=True, color=SLATE)
    box(gs, yr_x(2025 + 7 / 12) - 0.09, 2.84, 0.18, 0.26, fill=NAVY)
    add_text(gs, ax_l + 1.45, 2.82, 5.6, 0.32,
             [P(R("decision window closed at release — months before the set",
                  size=12.5, color=SLATE))])
    # swimlane 2: next modules (open window)
    add_text(gs, 0.55, 3.62, 2.15, 0.62,
             [P(R("Next modules in the queue", size=12.5, color=NAVY, bold=True),
                line=1.05)])
    ow = box(gs, 4.35, 3.64, 3.2, 0.34, fill=FOG, line=BLUE, lw=1.75, dash=True)
    for i in range(1, 10):
        box(gs, 4.35 + i * 0.32, 3.66, 0.012, 0.30, fill=BLUE)
    lock_glyph(gs, 7.65, 3.74, closed=False, color=BLUE)
    add_text(gs, 8.1, 3.66, 4.4, 0.32,
             [P(R("window open — closes at release-for-fabrication",
                  size=12.5, color=BLUE, bold=True))])
    # swimlane 3: later modules (unpublished)
    add_text(gs, 0.55, 4.46, 2.15, 0.62,
             [P(R("Later modules", size=12.5, color=NAVY, bold=True))])
    box(gs, 5.95, 4.48, 4.7, 0.34, fill=None, line=SLATE, lw=1.5, dash=True)
    add_text(gs, 6.1, 4.52, 4.45, 0.3,
             [P(R("schedule unpublished — D0/D2 disclosure", size=12.5,
                  color=SLATE, italic=True))])
    set_alt(g, "Timeline from 2025 to 2030 showing that each Terminal F "
               "module's design decisions close at its release-for-fabrication "
               "date, months before concrete placement; the six modules set in "
               "August 2025 are already closed, future modules' windows close "
               "at unpublished release dates.")
    # callout
    box(s.shapes, 0.55, 5.42, 0.07, 1.32, fill=GOLD, name="slide-03-callout-rule")
    box(s.shapes, 0.62, 5.42, 12.16, 1.32, fill=FOG, name="slide-03-callout-bg")
    add_text(s.shapes, 0.85, 5.58, 11.7, 1.05,
             [P(R("Post-release, modular is more change-hostile than "
                  "stick-built: ", size=16, color=INK, bold=True),
                R("template changes propagate across the fabrication line. The "
                  "~30% cost and ~30% schedule savings is a repetition dividend "
                  "— DFW-supplied, not independently audited.",
                  size=16, color=INK), line=1.12)],
             name="slide-03-callout")
    add_source(s, 3, "Module set: Innovation Next+ JV / Walsh Group, Aug 2025; "
                     "Dallas Innovates (2025). ~30%/~30% savings are "
                     "DFW-supplied trade-press figures; no independent audit "
                     "exists. Per-module release dates unpublished; obtaining "
                     "them is Decision D0/D2.")
    add_page(s, 3)
    add_notes(s, "The savings are a repetition dividend — same module, same MEP "
                 "tree, over and over. Every bespoke feature either forces an "
                 "off-standard module or makes the base module carry provisions "
                 "most positions don't need. Once a module clears release, "
                 "changes propagate across the fabrication line. That inversion "
                 "is why the register is non-optional.")


# ---------------------------------------------------------------- slide 4
def slide_04():
    s = new_slide()
    add_title(s, 4, "American has told the airport, for now, what F is for — a "
                    "signal, not a settlement")
    g = s.shapes.add_group_shape()
    g.name = "slide-04-map"
    gs = g.shapes
    # Skylink loop
    loop = box(gs, 0.9, 1.75, 5.5, 3.85, fill=None, line=BLUE, lw=2.5,
               shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(gs, 2.35, 1.95, 2.6, 0.3,
             [P(R("Skylink", size=13, color=BLUE, bold=True),
                align=PP_ALIGN.CENTER)])
    # minor stations A, B, E
    for sx, sy, lab, lx, ly in [
        (2.15, 1.75, "A", 1.90, 1.94), (4.35, 1.75, "B", 4.10, 1.94),
        (2.15, 5.60, "E", 2.03, 5.90),
    ]:
        box(gs, sx - 0.13, sy - 0.13, 0.26, 0.26, fill=GRAY,
            shape=MSO_SHAPE.OVAL)
        add_text(gs, lx, ly, 0.5, 0.3,
                 [P(R(lab, size=12.5, color=SLATE, bold=True),
                    align=PP_ALIGN.CENTER)])
    # highlighted nodes C, D, F with leaders to annotation boxes
    nodes = [
        (6.4, 2.55, "C", "Terminal C — Admirals Club, 37,000 sq ft, American's "
                         "largest ever, near the Skylink station", 2.28),
        (6.4, 3.90, "D", "Terminal D — Flagship Check-In near gate D30", 3.68),
        (4.9, 5.60, "F", "Terminal F — Provisions by Admirals Club grab-and-go, "
                         "explicitly scaled down", 5.08),
    ]
    for nx, ny, lab, text, by in nodes:
        box(gs, nx - 0.19, ny - 0.19, 0.38, 0.38, fill=BLUE,
            shape=MSO_SHAPE.OVAL)
        add_text(gs, nx - 0.19, ny - 0.155, 0.38, 0.3,
                 [P(R(lab, size=13, color=WHITE, bold=True),
                    align=PP_ALIGN.CENTER)])
        box(gs, nx + 0.22, by + 0.36, 6.95 - nx - 0.25, 0.016, fill=SLATE)
        add_text(gs, 7.0, by, 3.15, 0.85,
                 [P(R(text, size=13, color=INK), line=1.08)])
    # right rail: three readings
    box(gs, 10.4, 1.75, 0.045, 3.85, fill=GOLD)
    add_text(gs, 10.6, 1.80, 2.5, 0.35,
             [P(R("Three readings", size=16, color=NAVY, bold=True))])
    for i, txt in enumerate([
        "1 — Revealed preference",
        "2 — Sequencing: F opens 2027",
        "3 — An Arrivals-premium node",
    ]):
        add_text(gs, 10.6, 2.35 + i * 1.05, 2.6, 0.9,
                 [P(R(txt, size=14, color=SLATE, bold=True), line=1.08)])
    set_alt(g, "Schematic of DFW terminals on the Skylink loop showing "
               "American's July 2026 premium placements: the largest-ever "
               "Admirals Club at Terminal C, Flagship Check-In at Terminal D, "
               "and a scaled-down Provisions grab-and-go at Terminal F.")
    add_text(s.shapes, 0.55, 6.20, 12.23, 0.5,
             [P(R("Convergence traces to one American press release plus "
                  "secondary reporting.", size=16, color=SLATE, italic=True))],
             name="slide-04-caveat")
    add_source(s, 4, "American Airlines newsroom, July 2026, corroborated by "
                     "Dallas Morning News and View from the Wing. Fleet "
                     "targets: AA newsroom, 2026 (narrowbody premium ~25% → "
                     "~40%; lie-flat count +50%+ by end of decade).")
    add_page(s, 4)
    add_notes(s, "The repositioning is real — but it is a fleet and network "
                 "story, executable inside cabins the airport does not build. "
                 "Its physical anchor at F, on American's public record, is a "
                 "grab-and-go and a widebody-feeder concourse. None of the "
                 "register's premium items should be conditioned on the July "
                 "2026 posture holding through 2043.")


# ---------------------------------------------------------------- slide 5
def slide_05():
    s = new_slide()
    add_title(s, 5, "Size to the denser-peak reading: the rebank's flattening "
                    "effect is claimed, not proven")
    g = s.shapes.add_group_shape()
    g.name = "slide-05-panel"
    gs = g.shapes
    cols = [
        (0.55, BLUE, "American's case", [
            ("9 → 13 daily banks, effective April 2026", False),
            ("Missed connections down 50% in the first 17 days", False),
            ("AA internal data — short window, not independently audited",
             True),
        ]),
        (6.83, SLATE, "The Q2 2026 audit", [
            ("Departure rate up ~1% year over year — second-worst among busy "
             "US airports", False),
            ("Block times +6 minutes average across 145 network markets "
             "(CLT +17, DCA +10, MIA +9)", False),
            ("Network-wide measure — not DFW-specific", True),
        ]),
    ]
    for cx, hdr_fill, hdr, items in cols:
        box(gs, cx, 1.60, 5.95, 3.95, fill=FOG)
        box(gs, cx, 1.60, 5.95, 0.52, fill=hdr_fill)
        add_text(gs, cx + 0.25, 1.70, 5.5, 0.35,
                 [P(R(hdr, size=16, color=WHITE, bold=True))])
        yy = 2.40
        for txt, is_flag in items:
            if is_flag:
                add_text(gs, cx + 0.25, yy + 0.15, 5.45, 0.8,
                         [P(R("Flag — " + txt, size=12.5, color=SLATE,
                              italic=True), line=1.1)])
            else:
                box(gs, cx + 0.25, yy + 0.07, 0.12, 0.12, fill=hdr_fill)
                add_text(gs, cx + 0.55, yy, 5.15, 0.95,
                         [P(R(txt, size=14.5, color=INK), line=1.1)])
                yy += 1.05
    set_alt(g, "Two-column panel contrasting American's rebank claim — 13 "
               "banks and a 50 percent missed-connection reduction from "
               "internal data — with the Q2 2026 audit showing departure rate "
               "up about one percent and block times lengthened six minutes "
               "across 145 network markets.")
    box(s.shapes, 0.55, 5.85, 12.23, 0.72, fill=NAVY, name="slide-05-footer-bg")
    add_text(s.shapes, 0.85, 6.03, 11.6, 0.4,
             [P(R("The building must serve ~100,000 peak-day customers on "
                  "~930 peak-day departures.", size=16, color=WHITE,
                  bold=True), align=PP_ALIGN.CENTER)],
             name="slide-05-footer")
    add_source(s, 5, "AA newsroom Dec 2025 (rebank) and Apr 2026; View from "
                     "the Wing (2026) reporting AA internal briefing; Airline "
                     "Geeks, July 31, 2026 (Q2 2026 performance); Cirium "
                     "schedule-data analysis via trade press (block times).")
    add_page(s, 5)
    add_notes(s, "This is where a design decision has to be made, not assumed. "
                 "Decision D4 adopts the denser-peak reading and re-tests it "
                 "annually; two consecutive corroborations of the flatter "
                 "profile unwind the affected geometry.")


# ---------------------------------------------------------------- slide 6
def slide_06():
    s = new_slide()
    add_title(s, 6, "Seven purpose-built hub bets since deregulation: the "
                    "mixed cases carry the lesson")
    rows = [
        ("Pittsburgh Midfield (1992)",
         "US Airways dehubbed 2004 over a fee dispute; CPE nearly doubled to "
         "~$14.97 by 2011",
         "Carrier-specific bet; debt outlived the hub"),
        ("Cincinnati CVG (2005 peak)",
         "22.7M passengers, 600+ daily flights → under 6M, under 180 daily by "
         "2013",
         "Unwind can take a decade and still be total"),
        ("St. Louis (2001 peak)",
         "500+ daily flights July 2001 → 207 by 2003",
         "The fastest dehub in the set: two years"),
        ("Cleveland Concourse D",
         "Went dark May 2014; daily departures ~200 → 72",
         "A purpose-built concourse can simply go dark"),
        ("Detroit McNamara (2002)",
         "Absorbed the 2008 Delta–Northwest merger without material retrofit",
         "A hub form, not a Northwest form — generic geometry cost nothing"),
        ("JetBlue JFK Terminal 5 (2008)",
         "Required $200M international-arrivals extension (2014); 2025–26 "
         "premium refresh",
         "Omitted optionality was paid for inside six years"),
        ("Atlanta Concourse F (2012)",
         "Purpose-built premium international floor plate; functioned as "
         "designed",
         "Premium floor plates work when priced and sited deliberately"),
    ]
    gf = make_table(s, 0.55, 1.55, 12.23, 4.45, [3.15, 4.85, 4.23],
                    ["Terminal (opened)", "Outcome", "Design lesson"],
                    rows, "slide-06-table")
    set_alt(gf, "Table of seven purpose-built hub terminals since "
                "deregulation: four dehubbed, Detroit's generic geometry "
                "absorbed a merger, JFK Terminal 5 paid 200 million dollars "
                "for omitted optionality, and Atlanta's Concourse F premium "
                "floor plate functioned as designed.")
    box(s.shapes, 0.55, 6.22, 12.23, 0.58, fill=GOLD, name="slide-06-bar-bg")
    add_text(s.shapes, 0.85, 6.34, 11.6, 0.38,
             [P(R("Hub-generic bones plus priced, conditioned premium shells "
                  "— not either alone.", size=17, font=DISPLAY, color=NAVY,
                  bold=True), align=PP_ALIGN.CENTER)],
             name="slide-06-takeaway")
    add_source(s, 6, "Cranky Flier (PIT, 2012); Simple Flying (CVG, STL, CLE); "
                     "DTW McNamara record; JetBlue investor release 2014 and "
                     "World Construction Network (T5/T5i); "
                     "airport-technology.com (ATL Maynard Jackson, 2012). "
                     "Small sample; DFW is not Pittsburgh — AA carried 82.6% "
                     "of DFW passengers in 2025 on an ~8M-person O&D economy.")
    add_page(s, 6)
    add_notes(s, "DTW's discipline — a hub form, not a Northwest form — cost "
                 "nothing at design time. T5's warning — no premium "
                 "optionality in original scope — cost $200M inside six "
                 "years. ATL Concourse F shows a purpose-built premium floor "
                 "plate can work at a mega-connecting hub. Read together: "
                 "generic bones, priced shells.")


# ---------------------------------------------------------------- slide 7
def slide_07():
    s = new_slide()
    add_title(s, 7, "The 6:15 a.m. case: size the pour-locked geometry for "
                    "the morning the building is worst")
    g = s.shapes.add_group_shape()
    g.name = "slide-07-flow"
    gs = g.shapes
    nodes = [
        ("Skylink overnight",
         "Single loop 22:00–06:00 maintenance window; transit ~doubles, "
         "9 → ~15 minutes"),
        ("Terminal F Skylink node",
         "Platform width, circulation cores, concourse-side queuing depth"),
        ("Checkpoint",
         "Next-generation CT lane geometry; DHS target 300 passengers/hour "
         "per lane"),
        ("FIS primary hall",
         "CBP ~100 passengers/hour per double booth; 50–75 ft queue depth"),
    ]
    for i, (hdr, body_txt) in enumerate(nodes):
        nx = 0.55 + i * 3.16
        box(gs, nx, 1.75, 2.78, 2.55, fill=FOG)
        box(gs, nx, 1.75, 2.78, 0.52, fill=NAVY)
        add_text(gs, nx + 0.16, 1.83, 2.5, 0.38,
                 [P(R(hdr, size=14.5, color=WHITE, bold=True))])
        add_text(gs, nx + 0.16, 2.48, 2.48, 1.7,
                 [P(R(body_txt, size=13, color=INK), line=1.12)])
        if i < 3:
            box(gs, nx + 2.82, 2.78, 0.30, 0.42, fill=SLATE,
                shape=MSO_SHAPE.RIGHT_ARROW)
    box(gs, 0.55, 4.62, 12.23, 0.60, fill=NAVY)
    add_text(gs, 0.85, 4.76, 11.6, 0.36,
             [P(R("Stress overlay — January 2023 winter event: ground stop, "
                  "1,100+ peak-day cancellations, 600+ the next day",
                  size=14, color=WHITE, bold=True), align=PP_ALIGN.CENTER)])
    set_alt(g, "Left-to-right passenger flow for the degraded morning: "
               "Skylink on its overnight single loop with transit near 15 "
               "minutes, the Terminal F Skylink node, the checkpoint at the "
               "300-passenger-per-hour-per-lane target, and the federal "
               "inspection hall at about 100 passengers per hour per double "
               "booth, overlaid with the January 2023 winter event.")
    add_text(s.shapes, 0.55, 5.55, 12.23, 0.85,
             [P(R("Constructed stress case — professional judgment, not a "
                  "documented daily event; adopted as the formal sizing "
                  "scenario for exactly that reason (D4).",
                  size=16, color=SLATE, italic=True), line=1.12)],
             name="slide-07-caveat")
    add_source(s, 7, "Skylink headway/transit and overnight single-loop window "
                     "per operational guidance (see appendix); Spectrum News, "
                     "Jan 2023 (storm); PARAS 0052 / CBP Airport Technical "
                     "Design Standard 2021; DHS S&T '300 People Per Hour Per "
                     "Lane'.")
    add_page(s, 7)
    add_notes(s, "The scenario has a geometry test and an operating test. The "
                 "operating test is a tabletop with AA Hub Ops, TSA FSD, CBP "
                 "Port Director, Guest Experience, ARFF, police, and Skylink "
                 "ops at the same table. If CBP will not staff the sized hall, "
                 "the shell is a moot preservation.")


# ---------------------------------------------------------------- slide 8
def slide_08():
    s = new_slide()
    add_title(s, 8, "Eight pour-locked items, none of them a Flagship Lounge")
    rows = [
        ("MARS-capable stands (minority of widebody positions)",
         "Fuel pit, PCA/400 Hz risers, jetbridge foundation, apron lead-in "
         "geometry",
         "FAA modification-of-standard letter"),
        ("FIS floor plate & sterile corridor",
         "Primary-hall dimensions; CBP ~100 pax/hr per double booth; 50–75 ft "
         "queues",
         "CBP Reimbursable Services agreement — a two-year negotiation"),
        ("Checkpoint rooms",
         "Conveyor run, search depth, ceiling height; 300 pax/hr/lane CT "
         "geometry",
         "Written TSA concurrence — on the critical path"),
        ("Skylink node",
         "Platform width, cores, queuing depth for the ~15-minute overnight "
         "loop",
         "Adopt the D4 sizing scenario"),
        ("Baggage-handling trunk cross-section",
         "Trunk dimension through the concourse at next-decade bank volumes",
         "None federal; disruptive and costlier as an operating retrofit"),
        ("Lounge-capable shells (head-of-pier)",
         "MEP risers, grease waste, structural allowance, curtain wall — "
         "deliberately not built out",
         "Terminal F Reversibility Side Letter with American"),
        ("Instrumentation baseline",
         "Sensor conduit, PoE++ drops, fiber diversity — through the factory "
         "once",
         "D0 — template changes at repeat-part cost"),
        ("CUPPS-capable gate hardware",
         "Common-use IT provision across the exclusive-use envelope",
         "American IT participation in factory witness testing"),
    ]
    gf = make_table(s, 0.55, 1.55, 12.23, 4.75, [3.35, 4.95, 3.93],
                    ["Register item", "What locks at pour or release",
                     "Dependency to clear first"],
                    rows, "slide-08-table")
    set_alt(gf, "Register table of eight pour-locked Terminal F items — MARS "
                "stands, federal inspection floor plate, checkpoint rooms, "
                "Skylink node, baggage trunk, lounge-capable shells, "
                "instrumentation baseline, and common-use gate hardware — "
                "each with its locking geometry and the federal or airline "
                "dependency that must clear first.")
    add_text(s.shapes, 0.55, 6.48, 12.23, 0.42,
             [P(R("Costs are register-grade, not audit-grade — confirm before "
                  "the Finance & Audit Committee.", size=16, color=SLATE,
                  bold=True, italic=True))],
             name="slide-08-flag")
    add_source(s, 8, "ACI-NA MARS gates (Sept 2024); FAA AC 150/5300-13B; "
                     "PARAS 0052 / CBP ATDS 2021; DHS S&T CT-lane target; "
                     "Skylink operational guidance. Order-of-magnitude cost "
                     "characterizations are analyst-constructed and must be "
                     "confirmed before the register goes to committee.")
    add_page(s, 8)
    add_notes(s, "All of it works for a tenant DFW has not met. The airport "
                 "pays an option premium now — most of it inside the module "
                 "template at repeat-part cost if D0 clears — and preserves "
                 "the ability for American, or a future signatory, to fit out "
                 "any envelope at the tenant's expense.")


# ---------------------------------------------------------------- slide 9
def slide_09():
    s = new_slide()
    add_title(s, 9, "The affordability arbitrage is being spent — every "
                    "option needs a number and a ceiling")
    g = s.shapes.add_group_shape()
    g.name = "slide-09-chart"
    gs = g.shapes
    # Panel A — CPE
    add_text(gs, 0.55, 1.62, 5.9, 0.35,
             [P(R("Cost per enplanement — signatory CPE, USD",
                  size=16, color=NAVY, bold=True))])
    base_y = 5.05
    for bx, val, h, lab, fill_solid in [
        (1.55, "$13.59", 13.59 / 16.99 * 2.45, "FY25 actual", True),
        (3.75, "$16.99", 2.45, "FY26 projected", False),
    ]:
        by = base_y - h
        if fill_solid:
            box(gs, bx, by, 0.95, h, fill=NAVY)
        else:
            box(gs, bx, by, 0.95, h, fill=FOG, line=NAVY, lw=1.75, dash=True)
        add_text(gs, bx - 0.35, by - 0.34, 1.65, 0.3,
                 [P(R(val, size=14, color=NAVY, bold=True),
                    align=PP_ALIGN.CENTER)])
        add_text(gs, bx - 0.45, base_y + 0.08, 1.85, 0.3,
                 [P(R(lab, size=12.5, color=SLATE), align=PP_ALIGN.CENTER)])
    box(gs, 0.85, base_y, 4.9, 0.018, fill=SLATE)
    add_text(gs, 4.9, 3.05, 1.55, 0.55,
             [P(R("+25% in one year", size=12.5, color=NAVY, bold=True),
                line=1.05)])
    add_text(gs, 0.55, 5.55, 5.9, 0.35,
             [P(R("Flag — no public CPE projection exists for FY27–FY30",
                  size=12, color=SLATE, italic=True))])
    # Panel B — debt
    add_text(gs, 6.85, 1.62, 5.9, 0.35,
             [P(R("Outstanding debt — USD billions", size=16, color=NAVY,
                  bold=True))])
    for bx, val, h, lab, fill_solid in [
        (7.85, "$7.2B", 7.224 / 12.4 * 2.45, "2024, at upgrade", True),
        (10.05, "$12.4B", 2.45, "projected FY29", False),
    ]:
        by = base_y - h
        if fill_solid:
            box(gs, bx, by, 0.95, h, fill=NAVY)
        else:
            box(gs, bx, by, 0.95, h, fill=FOG, line=NAVY, lw=1.75, dash=True)
        add_text(gs, bx - 0.35, by - 0.34, 1.65, 0.3,
                 [P(R(val, size=14, color=NAVY, bold=True),
                    align=PP_ALIGN.CENTER)])
        add_text(gs, bx - 0.45, base_y + 0.08, 1.85, 0.3,
                 [P(R(lab, size=12.5, color=SLATE), align=PP_ALIGN.CENTER)])
    box(gs, 7.15, base_y, 4.9, 0.018, fill=SLATE)
    add_text(gs, 6.85, 5.55, 6.0, 0.6,
             [P(R("Flag — projection pre-dates the May 2025 +$2.4B scope expansion "
                  "— post-expansion path unpublished and higher",
                  size=12, color=SLATE, italic=True), line=1.08)])
    set_alt(g, "Two-panel chart: DFW signatory cost per enplanement rises "
               "from 13.59 dollars actual in fiscal 2025 to a projected 16.99 "
               "dollars in fiscal 2026, and outstanding debt rises from 7.2 "
               "billion dollars in 2024 to a projected 12.4 billion by fiscal "
               "2029, a projection that pre-dates the May 2025 scope "
               "expansion.")
    add_text(s.shapes, 0.55, 6.28, 12.23, 0.62,
             [P(R("The 70th Supplemental Bond Ordinance authorizes $3.0B in "
                  "new debt through February 2026; beyond it, both the Dallas "
                  "and Fort Worth city councils must approve.",
                  size=16, color=INK), line=1.1)],
             name="slide-09-footer")
    add_source(s, 9, "Cost per Enplanement by Airport — CPE Data 2026 "
                     "(industry aggregation of the 2025 A/B bond Official "
                     "Statement); The Bond Buyer, Aug 2024 (S&P upgrade to "
                     "AA-; debt projection); DFW 70th Supplemental Bond "
                     "Ordinance.")
    add_page(s, 9)
    add_notes(s, "DFW remains well below JFK, LAX, EWR, and ORD on CPE — the "
                 "peer table is in the read-ahead. The register's answer to "
                 "affordability: a Finance & Audit Committee forward-CPE "
                 "ceiling the aggregate register premium may not breach, with "
                 "unwind on breach.")


# ---------------------------------------------------------------- slide 10
def slide_10():
    s = new_slide()
    add_title(s, 10, "Six decisions before the fabrication window closes in "
                     "2030 — D0 precedes everything")
    g = s.shapes.add_group_shape()
    g.name = "slide-10-decisions"
    gs = g.shapes
    decisions = [
        ("D0", "Confirm the JV template-change path at repeat-part cost",
         "Chief Procurement Officer", True),
        ("D1", "Close the two U&L contract questions",
         "General Counsel", False),
        ("D2", "Publish the reversibility register, module by module",
         "VP Terminal Development", False),
        ("D3", "Execute the Reversibility Side Letter",
         "Chief Development Officer", False),
        ("D4", "Adopt the degraded-morning sizing case",
         "Chief Operating Officer", False),
        ("D5", "Fix the instrumentation template",
         "Chief Information Officer", False),
    ]
    for i, (did, action, owner, is_gate) in enumerate(decisions):
        nx = 0.55 + i * 2.045
        box(gs, nx, 1.65, 1.92, 2.65, fill=FOG)
        band = GOLD if is_gate else NAVY
        band_txt = NAVY if is_gate else WHITE
        box(gs, nx, 1.65, 1.92, 0.46, fill=band)
        add_text(gs, nx + 0.12, 1.72, 1.7, 0.32,
                 [P(R(did, size=15, font=DISPLAY, color=band_txt, bold=True))])
        add_text(gs, nx + 0.12, 2.25, 1.7, 1.35,
                 [P(R(action, size=12.5, color=INK), line=1.08)])
        add_text(gs, nx + 0.12, 3.62, 1.7, 0.62,
                 [P(R(owner, size=11.5, color=SLATE, italic=True),
                    line=1.05)])
    box(gs, 0.62, 4.48, 11.5, 0.03, fill=GOLD)
    box(gs, 12.12, 4.36, 0.24, 0.27, fill=GOLD, shape=MSO_SHAPE.RIGHT_ARROW)
    add_text(gs, 0.62, 4.58, 11.7, 0.3,
             [P(R("D0 gates the register's cadence and the instrumentation "
                  "rollout", size=12.5, color=NAVY, bold=True))])
    set_alt(g, "Six sequenced decisions with owners: D0 confirms the joint "
               "venture template-change path and gates the rest; D1 through "
               "D5 cover the lease questions, the register, the Side Letter, "
               "the sizing case, and the instrumentation template — all "
               "before the fabrication-release window closes in 2030.")
    box(s.shapes, 0.55, 5.10, 12.23, 1.62, fill=FOG, name="slide-10-stop-bg")
    stop_rows = [
        "American announces a Flagship-tier build-out on a published schedule "
        "→ a fit-out program replaces the register",
        "Connecting share falls below a connecting-majority floor for two "
        "consecutive cycles → re-test the affected geometry",
        "U&L renegotiation opens before 2043 → the Side Letter folds into "
        "the amendment",
    ]
    stop_paras = []
    for row in stop_rows:
        stop_paras.append(P(R("STOP  ", size=13, color=RED, bold=True),
                            R(row, size=13, color=INK),
                            space_after=6, line=1.05))
    gstop = s.shapes.add_group_shape()
    gstop.name = "slide-10-stop-band"
    add_text(gstop.shapes, 0.85, 5.28, 11.7, 1.3, stop_paras)
    box(gstop.shapes, 0.62, 5.10, 0.07, 1.62, fill=RED)
    set_alt(gstop, "Three stop conditions that halt or convert the register "
                   "program.")
    add_source(s, 10, "Decision program per this report; fabrication window "
                      "bounded by program completion 2030 (DFW/AA joint "
                      "release, May 2025).")
    add_page(s, 10)
    add_notes(s, "The register does not tell DFW what to build. It tells DFW "
                 "what remains to be decided, before which date, at what "
                 "price, from which funding source, against which federal "
                 "dependency, inside which contract instrument. The customer "
                 "it protects is the connecting passenger in the first bank "
                 "of a bad morning.")


# ---------------------------------------------------------------- slide 11
def slide_11():
    s = new_slide()
    add_title(s, 11, "The counter-case, honestly presented: four objections, "
                     "four constraints", w=10.4)
    add_chip(s, 11)
    rows = [
        ("American signed to 2043 — why insure a committed relationship?",
         "Generic geometry insures against strategy shifts inside a 17-year "
         "term — and protects the tenant as much as the airport"),
        ("Premium is the airline's side of the line: Delta's JFK T4 Delta One "
         "Lounge (~40,000 sq ft, 515 seats) was carrier-funded on someone "
         "else's base building",
         "Exactly — hence Side Letter conditioning: the airport builds shells "
         "only; American exercises at American's expense"),
        ("Options accumulate into an unfunded shopping list",
         "A three-to-four-line register that refuses everything else is the "
         "discipline, not the abandonment"),
        ("Affordability: signatory CPE is already rising ~25% in a single "
         "year",
         "The Finance & Audit forward-CPE ceiling binds the aggregate premium "
         "— and under-scoping is often the more expensive error"),
    ]
    gf = make_table(s, 0.55, 1.55, 12.23, 5.15, [6.1, 6.13],
                    ["Objection, at full strength",
                     "Disposition — as a constraint, not a defeat"],
                    rows, "slide-11-table", cell_size=12.5, first_col_bold=False)
    set_alt(gf, "Table pairing the four strongest objections to the "
                "reversibility register — the 2043 commitment, "
                "carrier-funded premium precedent at JFK Terminal 4, "
                "shopping-list risk, and affordability — with the constraint "
                "each imposes on the program.")
    add_source(s, 11, "Delta News Hub (JFK T4 Delta One Lounge); DFW/AA joint "
                      "release May 2025; CPE Data 2026; The Bond Buyer, "
                      "Aug 2024.")
    add_page(s, 11)
    add_notes(s, "Read-ahead only; the presenter references it if challenged. "
                 "The sharpest objection is the shopping list — and its "
                 "answer is the register's discipline, not its abandonment.")


# ---------------------------------------------------------------- slide 12
def slide_12():
    s = new_slide()
    add_title(s, 12, "Premium posture is a cycle, not a constant", w=10.4)
    add_chip(s, 12)
    g = s.shapes.add_group_shape()
    g.name = "slide-12-timeline"
    gs = g.shapes
    ax_l, ax_y, pitch = 1.0, 3.95, 1.45
    box(gs, ax_l, ax_y, 11.6, 0.022, fill=SLATE)
    for i in range(9):
        tx = ax_l + i * pitch
        box(gs, tx - 0.008, ax_y - 0.05, 0.016, 0.12, fill=SLATE)
        add_text(gs, tx - 0.35, ax_y + 0.10, 0.7, 0.24,
                 [P(R(str(2020 + i), size=12, color=SLATE),
                    align=PP_ALIGN.CENTER)])
    events = [
        (1.24, "above", 0.60, 2.55, "Mar 2020",
         "American closes every Flagship Lounge except JFK (LAX, DFW, MIA, "
         "ORD)"),
        (3.50, "below", 2.30, 2.55, "Fall 2021 – Apr 2022",
         "Reopening sequence restores the Flagship network"),
        (9.20, "above", 7.60, 2.55, "Q4 2025",
         "Delta premium ticket revenue crosses main cabin: $5.70B vs $5.62B "
         "(FY25 premium $22.1B)"),
        (10.42, "below", 8.90, 3.30, "2026",
         "American shifts long-haul into DFW: +6% Q3 YoY (LHR −13%); six new "
         "international routes summer 2026"),
        (12.55, "above", 10.30, 2.45, "2028",
         "ATL Delta One Lounge target — 'Delta stated the date, not the "
         "reason'"),
    ]
    for cx, side, bx, bw, date, txt in events:
        box(gs, cx - 0.06, ax_y - 0.05, 0.12, 0.12, fill=NAVY,
            shape=MSO_SHAPE.OVAL)
        if side == "above":
            by, bh = 1.90, 1.35
            box(gs, cx - 0.008, by + bh, 0.016, ax_y - (by + bh), fill=GRAY)
        else:
            by, bh = 4.62, 1.35
            box(gs, cx - 0.008, ax_y + 0.07, 0.016, by - ax_y - 0.07,
                fill=GRAY)
        box(gs, bx, by, bw, bh, fill=FOG)
        add_text(gs, bx + 0.14, by + 0.10, bw - 0.28, 0.28,
                 [P(R(date, size=12.5, color=NAVY, bold=True))])
        add_text(gs, bx + 0.14, by + 0.40, bw - 0.28, bh - 0.5,
                 [P(R(txt, size=12, color=INK), line=1.08)])
    set_alt(g, "Timeline 2020 to 2028: American closed every Flagship Lounge "
               "except JFK in March 2020 and reopened them through April "
               "2022; Delta premium revenue crossed main cabin in the fourth "
               "quarter of 2025; American shifted long-haul capacity into DFW "
               "in 2026; Delta targets a 2028 Delta One Lounge in Atlanta.")
    add_text(s.shapes, 0.55, 6.30, 12.23, 0.55,
             [P(R("A design bet that must survive to 2043 needs to survive "
                  "several repetitions of the 2020–2022 cycle.",
                  size=16, color=SLATE, italic=True))],
             name="slide-12-takeaway")
    add_source(s, 12, "Simple Flying (2020 Flagship closures and reopenings); "
                      "Delta Q4 2025 earnings release; Afar and industry "
                      "coverage (ATL 2028 target); Simple Flying / Travel and "
                      "Tour World (AA long-haul capacity); AA newsroom via "
                      "aviationa2z (summer 2026 routes).")
    add_page(s, 12)
    add_notes(s, "A design bet that must survive to 2043 needs to survive "
                 "several repetitions of the 2020–2022 cycle. That is the "
                 "argument for shells American exercises at American's "
                 "expense — not airport-side preload.")


# ---------------------------------------------------------------- slide 13
def slide_13():
    s = new_slide()
    add_title(s, 13, "DFW's cost position versus peer hubs, 2024", w=10.4)
    add_chip(s, 13)
    g = s.shapes.add_group_shape()
    g.name = "slide-13-bars"
    gs = g.shapes
    add_text(gs, 0.55, 1.46, 8.0, 0.3,
             [P(R("USD per enplaned passenger — 2024 peer table",
                  size=13, color=SLATE, bold=True))])
    data = [
        ("DTW", 9.20), ("IAH", 10.66), ("MSP", 11.06), ("DEN", 12.76),
        ("DFW", 13.44), ("PHL", 15.03), ("MIA", 16.83), ("SEA", 18.24),
        ("ORD", 29.56), ("LAX", 30.16), ("EWR", 31.67), ("JFK", 36.01),
    ]
    y0, bpitch, bh = 1.84, 0.360, 0.25
    for i, (code, val) in enumerate(data):
        yy = y0 + i * bpitch
        is_dfw = code == "DFW"
        add_text(gs, 0.55, yy + 0.015, 1.05, 0.26,
                 [P(R(code, size=12.5, color=NAVY if is_dfw else SLATE,
                      bold=True), align=PP_ALIGN.RIGHT)])
        bw = val / 36.01 * 9.0
        box(gs, 1.80, yy, bw, bh, fill=BLUE if is_dfw else SLATE)
        label = f"${val:,.2f}" + ("   — this airport" if is_dfw else "")
        add_text(gs, 1.90 + bw, yy + 0.015, 2.6 if is_dfw else 1.4, 0.26,
                 [P(R(label, size=12, color=NAVY if is_dfw else INK,
                      bold=is_dfw))])
    add_text(gs, 0.55, 6.30, 12.2, 0.6,
             [P(R("DWU Consulting peer table, 2024 vintage. DFW's "
                  "bond-covenanted signatory CPE — FY25 $13.59 actual, FY26 "
                  "$16.99 projected — uses a different scope; the two "
                  "measures are directionally consistent.",
                  size=12, color=SLATE, italic=True), line=1.08)])
    set_alt(g, "Horizontal bar chart of 2024 cost per enplanement across "
               "twelve US hubs, ascending from Detroit at 9.20 dollars to "
               "JFK at 36.01 dollars, with DFW highlighted at 13.44 dollars "
               "— well below JFK, LAX, Newark, and O'Hare.")
    add_source(s, 13, "Cost per Enplanement by Airport — CPE Data 2026; DWU "
                      "Consulting peer table (2024). Units: USD per enplaned "
                      "passenger.")
    add_page(s, 13)
    add_notes(s, "The arbitrage is real and closing; the ceiling converts it "
                 "from a talking point into a governance instrument.")


# ---------------------------------------------------------------- slide 14
def slide_14():
    s = new_slide()
    add_title(s, 14, "What the retrofit costs when the option was never "
                     "bought", w=10.4)
    add_chip(s, 14)
    rows = [
        ("IAD AeroTrain",
         "Concourse D stop omitted in the 2000s routing; 16 years on mobile "
         "lounges",
         "Extension budgeted ~$3.75B inside a $22.5B program — program "
         "direction, not committed contract"),
        ("DEN Great Hall",
         "P3 governance failure mid-build",
         "$184M termination payment plus ~$2.1B completion"),
        ("LHR Terminal 5",
         "Commissioning, not geometry — day-one systems failure",
         "£4.3B build on time and on budget; 42,000 bags stranded at "
         "opening"),
        ("JFK Terminal 5",
         "No premium or international optionality in the original scope",
         "$875M build; $200M T5i extension (2014); 2025–26 premium refresh, "
         "cost undisclosed"),
        ("LGA Terminal C",
         "The governance counter-example — disciplined single-airline "
         "delivery",
         "$4B, completed ~2 years ahead of schedule"),
    ]
    gf = make_table(s, 0.55, 1.55, 12.23, 4.95, [2.35, 4.7, 5.18],
                    ["Precedent", "What was omitted or failed",
                     "What it cost"],
                    rows, "slide-14-table", cell_size=12.5)
    set_alt(gf, "Five retrofit and delivery precedents: Dulles AeroTrain's "
                "omitted Concourse D stop, Denver's Great Hall P3 "
                "termination, Heathrow Terminal 5's opening-day systems "
                "failure, JFK Terminal 5's paid-for extension, and LaGuardia "
                "Terminal C as the disciplined governance counter-example.")
    add_source(s, 14, "NBC News and trade press (IAD); Colorado Public Radio, "
                      "Aug 2019 (DEN); IEEE Spectrum, 2008 (LHR T5); JetBlue "
                      "investor release 2014 / World Construction Network "
                      "(JFK T5); Delta News Hub 2022 and Aviation Week 2024 "
                      "(LGA).")
    add_page(s, 14)
    add_notes(s, "The affordability constraint does not repeal these lessons; "
                 "it sharpens them. Under-scoping is often the more expensive "
                 "error — and disciplined single-airline governance (LGA) is "
                 "the existence proof for the delivery model DFW is running.")


# ---------------------------------------------------------------- slide 15
def slide_15():
    s = new_slide()
    add_title(s, 15, "The program clock: 2025–2030", w=10.4, h=0.62)
    add_chip(s, 15)
    g = s.shapes.add_group_shape()
    g.name = "slide-15-milestones"
    gs = g.shapes
    ax_l, ax_y, pitch = 1.0, 3.60, 2.32
    box(gs, ax_l, ax_y, 11.6, 0.022, fill=SLATE)
    for i in range(6):
        tx = ax_l + i * pitch
        box(gs, tx - 0.008, ax_y - 0.05, 0.016, 0.12, fill=SLATE)
        add_text(gs, tx - 0.35, ax_y + 0.10, 0.7, 0.24,
                 [P(R(str(2025 + i), size=12, color=SLATE),
                    align=PP_ALIGN.CENTER)])

    def mx(year_frac):
        return ax_l + (year_frac - 2025.0) * pitch

    events = [
        (mx(2025 + 7 / 12), "above-low", 1.20, 2.30,
         "Aug 2025 — six modules set"),
        (mx(2025 + 11 / 12), "below-1", 2.00, 2.30,
         "Dec 2025 — rebank announced (9 → 13 banks)"),
        (mx(2026 + 3 / 12), "above-high", 2.85, 2.30,
         "Apr 2026 — 13-bank schedule effective"),
        (mx(2026 + 6 / 12), "below-2", 3.60, 2.35,
         "Jul 2026 — premium siting announced: C / D / F"),
        (mx(2027), "above-low", 4.70, 2.10, "2027 — Phase 1 opens"),
        (mx(2030), "above-low", 10.15, 2.45,
         "2030 — program completes; fabrication-release window closes"),
    ]
    tiers = {"above-high": (1.80, 0.66), "above-low": (2.70, 0.72),
             "below-1": (3.94, 0.66), "below-2": (4.78, 0.66)}
    for cx, tier, bx, bw, txt in events:
        by, bh = tiers[tier]
        box(gs, cx - 0.06, ax_y - 0.06, 0.12, 0.12, fill=NAVY,
            shape=MSO_SHAPE.OVAL)
        if tier.startswith("above"):
            box(gs, cx - 0.008, by + bh, 0.016, ax_y - (by + bh), fill=GRAY)
        else:
            box(gs, cx - 0.008, ax_y + 0.07, 0.016, by - ax_y - 0.07,
                fill=GRAY)
        box(gs, bx, by, bw, bh, fill=FOG)
        add_text(gs, bx + 0.12, by + 0.08, bw - 0.24, bh - 0.14,
                 [P(R(txt, size=12, color=INK, bold=False), line=1.06)])
    band = box(gs, 1.0, 5.72, 11.6, 0.62, fill=None, line=SLATE, lw=1.5,
               dash=True)
    add_text(gs, 1.2, 5.86, 11.2, 0.36,
             [P(R("Per-module release-for-fabrication dates — unpublished; "
                  "obtaining them is the register's first disclosure "
                  "obligation (D0 / D2)", size=13, color=SLATE, italic=True),
                align=PP_ALIGN.CENTER)])
    set_alt(g, "Program timeline 2025 to 2030 with six fixed public "
               "milestones from the August 2025 module set to program "
               "completion in 2030, and a dashed band marking the "
               "unpublished per-module release-for-fabrication dates.")
    add_source(s, 15, "Innovation Next+ JV / Walsh Group (Aug 2025); AA "
                      "newsroom Dec 2025 and Jul 2026; DFW/AA joint release "
                      "May 2025.")
    add_page(s, 15)
    add_notes(s, "Every recommendation in the report is timed to the next "
                 "module's fabrication release, not to a pour date. The "
                 "unpopulated band is the point: the board cannot currently "
                 "see when its own decisions close.")


# ---------------------------------------------------------------- slide 16
def slide_16():
    s = new_slide()
    add_title(s, 16, "What the board watches quarterly, and what stops the "
                     "program", w=10.4)
    add_chip(s, 16)
    g = s.shapes.add_group_shape()
    g.name = "slide-16-indicators"
    gs = g.shapes
    add_text(gs, 0.55, 1.55, 6.0, 0.35,
             [P(R("Quarterly indicators", size=17, font=DISPLAY, color=NAVY,
                  bold=True))])
    indicators = [
        "% of fabrication releases preceded by a written register finding — "
        "target 100%",
        "Options countersigned vs retired (two-cycle rule)",
        "Aggregate register premium vs the Finance & Audit forward-CPE "
        "ceiling",
        "American connecting share holds a connecting majority",
        "American premium build-out on schedule: narrowbody premium seats "
        "~25% → ~40%; lie-flat +50%+ by decade end",
        "Skylink overnight single-loop performance",
    ]
    yy = 2.02
    for txt in indicators:
        box(gs, 0.62, yy + 0.075, 0.12, 0.12, fill=NAVY)
        add_text(gs, 0.92, yy, 11.7, 0.42,
                 [P(R(txt, size=13, color=INK), line=1.05)])
        yy += 0.44
    add_text(gs, 0.55, 4.80, 6.0, 0.35,
             [P(R("Stop conditions", size=17, font=DISPLAY, color=RED,
                  bold=True))])
    stops = [
        "Flagship-tier build-out announced on a published schedule",
        "Connecting share below the majority floor, two consecutive annual "
        "cycles",
        "U&L renegotiation opens before 2043",
    ]
    yy = 5.26
    for txt in stops:
        add_text(gs, 0.62, yy, 11.9, 0.36,
                 [P(R("STOP  ", size=13, color=RED, bold=True),
                    R(txt, size=13, color=INK))])
        yy += 0.40
    set_alt(g, "Governance panel listing six quarterly board indicators for "
               "the reversibility register and three stop conditions that "
               "halt or convert the program.")
    add_text(s.shapes, 0.55, 6.50, 12.23, 0.42,
             [P(R("The VP Terminal Development convenes any unwind review; "
                  "the CEO's office signs the unwind.", size=16, color=NAVY,
                  bold=True))],
             name="slide-16-owners")
    add_source(s, 16, "Indicator definitions per this report's decision "
                      "program; AA fleet targets per AA newsroom 2026; rebank "
                      "monitoring per View from the Wing 2026 (AA internal "
                      "data caveat applies).")
    add_page(s, 16)
    add_notes(s, "This is the page that keeps the register from becoming "
                 "decoration on a strategy deck: a percentage, a ceiling, and "
                 "three tripwires with named conveners.")


# ---------------------------------------------------------------- build
for fn in [slide_01, slide_02, slide_03, slide_04, slide_05, slide_06,
           slide_07, slide_08, slide_09, slide_10, slide_11, slide_12,
           slide_13, slide_14, slide_15, slide_16]:
    fn()

OUT = ("/Users/christiankessleriv/Repos/ai-council-mwaa/outputs/stage4/"
       "designing-for-2043-what-terminal-f-should-absorb-before-the-"
       "concrete-cures.pptx")
prs.save(OUT)
print("saved", OUT, "slides:", len(prs.slides._sldIdLst))
