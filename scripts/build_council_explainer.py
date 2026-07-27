#!/usr/bin/env python3
"""Build the Council v2 "How It Works" board explainer.

The deck is intentionally generated from code so the two distributed copies
remain byte-identical and future Council changes can be reflected without
manual PowerPoint drift.
"""

from __future__ import annotations

import argparse
import json
import os
import shutil
from pathlib import Path
from typing import Iterable

from jsonschema import validate
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


REPO_ROOT = Path(__file__).resolve().parents[1]
FINAL_DECK = REPO_ROOT / "final" / "PowerPoint" / "AI Research Council — How It Works.pptx"
ROOT_DECK = REPO_ROOT / "AI Research Council — How It Works.pptx"
LOGO = REPO_ROOT / "assets" / "council-logo.png"
BRAND_SCHEMA = REPO_ROOT / "assets" / "brand" / "visual-brief.schema.json"

SLIDE_W = 13.333333
SLIDE_H = 7.5

COLORS = {
    "navy": "0B2D4D",
    "blue": "2E84A5",
    "gold": "D4A24C",
    "slate": "415669",
    "white": "FFFFFF",
    "fog": "EDF3F6",
    "green": "24745C",
    "red": "A6413A",
    "ink": "17232D",
}

DISPLAY = "Georgia"
BODY = "Aptos"


def rgb(name_or_hex: str) -> RGBColor:
    value = COLORS.get(name_or_hex, name_or_hex).lstrip("#")
    return RGBColor.from_string(value)


def _set_alt_text(shape, description: str) -> None:
    nodes = shape._element.xpath(".//p:cNvPr")
    if nodes:
        nodes[0].set("descr", description)


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 18,
    font: str = BODY,
    color: str = "ink",
    bold: bool = False,
    italic: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    name: str | None = None,
    margin: float = 0,
    line_spacing: float = 1.0,
) -> object:
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if name:
        shape.name = name
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    paragraphs = text.split("\n")
    for index, line in enumerate(paragraphs):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = rgb(color)
    return shape


def add_shape(
    slide,
    shape_type: MSO_SHAPE,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str | None = None,
    line: str | None = None,
    line_width: float = 1.0,
    name: str | None = None,
) -> object:
    shape = slide.shapes.add_shape(
        shape_type, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    if name:
        shape.name = name
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(line_width)
    return shape


def add_line(
    slide,
    x1: float,
    y1: float,
    x2: float,
    y2: float,
    *,
    color: str = "slate",
    width: float = 1.0,
    name: str | None = None,
) -> object:
    shape = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    if name:
        shape.name = name
    shape.line.color.rgb = rgb(color)
    shape.line.width = Pt(width)
    return shape


def add_arrow(
    slide,
    x: float,
    y: float,
    w: float,
    h: float = 0.15,
    *,
    color: str = "blue",
) -> object:
    return add_shape(
        slide,
        MSO_SHAPE.RIGHT_ARROW,
        x,
        y,
        w,
        h,
        fill=color,
        line=None,
    )


def add_node(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    label: str,
    detail: str,
    *,
    accent: str = "navy",
    fill: str = "white",
    label_size: float = 18,
    detail_size: float = 16,
) -> None:
    add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        x,
        y,
        w,
        h,
        fill=fill,
        line=accent,
        line_width=1.5,
    )
    add_shape(
        slide,
        MSO_SHAPE.RECTANGLE,
        x,
        y,
        0.08,
        h,
        fill=accent,
        line=None,
    )
    add_text(
        slide,
        label,
        x + 0.22,
        y + 0.15,
        w - 0.35,
        0.35,
        size=label_size,
        font=DISPLAY,
        color=accent,
        bold=True,
    )
    add_text(
        slide,
        detail,
        x + 0.22,
        y + 0.58,
        w - 0.35,
        h - 0.7,
        size=detail_size,
        color="slate",
    )


def add_title(slide, slide_number: int, title: str, source: str, notes: str) -> None:
    add_text(
        slide,
        title,
        0.55,
        0.30,
        12.05,
        0.68,
        size=36,
        font=DISPLAY,
        color="ink",
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
        name=f"slide-{slide_number}-title",
    )
    add_line(slide, 0.55, 1.16, 12.78, 1.16, color="slate", width=0.7)
    add_line(slide, 0.55, 6.92, 12.78, 6.92, color="slate", width=0.55)
    add_text(
        slide,
        source,
        0.55,
        7.01,
        11.45,
        0.22,
        size=9.2,
        color="slate",
        name=f"slide-{slide_number}-source",
    )
    add_text(
        slide,
        f"{slide_number:02d}",
        12.25,
        7.00,
        0.52,
        0.24,
        size=9.2,
        color="slate",
        align=PP_ALIGN.RIGHT,
        name=f"slide-{slide_number}-number",
    )
    slide.notes_slide.notes_text_frame.text = notes


def add_kicker(
    slide, text: str, x: float, y: float, w: float, *, color: str = "navy"
) -> None:
    add_text(
        slide,
        text.upper(),
        x,
        y,
        w,
        0.28,
        size=16,
        color=color,
        bold=True,
    )


def _visual_brief() -> dict:
    communication_job = (
        "By the end, airport executives and board members should understand how "
        "the Council produces a defensible decision package because independent "
        "research becomes curated evidence, survives two distinct challenges, is "
        "verified in fresh context, and cannot release without traceable QA."
    )
    slides = [
        (
            "Orient the audience to the system's promise.",
            "How the AI Research Council works",
            [],
            "Minimal branded cover",
            "Large editorial title, current-system subtitle, MWAA Strat Ops mark.",
            "Current Council v2 repository and operating contract.",
        ),
        (
            "Show the complete governed production line before explaining its parts.",
            "One question moves through a governed line",
            ["REF-PIPELINE"],
            "Cumulative process flow",
            "Six work handoffs with two visible human approvals.",
            "Council v2 operating contract and executable pipeline.",
        ),
        (
            "Establish the roster's scale and separation of responsibilities.",
            "The 54-agent roster separates depth from control",
            ["REF-REGISTRY"],
            "Proportional stacked bar",
            "20 airport research, 18 supplemental, and 16 process roles; note that runs seat a subset.",
            "Executable agent registry reviewed 23 July 2026.",
        ),
        (
            "Explain why airport context precedes research.",
            "Airport context comes before analysis",
            ["REF-CONTEXT"],
            "Context-to-swarm flow",
            "Six airport decision domains form one shared starting record.",
            "Airport Context Builder charter and operating contract.",
        ),
        (
            "Explain parallel independence and the Curator's role.",
            "A parallel research swarm protects disagreement",
            ["REF-RESEARCH", "REF-CURATION"],
            "Parallel lanes converging after the work",
            "Selected lenses work without cross-reading, then feed the Evidence Curator.",
            "Research contract, evidence schema, and Curator charter.",
        ),
        (
            "Make evidence custody memorable and concrete.",
            "Every released claim keeps a chain of custody",
            ["REF-EVIDENCE", "REF-LINEAGE", "REF-GATE"],
            "Signature claim-chain diagram",
            "Primary source → evidence ledger → reader claim → claim lineage → release gate.",
            "Evidence, claim-lineage, and publication-gate contracts.",
        ),
        (
            "Show how creative range enters without relaxing evidence standards.",
            "Creative framing keeps the story open",
            ["REF-CREATIVE", "REF-STRATEGY"],
            "Three narrative lines resolving to one draft",
            "Board-decision, counterintuitive, and operating frames converge at the Strategist.",
            "Creative Director and Strategist charters.",
        ),
        (
            "Differentiate the two adversarial review jobs.",
            "Two adversaries test different failure modes",
            ["REF-PROSECUTOR", "REF-EXEC-REVIEW"],
            "Dual-lane adversarial loop",
            "Evidence integrity above; airport executability below; Strategist revises after each.",
            "Evidence Prosecutor and Airport Executive Reviewer charters.",
        ),
        (
            "Separate editorial craft from independent source verification.",
            "Fresh eyes verify facts after prose is finished",
            ["REF-EDITORIAL", "REF-VERIFY", "REF-MODELS"],
            "Editorial-to-verification sequence",
            "Editor → Humanizer → Source Verifier in fresh Sonnet context → lineage gate → human approval.",
            "Editorial charters, verifier contract, and current model routing.",
        ),
        (
            "Explain routing and budget boundaries plainly.",
            "Model routing adds fresh eyes and hard cost limits",
            ["REF-MODELS", "REF-BUDGET", "REF-OPENAI"],
            "Separated budget lanes",
            "Claude run ceiling and per-call caps on one side; OpenAI Deep Research billing on the other.",
            "Current model routing, runtime budget implementation, and tests.",
        ),
        (
            "Show how an operator observes and controls a long run.",
            "Live agent telemetry makes a run observable",
            ["REF-EVENTS", "REF-WEBAPP"],
            "Event rail",
            "Role, artifact, evidence, cost, gate, and approval events flow to one operator.",
            "Event stream and run-screen reconnect behavior.",
        ),
        (
            "Close on the release standard and accountable action.",
            "Release requires provenance and rendered QA",
            ["REF-ART", "REF-PUBLISH", "REF-MANIFEST"],
            "Production and release gate",
            "Art direction → Word/PPT → structural QA → full render inspection; four provenance records travel with release.",
            "Art direction, Office QA, manifest, archive, and publishing contracts.",
        ),
    ]
    return {
        "schema_version": "2.0",
        "communication_job": communication_job,
        "audience": "Airport executives, board members, and accountable report sponsors",
        "decision": (
            "Use the provenance packet, human checkpoints, and rendered QA—not "
            "fluent prose alone—as the basis for approving a Council release."
        ),
        "decision_owner": (
            "The airport executive or board sponsor accountable for the decision"
        ),
        "approval_path": [
            "Council sponsor confirms the research question and decision",
            "Accountable airport executives approve the challenged argument",
            "Release owner accepts the provenance packet and rendered QA",
        ],
        "first_90_day_action": (
            "Use the Council v2 workflow for one consequential airport decision "
            "and review the complete release record at the final checkpoint."
        ),
        "success_measures": [
            "Every material released claim resolves to verified evidence",
            "Both adversarial reviews close without an unresolved critical issue",
            "Every Word page and PowerPoint slide passes full-size visual inspection",
        ],
        "deck_mode": "board_decision",
        "visual_thesis": (
            "The Council is a governed chain of custody for airport decisions: "
            "independent inquiry becomes evidence, evidence becomes a challenged "
            "claim, and only verified claims enter inspected deliverables."
        ),
        "signature_visual": {
            "slide_number": 6,
            "concept": "Claim chain of custody",
            "visual_type": "five-link process diagram",
            "visual_spec": (
                "A single horizontal chain connecting primary source, evidence "
                "ledger, reader claim, claim lineage, and deterministic release gate."
            ),
            "evidence_ids": ["REF-EVIDENCE", "REF-LINEAGE", "REF-GATE"],
            "source_note": "Council v2 evidence and release contracts.",
            "asset_source": "Native PowerPoint shapes; no external imagery.",
        },
        "brand_profile": {
            "name": "Transform Airports AI Research Council",
            "version": "2.0",
            "palette": list(COLORS.values()),
            "display_font": DISPLAY,
            "body_font": BODY,
        },
        "slides": [
            {
                "slide_number": number,
                "narrative_job": narrative_job,
                "headline": headline,
                "evidence_ids": evidence_ids,
                "visual_type": visual_type,
                "visual_spec": visual_spec,
                "source_note": source_note,
                "density_budget": "70 visible words maximum",
                "speaker_note": source_note,
            }
            for number, (
                narrative_job,
                headline,
                evidence_ids,
                visual_type,
                visual_spec,
                source_note,
            ) in enumerate(slides, 1)
        ],
        "report_visuals": [],
        "source_appendix": {
            "treatment": "Readable 9.2 pt source footer on every evidence-bearing slide; exact repository references in speaker notes."
        },
        "accessibility_checks": [
            "16:9 canvas",
            "Georgia display and Aptos body",
            "Body copy 16 pt or larger",
            "Source notes 9.2 pt",
            "Meaning is labeled in words, not color alone",
            "Material image has alt text",
        ],
        "asset_requests": [
            {
                "asset": "MWAA Strategy and Operational Performance logo",
                "source": "repository brand asset",
                "status": "reused",
            }
        ],
    }


def _build_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    prs.core_properties.title = "AI Research Council — How It Works"
    prs.core_properties.subject = "Current Council v2 process explainer"
    prs.core_properties.author = "Office of Strategy and Operational Performance"
    prs.core_properties.keywords = (
        "airport, research council, evidence, claim lineage, quality assurance"
    )
    prs.core_properties.comments = (
        "Board-deck explainer generated from the current Council v2 repository."
    )
    blank = prs.slide_layouts[6]

    # 1 — cover
    slide = prs.slides.add_slide(blank)
    add_text(
        slide,
        "TRANSFORM AIRPORTS / COUNCIL V2",
        0.55,
        0.44,
        5.7,
        0.30,
        size=16,
        color="navy",
        bold=True,
        name="cover-eyebrow",
    )
    picture = slide.shapes.add_picture(str(LOGO), Inches(8.65), Inches(0.34), width=Inches(4.05))
    picture.name = "MWAA Strategy and Operational Performance logo"
    _set_alt_text(
        picture,
        "Office of Strategy and Operational Performance, Metropolitan Washington Airports Authority.",
    )
    add_line(slide, 0.55, 1.96, 12.78, 1.96, color="slate", width=0.8)
    add_text(
        slide,
        "How the AI Research\nCouncil works",
        0.55,
        2.42,
        8.3,
        1.78,
        size=56,
        font=DISPLAY,
        color="ink",
        bold=True,
        line_spacing=0.92,
        name="cover-title",
    )
    add_text(
        slide,
        (
            "One contested thesis in. Airport context, independent evidence, "
            "two distinct challenges, source verification and a governed "
            "executive packet out."
        ),
        0.58,
        5.12,
        11.6,
        0.84,
        size=24,
        color="slate",
        name="cover-subtitle",
    )
    add_line(slide, 0.58, 6.56, 1.18, 6.56, color="blue", width=3)
    add_text(
        slide,
        "CODE-REVIEWED CURRENT SYSTEM",
        1.34,
        6.40,
        4.25,
        0.34,
        size=16,
        color="navy",
        bold=True,
    )
    add_text(
        slide,
        "23 July 2026",
        10.67,
        6.40,
        2.05,
        0.34,
        size=16,
        color="slate",
        align=PP_ALIGN.RIGHT,
    )
    slide.notes_slide.notes_text_frame.text = (
        "Communication job: By the end, airport executives and board members "
        "should understand how the Council produces a defensible decision package "
        "because independent research becomes curated evidence, survives two "
        "distinct challenges, is verified in fresh context, and cannot release "
        "without traceable QA.\n\nSources: README.md; prompts/orchestration.md; "
        "cli/orchestrator.py; cli/agents.py; council.toml. Repository reviewed "
        "23 July 2026."
    )

    # 2 — governed line
    slide = prs.slides.add_slide(blank)
    add_title(
        slide,
        2,
        "One question moves through a governed line",
        "Source: Council v2 operating contract and executable pipeline · repository review, 23 Jul 2026",
        "Sources: README.md “How a run works”; prompts/orchestration.md; "
        "cli/orchestrator.py PIPELINE_DEFINITION and run_pipeline().",
    )
    add_text(
        slide,
        "No model researches, writes and approves its own answer.",
        0.55,
        1.35,
        9.8,
        0.42,
        size=20,
        color="slate",
    )
    y = 3.0
    node_xs = [0.55, 2.64, 4.73, 6.82, 8.91, 11.00]
    for x in (2.42, 4.51, 6.60, 8.69, 10.78):
        add_arrow(slide, x, y + 0.67, 0.34, 0.14, color="blue")
    nodes = [
        ("Airport context", "shared facts", "navy"),
        ("Independent research", "separate briefs", "blue"),
        ("Curate + frame", "ledger + options", "green"),
        ("Draft + challenge", "two reviews", "red"),
        ("Edit + verify", "lineage + gate", "navy"),
        ("Produce + inspect", "Word / PPT", "green"),
    ]
    for x, (label, detail, accent) in zip(node_xs, nodes):
        add_shape(
            slide,
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x,
            y,
            1.78,
            1.55,
            fill="white",
            line=accent,
            line_width=2,
        )
        add_text(
            slide,
            label,
            x + 0.12,
            y + 0.27,
            1.54,
            0.52,
            size=18,
            font=DISPLAY,
            color=accent,
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
        add_text(
            slide,
            detail,
            x + 0.08,
            y + 1.04,
            1.62,
            0.28,
            size=16,
            color="slate",
            align=PP_ALIGN.CENTER,
        )
    for x in (8.35, 10.44):
        add_shape(
            slide,
            MSO_SHAPE.DIAMOND,
            x,
            4.78,
            0.36,
            0.36,
            fill="gold",
            line="navy",
            line_width=1,
        )
    add_text(
        slide,
        "HUMAN REVIEW",
        7.31,
        5.19,
        1.72,
        0.30,
        size=16,
        color="navy",
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "CHECKPOINT 2",
        9.42,
        5.19,
        1.70,
        0.30,
        size=16,
        color="navy",
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        0.55,
        5.82,
        12.23,
        0.62,
        fill="fog",
        line=None,
    )
    add_text(
        slide,
        "Typed artifacts connect each handoff. Failed contracts stop the line; valid completed work resumes.",
        0.80,
        5.97,
        11.72,
        0.31,
        size=17,
        color="navy",
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    # 3 — roster
    slide = prs.slides.add_slide(blank)
    add_title(
        slide,
        3,
        "The 54-agent roster separates depth from control",
        "Source: Executable agent registry · 20 airport research + 18 supplemental + 16 process · 23 Jul 2026",
        "Sources: cli/agents.py RESEARCH_AGENT_NAMES, SUPPLEMENTAL_AGENT_NAMES, "
        "PROCESS_AGENT_NAMES; .claude/agents/ (54 current definitions); "
        "tests/test_council_v2_registry.py.",
    )
    add_text(
        slide,
        "54",
        0.58,
        1.55,
        2.15,
        1.05,
        size=76,
        font=DISPLAY,
        color="navy",
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "defined agents",
        0.58,
        2.57,
        2.15,
        0.34,
        size=20,
        color="slate",
        align=PP_ALIGN.CENTER,
    )
    add_kicker(slide, "Defined agents · current repository · 23 Jul 2026", 3.05, 1.55, 8.9)
    bar_x, bar_y, bar_w, bar_h = 3.05, 2.08, 9.23, 0.92
    segments = [
        (20, "navy", "20  AIRPORT RESEARCH"),
        (18, "blue", "18  SUPPLEMENTAL"),
        (16, "gold", "16  PROCESS"),
    ]
    cursor = bar_x
    for value, color, label in segments:
        width = bar_w * value / 54
        add_shape(
            slide,
            MSO_SHAPE.RECTANGLE,
            cursor,
            bar_y,
            width,
            bar_h,
            fill=color,
            line="white",
            line_width=1.2,
        )
        add_text(
            slide,
            label,
            cursor + 0.08,
            bar_y + 0.27,
            width - 0.16,
            0.32,
            size=16,
            color="navy" if color == "gold" else "white",
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        cursor += width
    legend = [
        (
            "AIRPORT RESEARCH",
            "Economics, operations, engineering, airlines, safety and governance.",
            "navy",
        ),
        (
            "SUPPLEMENTAL",
            "Outside thinkers, seated only when the question benefits.",
            "blue",
        ),
        (
            "PROCESS",
            "Context, curation, writing, review, verification, design and production.",
            "gold",
        ),
    ]
    for i, (label, detail, color) in enumerate(legend):
        y0 = 3.55 + i * 0.77
        add_line(slide, 3.05, y0 + 0.13, 3.45, y0 + 0.13, color=color, width=4)
        add_text(
            slide,
            label,
            3.64,
            y0,
            2.25,
            0.32,
            size=17,
            color=color if color != "gold" else "navy",
            bold=True,
        )
        add_text(
            slide,
            detail,
            5.92,
            y0 - 0.02,
            6.20,
            0.46,
            size=17,
            color="slate",
        )
    add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        0.58,
        6.08,
        12.03,
        0.50,
        fill="fog",
        line=None,
    )
    add_text(
        slide,
        "Each run seats only the lenses it needs; the manifest records what actually ran.",
        0.88,
        6.20,
        11.43,
        0.25,
        size=17,
        color="navy",
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    # 4 — context builder
    slide = prs.slides.add_slide(blank)
    add_title(
        slide,
        4,
        "Airport context comes before analysis",
        "Source: Airport Context Builder charter and Council v2 operating contract · 23 Jul 2026",
        "Sources: .claude/agents/airport-context-builder.md; "
        "prompts/orchestration.md “Public stage 1”; cli/orchestrator.py "
        "airport-context pipeline step.",
    )
    add_text(
        slide,
        "It does not argue the thesis. It establishes the decision environment every researcher may read.",
        0.55,
        1.35,
        11.8,
        0.50,
        size=20,
        color="slate",
    )
    add_line(slide, 2.35, 3.98, 10.62, 3.98, color="blue", width=2.2)
    add_arrow(slide, 9.88, 3.90, 0.76, 0.16, color="blue")
    add_shape(
        slide,
        MSO_SHAPE.OVAL,
        0.67,
        2.76,
        2.10,
        2.10,
        fill="navy",
        line=None,
    )
    add_text(
        slide,
        "AIRPORT\nCONTEXT\nBUILDER",
        0.85,
        3.18,
        1.74,
        0.98,
        size=21,
        font=DISPLAY,
        color="white",
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    domains = [
        ("Governance", 3.10, 2.25),
        ("Finance", 5.18, 2.25),
        ("Capital", 7.26, 2.25),
        ("Airlines", 3.10, 4.60),
        ("Regulation", 5.18, 4.60),
        ("Operations", 7.26, 4.60),
    ]
    for label, x, y0 in domains:
        add_shape(
            slide,
            MSO_SHAPE.OVAL,
            x,
            y0,
            1.63,
            0.72,
            fill="white",
            line="blue",
            line_width=1.5,
        )
        add_text(
            slide,
            label,
            x + 0.08,
            y0 + 0.20,
            1.47,
            0.28,
            size=17,
            color="navy",
            bold=True,
            align=PP_ALIGN.CENTER,
        )
    add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        9.91,
        2.61,
        2.78,
        2.75,
        fill="fog",
        line="green",
        line_width=1.5,
    )
    add_text(
        slide,
        "SHARED RECORD",
        10.17,
        2.90,
        2.26,
        0.55,
        size=18,
        color="green",
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "Supplied facts\nAuthoritative records\nNamed gaps\nCurrent constraints",
        10.17,
        3.72,
        2.26,
        1.20,
        size=17,
        color="slate",
        align=PP_ALIGN.CENTER,
        line_spacing=1.15,
    )
    add_text(
        slide,
        "Context removes false assumptions without forcing consensus.",
        2.98,
        5.82,
        8.45,
        0.54,
        size=18,
        color="navy",
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    # 5 — swarm and curation
    slide = prs.slides.add_slide(blank)
    add_title(
        slide,
        5,
        "A parallel research swarm protects disagreement",
        "Source: Research contract, evidence schema, and Evidence Curator charter · 23 Jul 2026",
        "Sources: prompts/research-contract.md; prompts/orchestration.md; "
        "cli/orchestrator.py run_stage1(); cli/evidence.py; "
        ".claude/agents/evidence-curator.md.",
    )
    add_text(
        slide,
        "Plain English: specialists investigate the same question at the same time without reading one another’s conclusions.",
        0.55,
        1.34,
        12.0,
        0.48,
        size=20,
        color="slate",
    )
    add_arrow(slide, 2.13, 3.75, 0.72, 0.18, color="blue")
    add_arrow(slide, 8.86, 3.75, 0.78, 0.18, color="green")
    add_shape(
        slide,
        MSO_SHAPE.OVAL,
        0.58,
        3.05,
        1.54,
        1.54,
        fill="navy",
        line=None,
    )
    add_text(
        slide,
        "THESIS",
        0.78,
        3.43,
        1.14,
        0.64,
        size=20,
        font=DISPLAY,
        color="white",
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        2.84,
        2.18,
        5.98,
        3.58,
        fill="fog",
        line="blue",
        line_width=1.2,
    )
    add_kicker(slide, "Selected independent lenses", 3.10, 2.43, 5.2, color="navy")
    labels = [
        ("Economics", 3.12, 3.09),
        ("Operations", 4.90, 3.09),
        ("Airfield", 6.68, 3.09),
        ("Airlines", 3.12, 4.38),
        ("Safety", 4.90, 4.38),
        ("Contrarian", 6.68, 4.38),
    ]
    for label, x, y0 in labels:
        add_shape(
            slide,
            MSO_SHAPE.OVAL,
            x,
            y0,
            1.36,
            0.78,
            fill="white",
            line="blue" if label != "Contrarian" else "red",
            line_width=1.5,
        )
        add_text(
            slide,
            label,
            x + 0.08,
            y0 + 0.22,
            1.20,
            0.30,
            size=16,
            color="navy" if label != "Contrarian" else "red",
            bold=True,
            align=PP_ALIGN.CENTER,
        )
    add_text(
        slide,
        "NO CROSS-READING",
        4.39,
        5.24,
        2.72,
        0.30,
        size=16,
        color="red",
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_node(
        slide,
        9.65,
        2.62,
        3.03,
        2.70,
        "Evidence Curator",
        "Deduplicates\nPreserves conflict\nRanks what matters\nCloses narrow gaps",
        accent="green",
        fill="white",
        label_size=20,
        detail_size=17,
    )
    add_text(
        slide,
        "Independence ends at curation—after each lens has left a separate record.",
        2.99,
        6.08,
        8.66,
        0.38,
        size=18,
        color="navy",
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    # 6 — signature claim chain
    slide = prs.slides.add_slide(blank)
    add_title(
        slide,
        6,
        "Every released claim keeps a chain of custody",
        "Source: Evidence, claim-lineage, and deterministic publication-gate contracts · 23 Jul 2026",
        "Sources: cli/evidence.py; cli/quality_gate.py; prompts/orchestration.md "
        "“Edit, verify, and release gate”; docs/methodology.md “Evidence is a "
        "typed artifact”; cli/orchestrator.py run_quality_gate_with_remediation().",
    )
    add_text(
        slide,
        "A reader-facing claim must stay connected to the source that proves it.",
        0.55,
        1.35,
        11.8,
        0.46,
        size=20,
        color="slate",
    )
    add_kicker(slide, "Signature visual · claim chain of custody", 0.55, 2.04, 5.2, color="navy")
    chain_y = 2.74
    node_w = 2.05
    node_h = 2.20
    xs = [0.55, 3.06, 5.57, 8.08, 10.59]
    colors = ["navy", "blue", "slate", "green", "red"]
    labels = [
        ("PRIMARY SOURCE", "What it says"),
        ("EVIDENCE LEDGER", "ID + caveats"),
        ("READER CLAIM", "Exact words + footnote"),
        ("CLAIM LINEAGE", "Source checked + draft hash"),
        ("RELEASE GATE", "Pass—or stop"),
    ]
    for i in range(4):
        add_arrow(slide, xs[i] + node_w + 0.10, chain_y + 0.99, 0.35, 0.16, color="blue")
    for x, color, (label, detail) in zip(xs, colors, labels):
        add_shape(
            slide,
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x,
            chain_y,
            node_w,
            node_h,
            fill="white",
            line=color,
            line_width=2,
        )
        add_shape(
            slide,
            MSO_SHAPE.OVAL,
            x + 0.79,
            chain_y + 0.20,
            0.47,
            0.47,
            fill=color,
            line=None,
        )
        add_text(
            slide,
            label,
            x + 0.16,
            chain_y + 0.86,
            node_w - 0.32,
            0.49,
            size=17,
            color=color,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            detail,
            x + 0.16,
            chain_y + 1.38,
            node_w - 0.32,
            0.60,
            size=16,
            color="slate",
            align=PP_ALIGN.CENTER,
        )
    add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        1.38,
        5.56,
        10.58,
        0.70,
        fill="navy",
        line=None,
        name="SIGNATURE VISUAL — claim chain of custody",
    )
    add_text(
        slide,
        "Every link must match: claim • footnote • evidence • source • draft.",
        1.72,
        5.75,
        9.90,
        0.32,
        size=18,
        color="white",
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    # 7 — creative framing
    slide = prs.slides.add_slide(blank)
    add_title(
        slide,
        7,
        "Creative framing keeps the story open",
        "Source: Creative Director and Strategist charters; Council v2 synthesis contract · 23 Jul 2026",
        "Sources: .claude/agents/creative-director.md; "
        ".claude/agents/strategist.md; cli/orchestrator.py _stage2_prompts() and "
        "PIPELINE_DEFINITION.",
    )
    add_text(
        slide,
        "The Creative Director proposes truthful ways to tell the same evidence before the report’s structure hardens.",
        0.55,
        1.34,
        11.9,
        0.48,
        size=20,
        color="slate",
    )
    add_line(slide, 4.50, 2.85, 7.00, 3.79, color="blue", width=1.8)
    add_line(slide, 4.50, 3.79, 7.00, 3.79, color="red", width=1.8)
    add_line(slide, 4.50, 4.73, 7.00, 3.79, color="green", width=1.8)
    frames = [
        ("BOARD DECISION", "What must be approved?", 2.25, "blue"),
        ("COUNTERINTUITIVE", "What changes the reader’s mind?", 3.43, "red"),
        ("OPERATIONAL", "What happens at the airport?", 4.61, "green"),
    ]
    for label, detail, y0, color in frames:
        add_line(slide, 0.73, y0 + 0.23, 1.18, y0 + 0.23, color=color, width=4)
        add_text(
            slide,
            label,
            1.42,
            y0,
            2.45,
            0.35,
            size=18,
            color=color,
            bold=True,
        )
        add_text(
            slide,
            detail,
            1.42,
            y0 + 0.43,
            3.10,
            0.36,
            size=17,
            color="slate",
        )
    add_shape(
        slide,
        MSO_SHAPE.OVAL,
        7.00,
        2.64,
        2.30,
        2.30,
        fill="navy",
        line=None,
    )
    add_text(
        slide,
        "STRATEGIST",
        7.25,
        3.31,
        1.80,
        0.39,
        size=21,
        font=DISPLAY,
        color="white",
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "chooses the frame the evidence can carry",
        7.25,
        3.76,
        1.80,
        0.63,
        size=16,
        color="white",
        align=PP_ALIGN.CENTER,
    )
    add_arrow(slide, 9.45, 3.68, 0.88, 0.20, color="blue")
    add_node(
        slide,
        10.43,
        2.73,
        2.26,
        2.15,
        "Argument",
        "One narrative spine\nCurated evidence\nExplicit decision",
        accent="blue",
        fill="white",
        label_size=21,
        detail_size=17,
    )
    add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        2.29,
        5.74,
        8.75,
        0.62,
        fill="fog",
        line=None,
    )
    add_text(
        slide,
        "Creativity changes the route through evidence—not the evidence itself.",
        2.58,
        5.90,
        8.17,
        0.31,
        size=18,
        color="navy",
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    # 8 — adversarial loop
    slide = prs.slides.add_slide(blank)
    add_title(
        slide,
        8,
        "Two adversaries test different failure modes",
        "Source: Evidence Prosecutor and Airport Executive Reviewer charters · 23 Jul 2026",
        "Sources: .claude/agents/evidence-prosecutor.md; "
        ".claude/agents/airport-executive-reviewer.md; prompts/orchestration.md "
        "“Creative framing and adversarial synthesis”; cli/orchestrator.py "
        "PIPELINE_DEFINITION.",
    )
    add_text(
        slide,
        "Adversarial synthesis loop: the writer must answer each critic in a new draft.",
        0.55,
        1.34,
        11.9,
        0.48,
        size=20,
        color="slate",
    )
    add_line(slide, 6.67, 2.18, 6.67, 5.84, color="slate", width=0.8)
    add_kicker(slide, "Evidence integrity", 0.72, 2.15, 2.4, color="red")
    add_text(
        slide,
        "EVIDENCE PROSECUTOR",
        0.72,
        2.63,
        4.86,
        0.45,
        size=24,
        font=DISPLAY,
        color="red",
        bold=True,
    )
    add_text(
        slide,
        "Sources\nArithmetic\nCausality\nCounterevidence",
        0.72,
        3.33,
        2.15,
        1.55,
        size=18,
        color="slate",
        line_spacing=1.22,
    )
    add_text(
        slide,
        "Finds where the conclusion outruns the record.",
        3.06,
        3.37,
        2.86,
        1.15,
        size=20,
        color="ink",
        bold=True,
    )
    add_kicker(slide, "Airport executability", 7.23, 2.15, 2.8, color="navy")
    add_text(
        slide,
        "AIRPORT EXECUTIVE REVIEWER",
        7.23,
        2.63,
        5.34,
        0.45,
        size=24,
        font=DISPLAY,
        color="navy",
        bold=True,
    )
    add_text(
        slide,
        "Authority\nAirlines + funding\nProcurement + staffing\nPeak-hour operations",
        7.23,
        3.33,
        2.45,
        1.55,
        size=18,
        color="slate",
        line_spacing=1.22,
    )
    add_text(
        slide,
        "Finds where a sound idea fails in airport reality.",
        9.86,
        3.37,
        2.48,
        1.15,
        size=20,
        color="ink",
        bold=True,
    )
    add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        2.45,
        5.38,
        8.44,
        0.80,
        fill="navy",
        line=None,
    )
    add_text(
        slide,
        "STRATEGIST DRAFT → CHALLENGE → REVISION → CHALLENGE → REVISION → HUMAN REVIEW",
        2.72,
        5.62,
        7.90,
        0.31,
        size=17,
        color="white",
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    # 9 — editorial and verifier
    slide = prs.slides.add_slide(blank)
    add_title(
        slide,
        9,
        "Fresh eyes verify facts after prose is finished",
        "Source: Editorial and Source Verifier charters; current model routing · 23 Jul 2026",
        "Sources: .claude/agents/editor.md; .claude/agents/humanizer.md; "
        ".claude/agents/fact-checker.md; cli/config.py DEFAULT_MODELS; "
        "council.toml; cli/orchestrator.py _stage3_prompts() and run_stage3().",
    )
    add_text(
        slide,
        "Editor cuts; Humanizer improves voice. Neither may change facts.",
        0.55,
        1.34,
        11.9,
        0.48,
        size=20,
        color="slate",
    )
    xs = [0.60, 3.03, 5.46, 8.54, 11.00]
    widths = [1.93, 1.93, 2.58, 1.96, 1.70]
    for x in (2.54, 4.97, 8.04, 10.51):
        add_arrow(slide, x, 3.48, 0.40, 0.16, color="blue")
    add_node(slide, xs[0], 2.52, widths[0], 2.20, "Editor", "Cuts repetition\nTightens logic", accent="slate")
    add_node(slide, xs[1], 2.52, widths[1], 2.20, "Humanizer", "Improves voice\nAdds no facts", accent="blue")
    add_node(
        slide,
        xs[2],
        2.18,
        widths[2],
        2.88,
        "Source Verifier",
        "Opens sources\nChecks exact wording\nVerifies · qualifies\ncorrects · removes",
        accent="navy",
        fill="fog",
        label_size=22,
        detail_size=17,
    )
    add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        5.66,
        4.64,
        2.18,
        0.31,
        fill="gold",
        line=None,
    )
    add_text(
        slide,
        "FRESH SONNET CONTEXT",
        5.73,
        4.70,
        2.04,
        0.18,
        size=9.2,
        color="navy",
        bold=True,
        align=PP_ALIGN.CENTER,
        name="slide-9-verifier-source",
    )
    add_node(
        slide,
        xs[3],
        2.52,
        widths[3],
        2.20,
        "Claim lineage",
        "Binds claim to source\nGate blocks defects",
        accent="red",
        label_size=17,
    )
    add_node(slide, xs[4], 2.52, widths[4], 2.20, "Human", "Approves\nrelease", accent="green")
    add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        1.28,
        5.63,
        10.76,
        0.64,
        fill="navy",
        line=None,
    )
    add_text(
        slide,
        "Fresh Sonnet verifies the final draft—not the writer or polisher.",
        1.61,
        5.79,
        10.10,
        0.32,
        size=18,
        color="white",
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    # 10 — model routing and budgets
    slide = prs.slides.add_slide(blank)
    add_title(
        slide,
        10,
        "Model routing adds fresh eyes and hard cost limits",
        "Source: Current model routing and runtime budget implementation/tests · 23 Jul 2026",
        "Sources: council.toml; cli/config.py; cli/orchestrator.py CostTally, "
        "_run_agent(), _run_openai_deep_research(); "
        "tests/test_orchestrator_runtime.py; docs/how-to-run.md budget note.",
    )
    add_text(
        slide,
        "Multi-model orchestration gives different models different jobs, so the model that writes the argument does not grade all of its own work.",
        0.55,
        1.27,
        12.0,
        0.68,
        size=19,
        color="slate",
    )
    add_line(slide, 8.64, 2.10, 8.64, 6.19, color="slate", width=0.8)
    add_kicker(slide, "Claude calls · one shared ceiling", 0.67, 2.12, 4.6, color="navy")
    add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        0.67,
        2.67,
        2.12,
        2.85,
        fill="navy",
        line=None,
    )
    add_text(
        slide,
        "RUN\nCEILING",
        0.94,
        3.24,
        1.58,
        0.80,
        size=27,
        font=DISPLAY,
        color="white",
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "budget reserved before each call",
        0.94,
        4.28,
        1.58,
        0.70,
        size=16,
        color="white",
        align=PP_ALIGN.CENTER,
    )
    add_arrow(slide, 2.90, 3.93, 0.73, 0.18, color="blue")
    calls = [
        ("RESEARCH", 3.70, 2.66, "blue"),
        ("REVIEW", 3.70, 3.77, "red"),
        ("VERIFY", 3.70, 4.88, "green"),
    ]
    for label, x, y0, color in calls:
        add_shape(
            slide,
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x,
            y0,
            3.86,
            0.82,
            fill="white",
            line=color,
            line_width=1.6,
        )
        add_text(
            slide,
            label,
            x + 0.20,
            y0 + 0.19,
            1.28,
            0.31,
            size=17,
            color=color,
            bold=True,
        )
        add_text(
            slide,
            "PER-CALL CAP",
            x + 1.46,
            y0 + 0.18,
            2.18,
            0.35,
            size=16,
            color="slate",
            bold=True,
            align=PP_ALIGN.RIGHT,
        )
    add_text(
        slide,
        "Ceiling reached → stop between calls; keep valid completed work.",
        0.70,
        5.95,
        7.42,
        0.45,
        size=17,
        color="navy",
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_kicker(slide, "Optional second-provider lens", 9.08, 2.12, 3.8, color="gold")
    add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        9.08,
        2.67,
        3.64,
        2.85,
        fill="fog",
        line="gold",
        line_width=2,
    )
    add_text(
        slide,
        "OPENAI\nDEEP RESEARCH",
        9.41,
        3.15,
        2.98,
        0.76,
        size=25,
        font=DISPLAY,
        color="navy",
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "Separate account\nSeparate bill\nOutside Claude tally",
        9.41,
        4.15,
        2.98,
        0.96,
        size=17,
        color="slate",
        align=PP_ALIGN.CENTER,
        line_spacing=1.15,
    )

    # 11 — live telemetry
    slide = prs.slides.add_slide(blank)
    add_title(
        slide,
        11,
        "Live agent telemetry makes a run observable",
        "Source: Event stream, run screen, and reconnect behavior · repository review, 23 Jul 2026",
        "Sources: cli/events.py WebSink; cli/server.py WebSocket stream and active "
        "run registry; cli/webapp/app.js event handling; docs/how-to-run.md "
        "“Live agent telemetry”; tests/test_events.py.",
    )
    add_text(
        slide,
        "Plain English: the system streams operating state and receipts while the Council is working.",
        0.55,
        1.34,
        11.9,
        0.48,
        size=20,
        color="slate",
    )
    rail_y = 3.48
    add_line(slide, 1.16, rail_y, 10.44, rail_y, color="blue", width=2.4)
    add_arrow(slide, 10.03, rail_y - 0.09, 0.52, 0.18, color="blue")
    events = [
        ("ROLE", "who is active", "blue"),
        ("ARTIFACT", "what is being made", "navy"),
        ("EVIDENCE", "what changed", "green"),
        ("COST", "call + running total", "gold"),
        ("GATE", "pass, stop or review", "red"),
    ]
    for index, (label, detail, color) in enumerate(events):
        x = 0.72 + index * 2.16
        add_shape(
            slide,
            MSO_SHAPE.OVAL,
            x,
            rail_y - 0.39,
            0.78,
            0.78,
            fill=color,
            line="white",
            line_width=1.2,
        )
        add_text(
            slide,
            label,
            x - 0.43,
            rail_y - 1.03,
            1.64,
            0.33,
            size=17,
            color=color if color != "gold" else "navy",
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            detail,
            x - 0.61,
            rail_y + 0.57,
            2.00,
            0.56,
            size=16,
            color="slate",
            align=PP_ALIGN.CENTER,
        )
    add_shape(
        slide,
        MSO_SHAPE.OVAL,
        11.13,
        2.46,
        1.51,
        1.51,
        fill="navy",
        line=None,
    )
    add_text(
        slide,
        "OPERATOR",
        11.17,
        3.04,
        1.42,
        0.32,
        size=16,
        font=DISPLAY,
        color="white",
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "The screen shows artifacts and decisions—not hidden model reasoning.",
        1.29,
        5.26,
        7.18,
        0.45,
        size=18,
        color="navy",
        bold=True,
    )
    add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        8.78,
        5.03,
        3.86,
        0.82,
        fill="fog",
        line=None,
    )
    add_text(
        slide,
        "Browser refresh → reconnect to the active run",
        9.07,
        5.25,
        3.28,
        0.34,
        size=17,
        color="green",
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "Completed valid work does not restart.",
        4.14,
        6.06,
        5.04,
        0.34,
        size=17,
        color="slate",
        italic=True,
        align=PP_ALIGN.CENTER,
    )

    # 12 — release standard
    slide = prs.slides.add_slide(blank)
    add_title(
        slide,
        12,
        "Release requires provenance and rendered QA",
        "Source: Art direction, Office QA, manifest, archive, and publishing contracts · 23 Jul 2026",
        "Sources: .claude/agents/art-director.md; "
        ".claude/agents/presentation-designer.md; cli/docx_builder.py; "
        "cli/presentation_qa.py; cli/publishing_quality.py; cli/run_manifest.py; "
        "cli/archive.py; cli/publish.py; prompts/orchestration.md.",
    )
    add_text(
        slide,
        "Art Director sets the visual contract before Word or PowerPoint begins.",
        0.55,
        1.34,
        12.0,
        0.48,
        size=20,
        color="slate",
    )
    flow_y = 2.45
    flow = [
        ("ART DIRECT", "visual contract", "navy", 0.58, 1.93),
        ("WORD + PPT", "build + reopen", "blue", 3.10, 1.93),
        ("STRUCTURE", "package checks", "slate", 5.62, 1.93),
        ("RENDER + REVIEW", "full-size + montage", "green", 8.14, 1.93),
        ("RELEASE GATE", "only after pass", "red", 10.66, 2.04),
    ]
    for x in (2.58, 5.10, 7.62, 10.14):
        add_arrow(slide, x, flow_y + 0.78, 0.42, 0.16, color="blue")
    for label, detail, color, x, w in flow:
        add_shape(
            slide,
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x,
            flow_y,
            w,
            1.72,
            fill="white",
            line=color,
            line_width=1.7,
        )
        add_text(
            slide,
            label,
            x + 0.12,
            flow_y + 0.37,
            w - 0.24,
            0.38,
            size=17,
            color=color,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            detail,
            x + 0.12,
            flow_y + 0.93,
            w - 0.24,
            0.40,
            size=16,
            color="slate",
            align=PP_ALIGN.CENTER,
        )
    add_kicker(slide, "Provenance that travels", 0.58, 4.67, 5.5, color="navy")
    provenance = [
        ("RUN MANIFEST", "what ran"),
        ("EVIDENCE LEDGER", "what supports it"),
        ("CLAIM LINEAGE", "what reached the reader"),
        ("INSPECTION RECEIPT", "what was seen"),
    ]
    for index, (label, detail) in enumerate(provenance):
        x = 0.58 + index * 3.07
        add_line(slide, x, 5.36, x + 0.43, 5.36, color=("navy", "blue", "green", "red")[index], width=4)
        add_text(
            slide,
            label,
            x,
            5.52,
            2.76,
            0.31,
            size=17,
            color="navy",
            bold=True,
        )
        add_text(
            slide,
            detail,
            x,
            5.91,
            2.76,
            0.31,
            size=16,
            color="slate",
        )
    add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        7.43,
        6.26,
        5.25,
        0.46,
        fill="navy",
        line=None,
    )
    add_text(
        slide,
        "APPROVE THE RECORD, NOT THE FLUENCY",
        7.66,
        6.37,
        4.79,
        0.22,
        size=16,
        color="white",
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    return prs


def _write_visual_brief(records_dir: Path) -> Path:
    records_dir.mkdir(parents=True, exist_ok=True)
    brief = _visual_brief()
    schema = json.loads(BRAND_SCHEMA.read_text(encoding="utf-8"))
    validate(instance=brief, schema=schema)
    path = records_dir / "AI Research Council — How It Works-visual-brief.json"
    path.write_text(json.dumps(brief, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")
    return path


def _save_identical(prs: Presentation, output: Path, mirror: Path | None) -> None:
    output.parent.mkdir(parents=True, exist_ok=True)
    temporary = output.with_name(f".{output.name}.building")
    prs.save(temporary)
    os.replace(temporary, output)
    if mirror is not None:
        mirror.parent.mkdir(parents=True, exist_ok=True)
        shutil.copyfile(output, mirror)


def build(output: Path, mirror: Path | None, records_dir: Path | None) -> tuple[Path, Path | None]:
    if records_dir is not None:
        _write_visual_brief(records_dir)
    prs = _build_deck()
    _save_identical(prs, output, mirror)
    # Reopen both packages before returning.
    if len(Presentation(output).slides) != 12:
        raise RuntimeError("Primary deck did not reopen with 12 slides.")
    if mirror is not None and len(Presentation(mirror).slides) != 12:
        raise RuntimeError("Mirror deck did not reopen with 12 slides.")
    return output, mirror


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--output", type=Path, default=FINAL_DECK)
    parser.add_argument("--mirror", type=Path, default=ROOT_DECK)
    parser.add_argument(
        "--records-dir",
        type=Path,
        help="Optional external directory for the validated visual brief.",
    )
    args = parser.parse_args()
    output, mirror = build(args.output.resolve(), args.mirror.resolve(), args.records_dir)
    print(output)
    if mirror is not None:
        print(mirror)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
