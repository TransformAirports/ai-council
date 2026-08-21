from __future__ import annotations

import json
import tempfile
import unittest
from pathlib import Path
from unittest.mock import patch

from docx import Document
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pypdf import PdfWriter
from jsonschema import Draft202012Validator

from cli.docx_builder import (
    _build_decision_brief,
    _configure_document,
    _markdown_to_docx,
    _renderable_exhibits,
    build_documents,
)
from cli.evidence import file_sha256
from cli.presentation_qa import (
    PresentationQAConfig,
    prepare_visual_inspection_receipt,
    qa_presentation,
    qa_visual_inspection_receipt,
)
from cli.publishing_quality import (
    lint_reader_text,
    prepare_word_visual_inspection_receipt,
    qa_word_visual_inspection_receipt,
    render_office_artifact,
)


FINAL_DRAFT = """# Checkpoint access

The airport must decide whether to authorize a controlled pilot.[^1]

---

## Executive summary

1. **The decision is small but time-sensitive.** Demand reaches 12 million passengers a year.[^1]

2. **The operating window determines the risk.** Peak waits reach 42 minutes.[^2]

3. **Peer programs are deliberately capped.** The largest admits 300 visitors a day.[^3]

4. **The financial exposure is bounded.** Annual cost is below $1 million.[^4]

**The recommendation.** Authorize a 90-day off-peak pilot at one checkpoint. Give the operations center a pause switch. Stop if waits exceed the published threshold.

## The counter-case, honestly presented

The checkpoint is already constrained, and discretionary users could displace ticketed passengers during irregular operations.[^2]

The benefit is primarily community value, not a measurable operating saving.[^4]

[^1]: Airport traffic report.
[^2]: Airport checkpoint study.
[^3]: Peer airport program record.
[^4]: Airport budget.
"""


class PublishingV2Tests(unittest.TestCase):
    def test_report_without_decision_frame_does_not_invent_decision_package(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            final = root / "final.md"
            method = root / "method.md"
            final.write_text(FINAL_DRAFT, encoding="utf-8")
            method.write_text("# Method\n\nIndependent review.", encoding="utf-8")
            with patch(
                "cli.docx_builder.render_office_artifact", return_value=([], [])
            ), patch("cli.docx_builder.prepare_word_visual_inspection_receipt"):
                report, summary = build_documents(
                    slug="narrative-report",
                    title="Narrative report",
                    final_draft=final,
                    methodology=method,
                    out_dir=root / "stage4",
                    output_format="report",
                    decision_frame_enabled=False,
                )

            text = "\n".join(paragraph.text for paragraph in Document(report).paragraphs)
            self.assertIsNone(summary)
            self.assertNotIn("Executive decision brief", text)
            self.assertIn("Executive summary", text)

    def test_unverified_tag_is_release_blocking(self) -> None:
        issues = lint_reader_text(
            "The project costs $4 million. [UNVERIFIED — HUMAN REVIEW]"
        )
        self.assertTrue(
            any(issue.code == "unverified_release_tag" for issue in issues)
        )

    def test_decision_brief_prefers_explicit_verified_summary(self) -> None:
        brief = _build_decision_brief(FINAL_DRAFT)
        self.assertIn("Authorize a 90-day", brief.bottom_line)
        self.assertTrue(any("pause switch" in item for item in brief.recommendations))
        self.assertTrue(any("already constrained" in item for item in brief.risks))
        self.assertFalse(any("Hold that definition" in item for item in brief.evidence))

    def test_visual_brief_can_supply_the_canonical_decision(self) -> None:
        brief = _build_decision_brief(
            FINAL_DRAFT,
            visual_brief={
                "decision": "Approve a 90-day controlled pilot at one checkpoint.",
                "decision_owner": "Chief Operating Officer",
                "approval_route": ["Security Committee", "Director of Aviation"],
                "first_90_day_action": "Open one off-peak checkpoint pilot.",
                "success_measures": [
                    "Keep peak waits below 42 minutes.",
                    "Record no visitor-related security incidents.",
                ],
            },
        )
        self.assertEqual(
            brief.bottom_line,
            "Approve a 90-day controlled pilot at one checkpoint.",
        )
        self.assertEqual(brief.decision_owner, "Chief Operating Officer")
        self.assertEqual(
            brief.approval_route,
            "Security Committee → Director of Aviation",
        )
        self.assertEqual(
            brief.first_90_day_action,
            "Open one off-peak checkpoint pilot.",
        )
        self.assertEqual(
            brief.success_measures,
            (
                "Keep peak waits below 42 minutes.",
                "Record no visitor-related security incidents.",
            ),
        )

    def test_structured_decision_context_does_not_depend_on_prose_inference(
        self,
    ) -> None:
        brief = _build_decision_brief(
            """# Pilot brief

## Decision and bottom line

Approve the controlled pilot.

## Recommended action and guardrails

The success measure is stated here, but it is not the first action.
""",
            decision_context={
                "decision": "Approve one bounded pilot.",
                "decision_owner": "Chief Operating Officer",
                "approval_path": "Security Committee then Director of Aviation",
                "first_action": "Open one off-peak operating test.",
                "time_horizon": "Authorize within 30 days; test for 90 days.",
                "success_measure": "Preserve the passenger-service threshold.",
            },
        )
        self.assertEqual(brief.bottom_line, "Approve one bounded pilot.")
        self.assertEqual(brief.decision_owner, "Chief Operating Officer")
        self.assertEqual(
            brief.approval_route,
            "Security Committee then Director of Aviation",
        )
        self.assertEqual(
            brief.first_90_day_action,
            "Open one off-peak operating test.",
        )
        self.assertEqual(
            brief.time_horizon,
            "Authorize within 30 days; test for 90 days.",
        )
        self.assertEqual(
            brief.success_measures,
            ("Preserve the passenger-service threshold.",),
        )

    def test_horizontal_rule_is_not_emitted_as_literal_text(self) -> None:
        doc = Document()
        _configure_document(doc)
        _markdown_to_docx(doc, "# Decision\n\nBefore.\n\n---\n\nAfter.")
        self.assertNotIn("---", [paragraph.text for paragraph in doc.paragraphs])

    def test_word_exhibits_fail_closed_and_ignore_production_specs(self) -> None:
        self.assertEqual(
            _renderable_exhibits(
                {
                    "signature_visual": {
                        "concept": "Terminal curb allocation map",
                        "visual_type": "Annotated map",
                        "visual_spec": "Color curb zones by operating use.",
                    },
                    "report_visuals": [
                        {
                            "title": "Terminal curb allocation map",
                            "format": "Annotated map",
                            "visual_spec": "Color curb zones by operating use.",
                        }
                    ],
                }
            ),
            (),
        )
        with self.assertRaisesRegex(ValueError, "values for 2 columns"):
            _renderable_exhibits(
                {
                    "report_visuals": [
                        {
                            "title": "Operating comparison",
                            "exhibit_type": "comparison",
                            "takeaway": "The supplied options have different limits.",
                            "evidence_ids": ["EV-001"],
                            "source_note": "Airport operating record.",
                            "row_header": "Criterion",
                            "columns": [
                                {"label": "Current"},
                                {"label": "Pilot"},
                            ],
                            "rows": [
                                {
                                    "label": "Operating window",
                                    "values": ["Off-peak only"],
                                }
                            ],
                        }
                    ]
                }
            )

    def test_build_documents_accepts_visual_brief_and_writes_bundle_qa(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            stage4 = root / "stage4"
            final = root / "final.md"
            method = root / "method.md"
            visual = stage4 / "visual-brief.json"
            final.write_text(FINAL_DRAFT, encoding="utf-8")
            method.write_text("# Method\n\nIndependent review.", encoding="utf-8")
            stage4.mkdir()
            visual_payload = {
                "schema_version": "2.0",
                "communication_job": (
                    "Enable an executive decision on a controlled checkpoint pilot."
                ),
                "audience": "Airport executive committee",
                "decision": "Approve a 90-day controlled pilot.",
                "decision_owner": "Chief Operating Officer",
                "approval_path": "Security Committee then Director of Aviation",
                "first_90_day_action": "Launch one off-peak operating test.",
                "success_measures": ["Hold peak waits below 42 minutes."],
                "deck_mode": "board_decision",
                "visual_thesis": (
                    "A bounded operating envelope makes the pilot reversible."
                ),
                "signature_visual": {
                    "slide_number": 2,
                    "concept": "The controlled checkpoint operating envelope",
                    "visual_type": "Annotated terminal map",
                    "visual_spec": "Locate the proposed checkpoint operating zone.",
                    "evidence_ids": ["EV-001", "EV-002"],
                    "source_note": "Airport traffic report and checkpoint study.",
                },
                "brand_profile": "Transform Airports executive",
                "slides": [
                    {
                        "slide_number": 1,
                        "narrative_job": "Frame the decision",
                        "headline": "A controlled pilot keeps the risk bounded",
                        "evidence_ids": ["EV-001", "EV-002"],
                        "visual_type": "Decision flow",
                        "visual_spec": "Show launch, pause, and stop gates.",
                        "source_note": "Airport reports.",
                        "density_budget": "One claim and three gates",
                        "speaker_note": "Ask for authorization.",
                    }
                ],
                "report_visuals": [
                    {
                        "title": "Current operation versus controlled pilot",
                        "exhibit_type": "comparison",
                        "takeaway": (
                            "The pilot adds a bounded off-peak access path while "
                            "preserving a measurable stop condition."
                        ),
                        "evidence_ids": ["EV-001", "EV-002"],
                        "source_note": "Airport checkpoint study.",
                        "row_header": "Criterion",
                        "columns": [
                            {"label": "Current operation"},
                            {"label": "Controlled pilot"},
                        ],
                        "rows": [
                            {
                                "label": "Operating window",
                                "values": [
                                    "Ticketed passengers only",
                                    "Visitors during off-peak periods",
                                ],
                            },
                            {
                                "label": "Peak wait threshold",
                                "unit": "minutes",
                                "values": ["No pilot threshold", 42],
                            },
                        ],
                    },
                    {
                        "title": "Pilot authorization and stop flow",
                        "exhibit_type": "flow",
                        "takeaway": (
                            "Named gates keep launch, monitoring, and pause "
                            "authority explicit."
                        ),
                        "evidence_ids": ["EV-002"],
                        "source_note": "Airport checkpoint study.",
                        "steps": [
                            {
                                "label": "Authorize",
                                "detail": "Approve the bounded 90-day test.",
                                "owner": "Director of Aviation",
                            },
                            {
                                "label": "Launch",
                                "detail": "Open the off-peak operating window.",
                                "owner": "Chief Operating Officer",
                            },
                            {
                                "label": "Pause",
                                "detail": "Close visitor access before displacement.",
                                "trigger": "Peak wait exceeds 42 minutes",
                            },
                        ],
                    },
                    {
                        "title": "The first 90 days",
                        "exhibit_type": "timeline",
                        "takeaway": (
                            "The decision advances through authorization, "
                            "operation, and a measured continuation decision."
                        ),
                        "evidence_ids": ["EV-001", "EV-002"],
                        "source_note": "Airport pilot operating record.",
                        "milestones": [
                            {
                                "period": "Days 1–30",
                                "action": "Authorize and configure",
                                "owner": "Director of Aviation",
                            },
                            {
                                "period": "Days 31–60",
                                "action": "Operate and monitor",
                                "owner": "Chief Operating Officer",
                                "success_measure": (
                                    "Hold peak waits below 42 minutes"
                                ),
                            },
                            {
                                "period": "Days 61–90",
                                "action": "Review and decide",
                                "success_measure": (
                                    "Document a continuation or stop decision"
                                ),
                            },
                        ],
                    }
                ],
                "source_appendix": [
                    {
                        "evidence_id": "EV-001",
                        "source": "Airport traffic report",
                    }
                ],
                "accessibility_checks": [
                    "Do not rely on color alone to show a stop condition."
                ],
                "asset_requests": [],
            }
            first_slide = visual_payload["slides"][0]
            visual_payload["slides"] = [
                {
                    **first_slide,
                    "slide_number": number,
                    "headline": (
                        "A controlled pilot keeps the risk bounded"
                        if number == 1
                        else f"Decision guardrail {number} keeps the pilot bounded"
                    ),
                }
                for number in range(1, 9)
            ]
            schema = json.loads(
                (
                    Path(__file__).parents[1]
                    / "assets"
                    / "brand"
                    / "visual-brief.schema.json"
                ).read_text(encoding="utf-8")
            )
            self.assertEqual(
                list(Draft202012Validator(schema).iter_errors(visual_payload)),
                [],
            )
            unbound_signature = json.loads(json.dumps(visual_payload))
            unbound_signature["signature_visual"].pop("slide_number")
            self.assertTrue(
                list(
                    Draft202012Validator(schema).iter_errors(
                        unbound_signature
                    )
                )
            )
            unsupported_payload = json.loads(json.dumps(visual_payload))
            unsupported_payload["report_visuals"] = [
                {
                    "title": "Terminal curb allocation",
                    "exhibit_type": "map",
                    "takeaway": "The curb zones support different operating uses.",
                    "evidence_ids": ["EV-001"],
                    "source_note": "Airport curb plan.",
                }
            ]
            self.assertTrue(
                list(
                    Draft202012Validator(schema).iter_errors(
                        unsupported_payload
                    )
                )
            )
            visual.write_text(
                json.dumps(visual_payload),
                encoding="utf-8",
            )
            with patch(
                "cli.docx_builder.render_office_artifact",
                return_value=([], []),
            ), patch(
                "cli.docx_builder.prepare_word_visual_inspection_receipt",
            ):
                report, summary = build_documents(
                    slug="pilot",
                    title="Controlled checkpoint pilot",
                    final_draft=final,
                    methodology=method,
                    out_dir=stage4,
                    visual_brief=visual,
                    revision_label="Revised — Version 2",
                )

            self.assertTrue(report.is_file())
            self.assertIsNotNone(summary)
            self.assertIn(
                "REVISED — VERSION 2",
                "\n".join(paragraph.text for paragraph in Document(report).paragraphs),
            )
            assert summary is not None
            self.assertIn(
                "REVISED — VERSION 2",
                "\n".join(paragraph.text for paragraph in Document(summary).paragraphs),
            )
            quality = json.loads(
                (root / "publishing-quality.json").read_text(encoding="utf-8")
            )
            self.assertTrue(quality["ok"])
            self.assertEqual(quality["metadata"]["visual_brief"], str(visual))

            report_doc = Document(report)
            report_text = "\n".join(
                [paragraph.text for paragraph in report_doc.paragraphs]
                + [
                    cell.text
                    for table in report_doc.tables
                    for row in table.rows
                    for cell in row.cells
                ]
            )
            table_rows = [
                tuple(cell.text for cell in row.cells)
                for table in report_doc.tables
                for row in table.rows
            ]
            self.assertIn(
                ("Criterion", "Current operation", "Controlled pilot"),
                table_rows,
            )
            self.assertIn(
                ("When", "Action", "Accountability / success test"),
                table_rows,
            )
            self.assertTrue(
                any(
                    "Authorize" in " ".join(row)
                    and any("1" == cell.strip() for cell in row)
                    for row in table_rows
                )
            )
            self.assertTrue(
                any("↓" in cell for row in table_rows for cell in row)
            )
            self.assertIn(
                "generated with assistance from a multi-model AI research system",
                report_text,
            )
            self.assertIn("Chief Operating Officer", report_text)
            self.assertIn("Launch one off-peak operating test.", report_text)
            self.assertIn("Hold peak waits below 42 minutes.", report_text)
            self.assertIn("Decision exhibits", report_text)
            self.assertIn("Current operation versus controlled pilot", report_text)
            self.assertIn("Current operation", report_text)
            self.assertIn("Controlled pilot", report_text)
            self.assertIn("Peak wait threshold (minutes)", report_text)
            self.assertIn("Pilot authorization and stop flow", report_text)
            self.assertIn("Authorize", report_text)
            self.assertIn("Launch", report_text)
            self.assertIn("Pause", report_text)
            self.assertIn("The first 90 days", report_text)
            self.assertIn("Days 31–60", report_text)
            self.assertIn("Operate and monitor", report_text)
            self.assertIn("Airport checkpoint study.", report_text)
            self.assertNotIn(
                "The controlled checkpoint operating envelope",
                report_text,
            )
            self.assertNotIn("Exhibit form", report_text)
            self.assertNotIn("What it shows", report_text)
            self.assertIn("Exhibit source notes", report_text)
            self.assertNotIn("production plan", report_text.lower())
            self.assertNotIn("evidence identifiers", report_text.lower())
            self.assertNotIn("Evidence IDs", report_text)
            self.assertNotIn("evidence_id", report_text)
            self.assertNotIn("EV-001", report_text)
            self.assertNotIn("EV-002", report_text)
            self.assertIn("Technical appendix: Evidence register", report_text)
            self.assertIn("Airport traffic report.", report_text)
            self.assertIn("Additional exhibit sources", report_text)
            self.assertIn("Technical appendix: Methodology", report_text)

            self.assertIsNotNone(summary)
            summary_doc = Document(summary)
            summary_text = "\n".join(
                [paragraph.text for paragraph in summary_doc.paragraphs]
                + [
                    cell.text
                    for table in summary_doc.tables
                    for row in table.rows
                    for cell in row.cells
                ]
            )
            self.assertIn(
                "generated with assistance from a multi-model AI research system",
                summary_text,
            )
            self.assertIn("Chief Operating Officer", summary_text)
            self.assertIn("Security Committee then Director of Aviation", summary_text)
            self.assertIn("Launch one off-peak operating test.", summary_text)
            self.assertIn("Hold peak waits below 42 minutes.", summary_text)
            self.assertNotIn("EV-001", summary_text)
            self.assertNotIn("EV-002", summary_text)

    def test_word_output_formats_use_distinct_canonical_structures(self) -> None:
        drafts = {
            "report": FINAL_DRAFT,
            "article": """# The operating window

At 6 a.m., the checkpoint decision is already visible in the queue.

## The constraint

The airport can test the operating idea without committing to a permanent
program.

## The consequence

Airport leaders should authorize only the bounded test.
""",
            "brief": """# Controlled pilot brief

## Decision and bottom line

Approve a bounded checkpoint pilot.

## Three findings that carry the decision

- The operating window is reversible.
- The approval route is known.
- The stop condition is measurable.

## Strongest counter-case

The checkpoint may already be too constrained.

## Recommended action and guardrails

Open one off-peak test and stop if the approved threshold is crossed.
""",
            "recommendations": """# Controlled pilot recommendations

The airport can answer the decision with a bounded operating test.

1. Authorize one off-peak pilot with a named owner and stop condition.
2. Report results to the approval committee before any expansion.

## Do not proceed if

The airport cannot preserve the published passenger-service threshold.
""",
        }
        decision_context = {
            "decision": "Approve a bounded checkpoint pilot.",
            "decision_owner": "Chief Operating Officer",
            "approval_path": "Security Committee then Director of Aviation",
            "time_horizon": "Authorize within 30 days; test for 90 days.",
            "success_measure": "Preserve the approved passenger-service threshold.",
        }

        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            methodology = root / "methodology.md"
            methodology.write_text(
                "# Methodology\n\nIndependent evidence review.",
                encoding="utf-8",
            )
            built: dict[str, tuple[Path, Path | None]] = {}
            with patch(
                "cli.docx_builder.render_office_artifact",
                return_value=([], []),
            ), patch(
                "cli.docx_builder.prepare_word_visual_inspection_receipt",
            ):
                for output_format, markdown in drafts.items():
                    format_root = root / output_format
                    format_root.mkdir()
                    final = format_root / "final.md"
                    final.write_text(markdown, encoding="utf-8")
                    built[output_format] = build_documents(
                        slug=output_format,
                        title=f"{output_format.title()} output",
                        final_draft=final,
                        methodology=methodology,
                        out_dir=format_root / "stage4",
                        output_format=output_format,
                        decision_context=decision_context,
                    )

            def document_text(path: Path) -> str:
                doc = Document(path)
                return "\n".join(
                    [paragraph.text for paragraph in doc.paragraphs]
                    + [
                        cell.text
                        for table in doc.tables
                        for row in table.rows
                        for cell in row.cells
                    ]
                )

            report_path, report_summary = built["report"]
            report_text = document_text(report_path)
            self.assertIsNotNone(report_summary)
            self.assertIn("Executive decision brief", report_text)
            self.assertIn("Contents", report_text)
            self.assertIn("Technical appendix: Evidence register", report_text)
            self.assertIn("Technical appendix: Methodology", report_text)

            article_path, article_summary = built["article"]
            article_text = document_text(article_path)
            self.assertIsNone(article_summary)
            self.assertIn("At 6 a.m.", article_text)
            self.assertNotIn("Executive decision brief", article_text)
            self.assertNotIn("Contents", article_text)
            self.assertNotIn("Technical appendix:", article_text)

            brief_path, brief_summary = built["brief"]
            brief_text = document_text(brief_path)
            self.assertIsNone(brief_summary)
            self.assertIn("Executive brief", brief_text)
            self.assertIn("DECISION FRAME", brief_text)
            self.assertNotIn("Contents", brief_text)
            self.assertNotIn("Technical appendix:", brief_text)

            recommendations_path, recommendations_summary = built[
                "recommendations"
            ]
            recommendations_text = document_text(recommendations_path)
            self.assertIsNone(recommendations_summary)
            self.assertIn("Action recommendations", recommendations_text)
            self.assertIn("DECISION MANDATE", recommendations_text)
            self.assertNotIn("Contents", recommendations_text)
            self.assertNotIn("Technical appendix:", recommendations_text)

            for short_text in (brief_text, recommendations_text):
                self.assertIn("Chief Operating Officer", short_text)
                self.assertIn(
                    "Security Committee then Director of Aviation",
                    short_text,
                )
                self.assertIn(
                    "Authorize within 30 days; test for 90 days.",
                    short_text,
                )
                self.assertIn(
                    "Preserve the approved passenger-service threshold.",
                    short_text,
                )
                self.assertNotIn(
                    "Assign the accountable decision owner",
                    short_text,
                )

    def _manual_title_deck(self, path: Path, *, include_source: bool = True) -> None:
        deck = Presentation()
        deck.slide_width = Inches(13.333)
        deck.slide_height = Inches(7.5)
        slide = deck.slides.add_slide(deck.slide_layouts[6])

        title = slide.shapes.add_textbox(
            Inches(0.7), Inches(0.5), Inches(11.8), Inches(0.8)
        )
        run = title.text_frame.paragraphs[0].add_run()
        run.text = "Stage 1 launches a parallel research swarm"
        run.font.size = Pt(42)

        body = slide.shapes.add_textbox(
            Inches(0.7), Inches(1.7), Inches(10.5), Inches(1.2)
        )
        run = body.text_frame.paragraphs[0].add_run()
        run.text = "The Strategist combines briefs covering 42 million passengers."
        run.font.size = Pt(15.75)

        if include_source:
            source = slide.shapes.add_textbox(
                Inches(0.7), Inches(6.8), Inches(10.5), Inches(0.3)
            )
            source.name = "Source footer"
            run = source.text_frame.paragraphs[0].add_run()
            run.text = "Source: Council architecture and airport traffic record"
            run.font.size = Pt(9)
        deck.save(path)

    def test_process_explainer_profile_accepts_manual_title_and_internal_terms(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            path = Path(directory) / "process.pptx"
            self._manual_title_deck(path)
            report = qa_presentation(
                path,
                config=PresentationQAConfig.process_explainer(),
            )
            self.assertEqual(report.errors, [])
            self.assertEqual(report.metadata["profile"], "internal_process_explainer")

    def test_numeric_units_trigger_source_check_and_default_profile_stays_strict(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            path = Path(directory) / "process.pptx"
            self._manual_title_deck(path, include_source=False)
            report = qa_presentation(
                path,
                config=PresentationQAConfig.process_explainer(),
            )
            self.assertTrue(
                any(issue.code == "unsourced_numeric_slide" for issue in report.warnings)
            )

            strict = qa_presentation(
                path,
                config=PresentationQAConfig(
                    min_slide_count=1,
                    max_slide_count=18,
                ),
            )
            codes = {issue.code for issue in strict.errors}
            self.assertIn("internal_process_language", codes)
            self.assertIn("small_body_type", codes)
            self.assertIn("unsourced_numeric_slide", codes)
            self.assertIn("no_evidence_visuals", codes)

    def test_deck_mode_enforces_hard_slide_count_contracts(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            path = root / "mode.pptx"
            deck = Presentation()
            for index in range(13):
                slide = deck.slides.add_slide(deck.slide_layouts[5])
                title = slide.shapes.title
                title.text = f"Decision assertion {index + 1}"
                for paragraph in title.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(44 if index == 0 else 36)
            deck.save(path)

            board = qa_presentation(path, deck_mode="board_decision")
            self.assertIn(
                "long_deck",
                {issue.code for issue in board.errors},
            )

            technical = qa_presentation(
                path,
                deck_mode="technical_read_ahead",
            )
            technical_codes = {issue.code for issue in technical.errors}
            self.assertIn("short_deck", technical_codes)
            self.assertNotIn("long_deck", technical_codes)

    def test_signature_exhibit_must_be_on_the_canonical_slide(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            path = Path(directory) / "signature.pptx"
            deck = Presentation()
            slide = deck.slides.add_slide(deck.slide_layouts[5])
            slide.shapes.title.text = "Approve the bounded operating test"
            for run in slide.shapes.title.text_frame.paragraphs[0].runs:
                run.font.size = Pt(44)
            deck.save(path)
            brief = {
                "deck_mode": "board_decision",
                "signature_visual": {
                    "slide_number": 1,
                    "concept": "Operating envelope",
                    "visual_type": "Decision flow",
                },
                "slides": [
                    {
                        "slide_number": 1,
                        "headline": "Approve the bounded operating test",
                        "evidence_ids": [],
                        "source_note": "",
                        "visual_type": "text",
                    }
                ],
            }
            config = PresentationQAConfig(
                min_slide_count=1,
                max_slide_count=1,
                require_visual_evidence=False,
            )
            missing = qa_presentation(
                path,
                config=config,
                visual_brief=brief,
            )
            self.assertIn(
                "signature_visual_exhibit_missing",
                {issue.code for issue in missing.errors},
            )

            exhibit = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(1),
                Inches(2),
                Inches(4),
                Inches(2),
            )
            exhibit.name = "SIGNATURE VISUAL — Operating envelope"
            deck.save(path)
            present = qa_presentation(
                path,
                config=config,
                visual_brief=brief,
            )
            self.assertNotIn(
                "signature_visual_exhibit_missing",
                {issue.code for issue in present.errors},
            )

    def test_explicit_native_shape_title_outranks_a_larger_metric(
        self,
    ) -> None:
        with tempfile.TemporaryDirectory() as directory:
            path = Path(directory) / "metric-title.pptx"
            deck = Presentation()
            slide = deck.slides.add_slide(deck.slide_layouts[6])
            title = slide.shapes.add_textbox(
                Inches(0.5),
                Inches(0.3),
                Inches(10),
                Inches(0.7),
            )
            title.name = "slide-1-title"
            title.text = "The roster separates depth from control"
            for run in title.text_frame.paragraphs[0].runs:
                run.font.size = Pt(36)
            metric = slide.shapes.add_textbox(
                Inches(0.5),
                Inches(1.5),
                Inches(2),
                Inches(1),
            )
            metric.text = "54"
            for run in metric.text_frame.paragraphs[0].runs:
                run.font.size = Pt(76)
            signature = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(3),
                Inches(2),
                Inches(4),
                Inches(2),
            )
            signature.name = "SIGNATURE VISUAL — roster"
            deck.save(path)

            report = qa_presentation(
                path,
                config=PresentationQAConfig(
                    min_slide_count=1,
                    max_slide_count=1,
                    require_visual_evidence=False,
                ),
                visual_brief={
                    "deck_mode": "board_decision",
                    "signature_visual": {
                        "slide_number": 1,
                        "concept": "Roster",
                        "visual_type": "Proportional exhibit",
                    },
                    "slides": [
                        {
                            "slide_number": 1,
                            "headline": (
                                "The roster separates depth from control"
                            ),
                            "evidence_ids": [],
                            "source_note": "",
                            "visual_type": "proportional exhibit",
                        }
                    ],
                },
            )
            self.assertNotIn(
                "visual_brief_headline_mismatch",
                {issue.code for issue in report.errors},
            )

    def test_visual_inspection_receipt_is_bound_to_exact_deck_bytes(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            deck_path = root / "decision.pptx"
            brief_path = root / "visual-brief.json"
            render_dir = root / "inspection" / "decision"
            render_dir.mkdir(parents=True)

            deck = Presentation()
            slide = deck.slides.add_slide(deck.slide_layouts[5])
            slide.shapes.title.text = "Authorize the bounded pilot"
            deck.save(deck_path)
            brief_path.write_text(
                json.dumps(
                    {
                        "deck_mode": "board_decision",
                        "signature_visual": {
                            "slide_number": 1,
                            "concept": "Bounded pilot decision",
                            "visual_type": "Decision exhibit",
                        },
                    }
                ),
                encoding="utf-8",
            )
            slide_png = render_dir / "decision-1.png"
            Image.new("RGB", (640, 360), "white").save(slide_png)
            receipt_path = root / "decision-visual-inspection.json"
            prepare_visual_inspection_receipt(
                artifact=deck_path,
                visual_brief=brief_path,
                deck_mode="board_decision",
                rendered_files=[slide_png],
                receipt_path=receipt_path,
            )

            pending = qa_visual_inspection_receipt(
                receipt_path,
                artifact=deck_path,
                visual_brief=brief_path,
                deck_mode="board_decision",
            )
            self.assertFalse(pending.ok)
            self.assertIn(
                "inspection_not_passed",
                {issue.code for issue in pending.errors},
            )

            receipt = json.loads(receipt_path.read_text(encoding="utf-8"))
            receipt["inspection"].update(
                {
                    "full_size_each_slide_inspected": True,
                    "montage_inspected": True,
                    "signature_exhibit_present": True,
                    "signature_exhibit_matches_brief": True,
                    "findings_resolved": True,
                    "status": "pass",
                }
            )
            receipt_path.write_text(json.dumps(receipt), encoding="utf-8")
            passed = qa_visual_inspection_receipt(
                receipt_path,
                artifact=deck_path,
                visual_brief=brief_path,
                deck_mode="board_decision",
            )
            self.assertTrue(passed.ok)

            deck.core_properties.title = "Changed after inspection"
            deck.save(deck_path)
            mutated = qa_visual_inspection_receipt(
                receipt_path,
                artifact=deck_path,
                visual_brief=brief_path,
                deck_mode="board_decision",
            )
            self.assertIn(
                "inspection_hash_mismatch",
                {issue.code for issue in mutated.errors},
            )

    def test_word_page_inspection_binds_every_render_and_exact_docx(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            artifact = root / "decision.docx"
            artifact.write_bytes(b"exact Word package bytes")
            render_dir = root / "qa" / "decision"
            render_dir.mkdir(parents=True)
            pdf = render_dir / "decision.pdf"
            writer = PdfWriter()
            writer.add_blank_page(width=612, height=792)
            writer.add_blank_page(width=612, height=792)
            with pdf.open("wb") as handle:
                writer.write(handle)
            pages = [
                render_dir / "decision-1.png",
                render_dir / "decision-2.png",
            ]
            for index, page in enumerate(pages, 1):
                Image.new(
                    "RGB",
                    (480, 620),
                    "white" if index == 1 else "#edf3f6",
                ).save(page)
            receipt_path = root / "decision-word-visual-inspection.json"
            prepare_word_visual_inspection_receipt(
                artifact=artifact,
                rendered_files=[pdf, *pages],
                receipt_path=receipt_path,
            )

            pending = qa_word_visual_inspection_receipt(
                receipt_path,
                artifact=artifact,
            )
            self.assertFalse(pending.ok)
            self.assertIn(
                "word_inspection_not_passed",
                {issue.code for issue in pending.errors},
            )

            receipt = json.loads(receipt_path.read_text(encoding="utf-8"))
            receipt["inspection"].update(
                {
                    "full_size_each_page_inspected": True,
                    "montage_inspected": True,
                    "findings_resolved": True,
                    "status": "pass",
                }
            )
            receipt_path.write_text(json.dumps(receipt), encoding="utf-8")
            passed = qa_word_visual_inspection_receipt(
                receipt_path,
                artifact=artifact,
            )
            self.assertTrue(passed.ok)

            pages[1].write_bytes(b"changed after inspection")
            mutated = qa_word_visual_inspection_receipt(
                receipt_path,
                artifact=artifact,
            )
            self.assertIn(
                "word_inspection_page_hash_mismatch",
                {issue.code for issue in mutated.errors},
            )

    def test_word_inspection_rejects_truncated_page_render(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            artifact = root / "two-pages.docx"
            artifact.write_bytes(b"synthetic Word package")
            render_dir = root / "qa" / "two-pages"
            render_dir.mkdir(parents=True)
            pdf = render_dir / "two-pages.pdf"
            writer = PdfWriter()
            writer.add_blank_page(width=612, height=792)
            writer.add_blank_page(width=612, height=792)
            with pdf.open("wb") as handle:
                writer.write(handle)
            only_page = render_dir / "two-pages-1.png"
            Image.new("RGB", (480, 620), "white").save(only_page)

            with self.assertRaisesRegex(
                ValueError,
                "one sequential PNG for every PDF page",
            ):
                prepare_word_visual_inspection_receipt(
                    artifact=artifact,
                    rendered_files=[pdf, only_page],
                    receipt_path=(
                        root / "two-pages-word-visual-inspection.json"
                    ),
                )

    def test_word_receipt_qa_rechecks_bound_pdf_page_count(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            artifact = root / "decision.docx"
            artifact.write_bytes(b"synthetic Word package")
            render_dir = root / "qa" / "decision"
            render_dir.mkdir(parents=True)
            pdf = render_dir / "decision.pdf"
            writer = PdfWriter()
            writer.add_blank_page(width=612, height=792)
            with pdf.open("wb") as handle:
                writer.write(handle)
            page = render_dir / "decision-1.png"
            Image.new("RGB", (480, 620), "white").save(page)
            receipt_path = (
                root / "decision-word-visual-inspection.json"
            )
            prepare_word_visual_inspection_receipt(
                artifact=artifact,
                rendered_files=[pdf, page],
                receipt_path=receipt_path,
            )

            writer = PdfWriter()
            writer.add_blank_page(width=612, height=792)
            writer.add_blank_page(width=612, height=792)
            with pdf.open("wb") as handle:
                writer.write(handle)
            receipt = json.loads(receipt_path.read_text(encoding="utf-8"))
            receipt["pdf"]["sha256"] = file_sha256(pdf)
            receipt["pdf"]["size_bytes"] = pdf.stat().st_size
            receipt["inspection"].update(
                {
                    "full_size_each_page_inspected": True,
                    "montage_inspected": True,
                    "findings_resolved": True,
                    "status": "pass",
                }
            )
            receipt_path.write_text(
                json.dumps(receipt),
                encoding="utf-8",
            )

            report = qa_word_visual_inspection_receipt(
                receipt_path,
                artifact=artifact,
            )
            self.assertIn(
                "word_inspection_pdf_page_count_mismatch",
                {issue.code for issue in report.errors},
            )

    def test_office_render_retry_removes_only_owned_stale_outputs(
        self,
    ) -> None:
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            artifact = root / "decision.docx"
            artifact.write_bytes(b"synthetic Word package")
            render_dir = root / "qa"
            render_dir.mkdir()
            owned = [
                render_dir / "decision.pdf",
                render_dir / "decision-1.png",
                render_dir / "decision-4.png",
                render_dir / "montage.png",
            ]
            for path in owned:
                path.write_bytes(b"stale")
            unrelated = render_dir / "keep-this.json"
            unrelated.write_text("{}", encoding="utf-8")

            with patch(
                "cli.publishing_quality.shutil.which",
                return_value=None,
            ):
                rendered, issues = render_office_artifact(
                    artifact,
                    render_dir,
                    required=True,
                )

            self.assertEqual(rendered, [])
            self.assertTrue(all(not path.exists() for path in owned))
            self.assertTrue(unrelated.is_file())
            self.assertIn(
                "render_unavailable",
                {issue.code for issue in issues},
            )


if __name__ == "__main__":
    unittest.main()
