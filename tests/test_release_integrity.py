from __future__ import annotations

import hashlib
import json
import shutil
import tempfile
import unittest
from datetime import date
from pathlib import Path
from unittest.mock import patch

from PIL import Image
from pptx import Presentation
from pypdf import PdfWriter

from cli.archive import archive_run
from cli.orchestrator import CostTally
from cli.publish import (
    _promote_release_with_deck_supplement,
    discover_reports,
    promote_release,
    publish_all,
    stage_release_artifacts,
    verify_release_bundle,
)
from cli.presentation_qa import prepare_visual_inspection_receipt
from cli.publishing_quality import (
    QualityReport,
    prepare_word_visual_inspection_receipt,
)
from cli.run_manifest import assert_manifest_complete
from cli.sources import attach_sources, discover_dropzone


class ReleaseIntegrityTests(unittest.TestCase):
    @staticmethod
    def _passing_docx_report(path: Path) -> QualityReport:
        return QualityReport(artifact=str(path), kind="docx")

    @staticmethod
    def _passing_pptx_report(path: Path, **_kwargs) -> QualityReport:
        return QualityReport(artifact=str(path), kind="pptx")

    def _stage_synthetic_release(
        self, root: Path, *, slug: str = "exact-release"
    ) -> tuple[Path, Path, bytes, dict]:
        stage4 = root / "stage4"
        release = root / "release"
        stage4.mkdir(parents=True)
        report_bytes = b"synthetic-stage-4-office-package\x00\x01\x02"
        report = stage4 / f"{slug}.docx"
        report.write_bytes(report_bytes)

        with (
            patch(
                "cli.publish.qa_docx",
                side_effect=self._passing_docx_report,
            ),
            patch(
                "cli.publish.render_office_artifact",
                return_value=([], []),
            ),
        ):
            payload = stage_release_artifacts(
                stage4_dir=stage4,
                slug=slug,
                release_dir=release,
            )
        return stage4, release, report_bytes, payload

    def test_stage_release_copies_exact_bytes_and_records_passing_qa(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            _, release, expected_bytes, payload = self._stage_synthetic_release(
                root
            )

            released = release / "exact-release.docx"
            self.assertEqual(released.read_bytes(), expected_bytes)
            artifact = payload["artifacts"][0]
            expected_hash = hashlib.sha256(expected_bytes).hexdigest()
            self.assertEqual(artifact["role"], "word_report")
            self.assertEqual(artifact["source_sha256"], expected_hash)
            self.assertEqual(artifact["sha256"], expected_hash)
            self.assertTrue(artifact["qa_ok"])

            qa_path = release / artifact["qa_path"]
            qa_payload = json.loads(qa_path.read_text(encoding="utf-8"))
            self.assertTrue(qa_payload["ok"])
            self.assertEqual(
                artifact["qa_sha256"],
                hashlib.sha256(qa_path.read_bytes()).hexdigest(),
            )
            persisted = json.loads(
                (release / "release-manifest.json").read_text(encoding="utf-8")
            )
            self.assertEqual(persisted["artifacts"], payload["artifacts"])

    def test_presentation_release_requires_hash_bound_visual_inspection(
        self,
    ) -> None:
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            slug = "inspected-deck"
            stage4 = root / "stage4"
            inspection = stage4 / "inspection" / slug
            inspection.mkdir(parents=True)
            deck_path = stage4 / f"{slug}.pptx"
            brief_path = stage4 / "visual-brief.json"
            receipt_path = stage4 / f"{slug}-visual-inspection.json"

            deck = Presentation()
            slide = deck.slides.add_slide(deck.slide_layouts[5])
            slide.shapes.title.text = "Approve the bounded operating pilot"
            deck.save(deck_path)
            brief_path.write_text(
                json.dumps(
                    {
                        "deck_mode": "board_decision",
                        "signature_visual": {
                            "slide_number": 1,
                            "concept": "Bounded operating pilot",
                            "visual_type": "Decision exhibit",
                        },
                    }
                ),
                encoding="utf-8",
            )
            slide_png = inspection / f"{slug}-1.png"
            Image.new("RGB", (640, 360), "white").save(slide_png)
            prepare_visual_inspection_receipt(
                artifact=deck_path,
                visual_brief=brief_path,
                deck_mode="board_decision",
                rendered_files=[slide_png],
                receipt_path=receipt_path,
            )
            receipt = json.loads(receipt_path.read_text(encoding="utf-8"))
            receipt["inspection"].update(
                {
                    "full_size_each_slide_inspected": True,
                    "montage_inspected": True,
                    "signature_exhibit_present": True,
                    "signature_exhibit_matches_brief": True,
                    "findings_resolved": True,
                    "status": "pass",
                }
            )
            receipt_path.write_text(json.dumps(receipt), encoding="utf-8")

            with patch(
                "cli.presentation_qa.qa_presentation",
                side_effect=self._passing_pptx_report,
            ):
                payload = stage_release_artifacts(
                    stage4_dir=stage4,
                    slug=slug,
                    release_dir=root / "release",
                    require_presentation=True,
                    include_roles={"presentation"},
                    presentation_mode="board_decision",
                    visual_brief=brief_path,
                    require_visual_inspection=True,
                )

            artifact = payload["artifacts"][0]
            self.assertEqual(artifact["role"], "presentation")
            self.assertIn("visual_inspection", artifact)
            verify_release_bundle(
                root / "release",
                require_word_report=False,
            )

            (root / "release" / f"{slug}.pptx").write_bytes(
                b"changed after inspection"
            )
            with self.assertRaisesRegex(RuntimeError, "hash mismatch"):
                verify_release_bundle(
                    root / "release",
                    require_word_report=False,
                )

    def test_word_release_requires_exact_inspected_page_packet(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            slug = "inspected-report"
            stage4 = root / "stage4"
            render_dir = stage4 / "qa" / slug
            render_dir.mkdir(parents=True)
            report_path = stage4 / f"{slug}.docx"
            report_path.write_bytes(b"exact synthetic Word package")
            pdf = render_dir / f"{slug}.pdf"
            writer = PdfWriter()
            writer.add_blank_page(width=612, height=792)
            with pdf.open("wb") as handle:
                writer.write(handle)
            page = render_dir / f"{slug}-1.png"
            Image.new("RGB", (480, 620), "white").save(page)
            receipt_path = (
                stage4 / f"{slug}-word-visual-inspection.json"
            )
            prepare_word_visual_inspection_receipt(
                artifact=report_path,
                rendered_files=[pdf, page],
                receipt_path=receipt_path,
            )
            receipt = json.loads(receipt_path.read_text(encoding="utf-8"))
            receipt["inspection"].update(
                {
                    "full_size_each_page_inspected": True,
                    "montage_inspected": True,
                    "findings_resolved": True,
                    "status": "pass",
                }
            )
            receipt_path.write_text(json.dumps(receipt), encoding="utf-8")

            with (
                patch(
                    "cli.publish.qa_docx",
                    side_effect=self._passing_docx_report,
                ),
                patch(
                    "cli.publish.render_office_artifact",
                    return_value=([], []),
                ),
            ):
                payload = stage_release_artifacts(
                    stage4_dir=stage4,
                    slug=slug,
                    release_dir=root / "release",
                    require_word_visual_inspection=True,
                )

            artifact = payload["artifacts"][0]
            self.assertEqual(
                artifact["visual_inspection"]["type"],
                "word_pages",
            )
            verify_release_bundle(root / "release")

            released_page = (
                root / "release" / "qa" / slug / f"{slug}-1.png"
            )
            released_page.write_bytes(b"changed after inspection")
            with self.assertRaisesRegex(
                RuntimeError,
                "visual-inspection file hash mismatch",
            ):
                verify_release_bundle(root / "release")

    def test_promote_release_rejects_a_post_qa_byte_change(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            _, release, _, _ = self._stage_synthetic_release(root)
            (release / "exact-release.docx").write_bytes(
                b"tampered-after-release-qa"
            )
            reports = root / "reports"

            with self.assertRaisesRegex(RuntimeError, "hash mismatch"):
                promote_release(release_dir=release, out_dir=reports)

            self.assertFalse((reports / "exact-release.docx").exists())

    def test_promotion_commits_an_immutable_bundle_and_retires_stale_roles(
        self,
    ) -> None:
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            _, release, expected_bytes, _ = self._stage_synthetic_release(root)
            reports = root / "reports"
            reports.mkdir()
            (reports / "exact-release.pptx").write_bytes(b"stale deck")
            (reports / "exact-release-executive-summary.docx").write_bytes(
                b"stale summary"
            )
            (
                reports / "exact-release-deck-release-manifest.json"
            ).write_text('{"status":"current"}', encoding="utf-8")

            published = promote_release(
                release_dir=release,
                out_dir=reports,
            )

            self.assertEqual(
                published["word_report"].read_bytes(),
                expected_bytes,
            )
            self.assertFalse((reports / "exact-release.pptx").exists())
            self.assertFalse(
                (reports / "exact-release-executive-summary.docx").exists()
            )
            self.assertFalse(
                (
                    reports
                    / "exact-release-deck-release-manifest.json"
                ).exists()
            )
            pointer = json.loads(
                (
                    reports / "exact-release-release-manifest.json"
                ).read_text(encoding="utf-8")
            )
            self.assertEqual(pointer["status"], "current")
            artifact = pointer["artifacts"][0]
            bundled = reports / artifact["path"]
            self.assertTrue(bundled.is_file())
            self.assertEqual(
                hashlib.sha256(bundled.read_bytes()).hexdigest(),
                artifact["sha256"],
            )
            self.assertTrue((reports / pointer["bundle_path"]).is_file())

    def test_promotion_rejects_a_manifest_path_escape(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            _, release, _, _ = self._stage_synthetic_release(root)
            with self.assertRaisesRegex(
                ValueError, "Unsafe release manifest filename"
            ):
                promote_release(
                    release_dir=release,
                    out_dir=root / "reports",
                    release_manifest_name="../escaped.json",
                )
            self.assertFalse((root / "escaped.json").exists())

    def test_identical_republish_never_replaces_the_immutable_bundle(
        self,
    ) -> None:
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            _, release, _, _ = self._stage_synthetic_release(root)
            reports = root / "reports"
            first = promote_release(release_dir=release, out_dir=reports)
            bundle = first["release_bundle"]
            inode = bundle.stat().st_ino

            second = promote_release(release_dir=release, out_dir=reports)

            self.assertEqual(second["release_bundle"], bundle)
            self.assertEqual(bundle.stat().st_ino, inode)

    def test_supplement_replay_failure_restores_every_current_pointer(
        self,
    ) -> None:
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            _, canonical, _, _ = self._stage_synthetic_release(root)
            deck_stage = root / "deck-stage"
            deck_stage.mkdir()
            (deck_stage / "exact-release.pptx").write_bytes(b"new deck")
            supplement = root / "deck-supplement"
            with patch(
                "cli.presentation_qa.qa_presentation",
                side_effect=self._passing_pptx_report,
            ):
                stage_release_artifacts(
                    stage4_dir=deck_stage,
                    slug="exact-release",
                    release_dir=supplement,
                    include_roles={"presentation"},
                    require_presentation=True,
                )

            reports = root / "reports"
            first_publish = promote_release(
                release_dir=canonical,
                out_dir=reports,
            )
            immutable_bundle = first_publish["release_bundle"]
            immutable_inode = immutable_bundle.stat().st_ino
            prior = {
                "exact-release.docx": b"old word",
                "exact-release-executive-summary.docx": b"old summary",
                "exact-release.pptx": b"old deck",
                "exact-release-release-manifest.json": b"old canonical pointer",
                "exact-release-deck-release-manifest.json": b"old deck pointer",
            }
            for name, content in prior.items():
                (reports / name).write_bytes(content)

            real_promote = promote_release
            calls = 0

            def fail_second_promotion(**kwargs):
                nonlocal calls
                calls += 1
                if calls == 2:
                    raise OSError("simulated supplement replay failure")
                return real_promote(**kwargs)

            with patch(
                "cli.publish.promote_release",
                side_effect=fail_second_promotion,
            ):
                with self.assertRaisesRegex(
                    OSError,
                    "simulated supplement replay failure",
                ):
                    _promote_release_with_deck_supplement(
                        canonical_release_dir=canonical,
                        deck_supplement_dir=supplement,
                        out_dir=reports,
                        slug="exact-release",
                    )

            for name, content in prior.items():
                self.assertEqual((reports / name).read_bytes(), content)
            self.assertEqual(
                immutable_bundle.stat().st_ino,
                immutable_inode,
            )

    def test_legacy_republish_requires_explicit_opt_in(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            runs = root / "runs"
            archive_stage4 = runs / "2026-07-01-legacy" / "stage4"
            archive_stage4.mkdir(parents=True)
            (archive_stage4 / "legacy.docx").write_bytes(b"legacy bytes")

            refused = publish_all(
                runs_dir=runs,
                out_dir=root / "reports",
                only_slug="legacy",
            )
            self.assertEqual(len(refused), 1)
            self.assertIn("allow-legacy-publish", refused[0][2])

            with (
                patch(
                    "cli.publish.qa_docx",
                    side_effect=self._passing_docx_report,
                ),
                patch(
                    "cli.publish.render_office_artifact",
                    return_value=([], []),
                ),
            ):
                accepted = publish_all(
                    runs_dir=runs,
                    out_dir=root / "reports",
                    only_slug="legacy",
                    allow_legacy=True,
                )
            self.assertEqual(accepted[0][2], "ok (legacy re-QA)")
            self.assertEqual(
                accepted[0][1].read_bytes(),
                b"legacy bytes",
            )

    def test_v2_republish_replays_an_archived_deck_supplement(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            slug = "supplemented"
            archive = root / "runs" / f"2026-07-01-{slug}"
            stage4 = archive / "stage4"
            stage4.mkdir(parents=True)
            (stage4 / f"{slug}.docx").write_bytes(b"word release")

            with (
                patch(
                    "cli.publish.qa_docx",
                    side_effect=self._passing_docx_report,
                ),
                patch(
                    "cli.publish.render_office_artifact",
                    return_value=([], []),
                ),
            ):
                canonical = stage_release_artifacts(
                    stage4_dir=stage4,
                    slug=slug,
                    release_dir=archive / "release",
                )

            deck_stage = root / "deck-stage"
            deck_stage.mkdir()
            deck = Presentation()
            slide = deck.slides.add_slide(deck.slide_layouts[5])
            slide.shapes.title.text = "Approve the bounded supplement"
            deck_path = deck_stage / f"{slug}.pptx"
            deck.save(deck_path)
            expected_deck_bytes = deck_path.read_bytes()
            brief_path = deck_stage / "visual-brief.json"
            brief_path.write_text(
                json.dumps(
                    {
                        "deck_mode": "board_decision",
                        "signature_visual": {
                            "slide_number": 1,
                            "concept": "Bounded supplement",
                            "visual_type": "Decision exhibit",
                        },
                    }
                ),
                encoding="utf-8",
            )
            deck_inspection = deck_stage / "inspection" / slug
            deck_inspection.mkdir(parents=True)
            deck_slide = deck_inspection / f"{slug}-1.png"
            Image.new("RGB", (640, 360), "white").save(deck_slide)
            deck_receipt = (
                deck_stage / f"{slug}-visual-inspection.json"
            )
            prepare_visual_inspection_receipt(
                artifact=deck_path,
                visual_brief=brief_path,
                deck_mode="board_decision",
                rendered_files=[deck_slide],
                receipt_path=deck_receipt,
            )
            receipt_payload = json.loads(
                deck_receipt.read_text(encoding="utf-8")
            )
            receipt_payload["inspection"].update(
                {
                    "full_size_each_slide_inspected": True,
                    "montage_inspected": True,
                    "signature_exhibit_present": True,
                    "signature_exhibit_matches_brief": True,
                    "findings_resolved": True,
                    "status": "pass",
                }
            )
            deck_receipt.write_text(
                json.dumps(receipt_payload),
                encoding="utf-8",
            )
            with patch(
                "cli.presentation_qa.qa_presentation",
                side_effect=self._passing_pptx_report,
            ):
                stage_release_artifacts(
                    stage4_dir=deck_stage,
                    slug=slug,
                    release_dir=(
                        archive / "release-supplements" / "deck"
                    ),
                    include_roles={"presentation"},
                    require_presentation=True,
                    presentation_mode="board_decision",
                    visual_brief=brief_path,
                    require_visual_inspection=True,
                )

            # A durable backfill binds the reader-facing deck, its render QA,
            # and the exact release supplement into the archive transaction.
            (stage4 / f"{slug}.pptx").write_bytes(expected_deck_bytes)
            shutil.copy2(brief_path, stage4 / "visual-brief.json")
            shutil.copy2(
                deck_receipt,
                stage4 / f"{slug}-visual-inspection.json",
            )
            shutil.copytree(
                deck_inspection,
                stage4 / "inspection" / slug,
            )
            (stage4 / f"{slug}-qa.json").write_text(
                '{"ok":true}\n',
                encoding="utf-8",
            )
            render_dir = stage4 / "qa" / f"{slug}-presentation"
            render_dir.mkdir(parents=True)
            (render_dir / "slide-1.png").write_bytes(b"render")
            supplement_dir = archive / "release-supplements" / "deck"
            digest = lambda path: hashlib.sha256(path.read_bytes()).hexdigest()
            (stage4 / f"{slug}-deck-backfill.json").write_text(
                json.dumps(
                    {
                        "schema_version": "2.0",
                        "slug": slug,
                        "qa_ok": True,
                        "artifacts": {
                            "visual_brief": {
                                "path": "stage4/visual-brief.json",
                                "sha256": digest(
                                    stage4 / "visual-brief.json"
                                ),
                            },
                            "presentation": {
                                "path": f"stage4/{slug}.pptx",
                                "sha256": digest(stage4 / f"{slug}.pptx"),
                            },
                            "presentation_qa": {
                                "path": f"stage4/{slug}-qa.json",
                                "sha256": digest(
                                    stage4 / f"{slug}-qa.json"
                                ),
                            },
                            "visual_inspection": {
                                "path": (
                                    f"stage4/{slug}-visual-inspection.json"
                                ),
                                "sha256": digest(
                                    stage4
                                    / f"{slug}-visual-inspection.json"
                                ),
                            },
                            "qa_render": {
                                "path": (
                                    f"stage4/qa/{slug}-presentation"
                                ),
                                "files": {
                                    "slide-1.png": digest(
                                        render_dir / "slide-1.png"
                                    )
                                },
                            },
                            "inspection_render": {
                                "path": f"stage4/inspection/{slug}",
                                "files": {
                                    path.relative_to(
                                        stage4 / "inspection" / slug
                                    ).as_posix(): digest(path)
                                    for path in sorted(
                                        (
                                            stage4
                                            / "inspection"
                                            / slug
                                        ).rglob("*")
                                    )
                                    if path.is_file()
                                },
                            },
                            "release_supplement": {
                                "path": "release-supplements/deck",
                                "files": {
                                    path.relative_to(
                                        supplement_dir
                                    ).as_posix(): digest(path)
                                    for path in sorted(
                                        supplement_dir.rglob("*")
                                    )
                                    if path.is_file()
                                },
                            },
                        },
                    }
                ),
                encoding="utf-8",
            )

            manifest_artifacts = []
            for item in canonical["artifacts"]:
                for suffix, relative, digest in (
                    (
                        "artifact",
                        f"release/{item['path']}",
                        item["sha256"],
                    ),
                    (
                        "qa",
                        f"release/{item['qa_path']}",
                        item["qa_sha256"],
                    ),
                ):
                    manifest_artifacts.append(
                        {
                            "id": f"release/{item['role']}-{suffix}",
                            "path": relative,
                            "required": True,
                            "status": "complete",
                            "sha256": digest,
                        }
                    )
            release_manifest = archive / "release" / "release-manifest.json"
            manifest_artifacts.append(
                {
                    "id": "release/manifest",
                    "path": "release/release-manifest.json",
                    "required": True,
                    "status": "complete",
                    "sha256": hashlib.sha256(
                        release_manifest.read_bytes()
                    ).hexdigest(),
                }
            )
            (archive / "run-manifest.json").write_text(
                json.dumps({"artifacts": manifest_artifacts}),
                encoding="utf-8",
            )

            results = publish_all(
                runs_dir=root / "runs",
                out_dir=root / "reports",
                only_slug=slug,
            )

            self.assertEqual(results[0][2], "ok")
            self.assertEqual(
                (root / "reports" / f"{slug}.pptx").read_bytes(),
                expected_deck_bytes,
            )
            self.assertTrue(
                (
                    root
                    / "reports"
                    / f"{slug}-deck-release-manifest.json"
                ).is_file()
            )
            pointer = json.loads(
                (
                    root
                    / "reports"
                    / f"{slug}-deck-release-manifest.json"
                ).read_text(encoding="utf-8")
            )
            self.assertEqual(
                pointer["source_release_manifest_sha256"],
                hashlib.sha256(
                    (
                        supplement_dir / "release-manifest.json"
                    ).read_bytes()
                ).hexdigest(),
            )

    def test_discover_reports_requires_exact_slug_file_and_keeps_newest(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            runs = Path(directory) / "runs"

            old = runs / "2026-01-02-airfield-plan" / "stage4"
            old.mkdir(parents=True)
            (old / "airfield-plan.docx").write_bytes(b"old")

            newest = runs / "2026-06-02-airfield-plan" / "stage4"
            newest.mkdir(parents=True)
            (newest / "airfield-plan.docx").write_bytes(b"new")

            # A later directory with a loose filename match is not a valid
            # release for this slug.
            wrong = runs / "2026-07-02-airfield-plan" / "stage4"
            wrong.mkdir(parents=True)
            (wrong / "unrelated-airfield-plan.docx").write_bytes(b"wrong")

            # A neighboring slug must not collide with the exact target.
            neighbor = runs / "2026-07-03-airfield-plan-update" / "stage4"
            neighbor.mkdir(parents=True)
            (neighbor / "airfield-plan-update.docx").write_bytes(b"neighbor")

            discovered = {source.slug: source for source in discover_reports(runs)}
            self.assertEqual(set(discovered), {"airfield-plan", "airfield-plan-update"})
            selected = discovered["airfield-plan"]
            self.assertEqual(
                selected.archive_dir.name,
                "2026-06-02-airfield-plan",
            )
            self.assertEqual(selected.stage4_docx.read_bytes(), b"new")


class PersistentSourceTests(unittest.TestCase):
    def test_attached_sources_are_persistent_and_not_rediscovered(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            dropzone = root / "sources"
            persistent = dropzone / "runs" / "prior-run"
            persistent.mkdir(parents=True)
            (persistent / "already-attached.md").write_text(
                "prior evidence", encoding="utf-8"
            )
            incoming = dropzone / "airport-data.md"
            incoming.write_text("new airport evidence", encoding="utf-8")
            outputs = root / "outputs"
            outputs.mkdir()

            self.assertEqual(discover_dropzone(dropzone), [incoming])
            attached = attach_sources("new-run", [incoming], outputs)

            expected = (
                root / "sources" / "runs" / "new-run" / "airport-data.md"
            )
            self.assertEqual(len(attached), 1)
            self.assertEqual(attached[0].original, expected)
            self.assertEqual(attached[0].readable, expected)
            self.assertEqual(expected.read_text(encoding="utf-8"), "new airport evidence")
            self.assertFalse(incoming.exists())
            self.assertEqual(discover_dropzone(dropzone), [])


class ArchiveTransactionTests(unittest.TestCase):
    @staticmethod
    def _write_minimal_manifest(outputs: Path, artifact: Path) -> Path:
        digest = hashlib.sha256(artifact.read_bytes()).hexdigest()
        manifest = outputs / "run-manifest.json"
        manifest.write_text(
            json.dumps(
                {
                    "run": {},
                    "artifacts": [
                        {
                            "id": "stage4/audit-note",
                            "path": artifact.relative_to(outputs).as_posix(),
                            "role": "audit_note",
                            "required": True,
                            "status": "complete",
                            "sha256": digest,
                            "contract": {
                                "kind": "markdown",
                                "min_words": 0,
                                "min_records": 0,
                                "required_keys": [],
                                "required_any": [],
                                "optional": False,
                            },
                        }
                    ],
                }
            ),
            encoding="utf-8",
        )
        return manifest

    def test_archive_preserves_exact_run_prompt(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            outputs = root / "outputs"
            stage4 = outputs / "stage4"
            stage4.mkdir(parents=True)
            (stage4 / "airport-decision.docx").write_bytes(b"released report")
            run_file = root / "prompts" / "runs" / "airport-decision.md"
            run_file.parent.mkdir(parents=True)
            prompt_bytes = b"# Run: Airport decision\n\n## Thesis\nExact prompt.\n"
            run_file.write_bytes(prompt_bytes)

            archive = archive_run(
                repo_root=root,
                slug="airport-decision",
                tally=CostTally(),
                run_file=run_file,
            )

            self.assertEqual((archive / "run-prompt.md").read_bytes(), prompt_bytes)
            self.assertEqual(
                (archive / "stage4" / "airport-decision.docx").read_bytes(),
                b"released report",
            )
            self.assertFalse(stage4.exists())
            self.assertTrue((outputs / ".gitkeep").is_file())

    def test_archive_failure_leaves_outputs_resumable(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            outputs = root / "outputs"
            stage4 = outputs / "stage4"
            stage4.mkdir(parents=True)
            report = stage4 / "airport-decision.docx"
            report.write_bytes(b"work that must survive")
            run_file = root / "prompt.md"
            run_file.write_text("# Run: Airport decision", encoding="utf-8")

            with (
                patch(
                    "cli.sources.archive_sources",
                    side_effect=RuntimeError("synthetic archive failure"),
                ),
                self.assertRaisesRegex(RuntimeError, "synthetic archive failure"),
            ):
                archive_run(
                    repo_root=root,
                    slug="airport-decision",
                    tally=CostTally(),
                    run_file=run_file,
                )

            self.assertEqual(report.read_bytes(), b"work that must survive")
            self.assertTrue(run_file.is_file())
            archive_name = f"{date.today().isoformat()}-airport-decision"
            final_archive = root / "runs" / archive_name
            self.assertFalse(final_archive.exists())
            temporary_archives = list(
                (root / "runs").glob(f".{archive_name}-*")
            )
            self.assertEqual(temporary_archives, [])

    def test_archive_rechecks_manifest_artifacts_after_copy(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            artifact = root / "outputs" / "stage4" / "audit-note.md"
            artifact.parent.mkdir(parents=True)
            artifact.write_text("validated release bytes", encoding="utf-8")
            manifest = self._write_minimal_manifest(
                root / "outputs",
                artifact,
            )

            def mutate_archived_copy(
                slug,
                outputs_dir,
                archive_dir,
                source_material=None,
            ):
                del slug, outputs_dir, source_material
                (archive_dir / "stage4" / "audit-note.md").write_text(
                    "changed during copy",
                    encoding="utf-8",
                )
                return []

            with (
                patch(
                    "cli.sources.archive_sources",
                    side_effect=mutate_archived_copy,
                ),
                self.assertRaisesRegex(
                    RuntimeError,
                    "Archived generated artifact does not match",
                ),
            ):
                archive_run(
                    repo_root=root,
                    slug="copy-race",
                    tally=CostTally(),
                    manifest_path=manifest,
                )

            self.assertTrue(artifact.is_file())
            self.assertFalse(
                (
                    root
                    / "runs"
                    / f"{date.today().isoformat()}-copy-race"
                ).exists()
            )

    def test_cleanup_retry_recognises_an_already_committed_archive(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            artifact = root / "outputs" / "stage4" / "audit-note.md"
            artifact.parent.mkdir(parents=True)
            artifact.write_text("validated release bytes", encoding="utf-8")
            manifest = self._write_minimal_manifest(
                root / "outputs",
                artifact,
            )

            with (
                patch(
                    "cli.archive._clear_outputs",
                    side_effect=OSError("synthetic cleanup failure"),
                ),
                self.assertRaisesRegex(OSError, "synthetic cleanup failure"),
            ):
                archive_run(
                    repo_root=root,
                    slug="cleanup-retry",
                    tally=CostTally(),
                    manifest_path=manifest,
                )

            archive = archive_run(
                repo_root=root,
                slug="cleanup-retry",
                tally=CostTally(),
                manifest_path=manifest,
            )
            self.assertTrue(archive.is_dir())
            self.assertFalse(manifest.exists())
            self.assertTrue((root / "outputs" / ".gitkeep").is_file())

    def test_cleanup_retry_refuses_a_corrupt_committed_archive(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            artifact = root / "outputs" / "stage4" / "audit-note.md"
            artifact.parent.mkdir(parents=True)
            artifact.write_text("validated release bytes", encoding="utf-8")
            manifest = self._write_minimal_manifest(
                root / "outputs",
                artifact,
            )

            with (
                patch(
                    "cli.archive._clear_outputs",
                    side_effect=OSError("synthetic cleanup failure"),
                ),
                self.assertRaisesRegex(OSError, "synthetic cleanup failure"),
            ):
                archive_run(
                    repo_root=root,
                    slug="corrupt-cleanup-retry",
                    tally=CostTally(),
                    manifest_path=manifest,
                )

            archive = (
                root
                / "runs"
                / f"{date.today().isoformat()}-corrupt-cleanup-retry"
            )
            archived_artifact = archive / "stage4" / "audit-note.md"
            archived_artifact.write_text(
                "corrupt committed bytes",
                encoding="utf-8",
            )

            with (
                patch("cli.archive._clear_outputs") as clear_outputs,
                self.assertRaisesRegex(
                    RuntimeError,
                    "Archived generated artifact does not match",
                ),
            ):
                archive_run(
                    repo_root=root,
                    slug="corrupt-cleanup-retry",
                    tally=CostTally(),
                    manifest_path=manifest,
                )

            clear_outputs.assert_not_called()
            self.assertEqual(
                artifact.read_text(encoding="utf-8"),
                "validated release bytes",
            )
            self.assertTrue(manifest.is_file())


class ManifestCommitGateTests(unittest.TestCase):
    def test_manifest_commit_gate_detects_post_validation_mutation(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            outputs = Path(directory) / "outputs"
            artifact = outputs / "stage4" / "publishing-quality.json"
            artifact.parent.mkdir(parents=True)
            artifact.write_text(
                json.dumps(
                    {
                        "artifact": "report.docx",
                        "kind": "docx",
                        "ok": True,
                        "issues": [],
                    }
                ),
                encoding="utf-8",
            )
            expected_hash = hashlib.sha256(artifact.read_bytes()).hexdigest()
            manifest = outputs / "run-manifest.json"
            manifest.write_text(
                json.dumps(
                    {
                        "artifacts": [
                            {
                                "id": "stage4/publishing-quality",
                                "path": "stage4/publishing-quality.json",
                                "role": "publishing_quality",
                                "required": True,
                                "status": "complete",
                                "sha256": expected_hash,
                                "contract": {
                                    "kind": "json",
                                    "min_words": 0,
                                    "min_records": 0,
                                    "required_keys": [
                                        "artifact",
                                        "kind",
                                        "ok",
                                        "issues",
                                    ],
                                    "required_any": [],
                                    "optional": False,
                                },
                            }
                        ]
                    }
                ),
                encoding="utf-8",
            )

            assert_manifest_complete(manifest)
            artifact.write_text(
                json.dumps(
                    {
                        "artifact": "report.docx",
                        "kind": "docx",
                        "ok": True,
                        "issues": [],
                        "changed_after_validation": True,
                    }
                ),
                encoding="utf-8",
            )

            with self.assertRaisesRegex(
                RuntimeError, "current bytes do not match"
            ):
                assert_manifest_complete(manifest)


if __name__ == "__main__":
    unittest.main()
