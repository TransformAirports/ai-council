from __future__ import annotations

import asyncio
import json
import os
import shutil
import tempfile
import unittest
from pathlib import Path
from unittest.mock import AsyncMock, patch

from PIL import Image
from pptx import Presentation

from cli.agents import Agent
from cli.archive import _verify_committed_archive
from cli.evidence import file_sha256
from cli.orchestrator import (
    _build_presentation_backfill,
    _commit_presentation_backfill_archive,
    _presentation_backfill_identity,
    _promote_archive_backfill,
    _publish_presentation_backfill_release,
    _staged_presentation_release_matches_sources,
    run_presentation_for_archive,
)
from cli.presentation_qa import prepare_visual_inspection_receipt


class BackfillArchiveTransactionTests(unittest.TestCase):
    slug = "sample-run"

    def _staged_bundle(self, root: Path) -> Path:
        stage4 = root / "staged" / "stage4"
        render = stage4 / "qa" / f"{self.slug}-presentation"
        render.mkdir(parents=True)
        visual_brief = stage4 / "deck-backfill" / "visual-brief.json"
        visual_brief.parent.mkdir(parents=True)
        visual_brief.write_text(
            json.dumps(
                {
                    "title": "new visual",
                    "deck_mode": "board_decision",
                    "signature_visual": {
                        "slide_number": 1,
                        "concept": "Bounded decision",
                        "visual_type": "Decision exhibit",
                    },
                }
            )
            + "\n",
            encoding="utf-8",
        )
        deck = Presentation()
        slide = deck.slides.add_slide(deck.slide_layouts[5])
        slide.shapes.title.text = "A bounded decision"
        deck.save(stage4 / f"{self.slug}.pptx")
        deck_bytes = (stage4 / f"{self.slug}.pptx").read_bytes()
        (stage4 / f"{self.slug}-qa.json").write_text(
            '{"ok":true}\n', encoding="utf-8"
        )
        (render / "slide-1.png").write_bytes(b"new render")
        inspection = stage4 / "inspection" / self.slug
        inspection.mkdir(parents=True)
        Image.new("RGB", (320, 180), "white").save(
            inspection / "slide-1.png"
        )
        receipt = stage4 / f"{self.slug}-visual-inspection.json"
        prepare_visual_inspection_receipt(
            artifact=stage4 / f"{self.slug}.pptx",
            visual_brief=visual_brief,
            deck_mode="board_decision",
            rendered_files=[inspection / "slide-1.png"],
            receipt_path=receipt,
        )
        receipt_payload = json.loads(receipt.read_text(encoding="utf-8"))
        receipt_payload["inspection"].update(
            {
                "full_size_each_slide_inspected": True,
                "montage_inspected": True,
                "signature_exhibit_present": True,
                "signature_exhibit_matches_brief": True,
                "findings_resolved": True,
                "status": "pass",
            }
        )
        receipt.write_text(json.dumps(receipt_payload), encoding="utf-8")
        release = stage4.parent / "release"
        release_qa = release / "qa"
        release_qa.mkdir(parents=True)
        release_deck = release / f"{self.slug}.pptx"
        shutil.copy2(stage4 / f"{self.slug}.pptx", release_deck)
        release_visual_brief = (
            release / "deck-backfill" / "visual-brief.json"
        )
        release_visual_brief.parent.mkdir(parents=True)
        shutil.copy2(visual_brief, release_visual_brief)
        shutil.copy2(receipt, release / receipt.name)
        shutil.copytree(inspection, release / "inspection" / self.slug)
        release_qa_path = release_qa / f"{self.slug}.pptx.qa.json"
        release_qa_path.write_text(
            '{"artifact":"sample-run.pptx","kind":"pptx","ok":true,'
            '"issues":[],"rendered_files":[]}\n',
            encoding="utf-8",
        )
        release_hash = file_sha256(release_deck)
        release_qa_hash = file_sha256(release_qa_path)
        (release / "release-manifest.json").write_text(
            json.dumps(
                {
                    "schema_version": "1.0",
                    "slug": self.slug,
                    "status": "ready",
                    "artifacts": [
                        {
                            "role": "presentation",
                            "kind": "pptx",
                            "source_path": f"stage4/{self.slug}.pptx",
                            "path": release_deck.name,
                            "source_sha256": release_hash,
                            "sha256": release_hash,
                            "qa_path": (
                                release_qa_path.relative_to(release).as_posix()
                            ),
                            "qa_sha256": release_qa_hash,
                            "qa_ok": True,
                            "rendered_files": [],
                            "visual_inspection": {
                                "type": "presentation_slides",
                                "path": receipt.name,
                                "sha256": file_sha256(
                                    release / receipt.name
                                ),
                                "visual_brief_path": (
                                    "deck-backfill/visual-brief.json"
                                ),
                                "visual_brief_sha256": file_sha256(
                                    release_visual_brief
                                ),
                                "files": [
                                    {
                                        "path": path.relative_to(
                                            release
                                        ).as_posix(),
                                        "sha256": file_sha256(path),
                                        "size_bytes": path.stat().st_size,
                                    }
                                    for path in sorted(
                                        (release / "inspection" / self.slug).rglob("*")
                                    )
                                    if path.is_file()
                                ],
                            },
                        }
                    ],
                    "requirements": {
                        "executive_summary": False,
                        "presentation": True,
                        "visual_inspection": True,
                        "word_visual_inspection": False,
                    },
                }
            ),
            encoding="utf-8",
        )

        artifacts = {
            "visual_brief": {
                "path": "stage4/deck-backfill/visual-brief.json",
                "sha256": file_sha256(visual_brief),
            },
            "presentation": {
                "path": f"stage4/{self.slug}.pptx",
                "sha256": file_sha256(stage4 / f"{self.slug}.pptx")
            },
            "presentation_qa": {
                "path": f"stage4/{self.slug}-qa.json",
                "sha256": file_sha256(stage4 / f"{self.slug}-qa.json")
            },
            "visual_inspection": {
                "path": f"stage4/{self.slug}-visual-inspection.json",
                "sha256": file_sha256(receipt),
            },
            "qa_render": {
                "path": f"stage4/qa/{self.slug}-presentation",
                "files": {
                    "slide-1.png": file_sha256(render / "slide-1.png"),
                }
            },
            "inspection_render": {
                "path": f"stage4/inspection/{self.slug}",
                "files": {
                    path.relative_to(inspection).as_posix(): file_sha256(path)
                    for path in sorted(inspection.rglob("*"))
                    if path.is_file()
                },
            },
            "release_supplement": {
                "path": "release-supplements/deck",
                "files": {
                    path.relative_to(release).as_posix(): file_sha256(path)
                    for path in sorted(release.rglob("*"))
                    if path.is_file()
                }
            },
        }
        (stage4 / f"{self.slug}-deck-backfill.json").write_text(
            json.dumps(
                {
                    "schema_version": "3.0",
                    "slug": self.slug,
                    "deck_mode": "board_decision",
                    "artifacts": artifacts,
                    "qa_ok": True,
                }
            ),
            encoding="utf-8",
        )
        return stage4

    def test_valid_release_bundle_is_reused_only_for_exact_sources(
        self,
    ) -> None:
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            staged = self._staged_bundle(root)
            release = staged.parent / "release"
            manifest_hash = file_sha256(
                release / "release-manifest.json"
            )

            self.assertTrue(
                _staged_presentation_release_matches_sources(
                    release_dir=release,
                    staged_stage4=staged,
                    slug=self.slug,
                )
            )
            self.assertEqual(
                file_sha256(release / "release-manifest.json"),
                manifest_hash,
            )

            (
                staged / "deck-backfill" / "visual-brief.json"
            ).write_text(
                '{"changed":true}\n',
                encoding="utf-8",
            )
            self.assertFalse(
                _staged_presentation_release_matches_sources(
                    release_dir=release,
                    staged_stage4=staged,
                    slug=self.slug,
                )
            )

    def _old_archive(self, root: Path) -> Path:
        stage4 = root / "archive" / "stage4"
        render = stage4 / "qa" / f"{self.slug}-presentation"
        render.mkdir(parents=True)
        (stage4 / "visual-brief.json").write_bytes(b"old visual")
        (stage4 / f"{self.slug}.pptx").write_bytes(b"old deck")
        (stage4 / f"{self.slug}-qa.json").write_bytes(b"old qa")
        (stage4 / f"{self.slug}-deck-backfill.json").write_bytes(
            b"old manifest"
        )
        (render / "slide-1.png").write_bytes(b"old render")
        supplement = stage4.parent / "release-supplements" / "deck"
        supplement.mkdir(parents=True)
        (supplement / "old-release.txt").write_bytes(b"old release")
        (stage4 / f"{self.slug}.docx").write_bytes(b"word stays")
        return stage4

    def test_promotes_exact_backfill_set_and_preserves_word_artifact(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            staged = self._staged_bundle(root)
            archive = self._old_archive(root)
            deck_bytes = (staged / f"{self.slug}.pptx").read_bytes()

            promoted = _promote_archive_backfill(
                staged_stage4=staged,
                archive_stage4=archive,
                slug=self.slug,
            )

            self.assertEqual(
                (archive / f"{self.slug}.pptx").read_bytes(), deck_bytes
            )
            self.assertEqual(
                (archive / "visual-brief.json").read_bytes(), b"old visual"
            )
            self.assertEqual(
                (
                    archive
                    / "deck-backfill"
                    / "visual-brief.json"
                ).read_bytes(),
                (
                    staged
                    / "deck-backfill"
                    / "visual-brief.json"
                ).read_bytes(),
            )
            self.assertEqual(
                (
                    archive
                    / "qa"
                    / f"{self.slug}-presentation"
                    / "slide-1.png"
                ).read_bytes(),
                b"new render",
            )
            self.assertEqual(
                (archive / f"{self.slug}.docx").read_bytes(), b"word stays"
            )
            self.assertEqual(
                promoted["backfill_manifest"],
                archive / f"{self.slug}-deck-backfill.json",
            )
            self.assertEqual(
                (staged / f"{self.slug}.pptx").read_bytes(),
                deck_bytes,
            )
            self.assertEqual(
                (
                    archive.parent
                    / "release-supplements"
                    / "deck"
                    / f"{self.slug}.pptx"
                ).read_bytes(),
                deck_bytes,
            )

    def test_backfill_preserves_canonical_visual_and_run_manifest(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            archive_dir = (
                root / "runs" / "2026-07-23-sample-run"
            )
            archive_stage4 = archive_dir / "stage4"
            archive_stage4.mkdir(parents=True)
            canonical_visual = archive_stage4 / "visual-brief.json"
            canonical_visual.write_text(
                '{"canonical":"original","deck_mode":"board_decision"}\n',
                encoding="utf-8",
            )
            prompt = archive_dir / "run-prompt.md"
            prompt.write_text(
                "# Run: Sample\n\nCanonical archived prompt.\n",
                encoding="utf-8",
            )
            (archive_dir / "retrospective.md").write_text(
                "# Retrospective\n",
                encoding="utf-8",
            )
            canonical_record = {
                "id": "stage4/visual-brief",
                "path": "stage4/visual-brief.json",
                "role": "visual_brief",
                "required": True,
                "status": "complete",
                "sha256": file_sha256(canonical_visual),
                "contract": {
                    "kind": "json",
                    "min_words": 0,
                    "min_records": 0,
                    "required_keys": [],
                    "required_any": [],
                    "optional": False,
                },
            }
            manifest_payload = {
                "schema_version": "2.0",
                "run": {
                    "slug": self.slug,
                    "run_prompt_sha256": file_sha256(prompt),
                    "run_prompt_size": prompt.stat().st_size,
                    "source_material": [],
                    "source_library": [],
                },
                "artifacts": [canonical_record],
            }
            manifest = archive_dir / "run-manifest.json"
            manifest.write_text(
                json.dumps(manifest_payload, indent=2) + "\n",
                encoding="utf-8",
            )
            manifest_bytes = manifest.read_bytes()
            manifest_sha256 = file_sha256(manifest)
            canonical_bytes = canonical_visual.read_bytes()
            canonical_status = canonical_record["status"]
            canonical_sha256 = canonical_record["sha256"]

            _verify_committed_archive(
                archive_dir,
                manifest_payload,
                slug=self.slug,
                expected_manifest_sha256=manifest_sha256,
            )
            staged = self._staged_bundle(root / "backfill")
            staged_visual_bytes = (
                staged
                / "deck-backfill"
                / "visual-brief.json"
            ).read_bytes()

            _promote_archive_backfill(
                staged_stage4=staged,
                archive_stage4=archive_stage4,
                slug=self.slug,
            )

            self.assertEqual(canonical_visual.read_bytes(), canonical_bytes)
            self.assertEqual(manifest.read_bytes(), manifest_bytes)
            current_manifest = json.loads(
                manifest.read_text(encoding="utf-8")
            )
            current_record = current_manifest["artifacts"][0]
            self.assertEqual(current_record["status"], canonical_status)
            self.assertEqual(current_record["sha256"], canonical_sha256)
            self.assertEqual(
                (
                    archive_stage4
                    / "deck-backfill"
                    / "visual-brief.json"
                ).read_bytes(),
                staged_visual_bytes,
            )
            _verify_committed_archive(
                archive_dir,
                manifest_payload,
                slug=self.slug,
                expected_manifest_sha256=manifest_sha256,
            )

    def test_restores_every_prior_artifact_when_commit_fails(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            staged = self._staged_bundle(root)
            archive = self._old_archive(root)
            staged_deck = staged / f"{self.slug}.pptx"
            staged_deck_bytes = staged_deck.read_bytes()
            archive_deck = archive / f"{self.slug}.pptx"
            real_replace = os.replace

            def fail_during_deck_promotion(source, destination):
                if (
                    Path(source).parent.name == "staged"
                    and Path(source).name.endswith(staged_deck.name)
                    and Path(destination) == archive_deck
                ):
                    raise OSError("simulated archive commit failure")
                return real_replace(source, destination)

            with patch(
                "cli.orchestrator.os.replace",
                side_effect=fail_during_deck_promotion,
            ):
                with self.assertRaisesRegex(
                    OSError, "simulated archive commit failure"
                ):
                    _promote_archive_backfill(
                        staged_stage4=staged,
                        archive_stage4=archive,
                        slug=self.slug,
                    )

            self.assertEqual(
                (archive / "visual-brief.json").read_bytes(), b"old visual"
            )
            self.assertEqual(archive_deck.read_bytes(), b"old deck")
            self.assertEqual(
                (archive / f"{self.slug}-qa.json").read_bytes(), b"old qa"
            )
            self.assertEqual(
                (
                    archive
                    / "qa"
                    / f"{self.slug}-presentation"
                    / "slide-1.png"
                ).read_bytes(),
                b"old render",
            )
            self.assertEqual(
                (
                    archive / f"{self.slug}-deck-backfill.json"
                ).read_bytes(),
                b"old manifest",
            )
            self.assertEqual(
                (
                    archive.parent
                    / "release-supplements"
                    / "deck"
                    / "old-release.txt"
                ).read_bytes(),
                b"old release",
            )
            self.assertEqual(staged_deck.read_bytes(), staged_deck_bytes)
            self.assertTrue((staged.parent / "release").is_dir())

    def test_rejects_a_release_supplement_not_bound_to_the_manifest(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            staged = self._staged_bundle(root)
            archive = self._old_archive(root)
            (staged.parent / "release" / f"{self.slug}.pptx").write_bytes(
                b"substituted after backfill binding"
            )

            with self.assertRaisesRegex(
                RuntimeError,
                "release-supplement inventory mismatch",
            ):
                _promote_archive_backfill(
                    staged_stage4=staged,
                    archive_stage4=archive,
                    slug=self.slug,
                )

            self.assertEqual(
                (archive / f"{self.slug}.pptx").read_bytes(),
                b"old deck",
            )

    def test_staging_copy_failure_leaves_archive_and_sources_untouched(
        self,
    ) -> None:
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            staged = self._staged_bundle(root)
            archive = self._old_archive(root)
            real_copy2 = shutil.copy2
            staged_deck_bytes = (staged / f"{self.slug}.pptx").read_bytes()

            def fail_staging_copy(source, destination, *args, **kwargs):
                if Path(source) == staged / f"{self.slug}.pptx":
                    raise OSError("simulated staging copy failure")
                return real_copy2(source, destination, *args, **kwargs)

            with patch(
                "cli.orchestrator.shutil.copy2",
                side_effect=fail_staging_copy,
            ):
                with self.assertRaisesRegex(
                    OSError,
                    "simulated staging copy failure",
                ):
                    _promote_archive_backfill(
                        staged_stage4=staged,
                        archive_stage4=archive,
                        slug=self.slug,
                    )

            self.assertEqual(
                (archive / f"{self.slug}.pptx").read_bytes(),
                b"old deck",
            )
            self.assertEqual(
                (staged / f"{self.slug}.pptx").read_bytes(),
                staged_deck_bytes,
            )
            self.assertFalse(
                any(
                    archive.glob(
                        f".{self.slug}-backfill-commit-*"
                    )
                )
            )

    def test_completed_archive_commit_is_replayed_without_rebuilding(
        self,
    ) -> None:
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            staged = self._staged_bundle(root)
            archive_dir = root / "archive"
            archive_stage4 = archive_dir / "stage4"
            _promote_archive_backfill(
                staged_stage4=staged,
                archive_stage4=archive_stage4,
                slug=self.slug,
            )
            published = root / "reports" / f"{self.slug}.pptx"
            build = AsyncMock()

            with (
                patch(
                    "cli.publish.promote_release",
                    return_value={"presentation": published},
                ) as promote,
                patch(
                    "cli.orchestrator._build_presentation_backfill",
                    new=build,
                ),
            ):
                result = asyncio.run(
                    run_presentation_for_archive(
                        archive_dir=archive_dir,
                        slug=self.slug,
                        title="Sample run",
                        repo_root=root,
                    )
                )

            self.assertEqual(result, published)
            promote.assert_called_once()
            build.assert_not_awaited()

    def test_partial_archive_commit_resumes_from_valid_durable_staging(
        self,
    ) -> None:
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            seed = self._staged_bundle(root / "seed")
            staging_root = (
                root / "logs" / "deck-backfills" / self.slug
            )
            shutil.copytree(seed.parent, staging_root)
            archive_dir = root / "archive"
            archive_stage4 = archive_dir / "stage4"
            archive_stage4.mkdir(parents=True)
            shutil.copy2(
                staging_root / "stage4" / f"{self.slug}.pptx",
                archive_stage4 / f"{self.slug}.pptx",
            )
            expected = root / "reports" / f"{self.slug}.pptx"
            build = AsyncMock(return_value=expected)

            with patch(
                "cli.orchestrator._build_presentation_backfill",
                new=build,
            ):
                result = asyncio.run(
                    run_presentation_for_archive(
                        archive_dir=archive_dir,
                        slug=self.slug,
                        title="Sample run",
                        repo_root=root,
                    )
                )

            self.assertEqual(result, expected)
            build.assert_awaited_once()
            self.assertFalse(staging_root.exists())


class DurableBackfillStateTests(unittest.IsolatedAsyncioTestCase):
    async def test_source_mutation_blocks_publish_and_archive_commit(
        self,
    ) -> None:
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            archive = root / "runs" / "2026-07-23-sample-run"
            (archive / "stage3").mkdir(parents=True)
            final = archive / "stage3" / "final-draft.md"
            final.write_text(
                " ".join(["draft"] * 50), encoding="utf-8"
            )
            ledger = archive / "evidence-ledger.jsonl"
            ledger.write_text('{"evidence_id":"E1"}\n', encoding="utf-8")
            run_prompt = archive / "run-prompt.md"
            run_prompt.write_text("# Run: Sample", encoding="utf-8")
            charters = root / "agents"
            charters.mkdir()
            art_path = charters / "art-director.md"
            designer_path = charters / "presentation-designer.md"
            art_path.write_text("art charter", encoding="utf-8")
            designer_path.write_text("designer charter", encoding="utf-8")
            art = Agent(
                name="art-director",
                display_name="Art Director",
                description="test",
                tools=(),
                order=1,
                system_prompt="art charter",
                path=art_path,
            )
            designer = Agent(
                name="presentation-designer",
                display_name="Presentation Designer",
                description="test",
                tools=(),
                order=2,
                system_prompt="designer charter",
                path=designer_path,
            )
            stage4 = root / "staging" / "stage4"
            qa_render = stage4 / "qa" / "sample-run-presentation"
            inspection = stage4 / "inspection" / "sample-run"
            qa_render.mkdir(parents=True)
            inspection.mkdir(parents=True)
            visual = stage4 / "visual-brief.json"
            deck = stage4 / "sample-run.pptx"
            qa = stage4 / "sample-run-qa.json"
            receipt = stage4 / "sample-run-visual-inspection.json"
            visual.write_text('{"visual":"brief"}', encoding="utf-8")
            deck.write_bytes(b"deck")
            qa.write_text('{"ok":true}', encoding="utf-8")
            receipt.write_text('{"inspection":"pass"}', encoding="utf-8")
            (qa_render / "slide-1.png").write_bytes(b"qa render")
            (inspection / "slide-1.png").write_bytes(
                b"inspection render"
            )

            with patch(
                "cli.orchestrator.build_execution_contract_fingerprint",
                return_value=[{"path": "contract", "sha256": "stable"}],
            ):
                payload, identity, _, _ = (
                    _presentation_backfill_identity(
                        archive_dir=archive,
                        slug="sample-run",
                        title="Sample run",
                        repo_root=root,
                        art_director=art,
                        designer=designer,
                    )
                )
                state = {
                    "visual_brief_sha256": file_sha256(visual),
                    "presentation_sha256": file_sha256(deck),
                    "presentation_qa_sha256": file_sha256(qa),
                    "visual_inspection_sha256": file_sha256(receipt),
                    "qa_render_files": {
                        "slide-1.png": file_sha256(
                            qa_render / "slide-1.png"
                        )
                    },
                    "inspection_render_files": {
                        "slide-1.png": file_sha256(
                            inspection / "slide-1.png"
                        )
                    },
                }
                final.write_text(
                    " ".join(["mutated"] * 50), encoding="utf-8"
                )
                with (
                    patch("cli.publish.stage_release_artifacts") as stage,
                    patch("cli.publish.promote_release") as promote,
                    self.assertRaisesRegex(
                        RuntimeError,
                        "identity changed before publication",
                    ),
                ):
                    _publish_presentation_backfill_release(
                        archive_dir=archive,
                        slug="sample-run",
                        title="Sample run",
                        repo_root=root,
                        art_director=art,
                        designer=designer,
                        expected_payload=payload,
                        expected_identity=identity,
                        backfill_state=state,
                        stage4=stage4,
                        visual_path=visual,
                        out_path=deck,
                        qa_path=qa,
                        receipt_path=receipt,
                        qa_render_dir=qa_render,
                        inspection_dir=inspection,
                        release_dir=root / "staging" / "release",
                        deck_mode="board_decision",
                        out_dir=root / "reports",
                    )
                stage.assert_not_called()
                promote.assert_not_called()
                with (
                    patch(
                        "cli.orchestrator._promote_archive_backfill"
                    ) as archive_promote,
                    self.assertRaisesRegex(
                        RuntimeError,
                        "identity changed before publication",
                    ),
                ):
                    _commit_presentation_backfill_archive(
                        archive_dir=archive,
                        slug="sample-run",
                        title="Sample run",
                        repo_root=root,
                        art_director=art,
                        designer=designer,
                        expected_payload=payload,
                        expected_identity=identity,
                        backfill_state=state,
                        visual_path=visual,
                        out_path=deck,
                        qa_path=qa,
                        receipt_path=receipt,
                        qa_render_dir=qa_render,
                        inspection_dir=inspection,
                        staged_stage4=stage4,
                        archive_stage4=archive / "stage4",
                    )
                archive_promote.assert_not_called()

    async def test_execution_contract_change_quarantines_staged_work(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            archive = root / "runs" / "2026-07-23-sample-run"
            (archive / "stage3").mkdir(parents=True)
            (archive / "stage3" / "final-draft.md").write_text(
                " ".join(["draft"] * 50), encoding="utf-8"
            )
            (archive / "evidence-ledger.jsonl").write_text(
                '{"evidence_id":"E1","agent_id":"test","claim":"claim",'
                '"source_title":"source","source_type":"primary",'
                '"is_primary":true,"confidence":"high",'
                '"source_url":"https://example.com"}\n',
                encoding="utf-8",
            )
            run_prompt = root / "prompts" / "runs" / "sample-run.md"
            run_prompt.parent.mkdir(parents=True)
            run_prompt.write_text("# Run: Sample", encoding="utf-8")
            charters = root / "agents"
            charters.mkdir()
            art_path = charters / "art-director.md"
            design_path = charters / "presentation-designer.md"
            art_path.write_text("art charter", encoding="utf-8")
            design_path.write_text("design charter", encoding="utf-8")
            agents = [
                Agent(
                    name="art-director",
                    display_name="Art Director",
                    description="test",
                    tools=(),
                    order=1,
                    system_prompt="art",
                    path=art_path,
                ),
                Agent(
                    name="presentation-designer",
                    display_name="Presentation Designer",
                    description="test",
                    tools=(),
                    order=2,
                    system_prompt="design",
                    path=design_path,
                ),
            ]
            staging = (
                root / "logs" / "deck-backfills" / "sample-run"
            )
            stop = AsyncMock(side_effect=RuntimeError("stop after state"))

            with (
                patch(
                    "cli.orchestrator.load_all_agents",
                    return_value=agents,
                ),
                patch(
                    "cli.orchestrator.build_execution_contract_fingerprint",
                    return_value=[{"path": "contract", "sha256": "v1"}],
                ),
                patch("cli.orchestrator._run_agent", new=stop),
            ):
                with self.assertRaisesRegex(
                    RuntimeError, "stop after state"
                ):
                    await _build_presentation_backfill(
                        archive_dir=archive,
                        slug="sample-run",
                        title="Sample run",
                        repo_root=root,
                        staging_root=staging,
                    )

            first_state = json.loads(
                (staging / "state.json").read_text(encoding="utf-8")
            )
            sentinel = staging / "stage4" / "stale-deck.pptx"
            sentinel.write_bytes(b"prior generation")

            with (
                patch(
                    "cli.orchestrator.load_all_agents",
                    return_value=agents,
                ),
                patch(
                    "cli.orchestrator.build_execution_contract_fingerprint",
                    return_value=[{"path": "contract", "sha256": "v2"}],
                ),
                patch(
                    "cli.orchestrator._run_agent",
                    new=AsyncMock(
                        side_effect=RuntimeError("stop after state")
                    ),
                ),
            ):
                with self.assertRaisesRegex(
                    RuntimeError, "stop after state"
                ):
                    await _build_presentation_backfill(
                        archive_dir=archive,
                        slug="sample-run",
                        title="Sample run",
                        repo_root=root,
                        staging_root=staging,
                    )

            second_state = json.loads(
                (staging / "state.json").read_text(encoding="utf-8")
            )
            self.assertNotEqual(
                first_state["identity_sha256"],
                second_state["identity_sha256"],
            )
            quarantined = list(
                (staging / "stale").glob(
                    "*/stage4/stale-deck.pptx"
                )
            )
            self.assertEqual(len(quarantined), 1)

    async def test_refuses_to_replace_an_existing_archived_presentation(
        self,
    ) -> None:
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            archive = root / "runs" / "2026-07-23-sample-run"
            release = archive / "release"
            release.mkdir(parents=True)
            (release / "release-manifest.json").write_text(
                json.dumps(
                    {
                        "slug": "sample-run",
                        "status": "ready",
                        "artifacts": [{"role": "presentation"}],
                    }
                ),
                encoding="utf-8",
            )
            build = AsyncMock()

            with patch(
                "cli.orchestrator._build_presentation_backfill",
                new=build,
            ):
                with self.assertRaisesRegex(
                    RuntimeError,
                    "already has a canonical presentation",
                ):
                    await run_presentation_for_archive(
                        archive_dir=archive,
                        slug="sample-run",
                        title="Sample",
                        repo_root=root,
                        budget_usd=10,
                    )

            build.assert_not_awaited()
            self.assertFalse(
                (
                    root
                    / "logs"
                    / "deck-backfills"
                    / "sample-run"
                ).exists()
            )

    async def test_failure_preserves_resumable_state_and_budget(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            archive = root / "runs" / "2026-07-23-sample-run"
            archive.mkdir(parents=True)

            async def fail_build(**kwargs):
                state = kwargs["staging_root"] / "state.json"
                state.write_text(
                    json.dumps(
                        {
                            "identity_sha256": "bound-inputs",
                            "phase": "art_direction_complete",
                        }
                    ),
                    encoding="utf-8",
                )
                self.assertEqual(kwargs["budget_usd"], 12.5)
                raise RuntimeError("synthetic model failure")

            with patch(
                "cli.orchestrator._build_presentation_backfill",
                new=AsyncMock(side_effect=fail_build),
            ):
                with self.assertRaisesRegex(
                    RuntimeError, "synthetic model failure"
                ):
                    await run_presentation_for_archive(
                        archive_dir=archive,
                        slug="sample-run",
                        title="Sample",
                        repo_root=root,
                        budget_usd=12.5,
                    )

            state_path = (
                root
                / "logs"
                / "deck-backfills"
                / "sample-run"
                / "state.json"
            )
            state = json.loads(state_path.read_text(encoding="utf-8"))
            self.assertEqual(state["identity_sha256"], "bound-inputs")
            self.assertEqual(state["phase"], "art_direction_complete")
            self.assertEqual(state["status"], "interrupted")

    async def test_success_removes_durable_staging_tree(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            archive = root / "runs" / "2026-07-23-sample-run"
            archive.mkdir(parents=True)
            expected = root / "reports" / "sample-run.pptx"

            async def succeed(**kwargs):
                kwargs["staging_root"].mkdir(parents=True, exist_ok=True)
                (kwargs["staging_root"] / "state.json").write_text(
                    '{"status":"complete"}',
                    encoding="utf-8",
                )
                return expected

            with patch(
                "cli.orchestrator._build_presentation_backfill",
                new=AsyncMock(side_effect=succeed),
            ):
                result = await run_presentation_for_archive(
                    archive_dir=archive,
                    slug="sample-run",
                    title="Sample",
                    repo_root=root,
                    budget_usd=10,
                )

            self.assertEqual(result, expected)
            self.assertFalse(
                (
                    root
                    / "logs"
                    / "deck-backfills"
                    / "sample-run"
                ).exists()
            )


if __name__ == "__main__":
    unittest.main()
