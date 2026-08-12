from __future__ import annotations

import asyncio
import json
import hashlib
import tempfile
import unittest
import zipfile
from pathlib import Path
from types import SimpleNamespace
from urllib.parse import urlencode
from unittest.mock import AsyncMock, patch

from fastapi.testclient import TestClient
from PIL import Image
from pptx import Presentation
from starlette.websockets import WebSocketDisconnect

import cli.server as server
from cli.presentation_qa import prepare_visual_inspection_receipt


ORIGIN = "http://127.0.0.1:8723"
CLIENT_ID = "owner_client_1234567890"
WS_HEADERS = {"origin": ORIGIN, "host": "127.0.0.1:8723"}


class _FakeTask:
    def __init__(self) -> None:
        self.cancelled = False

    def done(self) -> bool:
        return False

    def cancel(self) -> None:
        self.cancelled = True


class _FakeSink:
    def __init__(self) -> None:
        self.resolved: list[tuple[object, object]] = []

    def resolve(self, checkpoint_id: object, value: object) -> None:
        self.resolved.append((checkpoint_id, value))


class ServerSecurityTests(unittest.TestCase):
    def setUp(self) -> None:
        self.client = TestClient(
            server.app,
            base_url=ORIGIN,
            client=("127.0.0.1", 51000),
        )
        self.previous = (
            server._active_task,
            server._active_sink,
            server._active_owner,
        )
        server._active_task = None
        server._active_sink = None
        server._active_owner = None
        server._live_clients.clear()

    def tearDown(self) -> None:
        (
            server._active_task,
            server._active_sink,
            server._active_owner,
        ) = self.previous
        server._live_clients.clear()
        self.client.close()

    def _ws_path(self, *, token: str = server._SESSION_TOKEN,
                 client_id: str = CLIENT_ID) -> str:
        return "/ws?" + urlencode({"token": token, "client_id": client_id})

    @staticmethod
    def _message(**values: object) -> dict[str, object]:
        return {
            "session_token": server._SESSION_TOKEN,
            "client_id": CLIENT_ID,
            **values,
        }

    @staticmethod
    def _write_release(
        reports: Path,
        slug: str = "verified-report",
        *,
        role: str = "word_report",
        extension: str = ".docx",
        pointer_name: str | None = None,
        bundle_name: str = "verified-bundle",
        with_visual_inspection: bool = False,
    ) -> Path:
        bundle = reports / "releases" / bundle_name
        qa = bundle / "qa" / f"{slug}{extension}.qa.json"
        artifact = bundle / f"{slug}{extension}"
        qa.parent.mkdir(parents=True)
        if with_visual_inspection:
            if role != "presentation" or extension != ".pptx":
                raise ValueError(
                    "Visual inspection fixture requires a presentation."
                )
            deck = Presentation()
            slide = deck.slides.add_slide(deck.slide_layouts[5])
            slide.shapes.title.text = "Authorize the bounded operating pilot"
            deck.save(artifact)
        else:
            artifact.write_bytes(
                b"hash-bound office package"
                if role == "word_report"
                else f"hash-bound {role}".encode()
            )
        qa.write_text(json.dumps({"ok": True}), encoding="utf-8")
        artifact_hash = hashlib.sha256(artifact.read_bytes()).hexdigest()
        qa_hash = hashlib.sha256(qa.read_bytes()).hexdigest()
        artifact_record = {
            "role": role,
            "kind": extension.lstrip("."),
            "path": artifact.name,
            "source_sha256": artifact_hash,
            "sha256": artifact_hash,
            "size_bytes": artifact.stat().st_size,
            "required": True,
            "qa_ok": True,
            "qa_path": f"qa/{qa.name}",
            "qa_sha256": qa_hash,
            "rendered_files": [],
        }
        requirements: dict[str, bool] = {}
        if with_visual_inspection:
            brief = bundle / "visual-brief.json"
            brief.write_text(
                json.dumps(
                    {
                        "deck_mode": "board_decision",
                        "signature_visual": {
                            "slide_number": 1,
                            "concept": "Bounded decision",
                            "visual_type": "Decision exhibit",
                        },
                    }
                ),
                encoding="utf-8",
            )
            inspection = bundle / "inspection" / slug
            inspection.mkdir(parents=True)
            slide_png = inspection / f"{slug}-1.png"
            Image.new("RGB", (640, 360), "white").save(slide_png)
            receipt_path = bundle / f"{slug}-visual-inspection.json"
            prepare_visual_inspection_receipt(
                artifact=artifact,
                visual_brief=brief,
                deck_mode="board_decision",
                rendered_files=[slide_png],
                receipt_path=receipt_path,
            )
            receipt = json.loads(receipt_path.read_text(encoding="utf-8"))
            receipt["inspection"].update(
                {
                    "full_size_each_slide_inspected": True,
                    "montage_inspected": True,
                    "signature_exhibit_present": True,
                    "signature_exhibit_matches_brief": True,
                    "findings_resolved": True,
                    "status": "pass",
                }
            )
            receipt_path.write_text(json.dumps(receipt), encoding="utf-8")
            artifact_record["visual_inspection"] = {
                "path": receipt_path.name,
                "sha256": hashlib.sha256(
                    receipt_path.read_bytes()
                ).hexdigest(),
                "visual_brief_path": brief.name,
                "visual_brief_sha256": hashlib.sha256(
                    brief.read_bytes()
                ).hexdigest(),
                "files": [
                    {
                        "path": path.relative_to(bundle).as_posix(),
                        "sha256": hashlib.sha256(
                            path.read_bytes()
                        ).hexdigest(),
                        "size_bytes": path.stat().st_size,
                    }
                    for path in sorted(inspection.rglob("*"))
                    if path.is_file()
                ],
            }
            requirements = {
                "presentation": True,
                "visual_inspection": True,
            }
        bundle_manifest = bundle / "release-manifest.json"
        bundle_manifest.write_text(
            json.dumps({
                "schema_version": "1.0",
                "status": "ready",
                "slug": slug,
                "requirements": requirements,
                "artifacts": [artifact_record],
            }),
            encoding="utf-8",
        )
        bundle_hash = hashlib.sha256(bundle_manifest.read_bytes()).hexdigest()
        pointer = reports / (
            pointer_name or f"{slug}-release-manifest.json"
        )
        pointer.write_text(
            json.dumps({
                "status": "current",
                "slug": slug,
                "source_release_manifest_sha256": bundle_hash,
                "bundle_path": bundle_manifest.relative_to(reports).as_posix(),
                "artifacts": [{
                    "role": role,
                    "path": artifact.relative_to(reports).as_posix(),
                    "sha256": artifact_hash,
                    "qa_ok": True,
                    "qa_path": qa.relative_to(reports).as_posix(),
                    "qa_sha256": qa_hash,
                }],
            }),
            encoding="utf-8",
        )
        return artifact

    @staticmethod
    def _write_scope_release(
        reports: Path,
        slug: str = "scope-engagement",
    ) -> Path:
        reports.mkdir(parents=True, exist_ok=True)
        package = reports / f"scope-{slug}"
        package.mkdir()
        (package / "deliverable.docx").write_bytes(b"scope deliverable")
        (package / "qa-report.md").write_text("Accepted.\n", encoding="utf-8")
        (package / "MANIFEST.md").write_text(
            "# Scope package\n", encoding="utf-8"
        )
        files = [
            {
                "path": path.relative_to(package).as_posix(),
                "sha256": hashlib.sha256(path.read_bytes()).hexdigest(),
                "size_bytes": path.stat().st_size,
            }
            for path in sorted(package.rglob("*"))
            if path.is_file()
        ]
        archive_path = reports / f"{slug}-deliverables.zip"
        with zipfile.ZipFile(
            archive_path, "w", zipfile.ZIP_DEFLATED
        ) as archive:
            for item in files:
                archive.write(package / item["path"], item["path"])
        archive_hash = hashlib.sha256(archive_path.read_bytes()).hexdigest()
        receipt = reports / f"scope-{slug}-package-receipt.json"
        receipt.write_text(
            json.dumps(
                {
                    "schema_version": "1.0",
                    "slug": slug,
                    "dependency_contract_sha256": "0" * 64,
                    "files": files,
                    "zip": {
                        "path": archive_path.name,
                        "sha256": archive_hash,
                        "size_bytes": archive_path.stat().st_size,
                    },
                }
            ),
            encoding="utf-8",
        )
        pointer = reports / f"scope-{slug}-package-manifest.json"
        pointer.write_text(
            json.dumps(
                {
                    "schema_version": "1.0",
                    "status": "current",
                    "mode": "scope",
                    "slug": slug,
                    "title": "Scope Engagement",
                    "date": "2026-07-23",
                    "receipt": {
                        "path": receipt.name,
                        "sha256": hashlib.sha256(
                            receipt.read_bytes()
                        ).hexdigest(),
                    },
                    "zip": {
                        "path": archive_path.name,
                        "sha256": archive_hash,
                        "size_bytes": archive_path.stat().st_size,
                    },
                    "package": {
                        "path": package.name,
                        "files": files,
                    },
                    "deliverables": [
                        {
                            "id": "deliverable",
                            "title": "Deliverable",
                            "filename": "deliverable.docx",
                        }
                    ],
                }
            ),
            encoding="utf-8",
        )
        return archive_path

    def test_origin_requires_exact_loopback_host_scheme_and_port(self) -> None:
        self.assertTrue(
            server._origin_matches_local_host(
                ORIGIN, "127.0.0.1:8723", "ws"
            )
        )
        self.assertTrue(
            server._origin_matches_local_host(
                "http://localhost:8723", "localhost:8723", "http"
            )
        )
        for origin, host, scheme in (
            ("https://127.0.0.1:8723", "127.0.0.1:8723", "ws"),
            ("http://127.0.0.1:9999", "127.0.0.1:8723", "ws"),
            ("http://evil.example", "127.0.0.1:8723", "ws"),
            ("http://council.example", "council.example", "ws"),
            (None, "127.0.0.1:8723", "ws"),
        ):
            with self.subTest(origin=origin, host=host, scheme=scheme):
                self.assertFalse(
                    server._origin_matches_local_host(origin, host, scheme)
                )

    def test_meta_does_not_disclose_token_for_non_loopback_host(self) -> None:
        response = self.client.get("/api/meta", headers={"host": "evil.example"})
        self.assertEqual(response.status_code, 403)
        self.assertNotIn("session_token", response.json())

    def test_meta_delivers_process_token_only_as_no_store_same_origin_data(self) -> None:
        config = SimpleNamespace(
            default_budget_usd=80,
            model=lambda _role: "test-model",
        )
        with (
            patch.object(server, "get_config", return_value=config),
            patch("cli.menu.check_claude_auth", return_value=(True, "ok")),
            patch("cli.sources.discover_dropzone", return_value=[]),
        ):
            response = self.client.get("/api/meta")
        self.assertEqual(response.status_code, 200)
        self.assertEqual(response.json()["session_token"], server._SESSION_TOKEN)
        self.assertEqual(
            response.headers["cross-origin-resource-policy"], "same-origin"
        )
        self.assertIn("no-store", response.headers["cache-control"])

    def test_websocket_rejects_bad_origin_and_bad_token_before_accept(self) -> None:
        attempts = (
            (self._ws_path(token="wrong-token"), ORIGIN),
            (self._ws_path(), "http://evil.example"),
        )
        for path, origin in attempts:
            with self.subTest(path=path, origin=origin):
                with self.assertRaises(WebSocketDisconnect) as raised:
                    with self.client.websocket_connect(
                        path,
                        headers={
                            "origin": origin,
                            "host": "127.0.0.1:8723",
                        },
                    ):
                        pass
                self.assertEqual(raised.exception.code, 1008)

    def test_websocket_rejects_unauthenticated_messages(self) -> None:
        with self.client.websocket_connect(
            self._ws_path(), headers=WS_HEADERS
        ) as socket:
            socket.send_json({"type": "cancel", "client_id": CLIENT_ID})
            response = socket.receive_json()
        self.assertEqual(response["type"], "control_error")
        self.assertIn("not authenticated", response["message"])

    def test_observer_cannot_cancel_owner_run(self) -> None:
        task = _FakeTask()
        server._active_task = task
        server._active_sink = _FakeSink()
        server._active_owner = CLIENT_ID
        # The owner has a socket open. A run whose owner has gone away is a
        # different case — control passes there, or a crashed browser would
        # strand the run forever (see tests/test_reattach.py).
        server._live_clients[CLIENT_ID] = 1
        observer = "observer_client_123456"
        path = self._ws_path(client_id=observer)
        with self.client.websocket_connect(path, headers=WS_HEADERS) as socket:
            socket.send_json({
                "type": "cancel",
                "session_token": server._SESSION_TOKEN,
                "client_id": observer,
            })
            response = socket.receive_json()
        self.assertEqual(response["type"], "control_error")
        self.assertIn("observing", response["message"])
        self.assertFalse(task.cancelled)

    def test_abandoned_run_can_be_cancelled_by_a_returning_tab(self) -> None:
        task = _FakeTask()
        server._active_task = task
        server._active_sink = _FakeSink()
        server._active_owner = "crashed_client_123456"
        returning = "returning_client_123456"
        with self.client.websocket_connect(
            self._ws_path(client_id=returning), headers=WS_HEADERS
        ) as socket:
            socket.send_json({
                "type": "cancel",
                "session_token": server._SESSION_TOKEN,
                "client_id": returning,
            })
            # Taking over announces the new control state first; the malformed
            # frame then gives a response boundary after the cancel is handled.
            self.assertEqual(socket.receive_json()["type"], "control_status")
            socket.send_text("{")
            self.assertEqual(socket.receive_json()["type"], "control_error")
        self.assertTrue(task.cancelled)
        self.assertEqual(server._active_owner, returning)

    def test_owner_can_cancel_its_run(self) -> None:
        task = _FakeTask()
        server._active_task = task
        server._active_sink = _FakeSink()
        server._active_owner = CLIENT_ID
        with self.client.websocket_connect(
            self._ws_path(), headers=WS_HEADERS
        ) as socket:
            socket.send_json(self._message(type="cancel"))
            # A second message gives the server a response boundary after the
            # cancellation message has been processed.
            socket.send_text("{")
            self.assertEqual(socket.receive_json()["type"], "control_error")
        self.assertTrue(task.cancelled)

    def test_review_post_requires_same_origin_token_and_client(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            archive = root / "runs" / "2026-07-23-security-test"
            archive.mkdir(parents=True)
            body = {"ratings": {"writing": 5}, "notes": "Clear."}
            authenticated_headers = {
                "origin": ORIGIN,
                server._SESSION_HEADER: server._SESSION_TOKEN,
                server._CLIENT_HEADER: CLIENT_ID,
            }
            with patch.object(server, "REPO_ROOT", root):
                missing = self.client.post(
                    "/api/review/security-test", json=body
                )
                cross_origin = self.client.post(
                    "/api/review/security-test",
                    json=body,
                    headers={**authenticated_headers, "origin": "http://evil.example"},
                )
                accepted = self.client.post(
                    "/api/review/security-test",
                    json=body,
                    headers=authenticated_headers,
                )

            self.assertEqual(missing.status_code, 403)
            self.assertEqual(cross_origin.status_code, 403)
            self.assertEqual(accepted.status_code, 200)
            saved = archive / accepted.json()["saved"]
            self.assertTrue(saved.is_file())
            self.assertEqual(
                json.loads(saved.read_text(encoding="utf-8"))["rubric"]["writing"]["score"],
                5.0,
            )

    def test_deck_backfill_receives_validated_budget(self) -> None:
        sink = SimpleNamespace(close=AsyncMock())
        drive_deck = AsyncMock()
        with patch.object(server, "_drive_deck", drive_deck):
            asyncio.run(server._drive_run(
                "deck",
                {"slug": "budgeted-deck", "budget": "42.50"},
                sink,
            ))
        drive_deck.assert_awaited_once_with(
            "budgeted-deck", sink, 42.5
        )
        sink.close.assert_awaited_once()

    def test_new_report_attaches_only_its_browser_selected_sources(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            (root / "outputs").mkdir()
            upload_dir = server._argument_upload_dir(
                CLIENT_ID, root, purpose="report"
            )
            upload_dir.mkdir(parents=True)
            (upload_dir / "selected.md").write_text(
                "selected in the report form", encoding="utf-8"
            )
            unselected = root / "sources" / "unselected.md"
            unselected.write_text("left in the terminal dropzone", encoding="utf-8")
            sink = SimpleNamespace(close=AsyncMock())
            drive_new = AsyncMock()

            with (
                patch.object(server, "REPO_ROOT", root),
                patch.object(server, "_drive_new", drive_new),
            ):
                asyncio.run(server._drive_run(
                    "new",
                    {
                        "client_id": CLIENT_ID,
                        "budget": 0,
                        "spec": {
                            "title": "Browser source isolation test",
                            "thesis": "Only explicitly selected files belong to this report.",
                            "source_tokens": ["selected.md"],
                        },
                    },
                    sink,
                ))

            spec = drive_new.await_args.args[0]
            self.assertEqual(
                spec.source_paths,
                ["sources/runs/browser-source-isolation-test/selected.md"],
            )
            self.assertTrue(unselected.is_file())
            self.assertTrue((root / spec.source_paths[0]).is_file())
            sink.close.assert_awaited_once()

    def test_current_downloads_are_manifest_driven_and_hash_verified(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            reports = Path(directory) / "reports"
            artifact = self._write_release(reports)

            downloads = server._downloads_for_slug(
                "verified-report", reports
            )
            self.assertEqual(len(downloads), 1)
            self.assertEqual(
                downloads[0]["path"],
                artifact.relative_to(reports).as_posix(),
            )

            artifact.write_bytes(b"tampered after release")
            self.assertEqual(
                server._downloads_for_slug("verified-report", reports),
                [],
            )

    def test_scope_zip_requires_an_exact_hash_bound_pointer_and_survives_reload(
        self,
    ) -> None:
        with tempfile.TemporaryDirectory() as directory:
            reports = Path(directory) / "reports"
            archive = self._write_scope_release(reports)
            unbound = reports / "unbound-deliverables.zip"
            unbound.write_bytes(b"not authorized")
            with (
                patch.object(server, "REPORTS_DIR", reports),
                patch("cli.publish.discover_reports", return_value=[]),
            ):
                download = self.client.get(f"/download/{archive.name}")
                rejected = self.client.get(f"/download/{unbound.name}")
                home = self.client.get("/api/home")
            self.assertEqual(download.status_code, 200)
            self.assertEqual(rejected.status_code, 404)
            scope_entries = [
                entry
                for entry in home.json()["archives"]
                if entry.get("mode") == "scope"
            ]
            self.assertEqual(len(scope_entries), 1)
            self.assertEqual(
                scope_entries[0]["downloads"][0]["url"],
                f"/download/{archive.name}",
            )
            self.assertFalse(scope_entries[0]["can_revise"])

            archive.write_bytes(b"tampered after pointer publication")
            with (
                patch.object(server, "REPORTS_DIR", reports),
                patch("cli.publish.discover_reports", return_value=[]),
            ):
                withdrawn = self.client.get(f"/download/{archive.name}")
                home_after_tamper = self.client.get("/api/home")
            self.assertEqual(withdrawn.status_code, 404)
            self.assertFalse(
                any(
                    entry.get("mode") == "scope"
                    for entry in home_after_tamper.json()["archives"]
                )
            )

    def test_revision_completion_and_library_resolve_the_revised_release(
        self,
    ) -> None:
        from cli.publish import ReportSource

        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            reports = root / "reports"
            released = self._write_release(
                reports,
                slug="airport-plan-revised-v1",
                bundle_name="revision-bundle",
            )
            archive = root / "runs" / "2026-07-23-airport-plan"
            (archive / "stage3").mkdir(parents=True)
            (archive / "stage4").mkdir()
            (archive / "stage3" / "final-draft.md").write_text(
                "# Original\n", encoding="utf-8"
            )
            original_docx = archive / "stage4" / "airport-plan.docx"
            original_docx.write_bytes(b"original")
            run_prompt = archive / "run-prompt.md"
            run_prompt.write_text(
                "# Run: Airport Plan\n\n## Output format\nreport\n",
                encoding="utf-8",
            )
            source = ReportSource(
                slug="airport-plan",
                archive_dir=archive,
                stage4_docx=original_docx,
                final_md=archive / "stage3" / "final-draft.md",
                run_file=run_prompt,
            )

            revision = archive / "revisions" / "v1"
            (revision / "stage4").mkdir(parents=True)
            (revision / "release").mkdir()
            final_draft = revision / "final-draft.md"
            final_draft.write_text("# Revised decision\n", encoding="utf-8")
            lineage = revision / "claim-lineage.jsonl"
            lineage.write_text("{}\n", encoding="utf-8")
            gate = revision / "quality-gate.json"
            gate.write_text("{}\n", encoding="utf-8")
            execution = revision / "revision-execution.json"
            execution.write_text("{}\n", encoding="utf-8")
            stage4_report = (
                revision / "stage4" / "airport-plan-revised-v1.docx"
            )
            stage4_report.write_bytes(released.read_bytes())
            release_manifest = revision / "release" / "release-manifest.json"
            release_manifest.write_text("{}\n", encoding="utf-8")

            def digest(path: Path) -> str:
                return hashlib.sha256(path.read_bytes()).hexdigest()

            (revision / "revision-manifest.json").write_text(
                json.dumps(
                    {
                        "schema_version": "1.0",
                        "slug": "airport-plan",
                        "revision": 1,
                        "created_at": "2026-07-23T12:00:00-04:00",
                        "final_draft_sha256": digest(final_draft),
                        "claim_lineage_sha256": digest(lineage),
                        "quality_gate_sha256": digest(gate),
                        "word_report_sha256": digest(stage4_report),
                        "release_manifest_sha256": digest(release_manifest),
                        "revision_execution_sha256": digest(execution),
                        "status": "released",
                    }
                ),
                encoding="utf-8",
            )

            with (
                patch.object(server, "REPORTS_DIR", reports),
                patch("cli.publish.discover_reports", return_value=[source]),
                patch("cli.revise._completed_revisions", return_value=[1]),
            ):
                report = self.client.get(
                    "/api/report/airport-plan-revised-v1"
                )
                home = self.client.get("/api/home")
                review = self.client.post(
                    "/api/review/airport-plan-revised-v1",
                    json={"ratings": {"writing": 5}},
                    headers={
                        "origin": ORIGIN,
                        server._SESSION_HEADER: server._SESSION_TOKEN,
                        server._CLIENT_HEADER: CLIENT_ID,
                    },
                )
            self.assertEqual(report.status_code, 200)
            self.assertIn("# Revised decision", report.json()["markdown"])
            self.assertEqual(report.json()["revise_slug"], "airport-plan")
            self.assertEqual(len(report.json()["downloads"]), 1)
            revision_entries = [
                entry
                for entry in home.json()["archives"]
                if entry.get("mode") == "revision"
            ]
            self.assertEqual(len(revision_entries), 1)
            self.assertEqual(
                revision_entries[0]["slug"],
                "airport-plan-revised-v1",
            )
            self.assertEqual(
                revision_entries[0]["revise_slug"], "airport-plan"
            )
            self.assertFalse(revision_entries[0]["can_build_deck"])
            self.assertEqual(review.status_code, 200)
            self.assertTrue(
                (
                    revision
                    / "evaluation"
                    / "reviews"
                    / "final-product.json"
                ).is_file()
            )

            final_draft.write_text(
                "# Tampered after revision release\n", encoding="utf-8"
            )
            with (
                patch.object(server, "REPORTS_DIR", reports),
                patch("cli.publish.discover_reports", return_value=[source]),
                patch("cli.revise._completed_revisions", return_value=[1]),
            ):
                withdrawn = self.client.get(
                    "/api/report/airport-plan-revised-v1"
                )
                home_after_tamper = self.client.get("/api/home")
            self.assertEqual(withdrawn.json()["markdown"], "")
            self.assertEqual(withdrawn.json()["downloads"], [])
            self.assertFalse(
                any(
                    entry.get("mode") == "revision"
                    for entry in home_after_tamper.json()["archives"]
                )
            )

    def test_revision_completion_event_names_the_revised_release(self) -> None:
        source = SimpleNamespace(
            slug="airport-plan",
            archive_dir=Path("/tmp/archive"),
        )
        sink = SimpleNamespace(emit=AsyncMock())
        tally = SimpleNamespace(total=4.25)
        with (
            patch("cli.revise.revisable_reports", return_value=[source]),
            patch("cli.revise.next_revision_version", return_value=2),
            patch(
                "cli.orchestrator.run_revision_pipeline",
                new=AsyncMock(return_value=(Path("/tmp/revised.docx"), tally)),
            ),
        ):
            asyncio.run(
                server._drive_revise(
                    "airport-plan", "Tighten the decision.", sink, True
                )
            )
        event, payload = sink.emit.await_args.args
        self.assertEqual(event, "run_complete")
        self.assertEqual(payload["slug"], "airport-plan-revised-v2")
        self.assertEqual(payload["revise_slug"], "airport-plan")
        self.assertEqual(payload["mode"], "revision")

    def test_tampered_visual_inspection_withdraws_presentation_download(
        self,
    ) -> None:
        with tempfile.TemporaryDirectory() as directory:
            reports = Path(directory) / "reports"
            artifact = self._write_release(
                reports,
                slug="inspection-bound",
                role="presentation",
                extension=".pptx",
                with_visual_inspection=True,
            )
            self.assertEqual(
                [item["role"] for item in server._downloads_for_slug(
                    "inspection-bound", reports
                )],
                ["presentation"],
            )
            relative = artifact.relative_to(reports).as_posix()

            montage = (
                artifact.parent
                / "inspection"
                / "inspection-bound"
                / "montage.png"
            )
            montage.write_bytes(b"altered after release")
            self.assertEqual(
                server._downloads_for_slug("inspection-bound", reports),
                [],
            )
            with patch.object(server, "REPORTS_DIR", reports):
                response = self.client.get(f"/download/{relative}")
            self.assertEqual(response.status_code, 404)

    def test_supplemental_deck_pointer_merges_with_canonical_release(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            reports = Path(directory) / "reports"
            word = self._write_release(reports, slug="with-backfill")
            deck = self._write_release(
                reports,
                slug="with-backfill",
                role="presentation",
                extension=".pptx",
                pointer_name="with-backfill-deck-release-manifest.json",
                bundle_name="verified-deck-bundle",
            )

            downloads = server._downloads_for_slug("with-backfill", reports)
            self.assertEqual(
                [item["role"] for item in downloads],
                ["word_report", "presentation"],
            )
            self.assertEqual(
                {item["path"] for item in downloads},
                {
                    word.relative_to(reports).as_posix(),
                    deck.relative_to(reports).as_posix(),
                },
            )

            # The supplemental pointer independently fails closed.
            deck.write_bytes(b"tampered deck")
            self.assertEqual(
                [item["role"] for item in server._downloads_for_slug(
                    "with-backfill", reports
                )],
                ["word_report"],
            )

            canonical_deck = self._write_release(
                reports,
                slug="canonical-wins",
                role="presentation",
                extension=".pptx",
                bundle_name="canonical-deck-bundle",
            )
            self._write_release(
                reports,
                slug="canonical-wins",
                role="presentation",
                extension=".pptx",
                pointer_name="canonical-wins-deck-release-manifest.json",
                bundle_name="supplemental-deck-bundle",
            )
            canonical_downloads = server._downloads_for_slug(
                "canonical-wins", reports
            )
            self.assertEqual(len(canonical_downloads), 1)
            self.assertEqual(
                canonical_downloads[0]["path"],
                canonical_deck.relative_to(reports).as_posix(),
            )

    def test_nested_download_route_enforces_manifest_and_legacy_boundary(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            reports = Path(directory) / "reports"
            artifact = self._write_release(reports)
            relative = artifact.relative_to(reports).as_posix()
            with patch.object(server, "REPORTS_DIR", reports):
                current = self.client.get(f"/download/{relative}")
                traversal = self.client.get(
                    "/download/releases/verified-bundle/../../outside.docx"
                )
            self.assertEqual(current.status_code, 200)
            self.assertEqual(current.content, b"hash-bound office package")
            self.assertEqual(traversal.status_code, 404)

            legacy = reports / "legacy-report.docx"
            legacy.write_bytes(b"legacy")
            with patch.object(server, "REPORTS_DIR", reports):
                accepted_legacy = self.client.get(
                    "/download/legacy-report.docx"
                )
            self.assertEqual(accepted_legacy.status_code, 200)

            (reports / "legacy-report-release-manifest.json").write_text(
                "{}", encoding="utf-8"
            )
            with patch.object(server, "REPORTS_DIR", reports):
                blocked_fallback = self.client.get(
                    "/download/legacy-report.docx"
                )
            self.assertEqual(blocked_fallback.status_code, 404)

            deck_pointer_only = reports / "deck-pointer-only.docx"
            deck_pointer_only.write_bytes(b"must not fall back")
            (
                reports / "deck-pointer-only-deck-release-manifest.json"
            ).write_text("{}", encoding="utf-8")
            with patch.object(server, "REPORTS_DIR", reports):
                blocked_by_supplement = self.client.get(
                    "/download/deck-pointer-only.docx"
                )
            self.assertEqual(blocked_by_supplement.status_code, 404)


if __name__ == "__main__":
    unittest.main()
