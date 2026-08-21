from __future__ import annotations

import json
import tempfile
import unittest
from pathlib import Path

from jsonschema import Draft202012Validator
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from cli.presentation_qa import PresentationQAConfig, qa_presentation
from cli.publishing_quality import executive_summary_word_target, lint_markdown


REPO_ROOT = Path(__file__).parents[1]


def _set_shape_text(shape, text: str, size: float) -> None:
    shape.text = text
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(size)


def _add_slide(deck: Presentation, title: str, body: str = ""):
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    title_shape = slide.shapes.add_textbox(
        Inches(0.7), Inches(0.45), Inches(11.9), Inches(0.8)
    )
    title_shape.name = f"slide-{len(deck.slides)}-title"
    _set_shape_text(title_shape, title, 46)
    if body:
        body_shape = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.6), Inches(11.5), Inches(4.8)
        )
        _set_shape_text(body_shape, body, 18)
    return slide


def _brief_slide(number: int, **overrides) -> dict:
    item = {
        "slide_number": number,
        "narrative_job": "Advance the decision",
        "headline": f"Decision evidence changes the choice on slide {number}",
        "evidence_ids": [],
        "visual_type": "text",
        "visual_spec": "Use one clear assertion.",
        "source_note": "",
        "density_budget": "No more than 70 visible words.",
        "speaker_note": "Explain the consequence.",
    }
    item.update(overrides)
    return item


def _visual_brief(slides: list[dict], *, asset_requests: list | None = None) -> dict:
    return {
        "schema_version": "2.0",
        "communication_job": (
            "Enable the airport board to choose a bounded operating decision."
        ),
        "audience": "Airport board",
        "decision": "Approve the bounded operating test.",
        "decision_owner": "Chief Operating Officer",
        "approval_path": "Executive committee then board",
        "first_90_day_action": "Open one controlled operating test.",
        "success_measures": ["Hold service inside the published threshold."],
        "deck_mode": "board_decision",
        "visual_thesis": "A bounded test makes the operating choice reversible.",
        "signature_visual": {
            "slide_number": 1,
            "concept": "The bounded operating envelope",
            "visual_type": "Decision frame",
            "visual_spec": "Show the launch, pause, and stop conditions.",
            "evidence_ids": [],
            "source_note": "",
        },
        "brand_profile": "Transform Airports executive",
        "slides": slides,
        "report_visuals": [],
        "source_appendix": [],
        "accessibility_checks": ["Do not encode the decision by color alone."],
        "asset_requests": asset_requests or [],
    }


class NarrativeQualityTests(unittest.TestCase):
    def test_generic_report_scaffold_and_duplicate_sections_are_rejected(self) -> None:
        repeated = " ".join(
            [
                "The operating record ties the decision to a specific morning peak."
                for _ in range(6)
            ]
        )
        report = lint_markdown(
            "# Decision\n\n"
            f"## The argument\n\n{repeated}\n\n"
            f"## Implications for the operator\n\n{repeated}\n"
        )
        warning_codes = {issue.code for issue in report.warnings}
        error_codes = {issue.code for issue in report.errors}
        self.assertIn("generic_section_heading", warning_codes)
        self.assertIn("generic_section_scaffold", warning_codes)
        self.assertIn("duplicate_section_body", error_codes)

    def test_density_checks_warn_without_blocking_specific_prose(self) -> None:
        long_paragraph = " ".join(
            "The terminal operating decision stays tied to evidence."
            for _ in range(30)
        )
        list_items = "\n".join(
            f"- Condition {number} "
            + "keeps the operating test bounded and gives the owner a clear stop rule "
            + "before passenger service deteriorates during the morning peak "
            + "or an irregular operation."
            for number in range(1, 7)
        )
        report = lint_markdown(
            "# Decision\n\n"
            "## The morning peak sets the operating limit\n\n"
            f"{long_paragraph}\n\n"
            "## Six conditions keep the test reversible\n\n"
            f"{list_items}\n"
        )
        warning_codes = {issue.code for issue in report.warnings}
        self.assertIn("oversized_paragraph", warning_codes)
        self.assertIn("list_heavy_prose", warning_codes)
        self.assertEqual(report.errors, [])

    def test_layered_summary_contract_flags_excess_claims_and_extreme_length(
        self,
    ) -> None:
        item_text = " ".join(
            "Verified evidence connects the choice to the operating threshold."
            for _ in range(15)
        )
        summary = "\n\n".join(
            f"{number}. **Claim {number}.** {item_text}"
            for number in range(1, 7)
        )
        report = lint_markdown(
            f"# Decision\n\n## Executive summary\n\n{summary}\n"
        )
        self.assertIn(
            "oversized_executive_summary",
            {issue.code for issue in report.warnings},
        )
        self.assertIn(
            "too_many_summary_claims",
            {issue.code for issue in report.warnings},
        )

    def test_explicit_summary_target_overrides_default_without_confusing_total_length(
        self,
    ) -> None:
        prompt = (
            "## Length\n\n4,000–6,000 words for the full report; "
            "~1,100-word executive summary.\n"
        )
        target = executive_summary_word_target(prompt)
        self.assertEqual(target, 1_100)
        summary = " ".join("Evidence supports the operating choice." for _ in range(180))
        report = lint_markdown(
            f"# Decision\n\n## Executive summary\n\n{summary}\n",
            executive_summary_target_words=target,
        )
        self.assertNotIn(
            "oversized_executive_summary",
            {issue.code for issue in report.issues},
        )

    def test_summary_target_parser_supports_summary_first_range_and_default(self) -> None:
        self.assertEqual(
            executive_summary_word_target(
                "## Length\n\n**Executive summary:** 900–1,100 words; "
                "report body: 6,000 words.\n"
            ),
            1_100,
        )
        self.assertIsNone(
            executive_summary_word_target(
                "## Thesis\n\nReplace the old 1,100-word executive summary.\n\n"
                "## Length\n\n4,000–6,000 words for the full report.\n"
            )
        )

    def test_summary_target_parser_accepts_natural_length_phrasings(self) -> None:
        bodies = (
            "1,100 words for the executive summary.",
            "The executive summary should be about 1,100 words.",
            "About 1,100 words in the executive summary.",
        )
        for body in bodies:
            with self.subTest(body=body):
                self.assertEqual(
                    executive_summary_word_target(f"## Length\n\n{body}\n"),
                    1_100,
                )

    def test_nested_summary_subsections_count_toward_advisory_target(self) -> None:
        nested = " ".join(
            "Evidence keeps the operating decision reversible."
            for _ in range(125)
        )
        report = lint_markdown(
            "# Decision\n\n## Executive summary\n\nThe decision is bounded.\n\n"
            f"### Evidence beneath the decision\n\n{nested}\n"
        )
        self.assertIn(
            "oversized_executive_summary",
            {issue.code for issue in report.warnings},
        )


class VisualBriefSchemaTests(unittest.TestCase):
    @classmethod
    def setUpClass(cls) -> None:
        cls.schema = json.loads(
            (REPO_ROOT / "assets/brand/visual-brief.schema.json").read_text(
                encoding="utf-8"
            )
        )
        cls.validator = Draft202012Validator(cls.schema)

    def test_legacy_visual_brief_remains_valid(self) -> None:
        slides = [_brief_slide(number) for number in range(1, 9)]
        payload = _visual_brief(
            slides,
            asset_requests=[
                {
                    "id": "AR-legacy",
                    "description": "Existing request without workflow fields.",
                    "rights": "Confirm before publication.",
                    "deadline": "Before deck production.",
                }
            ],
        )
        self.assertEqual(list(self.validator.iter_errors(payload)), [])

    def test_extended_layout_and_asset_contract_is_valid(self) -> None:
        families = (
            "cover",
            "split_media",
            "timeline",
            "decision",
            "chart",
            "map_plan",
            "comparison",
            "section_break",
        )
        slides = [
            _brief_slide(
                number,
                layout_family=families[number - 1],
                colorway="dark" if number in {1, 8} else "light",
                speaker_led=True,
                visual_priority="high" if number in {1, 8} else "medium",
                asset_request_ids=["AR-1"] if number == 2 else [],
                visible_word_budget=70,
            )
            for number in range(1, 9)
        ]
        payload = _visual_brief(
            slides,
            asset_requests=[
                {
                    "id": "AR-1",
                    "slide_numbers": [2],
                    "description": "Rights-cleared airport operating photograph.",
                    "media_role": "photograph",
                    "source": "Airport media library",
                    "rights": "Publication permission confirmed",
                    "credit": "Airport Authority",
                    "approval_status": "approved",
                    "fulfillment_status": "supplied",
                    "required": True,
                }
            ],
        )
        self.assertEqual(list(self.validator.iter_errors(payload)), [])


class PresentationVisualQualityTests(unittest.TestCase):
    def test_visible_word_budget_has_a_calibrated_hard_limit(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            path = Path(directory) / "dense.pptx"
            deck = Presentation()
            deck.slide_width = Inches(13.333)
            deck.slide_height = Inches(7.5)
            body = " ".join("operating evidence" for _ in range(60))
            _add_slide(deck, "The operating threshold controls the decision", body)
            deck.save(path)

            report = qa_presentation(
                path,
                config=PresentationQAConfig(
                    min_slide_count=1,
                    max_slide_count=1,
                    require_visual_evidence=False,
                    visible_word_warning=50,
                    visible_word_error=100,
                ),
            )
            self.assertIn(
                "visible_word_budget_exceeded",
                {issue.code for issue in report.errors},
            )
            self.assertGreater(report.metadata["maximum_visible_words"], 100)

    def test_visible_word_budget_counts_table_cells(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            path = Path(directory) / "dense-table.pptx"
            deck = Presentation()
            deck.slide_width = Inches(13.333)
            deck.slide_height = Inches(7.5)
            slide = _add_slide(deck, "The comparison must remain readable")
            table = slide.shapes.add_table(
                1, 1, Inches(1), Inches(2), Inches(10), Inches(3)
            ).table
            table.cell(0, 0).text = " ".join(
                "operating threshold" for _ in range(55)
            )
            deck.save(path)

            report = qa_presentation(
                path,
                config=PresentationQAConfig(
                    min_slide_count=1,
                    max_slide_count=1,
                    require_visual_evidence=False,
                    visible_word_warning=50,
                    visible_word_error=100,
                ),
            )
            self.assertIn(
                "visible_word_budget_exceeded",
                {issue.code for issue in report.errors},
            )
            self.assertGreater(report.metadata["maximum_visible_words"], 100)

    def test_layout_table_and_card_rhythm_produce_actionable_findings(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            path = Path(directory) / "repetitive.pptx"
            deck = Presentation()
            deck.slide_width = Inches(13.333)
            deck.slide_height = Inches(7.5)
            slides: list[dict] = []
            for number in range(1, 9):
                headline = f"Decision evidence changes the choice on slide {number}"
                slide = _add_slide(deck, headline)
                if number <= 5:
                    table = slide.shapes.add_table(
                        2, 2, Inches(1), Inches(2), Inches(8), Inches(2)
                    ).table
                    table.cell(0, 0).text = "Option"
                    table.cell(0, 1).text = "Threshold"
                    table.cell(1, 0).text = "Bounded test"
                    table.cell(1, 1).text = "Hold service"
                if number == 1:
                    for index in range(4):
                        card = slide.shapes.add_shape(
                            MSO_SHAPE.RECTANGLE,
                            Inches(0.8 + index * 2.2),
                            Inches(4.8),
                            Inches(1.8),
                            Inches(1.0),
                        )
                        card.name = f"decision-card-{index + 1}"
                    slide.shapes[-1].name = "SIGNATURE VISUAL — operating envelope"
                slides.append(
                    _brief_slide(
                        number,
                        headline=headline,
                        visual_type="table" if number <= 5 else "text",
                        layout_family="comparison",
                        colorway="light",
                    )
                )
            deck.save(path)

            report = qa_presentation(
                path,
                config=PresentationQAConfig(
                    min_slide_count=8,
                    max_slide_count=8,
                    require_visual_evidence=False,
                ),
                visual_brief=_visual_brief(slides),
            )
            warning_codes = {issue.code for issue in report.warnings}
            error_codes = {issue.code for issue in report.errors}
            self.assertIn("repeated_layout_family", warning_codes)
            self.assertIn("insufficient_layout_diversity", warning_codes)
            self.assertIn("flat_colorway_rhythm", warning_codes)
            self.assertIn("excessive_equal_card_grid", warning_codes)
            self.assertIn("table_heavy_deck", error_codes)

    def test_required_approved_asset_cannot_remain_pending(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            path = Path(directory) / "asset-pending.pptx"
            deck = Presentation()
            deck.slide_width = Inches(13.333)
            deck.slide_height = Inches(7.5)
            slide = _add_slide(
                deck,
                "The supplied airport plan makes the operating choice visible",
            )
            signature = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(1),
                Inches(2),
                Inches(5),
                Inches(2),
            )
            signature.name = "SIGNATURE VISUAL — operating envelope"
            deck.save(path)

            slide_brief = _brief_slide(
                1,
                headline=(
                    "The supplied airport plan makes the operating choice visible"
                ),
                asset_request_ids=["AR-1"],
            )
            brief = _visual_brief(
                [slide_brief],
                asset_requests=[
                    {
                        "id": "AR-1",
                        "slide_numbers": [1],
                        "description": "Official airport operating plan.",
                        "media_role": "official_plan",
                        "source": "Airport planning office",
                        "rights": "Publication permission requested",
                        "credit": "Airport Authority",
                        "approval_status": "approved",
                        "fulfillment_status": "pending",
                        "required": True,
                    }
                ],
            )
            report = qa_presentation(
                path,
                config=PresentationQAConfig(
                    min_slide_count=1,
                    max_slide_count=1,
                    require_visual_evidence=False,
                ),
                visual_brief=brief,
            )
            self.assertIn(
                "approved_asset_unfulfilled",
                {issue.code for issue in report.errors},
            )

            brief["asset_requests"][0]["fulfillment_status"] = "retrieved"
            placed_report = qa_presentation(
                path,
                config=PresentationQAConfig(
                    min_slide_count=1,
                    max_slide_count=1,
                    require_visual_evidence=False,
                ),
                visual_brief=brief,
            )
            self.assertIn(
                "fulfilled_asset_not_placed",
                {issue.code for issue in placed_report.errors},
            )


if __name__ == "__main__":
    unittest.main()
