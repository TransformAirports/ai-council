"""Structural and render QA for Council PowerPoint presentations.

The checks are deliberately deterministic and require no paid model calls.
They do not replace full-size visual inspection; they make that inspection
safer by catching canvas overflow, unreadable type, placeholders, unsourced
numeric claims, and decks that contain no evidence-bearing visual at all.
"""
from __future__ import annotations

import argparse
import hashlib
import json
import math
import re
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

from cli.publishing_quality import (
    QualityIssue,
    QualityReport,
    lint_reader_text,
    render_office_artifact,
)


NUMERIC_CLAIM_RE = re.compile(
    r"(?:"
    r"\$[\d,.]+"
    r"|\b\d+(?:\.\d+)?%"
    r"|\b\d{1,3}(?:,\d{3})+\b"
    r"|\b\d+(?:\.\d+)?\s*(?:"
    r"thousand|million|billion|passengers?|gates?|minutes?|hours?|days?|"
    r"years?|percent(?:age\s+points?)?|basis\s+points?|bps|"
    r"checkpoints?|flights?|operations?"
    r")\b"
    r"|\b\d+(?:\.\d+)?\s*:\s*\d+(?:\.\d+)?\b"
    r"|\b\d+(?:\.\d+)?\s*[xX]\b"
    r")",
    re.IGNORECASE,
)
SOURCE_RE = re.compile(r"\b(?:source|sources|footnote|notes?)\s*:", re.IGNORECASE)
INTERNAL_TOKEN_RE = re.compile(
    r"\b(?:stage[123]|strategist|red[- ]team|fact[- ]checker|"
    r"humanized[- ]draft|[-_]brief\.md)\b",
    re.IGNORECASE,
)
SIGNATURE_VISUAL_PREFIX = "SIGNATURE VISUAL —"


@dataclass(frozen=True)
class PresentationQAConfig:
    min_title_pt: float = 35.0
    min_cover_title_pt: float = 44.0
    min_body_pt: float = 16.0
    min_source_pt: float = 9.0
    max_slide_count: int = 18
    min_slide_count: int = 8
    require_numeric_sources: bool = True
    require_visual_evidence: bool = True
    allow_internal_terms: bool = False
    profile: str = "client_release"
    body_type_severity: str = "error"
    evidence_issue_severity: str = "error"
    slide_count_severity: str = "warning"
    deck_mode: str | None = None
    appendix_body_pt: float | None = None
    require_render: bool = False

    @classmethod
    def process_explainer(cls) -> "PresentationQAConfig":
        """QA profile for decks whose legitimate subject is Council mechanics."""
        return cls(
            min_title_pt=28.0,
            min_cover_title_pt=40.0,
            min_body_pt=15.5,
            min_source_pt=9.0,
            allow_internal_terms=True,
            profile="internal_process_explainer",
            body_type_severity="warning",
            evidence_issue_severity="warning",
        )

    @classmethod
    def for_mode(cls, deck_mode: str) -> "PresentationQAConfig":
        """Return the enforceable production contract for one deck mode."""

        if deck_mode == "board_decision":
            return cls(
                min_slide_count=8,
                max_slide_count=12,
                slide_count_severity="error",
                deck_mode=deck_mode,
                require_render=True,
            )
        if deck_mode == "executive_briefing":
            return cls(
                min_slide_count=12,
                max_slide_count=18,
                slide_count_severity="error",
                deck_mode=deck_mode,
                require_render=True,
            )
        if deck_mode == "technical_read_ahead":
            return cls(
                min_slide_count=15,
                max_slide_count=25,
                slide_count_severity="error",
                deck_mode=deck_mode,
                appendix_body_pt=12.0,
                require_render=True,
            )
        raise ValueError(f"Unknown deck mode: {deck_mode!r}")

    @classmethod
    def for_argument(cls, slide_count: int) -> "PresentationQAConfig":
        """Return the exact-length contract for Strengthen an argument decks."""

        if (
            isinstance(slide_count, bool)
            or not isinstance(slide_count, int)
            or not 3 <= slide_count <= 30
        ):
            raise ValueError("Argument deck slide count must be between 3 and 30.")
        return cls(
            min_slide_count=slide_count,
            max_slide_count=slide_count,
            slide_count_severity="error",
            deck_mode="argument_brief",
            require_render=True,
        )


def _normalise_headline(text: str) -> str:
    return " ".join(re.findall(r"[a-z0-9]+", text.casefold()))


def _load_visual_brief(
    visual_brief: Path | dict | None,
) -> dict | None:
    if isinstance(visual_brief, dict):
        return visual_brief
    if visual_brief is None:
        return None
    try:
        payload = json.loads(Path(visual_brief).read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError) as exc:
        raise ValueError(f"Cannot read visual brief {visual_brief}: {exc}") from exc
    if not isinstance(payload, dict):
        raise ValueError("Visual brief must contain a JSON object.")
    return payload


def _sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with Path(path).open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _safe_relative_file(root: Path, relative: str) -> Path:
    """Resolve one receipt path without permitting escape or symlink tricks."""

    if not relative or Path(relative).is_absolute():
        raise ValueError(f"Inspection receipt path must be relative: {relative!r}")
    root = root.resolve()
    candidate = root / relative
    resolved = candidate.resolve()
    try:
        resolved.relative_to(root)
    except ValueError as exc:
        raise ValueError(
            f"Inspection receipt path escapes its artifact root: {relative!r}"
        ) from exc
    cursor = candidate
    while cursor != root:
        if cursor.is_symlink():
            raise ValueError(
                f"Inspection receipt path uses a symlink: {relative!r}"
            )
        cursor = cursor.parent
    if not resolved.is_file():
        raise ValueError(f"Inspection receipt file is missing: {relative!r}")
    return resolved


def build_slide_montage(
    rendered_slides: list[Path],
    montage_path: Path,
    *,
    columns: int = 4,
    thumbnail_width: int = 420,
) -> Path:
    """Create a compact narrative-rhythm view from full-slide PNG renders."""

    if not rendered_slides:
        raise ValueError("Cannot create a presentation montage without slide PNGs.")
    columns = max(1, min(columns, len(rendered_slides)))
    opened = [Image.open(path).convert("RGB") for path in rendered_slides]
    try:
        ratio = opened[0].height / opened[0].width
        thumbnail_height = max(1, round(thumbnail_width * ratio))
        label_height = 34
        gutter = 18
        rows = math.ceil(len(opened) / columns)
        width = gutter + columns * (thumbnail_width + gutter)
        height = gutter + rows * (thumbnail_height + label_height + gutter)
        montage = Image.new("RGB", (width, height), "#E9EEF3")
        draw = ImageDraw.Draw(montage)
        for index, image in enumerate(opened):
            row, column = divmod(index, columns)
            x = gutter + column * (thumbnail_width + gutter)
            y = gutter + row * (thumbnail_height + label_height + gutter)
            thumbnail = image.copy()
            thumbnail.thumbnail(
                (thumbnail_width, thumbnail_height),
                Image.Resampling.LANCZOS,
            )
            card = Image.new("RGB", (thumbnail_width, thumbnail_height), "white")
            card.paste(
                thumbnail,
                (
                    (thumbnail_width - thumbnail.width) // 2,
                    (thumbnail_height - thumbnail.height) // 2,
                ),
            )
            montage.paste(card, (x, y))
            draw.text((x, y + thumbnail_height + 7), f"SLIDE {index + 1:02d}", fill="#152B45")
        montage_path.parent.mkdir(parents=True, exist_ok=True)
        montage.save(montage_path, format="PNG", optimize=True)
    finally:
        for image in opened:
            image.close()
    return montage_path


def prepare_visual_inspection_receipt(
    *,
    artifact: Path,
    visual_brief: Path,
    deck_mode: str,
    rendered_files: Iterable[Path | str],
    receipt_path: Path,
    inspector: str = "presentation-designer",
) -> Path:
    """Write a hash-bound inspection packet whose attestations start pending.

    The designer must inspect every full-size render and the montage, resolve
    findings, then deliberately change the four inspection fields to a passing
    state. The orchestrator independently validates every hash before release.
    """

    artifact = Path(artifact).resolve()
    visual_brief = Path(visual_brief).resolve()
    receipt_path = Path(receipt_path).resolve()
    root = receipt_path.parent.resolve()
    brief_payload = _load_visual_brief(visual_brief)
    signature = (
        brief_payload.get("signature_visual")
        if isinstance(brief_payload, dict)
        else None
    )
    signature_slide_number = (
        signature.get("slide_number")
        if isinstance(signature, dict)
        else None
    )
    if not isinstance(signature_slide_number, int):
        raise ValueError(
            "Visual brief must bind signature_visual to an integer slide_number."
        )

    def render_order(path: Path) -> tuple[int, str]:
        match = re.search(r"-(\d+)\.png$", path.name, re.IGNORECASE)
        return (
            int(match.group(1)) if match else 10**9,
            path.name.casefold(),
        )

    slides = sorted(
        (
            Path(item).resolve()
            for item in rendered_files
            if Path(item).suffix.lower() == ".png"
            and Path(item).name.casefold() != "montage.png"
        ),
        key=render_order,
    )
    expected_slides = len(Presentation(artifact).slides)
    if len(slides) != expected_slides:
        raise ValueError(
            f"Inspection packet needs one PNG per slide; found {len(slides)} "
            f"for {expected_slides} slides."
        )
    if not 1 <= signature_slide_number <= expected_slides:
        raise ValueError(
            "Visual brief signature_visual.slide_number is outside the deck."
        )
    for path in (artifact, visual_brief, *slides):
        try:
            path.relative_to(root)
        except ValueError as exc:
            raise ValueError(
                f"Inspection inputs must live under {root}: {path}"
            ) from exc

    montage_path = slides[0].parent / "montage.png"
    build_slide_montage(slides, montage_path)
    payload = {
        "schema_version": "1.0",
        "artifact": {
            "path": artifact.relative_to(root).as_posix(),
            "sha256": _sha256(artifact),
            "size_bytes": artifact.stat().st_size,
        },
        "visual_brief": {
            "path": visual_brief.relative_to(root).as_posix(),
            "sha256": _sha256(visual_brief),
        },
        "deck_mode": deck_mode,
        "slide_count": expected_slides,
        "rendered_slides": [
            {
                "slide_number": index,
                "path": slide.relative_to(root).as_posix(),
                "sha256": _sha256(slide),
                "size_bytes": slide.stat().st_size,
            }
            for index, slide in enumerate(slides, 1)
        ],
        "montage": {
            "path": montage_path.relative_to(root).as_posix(),
            "sha256": _sha256(montage_path),
            "size_bytes": montage_path.stat().st_size,
        },
        "signature_visual": {
            "slide_number": signature_slide_number,
            "concept": str(signature.get("concept") or ""),
            "visual_type": str(signature.get("visual_type") or ""),
            "rendered_path": slides[
                signature_slide_number - 1
            ].relative_to(root).as_posix(),
            "rendered_sha256": _sha256(
                slides[signature_slide_number - 1]
            ),
        },
        "inspection": {
            "inspector": inspector,
            "full_size_each_slide_inspected": False,
            "montage_inspected": False,
            "signature_exhibit_present": False,
            "signature_exhibit_matches_brief": False,
            "findings_resolved": False,
            "status": "pending",
            "resolved_findings": [],
            "unresolved_findings": [],
        },
    }
    receipt_path.parent.mkdir(parents=True, exist_ok=True)
    receipt_path.write_text(
        json.dumps(payload, indent=2, ensure_ascii=False) + "\n",
        encoding="utf-8",
    )
    return receipt_path


def qa_visual_inspection_receipt(
    receipt_path: Path,
    *,
    artifact: Path,
    visual_brief: Path,
    deck_mode: str,
) -> QualityReport:
    """Validate that inspection covers the exact deck, brief, and slide renders."""

    receipt_path = Path(receipt_path)
    artifact = Path(artifact)
    visual_brief = Path(visual_brief)
    issues: list[QualityIssue] = []
    payload: dict = {}
    try:
        brief_payload = _load_visual_brief(visual_brief)
    except ValueError as exc:
        brief_payload = None
        issues.append(
            QualityIssue(
                code="inspection_visual_brief_invalid",
                severity="error",
                message=str(exc),
                location=str(visual_brief),
            )
        )
    try:
        payload = json.loads(receipt_path.read_text(encoding="utf-8"))
        if not isinstance(payload, dict):
            raise ValueError("receipt root is not an object")
    except (OSError, json.JSONDecodeError, ValueError) as exc:
        return QualityReport(
            artifact=str(receipt_path),
            kind="visual_inspection",
            issues=[
                QualityIssue(
                    code="invalid_visual_inspection_receipt",
                    severity="error",
                    message=f"Visual inspection receipt cannot be read: {exc}",
                    location=str(receipt_path),
                )
            ],
        )

    root = receipt_path.parent

    def error(code: str, message: str, location: str = "") -> None:
        issues.append(
            QualityIssue(
                code=code,
                severity="error",
                message=message,
                location=location or str(receipt_path),
            )
        )

    if payload.get("schema_version") != "1.0":
        error("inspection_schema", "Visual inspection receipt schema must be 1.0.")
    if payload.get("deck_mode") != deck_mode:
        error(
            "inspection_mode_mismatch",
            "Visual inspection receipt does not match the requested deck mode.",
        )

    for key, expected in (("artifact", artifact), ("visual_brief", visual_brief)):
        record = payload.get(key)
        if not isinstance(record, dict):
            error("inspection_binding_missing", f"Receipt has no {key} binding.")
            continue
        try:
            bound_path = _safe_relative_file(root, str(record.get("path") or ""))
        except ValueError as exc:
            error("inspection_binding_invalid", str(exc))
            continue
        if bound_path.resolve() != expected.resolve():
            error(
                "inspection_binding_mismatch",
                f"Receipt {key} path is not the expected file.",
            )
        elif record.get("sha256") != _sha256(expected):
            error(
                "inspection_hash_mismatch",
                f"Receipt {key} hash does not match the current bytes.",
            )

    try:
        slide_count = len(Presentation(artifact).slides)
    except Exception as exc:  # noqa: BLE001
        error("inspection_artifact_unreadable", f"Cannot reopen inspected deck: {exc}")
        slide_count = 0
    if payload.get("slide_count") != slide_count:
        error(
            "inspection_slide_count_mismatch",
            "Inspection receipt slide count does not match the deck.",
        )

    rendered = payload.get("rendered_slides")
    if not isinstance(rendered, list) or len(rendered) != slide_count:
        error(
            "inspection_render_inventory",
            "Receipt must bind exactly one full-size PNG to every slide.",
        )
        rendered = []
    seen_numbers: set[int] = set()
    seen_paths: set[Path] = set()
    rendered_by_number: dict[int, tuple[Path, dict]] = {}
    for record in rendered:
        if not isinstance(record, dict):
            error("inspection_render_inventory", "Rendered-slide record is invalid.")
            continue
        number = record.get("slide_number")
        if not isinstance(number, int) or number in seen_numbers:
            error(
                "inspection_render_sequence",
                "Rendered slides need unique integer slide numbers.",
            )
        else:
            seen_numbers.add(number)
        try:
            rendered_path = _safe_relative_file(
                root, str(record.get("path") or "")
            )
        except ValueError as exc:
            error("inspection_render_missing", str(exc))
            continue
        if rendered_path in seen_paths:
            error(
                "inspection_render_inventory",
                "Every rendered-slide record must reference a distinct PNG.",
            )
        seen_paths.add(rendered_path)
        if isinstance(number, int):
            rendered_by_number[number] = (rendered_path, record)
        if rendered_path.suffix.lower() != ".png":
            error(
                "inspection_render_format",
                "Every inspected slide render must be a PNG.",
                str(rendered_path),
            )
        elif record.get("sha256") != _sha256(rendered_path):
            error(
                "inspection_render_hash_mismatch",
                "Rendered slide bytes changed after inspection.",
                str(rendered_path),
            )
    if seen_numbers and seen_numbers != set(range(1, slide_count + 1)):
        error(
            "inspection_render_sequence",
            "Rendered-slide numbering does not cover the complete deck.",
        )

    expected_signature = (
        brief_payload.get("signature_visual")
        if isinstance(brief_payload, dict)
        else None
    )
    expected_signature_slide = (
        expected_signature.get("slide_number")
        if isinstance(expected_signature, dict)
        else None
    )
    signature_record = payload.get("signature_visual")
    if (
        not isinstance(signature_record, dict)
        or not isinstance(expected_signature_slide, int)
    ):
        error(
            "inspection_signature_binding_missing",
            "Receipt has no canonical signature-slide binding.",
        )
    elif signature_record.get("slide_number") != expected_signature_slide:
        error(
            "inspection_signature_binding_mismatch",
            "Receipt signature slide does not match the canonical visual brief.",
        )
    else:
        rendered_signature = rendered_by_number.get(expected_signature_slide)
        if rendered_signature is None:
            error(
                "inspection_signature_render_missing",
                "Signature slide has no bound full-size render.",
            )
        else:
            rendered_path, rendered_record = rendered_signature
            if (
                signature_record.get("rendered_path")
                != rendered_record.get("path")
                or signature_record.get("rendered_sha256")
                != _sha256(rendered_path)
            ):
                error(
                    "inspection_signature_render_mismatch",
                    "Receipt signature visual is not bound to its exact slide render.",
                )
        for field in ("concept", "visual_type"):
            if str(signature_record.get(field) or "") != str(
                expected_signature.get(field) or ""
            ):
                error(
                    "inspection_signature_binding_mismatch",
                    f"Receipt signature {field} differs from the visual brief.",
                )

    montage = payload.get("montage")
    if not isinstance(montage, dict):
        error("inspection_montage_missing", "Receipt has no montage binding.")
    else:
        try:
            montage_path = _safe_relative_file(
                root, str(montage.get("path") or "")
            )
        except ValueError as exc:
            error("inspection_montage_missing", str(exc))
        else:
            if (
                montage_path.suffix.lower() != ".png"
                or montage.get("sha256") != _sha256(montage_path)
            ):
                error(
                    "inspection_montage_hash_mismatch",
                    "Montage bytes do not match the inspected montage.",
                )

    inspection = payload.get("inspection")
    if not isinstance(inspection, dict):
        error("inspection_attestation_missing", "Receipt has no inspection attestation.")
    else:
        if not str(inspection.get("inspector") or "").strip():
            error("inspection_identity_missing", "Receipt must name its inspector.")
        for field in (
            "full_size_each_slide_inspected",
            "montage_inspected",
            "signature_exhibit_present",
            "signature_exhibit_matches_brief",
            "findings_resolved",
        ):
            if inspection.get(field) is not True:
                error(
                    "inspection_incomplete",
                    f"Visual inspection attestation {field!r} is not true.",
                )
        if inspection.get("status") != "pass":
            error(
                "inspection_not_passed",
                "Visual inspection status must be 'pass' before release.",
            )
        unresolved = inspection.get("unresolved_findings")
        if not isinstance(unresolved, list) or unresolved:
            error(
                "inspection_findings_open",
                "Visual inspection has unresolved findings.",
            )

    return QualityReport(
        artifact=str(receipt_path),
        kind="visual_inspection",
        issues=issues,
        metadata={
            "deck_mode": deck_mode,
            "slides": slide_count,
            "inspector": (
                payload.get("inspection", {}).get("inspector")
                if isinstance(payload.get("inspection"), dict)
                else None
            ),
            "receipt_sha256": _sha256(receipt_path),
        },
    )


def _shape_text(shape) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    return shape.text_frame.text.strip()


def _runs_with_sizes(shape) -> Iterable[tuple[str, float | None]]:
    if not getattr(shape, "has_text_frame", False):
        return
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.text.strip():
                yield run.text, run.font.size.pt if run.font.size is not None else None


def _notes_text(slide) -> str:
    try:
        return slide.notes_slide.notes_text_frame.text or ""
    except (AttributeError, KeyError, ValueError):
        return ""


def _is_source_shape(shape, text: str) -> bool:
    if SOURCE_RE.search(text):
        return True
    name = getattr(shape, "name", "").lower()
    return "source" in name or "footnote" in name


def _shape_smallest_size(shape) -> float | None:
    sizes = [size for _, size in _runs_with_sizes(shape) if size is not None]
    return min(sizes) if sizes else None


def _find_title_shape(slide, slide_height: Emu):
    """Find a title placeholder or the strongest manual headline near the top."""
    if slide.shapes.title is not None and _shape_text(slide.shapes.title):
        return slide.shapes.title
    # Native-shape decks frequently pair an assertion headline with a larger
    # hero metric. An explicitly named ``*-title`` shape is the author's
    # unambiguous title contract and must outrank that metric.
    for shape in slide.shapes:
        name = str(getattr(shape, "name", "") or "").casefold()
        if (
            re.search(r"(?:^|[-_ ])title$", name)
            and _shape_text(shape)
            and getattr(shape, "top", slide_height + 1)
            <= int(slide_height * 0.40)
        ):
            return shape
    candidates: list[tuple[float, int, object]] = []
    top_limit = int(slide_height * 0.40)
    for shape in slide.shapes:
        text = _shape_text(shape)
        if not text or getattr(shape, "top", top_limit + 1) > top_limit:
            continue
        sizes = [size for _, size in _runs_with_sizes(shape) if size is not None]
        largest = max(sizes) if sizes else 0.0
        # Prefer the largest type, then the highest object on the canvas.
        candidates.append((largest, -int(getattr(shape, "top", 0)), shape))
    if not candidates:
        return None
    return max(candidates, key=lambda item: (item[0], item[1]))[2]


def _shape_outside_canvas(shape, slide_width: Emu, slide_height: Emu) -> bool:
    try:
        return (
            shape.left < 0
            or shape.top < 0
            or shape.left + shape.width > slide_width
            or shape.top + shape.height > slide_height
        )
    except (AttributeError, TypeError):
        return False


def _likely_text_overflow(shape, smallest_size: float | None) -> bool:
    """Conservative capacity estimate for text boxes.

    PowerPoint's exact autofit calculation is unavailable in python-pptx.  This
    only flags extreme cases and leaves normal wrapping for visual inspection.
    """
    text = _shape_text(shape)
    if not text or smallest_size is None or smallest_size <= 0:
        return False
    width_pt = shape.width.pt
    height_pt = shape.height.pt
    if width_pt <= 0 or height_pt <= 0:
        return True
    chars_per_line = max(8, int(width_pt / (smallest_size * 0.52)))
    explicit_lines = text.splitlines() or [text]
    estimated_lines = sum(
        max(1, math.ceil(len(line) / chars_per_line)) for line in explicit_lines
    )
    line_height = smallest_size * 1.22
    return estimated_lines * line_height > height_pt * 1.12


def qa_presentation(
    path: Path,
    *,
    config: PresentationQAConfig | None = None,
    render_dir: Path | None = None,
    deck_mode: str | None = None,
    visual_brief: Path | dict | None = None,
) -> QualityReport:
    path = Path(path)
    brief = _load_visual_brief(visual_brief)
    requested_mode = (
        deck_mode
        or (str(brief.get("deck_mode")) if brief else None)
    )
    config = config or (
        PresentationQAConfig.for_mode(requested_mode)
        if requested_mode
        else PresentationQAConfig()
    )
    if (
        requested_mode
        and config.deck_mode
        and requested_mode != config.deck_mode
    ):
        raise ValueError(
            f"QA config mode {config.deck_mode!r} does not match "
            f"{requested_mode!r}."
        )
    issues: list[QualityIssue] = []
    try:
        deck = Presentation(path)
    except Exception as exc:  # noqa: BLE001 - produce useful quality report
        return QualityReport(
            artifact=str(path),
            kind="pptx",
            issues=[
                QualityIssue(
                    code="unreadable_pptx",
                    severity="error",
                    message=f"Presentation cannot be reopened: {exc}",
                    location=str(path),
                )
            ],
        )

    slide_count = len(deck.slides)
    if slide_count < config.min_slide_count:
        issues.append(
            QualityIssue(
                code="short_deck",
                severity=config.slide_count_severity,
                message=(
                    f"{requested_mode or 'Presentation'} has only "
                    f"{slide_count} slides; minimum is "
                    f"{config.min_slide_count}."
                ),
                location=str(path),
            )
        )
    if slide_count > config.max_slide_count:
        issues.append(
            QualityIssue(
                code="long_deck",
                severity=config.slide_count_severity,
                message=(
                    f"{requested_mode or 'Presentation'} has "
                    f"{slide_count} slides; maximum is "
                    f"{config.max_slide_count}."
                ),
                location=str(path),
            )
        )

    ratio = deck.slide_width / deck.slide_height
    if abs(ratio - (16 / 9)) > 0.03:
        issues.append(
            QualityIssue(
                code="non_widescreen_canvas",
                severity="warning",
                message=f"Deck aspect ratio is {ratio:.3f}, not widescreen 16:9.",
                location=str(path),
            )
        )

    visual_shapes = 0
    signature_markers: list[int] = []
    numeric_slides = 0
    sourced_numeric_slides = 0
    title_sizes: list[float] = []
    body_sizes: list[float] = []
    source_sizes: list[float] = []
    brief_slides = (
        brief.get("slides", [])
        if isinstance(brief, dict)
        else []
    )
    expected_by_number = {
        int(item.get("slide_number")): item
        for item in brief_slides
        if isinstance(item, dict)
        and isinstance(item.get("slide_number"), int)
    }
    signature = (
        brief.get("signature_visual")
        if isinstance(brief, dict)
        else None
    )
    signature_slide_number = (
        signature.get("slide_number")
        if isinstance(signature, dict)
        else None
    )
    if brief is not None and (
        not isinstance(signature_slide_number, int)
        or signature_slide_number not in expected_by_number
    ):
        issues.append(
            QualityIssue(
                code="signature_visual_slide_invalid",
                severity="error",
                message=(
                    "Canonical visual brief must bind its signature visual "
                    "to one slide in the deck."
                ),
                location=str(path),
            )
        )
    if brief is not None and len(expected_by_number) != slide_count:
        issues.append(
            QualityIssue(
                code="visual_brief_slide_count_mismatch",
                severity="error",
                message=(
                    f"Deck has {slide_count} slides but its canonical visual "
                    f"brief specifies {len(expected_by_number)}."
                ),
                location=str(path),
            )
        )

    for slide_number, slide in enumerate(deck.slides, 1):
        location = f"slide {slide_number}"
        title_shape = _find_title_shape(slide, deck.slide_height)
        title_text = _shape_text(title_shape) if title_shape is not None else ""
        if not title_text:
            issues.append(
                QualityIssue(
                    code="missing_slide_title",
                    severity="error",
                    message="Every slide must have an audience-facing assertion title.",
                    location=location,
                )
            )

        slide_text_parts: list[str] = []
        source_present = False
        slide_has_numeric_claim = False
        slide_visuals = 0
        vector_shapes = 0
        expected_slide = expected_by_number.get(slide_number)
        if expected_slide is not None:
            expected_headline = str(expected_slide.get("headline") or "")
            if (
                expected_headline
                and _normalise_headline(title_text)
                != _normalise_headline(expected_headline)
            ):
                issues.append(
                    QualityIssue(
                        code="visual_brief_headline_mismatch",
                        severity="error",
                        message=(
                            "Slide title does not match the approved visual "
                            f"brief headline: {expected_headline!r}."
                        ),
                        location=location,
                    )
                )
        appendix_slide = bool(
            requested_mode == "technical_read_ahead"
            and expected_slide is not None
            and re.search(
                r"\b(?:appendix|sources?|technical detail|assumptions?|"
                r"methodology)\b",
                " ".join(
                    str(expected_slide.get(key) or "")
                    for key in (
                        "narrative_job",
                        "headline",
                        "density_budget",
                    )
                ),
                re.IGNORECASE,
            )
        )
        required_body_pt = (
            config.appendix_body_pt
            if appendix_slide and config.appendix_body_pt is not None
            else config.min_body_pt
        )

        for shape in slide.shapes:
            text = _shape_text(shape)
            if getattr(shape, "name", "").strip().casefold().startswith(
                SIGNATURE_VISUAL_PREFIX.casefold()
            ):
                signature_markers.append(slide_number)
            if text:
                slide_text_parts.append(text)
                copy_issues = lint_reader_text(text, location=location)
                if config.allow_internal_terms:
                    copy_issues = [
                        issue
                        for issue in copy_issues
                        if issue.code
                        not in {
                            "brief_reference",
                            "internal_filename",
                            "stage_header",
                        }
                    ]
                issues.extend(copy_issues)
                if not config.allow_internal_terms and INTERNAL_TOKEN_RE.search(text):
                    issues.append(
                        QualityIssue(
                            code="internal_process_language",
                            severity="error",
                            message=f"Internal production language appears in visible copy: {text[:120]!r}",
                            location=location,
                        )
                    )
                if NUMERIC_CLAIM_RE.search(text):
                    slide_has_numeric_claim = True

            if _shape_outside_canvas(shape, deck.slide_width, deck.slide_height):
                issues.append(
                    QualityIssue(
                        code="canvas_overflow",
                        severity="error",
                        message=f"Shape {getattr(shape, 'name', '<unnamed>')!r} extends beyond the slide canvas.",
                        location=location,
                    )
                )

            is_title = (
                title_shape is not None
                and getattr(shape, "shape_id", None)
                == getattr(title_shape, "shape_id", None)
            )
            is_source = bool(text and _is_source_shape(shape, text))
            if is_source:
                source_present = True
            sizes = [size for _, size in _runs_with_sizes(shape) if size is not None]
            smallest = min(sizes) if sizes else None
            if sizes:
                if is_title:
                    title_sizes.extend(sizes)
                    required = (
                        config.min_cover_title_pt
                        if slide_number == 1
                        else config.min_title_pt
                    )
                    if min(sizes) < required:
                        issues.append(
                            QualityIssue(
                                code="small_title_type",
                                severity="error",
                                message=f"Title type is {min(sizes):g} pt; minimum is {required:g} pt.",
                                location=location,
                            )
                        )
                elif is_source:
                    source_sizes.extend(sizes)
                    if min(sizes) < config.min_source_pt:
                        issues.append(
                            QualityIssue(
                                code="small_source_type",
                                severity="error",
                                message=f"Source note is {min(sizes):g} pt; minimum is {config.min_source_pt:g} pt.",
                                location=location,
                            )
                        )
                elif text and not re.fullmatch(r"\d{1,2}", text.strip()):
                    body_sizes.extend(sizes)
                    if min(sizes) < required_body_pt:
                        issues.append(
                            QualityIssue(
                                code="small_body_type",
                                severity=config.body_type_severity,
                                message=(
                                    f"Body type is {min(sizes):g} pt; minimum "
                                    f"is {required_body_pt:g} pt"
                                    + (
                                        " on a technical appendix slide."
                                        if appendix_slide
                                        else "."
                                    )
                                ),
                                location=location,
                            )
                        )
            elif is_title:
                issues.append(
                    QualityIssue(
                        code="unknown_title_type_size",
                        severity="error",
                        message="Slide title has no explicit, inspectable type size.",
                        location=location,
                    )
                )
            elif is_source:
                issues.append(
                    QualityIssue(
                        code="unknown_source_type_size",
                        severity="error",
                        message=(
                            "Source note inherits an unresolved theme size; "
                            f"set it explicitly to at least {config.min_source_pt:g} pt."
                        ),
                        location=location,
                    )
                )
            if text and _likely_text_overflow(shape, smallest):
                issues.append(
                    QualityIssue(
                        code="likely_text_overflow",
                        severity="warning",
                        message=f"Text may not fit shape {getattr(shape, 'name', '<unnamed>')!r}; render and inspect it.",
                        location=location,
                    )
                )

            if (
                getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE
                or getattr(shape, "has_chart", False)
                or getattr(shape, "has_table", False)
                or getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP
            ):
                slide_visuals += 1
            elif (
                getattr(shape, "shape_type", None)
                in {
                    MSO_SHAPE_TYPE.AUTO_SHAPE,
                    MSO_SHAPE_TYPE.LINE,
                    MSO_SHAPE_TYPE.FREEFORM,
                }
                and not is_title
                and not is_source
            ):
                vector_shapes += 1

        notes = _notes_text(slide)
        source_present = source_present or bool(SOURCE_RE.search(notes))
        if slide_has_numeric_claim:
            numeric_slides += 1
            if source_present:
                sourced_numeric_slides += 1
            elif config.require_numeric_sources:
                issues.append(
                    QualityIssue(
                        code="unsourced_numeric_slide",
                        severity=config.evidence_issue_severity,
                        message="Slide contains a material numeric claim but no visible or speaker-note source.",
                        location=location,
                    )
                )
        visual_shapes += slide_visuals
        if expected_slide is not None:
            expected_evidence = expected_slide.get("evidence_ids") or []
            expected_source = str(
                expected_slide.get("source_note") or ""
            ).strip()
            expected_visual_type = str(
                expected_slide.get("visual_type") or ""
            ).casefold()
            expects_visual = expected_visual_type not in {
                "",
                "text",
                "text only",
                "quote",
                "title",
            }
            if expects_visual and slide_visuals == 0 and vector_shapes < 3:
                issues.append(
                    QualityIssue(
                        code="visual_brief_exhibit_missing",
                        severity="error",
                        message=(
                            "Slide does not contain the evidence-bearing "
                            f"{expected_visual_type!r} exhibit required by the "
                            "visual brief."
                        ),
                        location=location,
                    )
                )
            if (expected_evidence or expected_source) and not source_present:
                issues.append(
                    QualityIssue(
                        code="visual_brief_source_missing",
                        severity="error",
                        message=(
                            "Visual brief binds this slide to evidence, but "
                            "the slide has no visible or speaker-note source."
                        ),
                        location=location,
                    )
                )

        if (
            isinstance(signature_slide_number, int)
            and slide_number == signature_slide_number
            and slide_number not in signature_markers
        ):
            issues.append(
                QualityIssue(
                    code="signature_visual_exhibit_missing",
                    severity="error",
                    message=(
                        "The canonical signature slide has no exhibit named "
                        f"'{SIGNATURE_VISUAL_PREFIX}<concept>'."
                    ),
                    location=location,
                )
            )

    misplaced_signature_markers = sorted(
        {
            number
            for number in signature_markers
            if number != signature_slide_number
        }
    )
    if misplaced_signature_markers:
        issues.append(
            QualityIssue(
                code="signature_visual_marker_misplaced",
                severity="error",
                message=(
                    "Signature-visual marker appears outside the canonical "
                    f"slide: {misplaced_signature_markers}."
                ),
                location=str(path),
            )
        )
    if len(signature_markers) > 1:
        issues.append(
            QualityIssue(
                code="signature_visual_marker_ambiguous",
                severity="error",
                message=(
                    "Deck must contain exactly one reserved signature-visual "
                    f"marker; found {len(signature_markers)}."
                ),
                location=str(path),
            )
        )

    if config.require_visual_evidence and visual_shapes == 0:
        issues.append(
            QualityIssue(
                code="no_evidence_visuals",
                severity=config.evidence_issue_severity,
                message=(
                    "Deck contains no pictures, charts, tables, or grouped diagrams. "
                    "A board deck needs at least one evidence-bearing signature visual."
                ),
                location=str(path),
            )
        )

    rendered_files: list[str] = []
    if render_dir is not None:
        rendered, render_issues = render_office_artifact(
            path,
            render_dir,
            required=config.require_render,
        )
        rendered_files = [str(item) for item in rendered]
        issues.extend(render_issues)
    if config.require_render and not any(
        Path(item).suffix.lower() == ".png"
        for item in rendered_files
    ):
        issues.append(
            QualityIssue(
                code="render_required",
                severity="error",
                message=(
                    "Client-release presentation QA requires a fresh PDF and "
                    "full-slide PNG render; structural inspection alone cannot "
                    "release the deck."
                ),
                location=str(path),
            )
        )

    return QualityReport(
        artifact=str(path),
        kind="pptx",
        issues=issues,
        metadata={
            "slides": slide_count,
            "aspect_ratio": round(ratio, 4),
            "visual_shapes": visual_shapes,
            "numeric_slides": numeric_slides,
            "sourced_numeric_slides": sourced_numeric_slides,
            "minimum_title_pt": min(title_sizes) if title_sizes else None,
            "minimum_body_pt": min(body_sizes) if body_sizes else None,
            "minimum_source_pt": min(source_sizes) if source_sizes else None,
            "profile": config.profile,
            "deck_mode": requested_mode,
            "visual_brief_slides": len(expected_by_number),
            "signature_visual_slide": signature_slide_number,
            "signature_visual_markers": signature_markers,
        },
        rendered_files=rendered_files,
    )


def main() -> int:
    parser = argparse.ArgumentParser(description="Quality-check a Council PowerPoint.")
    parser.add_argument("pptx", type=Path)
    parser.add_argument("--json", type=Path, dest="json_path")
    parser.add_argument("--render-dir", type=Path)
    parser.add_argument(
        "--mode",
        choices=(
            "board_decision",
            "executive_briefing",
            "technical_read_ahead",
            "argument_brief",
        ),
    )
    parser.add_argument(
        "--slide-count",
        type=int,
        help="Exact slide count required when --mode argument_brief is selected.",
    )
    parser.add_argument("--visual-brief", type=Path)
    parser.add_argument(
        "--prepare-inspection",
        type=Path,
        help=(
            "Write a hash-bound visual-inspection packet with pending "
            "attestations. Inspect the renders, then explicitly mark it passed."
        ),
    )
    parser.add_argument("--fail-on-warnings", action="store_true")
    parser.add_argument(
        "--profile",
        choices=("client-release", "process-explainer"),
        default="client-release",
    )
    args = parser.parse_args()

    if args.mode == "argument_brief" and args.slide_count is None:
        parser.error("--mode argument_brief requires --slide-count")
    if args.slide_count is not None and args.mode != "argument_brief":
        parser.error("--slide-count is only valid with --mode argument_brief")
    config = (
        PresentationQAConfig.process_explainer()
        if args.profile == "process-explainer"
        else (
            PresentationQAConfig.for_argument(args.slide_count)
            if args.mode == "argument_brief"
            else PresentationQAConfig.for_mode(args.mode)
            if args.mode
            else PresentationQAConfig()
        )
    )
    report = qa_presentation(
        args.pptx,
        config=config,
        render_dir=args.render_dir,
        deck_mode=args.mode,
        visual_brief=args.visual_brief,
    )
    if args.json_path:
        report.write_json(args.json_path)
    if args.prepare_inspection:
        if not args.mode or args.visual_brief is None or args.render_dir is None:
            parser.error(
                "--prepare-inspection requires --mode, --visual-brief, and "
                "--render-dir"
            )
        if not report.ok:
            print(json.dumps(report.to_dict(), indent=2))
            return 1
        prepare_visual_inspection_receipt(
            artifact=args.pptx,
            visual_brief=args.visual_brief,
            deck_mode=args.mode,
            rendered_files=report.rendered_files,
            receipt_path=args.prepare_inspection,
        )
    print(json.dumps(report.to_dict(), indent=2))
    if not report.ok or (args.fail_on_warnings and report.warnings):
        return 1
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
