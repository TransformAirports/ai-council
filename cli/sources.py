"""Operator-supplied source material — the file-drop workflow.

The repo's `sources/` folder is a drop zone. When the operator launches a
new run, the hub detects what's there, optionally moves the files into the
persistent `sources/runs/<slug>/` folder, and converts office formats to plain
markdown sidecars so every agent (including the OpenAI-hosted Deep Research
lens, which has no file tools) can read them. Persistent paths keep saved run
prompts reproducible after the transient `outputs/` workspace is cleared.

Supported formats:
  - .md, .txt, .csv, .json, .yaml → pass through; agents read directly
  - .pdf                          → pass through; Claude SDK's Read handles it
  - .docx, .pptx, .xlsx           → extracted to a .extracted.md sidecar
  - other                         → kept but flagged; agents see path only
"""
from __future__ import annotations

import shutil
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable

REPO_ROOT = Path(__file__).resolve().parent.parent
DROPZONE = REPO_ROOT / "sources"

EXTRACTABLE_OFFICE = {".docx", ".pptx", ".xlsx", ".pdf"}
NATIVELY_READABLE = {".md", ".txt", ".csv", ".json", ".yaml", ".yml"}
IGNORE_NAMES = {".ds_store", ".gitkeep"}

# Cap extracted text so a single tool-result can never exceed the SDK's stdout
# buffer or balloon token cost. 600k chars ≈ 150k words ≈ a 400-page document.
MAX_EXTRACT_CHARS = 600_000


class SourcePathError(ValueError):
    """Raised when a source path crosses the Council's trusted file boundary."""


@dataclass
class SourceFile:
    original: Path        # path of the file in sources/runs/<slug>/
    readable: Path        # path agents should read (sidecar or original)
    description: str
    size_bytes: int
    extracted: bool


# ----------------------------------------------------------------------------
# Discovery.
# ----------------------------------------------------------------------------

def _is_within(path: Path, root: Path) -> bool:
    try:
        path.relative_to(root)
    except ValueError:
        return False
    return True


def _is_lexically_within(path: Path, root: Path) -> bool:
    try:
        path.absolute().relative_to(root.absolute())
    except ValueError:
        return False
    return True


def _resolved_root(root: Path, *, label: str) -> Path:
    """Return a real approved root; the root itself may not be a symlink."""

    if root.is_symlink():
        raise SourcePathError(f"{label} may not be a symlink: {root}")
    try:
        resolved = root.resolve(strict=True)
    except OSError as exc:
        raise SourcePathError(f"{label} is unavailable: {root}") from exc
    if not resolved.is_dir():
        raise SourcePathError(f"{label} is not a directory: {root}")
    return resolved


def _has_symlink_below(path: Path, root: Path) -> bool:
    """Whether any lexical component from ``root`` through ``path`` is a link."""

    try:
        relative = path.absolute().relative_to(root.absolute())
    except ValueError:
        return False
    cursor = root.absolute()
    for part in relative.parts:
        cursor = cursor / part
        if cursor.is_symlink():
            return True
    return False


def _approved_source_path(
    path: Path,
    *,
    dropzone: Path,
    allow_external: bool = False,
    approved_external_roots: Iterable[Path] = (),
) -> tuple[Path, bool]:
    """Validate one input and return ``(resolved_path, is_dropzone_input)``.

    Normal runs may ingest only regular, non-symlink files under the repository
    source drop zone (excluding the durable ``sources/runs`` library). A caller
    that deliberately needs an external file must opt in *and* name the
    external root it approves. This keeps a boolean flag from becoming an
    unrestricted filesystem read.
    """

    dropzone_root = _resolved_root(dropzone, label="Source drop zone")
    candidate = Path(path)
    if candidate.is_symlink():
        raise SourcePathError(f"Source files may not be symlinks: {candidate}")
    try:
        resolved = candidate.resolve(strict=True)
    except OSError as exc:
        raise SourcePathError(f"Source file is unavailable: {candidate}") from exc
    if not resolved.is_file():
        raise SourcePathError(f"Source path is not a regular file: {candidate}")

    if _is_within(resolved, dropzone_root) and _is_lexically_within(
        candidate, dropzone
    ):
        if _has_symlink_below(candidate, dropzone):
            raise SourcePathError(
                f"Source paths may not traverse symlinks: {candidate}"
            )
        relative = resolved.relative_to(dropzone_root)
        if relative.parts and relative.parts[0] == "runs":
            raise SourcePathError(
                "Already-attached files in sources/runs/ cannot be ingested "
                f"again: {candidate}"
            )
        return resolved, True

    external_roots = tuple(Path(root) for root in approved_external_roots)
    if not allow_external:
        raise SourcePathError(
            f"Source file is outside the approved drop zone {dropzone_root}: "
            f"{candidate}. External sources require explicit opt-in."
        )
    if not external_roots:
        raise SourcePathError(
            "External source opt-in requires at least one approved external root."
        )
    for raw_root in external_roots:
        approved_root = _resolved_root(
            raw_root, label="Approved external source root"
        )
        if not _is_within(resolved, approved_root) or not _is_lexically_within(
            candidate, raw_root
        ):
            continue
        if _has_symlink_below(candidate, raw_root):
            raise SourcePathError(
                f"External source paths may not traverse symlinks: {candidate}"
            )
        return resolved, False
    raise SourcePathError(
        f"External source is not within an explicitly approved root: {candidate}"
    )


def discover_dropzone(dropzone: Path = DROPZONE) -> list[Path]:
    """Files currently in the drop zone, ready to attach to a new run."""
    if not dropzone.is_dir():
        return []
    dropzone_root = _resolved_root(dropzone, label="Source drop zone")
    discovered: list[Path] = []
    for path in dropzone.rglob("*"):
        # Do not follow or advertise links, even when a link happens to resolve
        # to a regular file inside the source tree.
        if path.is_symlink():
            continue
        if not path.is_file():
            continue
        try:
            resolved = path.resolve(strict=True)
        except OSError:
            continue
        if not _is_within(resolved, dropzone_root):
            continue
        relative = resolved.relative_to(dropzone_root)
        # `sources/runs/` is the durable per-run library, not part of the
        # operator drop zone. Never offer an already-attached source again.
        if relative.parts and relative.parts[0] == "runs":
            continue
        # Browser uploads are selected explicitly by opaque token. Never let a
        # staged file leak into the legacy global drop-zone workflow.
        if any(part.startswith(".") for part in relative.parts):
            continue
        if path.name.lower() in IGNORE_NAMES or path.name.startswith("."):
            continue
        # Preserve the caller's lexical root in the returned path; attachment
        # resolves and revalidates it before any move. This keeps UI labels and
        # existing call sites stable on systems where /var resolves to
        # /private/var.
        discovered.append(path)
    return sorted(discovered)


def format_size(n: int) -> str:
    if n < 1024:
        return f"{n} B"
    if n < 1024 ** 2:
        return f"{n / 1024:.1f} KB"
    return f"{n / 1024 ** 2:.1f} MB"


# ----------------------------------------------------------------------------
# Attach: move + extract.
# ----------------------------------------------------------------------------

def attach_sources(
    slug: str,
    files: list[Path],
    outputs_dir: Path,
    *,
    allow_external: bool = False,
    approved_external_roots: Iterable[Path] = (),
) -> list[SourceFile]:
    """Attach trusted source files to ``sources/runs/<slug>/`` and extract them.

    Drop-zone files are moved. Explicitly approved external files are copied so
    attaching a source never removes a file outside the Council workspace.
    """

    dropzone = outputs_dir.parent / "sources"
    dropzone.mkdir(parents=True, exist_ok=True)
    library_root = dropzone / "runs"
    library_root.mkdir(parents=True, exist_ok=True)
    if not slug or Path(slug).name != slug or slug in {".", ".."}:
        raise SourcePathError(f"Unsafe source-library slug: {slug!r}")
    target_dir = library_root / slug
    if target_dir.is_symlink():
        raise SourcePathError(
            f"Source-library destination may not be a symlink: {target_dir}"
        )
    target_dir.mkdir(parents=True, exist_ok=True)
    resolved_library = _resolved_root(library_root, label="Source library")
    resolved_target = target_dir.resolve(strict=True)
    if not _is_within(resolved_target, resolved_library):
        raise SourcePathError(
            f"Source-library destination escapes its approved root: {target_dir}"
        )

    out: list[SourceFile] = []
    for supplied in files:
        src, from_dropzone = _approved_source_path(
            supplied,
            dropzone=dropzone,
            allow_external=allow_external,
            approved_external_roots=approved_external_roots,
        )
        def destination_collides(candidate: Path) -> bool:
            if candidate.exists() or candidate.is_symlink():
                return True
            sidecar = candidate.with_suffix(candidate.suffix + ".extracted.md")
            return sidecar.exists() or sidecar.is_symlink()

        dst = target_dir / src.name
        if destination_collides(dst):
            # Don't clobber a previously-attached file with the same name.
            stem, suffix = src.stem, src.suffix
            n = 2
            while destination_collides(target_dir / f"{stem}-{n}{suffix}"):
                n += 1
            dst = target_dir / f"{stem}-{n}{suffix}"
        if not _is_within(dst.resolve(strict=False), resolved_target):
            raise SourcePathError(
                f"Source destination escapes its approved root: {dst}"
            )
        if from_dropzone:
            shutil.move(str(src), str(dst))
        else:
            shutil.copy2(src, dst)

        ext = dst.suffix.lower()
        readable = dst
        extracted = False
        description = src.name

        if ext in EXTRACTABLE_OFFICE:
            sidecar = dst.with_suffix(dst.suffix + ".extracted.md")
            try:
                sidecar.write_text(_extract_office(dst), encoding="utf-8")
                readable = sidecar
                extracted = True
            except Exception as e:  # noqa: BLE001 — extraction failure is recoverable
                description = f"{src.name} (extraction failed: {e}; agents see binary path only)"
        elif ext in NATIVELY_READABLE:
            pass
        else:
            description = f"{src.name} (binary; agents see path only)"

        out.append(SourceFile(
            original=dst,
            readable=readable,
            description=description,
            size_bytes=dst.stat().st_size,
            extracted=extracted,
        ))
    return out


# ----------------------------------------------------------------------------
# Format extraction.
# ----------------------------------------------------------------------------

def _extract_office(path: Path) -> str:
    ext = path.suffix.lower()
    if ext == ".docx":
        text = _extract_docx(path)
    elif ext == ".pptx":
        text = _extract_pptx(path)
    elif ext == ".xlsx":
        text = _extract_xlsx(path)
    elif ext == ".pdf":
        text = _extract_pdf(path)
    else:
        raise ValueError(f"unsupported format: {ext}")
    if len(text) > MAX_EXTRACT_CHARS:
        text = (
            text[:MAX_EXTRACT_CHARS]
            + f"\n\n[…extraction truncated at {MAX_EXTRACT_CHARS:,} characters. "
            f"The source file is longer; the remainder is omitted to keep the "
            f"document within processing limits.]"
        )
    return text


def _extract_pdf(path: Path) -> str:
    """Extract text from a PDF to markdown, page by page.

    Reading a large PDF through Claude's native Read tool returns the file as
    base64 — which blows the SDK's 1 MB stdout buffer and burns huge token
    cost. Extracting text first keeps the tool-result small, cheap, and
    readable by every agent (including the OpenAI Deep Research lens).
    """
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    out: list[str] = [f"# Extracted from {path.name}", ""]
    total_chars = 0
    for i, page in enumerate(reader.pages, 1):
        try:
            page_text = (page.extract_text() or "").strip()
        except Exception:  # noqa: BLE001 — a single bad page shouldn't kill extraction
            page_text = ""
        if not page_text:
            continue
        out.append(f"## Page {i}")
        out.append(page_text)
        out.append("")
        total_chars += len(page_text)

    if total_chars < 100:
        # Almost no extractable text — very likely a scanned/image-only PDF.
        out.append(
            "\n[NOTE: this PDF yielded almost no extractable text. It is likely "
            "a scanned or image-only document. OCR is not available in this "
            "pipeline — if its contents matter, supply a text-based version "
            "(searchable PDF, .docx, or pasted .md).]"
        )
    return "\n".join(out)


def _extract_docx(path: Path) -> str:
    from docx import Document
    doc = Document(str(path))
    out: list[str] = [f"# Extracted from {path.name}", ""]
    for para in doc.paragraphs:
        t = para.text.strip()
        if not t:
            continue
        style = (para.style.name if para.style else "") or ""
        if style.startswith("Heading "):
            level_str = "".join(c for c in style if c.isdigit()) or "1"
            depth = min(int(level_str) + 1, 6)
            out.append(f"{'#' * depth} {t}")
        else:
            out.append(t)
        out.append("")
    for i, table in enumerate(doc.tables, 1):
        out.append(f"## Table {i}")
        for row in table.rows:
            cells = " | ".join((c.text.strip() or "—") for c in row.cells)
            out.append(f"| {cells} |")
        out.append("")
    return "\n".join(out)


def _extract_pptx(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(path))
    out: list[str] = [f"# Extracted from {path.name}", ""]
    for i, slide in enumerate(prs.slides, 1):
        out.append(f"## Slide {i}")
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    out.append(t)
        try:
            notes = slide.notes_slide.notes_text_frame.text.strip() if slide.has_notes_slide else ""
        except Exception:  # noqa: BLE001 — some decks have malformed notes parts
            notes = ""
        if notes:
            out.append(f"*Speaker notes:* {notes}")
        out.append("")
    return "\n".join(out)


def _extract_xlsx(path: Path) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(str(path), data_only=True, read_only=True)
    out: list[str] = [f"# Extracted from {path.name}", ""]
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        out.append(f"## Sheet: {sheet_name}")
        rows_written = 0
        for row in ws.iter_rows(values_only=True):
            if all(c is None for c in row):
                continue
            cells = " | ".join(("" if c is None else str(c)) for c in row)
            out.append(f"| {cells} |")
            rows_written += 1
            if rows_written >= 1000:
                out.append(f"| … truncated at 1,000 rows |")
                break
        out.append("")
    return "\n".join(out)


# ----------------------------------------------------------------------------
# Run-prompt section + agent prompt preamble + OpenAI inlining.
# ----------------------------------------------------------------------------

def render_for_run_prompt(sources: list[SourceFile], repo_root: Path = REPO_ROOT) -> str:
    """Body of the `## Source material` section written into the run prompt."""
    if not sources:
        return ""
    lines = [
        "The following operator-supplied files are the starting point for "
        "this run. Read them before conducting your own research. Where the "
        "source material conflicts with your default evidence base, name "
        "the conflict in your brief.",
        "",
    ]
    for s in sources:
        rel = s.readable.relative_to(repo_root).as_posix()
        size = format_size(s.size_bytes)
        if s.extracted:
            orig = s.original.relative_to(repo_root).as_posix()
            lines.append(f"- `{rel}` — extracted from `{orig}` ({size})")
        else:
            lines.append(f"- `{rel}` ({size}) — {s.description}")
    return "\n".join(lines)


def stage1_preamble(source_paths: list[str]) -> str:
    """The instruction injected into every Stage 1 agent's prompt."""
    if not source_paths:
        return ""
    lines = [
        "",
        "**BEFORE you research anything:** the operator attached the following "
        "source material to this run. Read every file listed below first, "
        "treat it as the primary starting point, and quote and engage with it "
        "directly in your brief. Where the source material conflicts with "
        "your default evidence base, name the conflict openly:",
    ]
    for p in source_paths:
        lines.append(f"  - `{p}`")
    return "\n".join(lines)


def inline_for_openai(
    source_paths: list[str],
    repo_root: Path = REPO_ROOT,
    max_chars_per_file: int = 30000,
) -> str:
    """Inline source-file text for the OpenAI-hosted Deep Research agent.

    OpenAI agents have no file tools — they only see the prompt text. Each
    file's readable form is inlined; PDFs are referenced by name only because
    we can't extract their text without adding another dep.
    """
    if not source_paths:
        return ""
    chunks: list[str] = [
        "\n\n--- OPERATOR-SUPPLIED SOURCE MATERIAL (inlined; you cannot read files) ---",
    ]
    for rel in source_paths:
        path = (repo_root / rel) if not Path(rel).is_absolute() else Path(rel)
        chunks.append(f"\n\n=== FILE: {path.name} ===")
        if not path.is_file():
            chunks.append(f"\n[File not found at {rel}]")
            continue
        try:
            text = path.read_text(encoding="utf-8", errors="ignore")
        except OSError as e:
            chunks.append(f"\n[Read failed: {e}]")
            continue
        if len(text) > max_chars_per_file:
            text = text[:max_chars_per_file] + f"\n\n[…truncated at {max_chars_per_file:,} chars]"
        chunks.append("\n" + text)
    return "".join(chunks)


# ----------------------------------------------------------------------------
# Archive.
# ----------------------------------------------------------------------------

def archive_sources(
    slug: str,
    outputs_dir: Path,
    archive_dir: Path,
    *,
    source_material: list[dict] | None = None,
) -> Path | None:
    """Copy the persistent library and any external inputs into the archive."""
    src = outputs_dir.parent / "sources" / "runs" / slug
    dst = archive_dir / "sources"
    copied = False
    if src.is_dir():
        shutil.copytree(src, dst)
        copied = True
    for record in source_material or []:
        runtime = Path(str(record.get("runtime_path") or ""))
        if not runtime.is_absolute():
            runtime = outputs_dir.parent / runtime
        archive_path = Path(str(record.get("archive_path") or ""))
        if not runtime.is_file() or not archive_path.parts:
            continue
        destination = (archive_dir / archive_path).resolve()
        try:
            destination.relative_to(archive_dir.resolve())
        except ValueError as exc:
            raise SourcePathError(
                f"Archived source path escapes the run archive: {archive_path}"
            ) from exc
        if destination.is_file():
            continue
        destination.parent.mkdir(parents=True, exist_ok=True)
        shutil.copy2(runtime, destination)
        copied = True
    return dst if copied else None
